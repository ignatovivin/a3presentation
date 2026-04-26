from __future__ import annotations

import base64
import copy
import math
import re
import zipfile
from datetime import UTC, datetime
from io import BytesIO
from pathlib import Path
from copy import deepcopy
from dataclasses import replace

from pptx.chart.data import CategoryChartData
from pptx.chart.axis import ValueAxis
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.xmlchemy import OxmlElement
from pptx import Presentation
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Pt

from a3presentation.domain.chart import ChartSpec, ChartType
from a3presentation.domain.presentation import (
    PresentationPlan,
    SlideContentBlock,
    SlideContentBlockKind,
    SlideKind,
    SlideSpec,
    TableBlock,
)
from a3presentation.domain.template import (
    GenerationMode,
    LayoutSpec,
    PlaceholderSpec,
    PlaceholderKind,
    PrototypeSlideSpec,
    TemplateManifest,
    TemplateShapeStyleSpec,
    TemplateTextStyleSpec,
)
from a3presentation.services.chart_style import CHART_STYLE_CONFIG
from a3presentation.services.chart_render_contract import (
    PRIMARY_AXIS,
    SECONDARY_AXIS,
    chart_axis_number_format,
    chart_axis_number_format_for_axis,
    render_chart_spec,
    uses_secondary_value_axis,
)
from a3presentation.services.layout_capacity import (
    LayoutCapacityProfile,
    derive_capacity_profile_for_geometry,
    geometry_policy_for_layout,
    profile_for_layout,
    runtime_profile_key_for_target,
    spacing_policy_for_layout,
)


class PptxGenerator:
    TOKEN_PATTERN = re.compile(r"{{\s*([a-zA-Z0-9_]+)\s*}}")
    RELATIONSHIP_NAMESPACE = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    SKIPPED_RELATIONSHIP_TYPES = {
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide",
    }
    EMU_PER_PT = 12700
    TITLE_CONTENT_GAP_EMU = 180000
    TITLE_BODY_GAP_NO_SUBTITLE_EMU = 300000
    CONTENT_FOOTER_GAP_EMU = 180000
    FOOTER_FONT_PT = 12.0
    COVER_TITLE_TOP_EMU = 651176
    COVER_TITLE_LEFT_EMU = 444249
    COVER_TITLE_WIDTH_EMU = 10693901
    COVER_TITLE_MIN_HEIGHT_EMU = 1422646
    COVER_META_LEFT_EMU = 444249
    COVER_META_WIDTH_EMU = 8200000
    COVER_META_DEFAULT_TOP_EMU = 2438400
    COVER_META_MIN_HEIGHT_EMU = 700000
    COVER_META_GAP_EMU = 220000
    COVER_BOTTOM_LIMIT_EMU = 6200000
    KPI_TITLE_TOP_EMU = 651176
    KPI_TITLE_LEFT_EMU = 444249
    KPI_TITLE_WIDTH_EMU = 10693901
    KPI_TITLE_HEIGHT_EMU = 720000
    KPI_DESCRIPTION_TOP_EMU = 1460000
    KPI_DESCRIPTION_LEFT_EMU = 444249
    KPI_DESCRIPTION_WIDTH_EMU = 8200000
    KPI_DESCRIPTION_HEIGHT_EMU = 900000
    KPI_METRIC_START_TOP_EMU = 3350000
    KPI_METRIC_LEFT_EMU = 444249
    KPI_METRIC_RIGHT_EMU = 6900000
    KPI_METRIC_WIDTH_EMU = 4100000
    KPI_METRIC_HEIGHT_EMU = 1200000
    KPI_METRIC_ROW_GAP_EMU = 650000
    FULL_CONTENT_LEFT_EMU = 442913
    FULL_CONTENT_WIDTH_EMU = 11198224
    FOOTER_TOP_EMU = 6384626
    FOOTER_HEIGHT_EMU = 260000
    DEFAULT_TEXT_MARGIN_X_EMU = 91440
    DEFAULT_TEXT_MARGIN_Y_EMU = 45720
    GEOMETRY_PROFILE_TOLERANCE_EMU = 120000
    KPI_FOURTH_CARD_LEFT_EMU = 6980000
    KPI_FOURTH_CARD_TOP_EMU = 4950000
    KPI_FOURTH_CARD_WIDTH_EMU = 3600000
    KPI_FOURTH_CARD_HEIGHT_EMU = 1350000
    BUILTIN_LAYOUT_KEYS = {
        "text_full_width",
        "dense_text_full_width",
        "list_full_width",
        "table",
        "image_text",
        "cards_3",
        "cards_kpi",
        "list_with_icons",
        "contacts",
        "cover",
    }
    _CHART_STYLE_CONFIG = None

    def __init__(self) -> None:
        self._active_manifest: TemplateManifest | None = None
        self._active_presentation: Presentation | None = None

    @classmethod
    def _chart_style_config(cls) -> dict:
        if cls._CHART_STYLE_CONFIG is None:
            cls._CHART_STYLE_CONFIG = copy.deepcopy(CHART_STYLE_CONFIG)
        return cls._CHART_STYLE_CONFIG

    @classmethod
    def _style_rgb(cls, key: str) -> RGBColor:
        value = cls._chart_style_config()[key].lstrip("#")
        return RGBColor.from_string(value)

    @classmethod
    def _series_color(cls, index: int) -> RGBColor:
        palette = cls._chart_style_config()["palette"]
        return RGBColor.from_string(palette[index % len(palette)].lstrip("#"))

    def generate(self, template_path: Path, manifest: TemplateManifest, plan: PresentationPlan, output_dir: Path) -> Path:
        self._active_manifest = manifest
        if manifest.generation_mode == GenerationMode.PROTOTYPE and manifest.prototype_slides:
            presentation = self._generate_from_prototypes(template_path, manifest, plan)
        elif self._should_use_direct_shape_binding(manifest, plan):
            presentation = self._generate_from_direct_shape_bindings(template_path, manifest, plan)
        else:
            presentation = self._generate_from_layouts(template_path, manifest, plan)
        self._active_presentation = presentation

        self._apply_core_properties(presentation, plan)
        output_dir.mkdir(parents=True, exist_ok=True)
        timestamp = datetime.now(UTC).strftime("%Y%m%dT%H%M%S%fZ")
        output_stem = self._build_output_stem(plan.title or plan.template_id, timestamp)
        output_path = output_dir / f"{output_stem}.pptx"
        presentation.save(str(output_path))
        self._validate_output_file(output_path, expected_slide_count=len(plan.slides))
        return output_path

    def _should_use_direct_shape_binding(self, manifest: TemplateManifest, plan: PresentationPlan) -> bool:
        if manifest.inventory.degradation_mode == "direct_shape_binding" and manifest.inventory.slides:
            return True
        return any(self._target_type_for_slide(slide) == "direct_shape_binding" for slide in plan.slides)

    def _build_output_stem(self, title: str, timestamp: str) -> str:
        normalized = re.sub(r"\s+", " ", (title or "").strip())
        safe = re.sub(r"[^A-Za-z0-9А-Яа-яЁё _-]+", "", normalized).strip(" ._-")
        safe = safe.replace(" ", "_")
        if not safe:
            safe = "presentation"
        return f"{safe[:48]}_{timestamp}"

    def _generate_from_prototypes(self, template_path: Path, manifest: TemplateManifest, plan: PresentationPlan) -> Presentation:
        output_presentation = Presentation(str(template_path))
        source_slides = [output_presentation.slides[index] for index in range(len(output_presentation.slides))]
        self._remove_all_slides(output_presentation)

        for slide_spec in plan.slides:
            prototype = self._resolve_prototype_slide(manifest, slide_spec)
            source_slide = source_slides[prototype.source_slide_index]
            target_slide = self._clone_slide(output_presentation, source_slide)
            self._replace_tokens_in_slide(target_slide, prototype, slide_spec, plan.title)
            if self._should_apply_prototype_layout_flow(prototype):
                self._apply_layout_expansion_and_flow(
                    target_slide,
                    self._runtime_profile_key_for_slide(slide_spec, prototype=prototype),
                    slide_spec,
                )

        return output_presentation

    def _should_apply_prototype_layout_flow(self, prototype: PrototypeSlideSpec) -> bool:
        if prototype.key not in self.BUILTIN_LAYOUT_KEYS:
            return False
        return not any(token.shape_name for token in prototype.tokens)

    def _generate_from_layouts(self, template_path: Path, manifest: TemplateManifest, plan: PresentationPlan) -> Presentation:
        presentation = Presentation(str(template_path))
        self._active_presentation = presentation
        self._remove_all_slides(presentation)
        self._apply_core_properties(presentation, plan)

        for slide_spec in plan.slides:
            layout = self._resolve_layout(manifest, slide_spec)
            slide_layout = presentation.slide_masters[layout.slide_master_index].slide_layouts[layout.slide_layout_index]
            slide = presentation.slides.add_slide(slide_layout)
            self._fill_slide_from_layout(slide, slide_spec, layout, plan.title)

        return presentation

    def _generate_from_direct_shape_bindings(self, template_path: Path, manifest: TemplateManifest, plan: PresentationPlan) -> Presentation:
        output_presentation = Presentation(str(template_path))
        source_slides = [output_presentation.slides[index] for index in range(len(output_presentation.slides))]
        self._remove_all_slides(output_presentation)

        for slide_spec in plan.slides:
            direct_target = self._resolve_direct_shape_binding_target(manifest, slide_spec)
            if direct_target is None:
                continue
            source_slide = source_slides[direct_target["source_index"]]
            target_slide = self._clone_slide(output_presentation, source_slide)
            self._fill_direct_shape_binding_slide(target_slide, slide_spec, plan.title, direct_target["components"])

        return output_presentation

    def _resolve_direct_shape_binding_target(self, manifest: TemplateManifest, slide_spec: SlideSpec) -> dict | None:
        components_by_id = {component.component_id: component for component in manifest.inventory.components}
        slide_targets = []
        for slide_inventory in manifest.inventory.slides:
            key = f"direct_slide_{slide_inventory.source_index}"
            components = [
                components_by_id[component_id]
                for component_id in slide_inventory.component_ids
                if component_id in components_by_id
                and getattr(components_by_id[component_id], "binding", None)
                and getattr(components_by_id[component_id], "shape_name", None)
            ]
            if not components:
                continue
            slide_targets.append(
                {
                    "key": key,
                    "source_index": slide_inventory.source_index,
                    "supported_slide_kinds": list(slide_inventory.supported_slide_kinds),
                    "components": components,
                }
            )

        target_key = self._target_key_for_slide(slide_spec)
        if target_key:
            preferred = next((item for item in slide_targets if item["key"] == target_key), None)
            if preferred is not None:
                return preferred
        return next((item for item in slide_targets if slide_spec.kind.value in item["supported_slide_kinds"]), slide_targets[0] if slide_targets else None)

    def _fill_direct_shape_binding_slide(self, slide, slide_spec: SlideSpec, presentation_title: str, components: list[object]) -> None:
        runtime_profile_key = self._runtime_profile_key_for_slide(slide_spec)
        layout_profile = profile_for_layout(runtime_profile_key)
        used_shape_names: set[str] = set()
        for component in components:
            shape_name = getattr(component, "shape_name", None)
            binding = getattr(component, "binding", None)
            if not shape_name or not binding:
                continue
            target_shape = next((shape for shape in slide.shapes if getattr(shape, "name", None) == shape_name), None)
            if target_shape is None:
                continue
            shape_profile = self._capacity_profile_for_shape(runtime_profile_key, target_shape, layout_profile)
            self._fill_shape_by_binding(
                target_shape,
                binding,
                slide_spec,
                presentation_title,
                shape_profile,
            )
            used_shape_names.add(shape_name)

        for shape in slide.shapes:
            if getattr(shape, "name", None) in used_shape_names:
                continue
            if getattr(shape, "has_text_frame", False):
                self._clear_placeholder(shape)

    def _apply_core_properties(self, presentation: Presentation, plan: PresentationPlan) -> None:
        props = presentation.core_properties
        props.title = plan.title
        props.author = plan.author or "a3presentation"
        if plan.subject:
            props.subject = plan.subject

    def _remove_all_slides(self, presentation: Presentation) -> None:
        for index in range(len(presentation.slides) - 1, -1, -1):
            slide_id = presentation.slides._sldIdLst[index]
            relationship_id = slide_id.rId
            presentation.part.drop_rel(relationship_id)
            del presentation.slides._sldIdLst[index]

    def _clone_slide(self, presentation: Presentation, source_slide):
        blank_layout = presentation.slide_layouts[6] if len(presentation.slide_layouts) > 6 else presentation.slide_layouts[-1]
        target_slide = presentation.slides.add_slide(blank_layout)
        self._clear_slide_shapes(target_slide)

        relationship_map: dict[str, str] = {}
        for relationship in source_slide.part.rels.values():
            if relationship.reltype in self.SKIPPED_RELATIONSHIP_TYPES:
                continue
            new_relationship_id = target_slide.part.relate_to(
                relationship.target_ref if relationship.is_external else relationship.target_part,
                relationship.reltype,
                relationship.is_external,
            )
            relationship_map[relationship.rId] = new_relationship_id

        source_background = source_slide._element.cSld.bg
        if source_background is not None:
            target_background = copy.deepcopy(source_background)
            self._remap_relationship_ids(target_background, relationship_map)
            if target_slide._element.cSld.bg is not None:
                target_slide._element.cSld.remove(target_slide._element.cSld.bg)
            target_slide._element.cSld.insert(0, target_background)

        for shape_element in list(source_slide.shapes._spTree.iterchildren()):
            if shape_element.tag.endswith("nvGrpSpPr") or shape_element.tag.endswith("grpSpPr"):
                continue
            cloned_element = copy.deepcopy(shape_element)
            self._remap_relationship_ids(cloned_element, relationship_map)
            target_slide.shapes._spTree.insert_element_before(cloned_element, "p:extLst")

        return target_slide

    def _clear_slide_shapes(self, slide) -> None:
        for shape_element in list(slide.shapes._spTree.iterchildren()):
            if shape_element.tag.endswith("nvGrpSpPr") or shape_element.tag.endswith("grpSpPr"):
                continue
            slide.shapes._spTree.remove(shape_element)

    def _remap_relationship_ids(self, element, relationship_map: dict[str, str]) -> None:
        for current_element in element.iter():
            for attr_name, attr_value in list(current_element.attrib.items()):
                if attr_value in relationship_map and attr_name.startswith(f"{{{self.RELATIONSHIP_NAMESPACE}}}"):
                    current_element.set(attr_name, relationship_map[attr_value])

    def _resolve_prototype_slide(self, manifest: TemplateManifest, slide_spec: SlideSpec) -> PrototypeSlideSpec:
        if slide_spec.kind == SlideKind.CHART:
            for prototype_slide in manifest.prototype_slides:
                if self._prototype_supports_chart(prototype_slide):
                    return prototype_slide

        target_key = self._target_key_for_slide(slide_spec)
        target_type = self._target_type_for_slide(slide_spec)
        if target_key and target_type in {None, "prototype"}:
            for prototype_slide in manifest.prototype_slides:
                if prototype_slide.key == target_key:
                    return prototype_slide

        for prototype_slide in manifest.prototype_slides:
            if slide_spec.kind.value in prototype_slide.supported_slide_kinds:
                return prototype_slide

        return manifest.prototype_slides[0]

    def _prototype_supports_chart(self, prototype: PrototypeSlideSpec) -> bool:
        if prototype.key == "chart" or SlideKind.CHART.value in prototype.supported_slide_kinds:
            return True
        return any(token.binding in {"chart", "chart_image"} for token in prototype.tokens)

    def _runtime_profile_key_for_slide(
        self,
        slide_spec: SlideSpec,
        *,
        layout: LayoutSpec | None = None,
        prototype: PrototypeSlideSpec | None = None,
    ) -> str:
        target = layout or prototype
        target_key = self._target_key_for_slide(slide_spec)
        target_type = self._target_type_for_slide(slide_spec)
        if target is None and self._active_manifest is not None and target_key:
            if target_type in {None, "layout", "auto_layout", "direct_shape_binding"}:
                target = next(
                    (item for item in self._active_manifest.layouts if item.key == target_key),
                    None,
                )
            if target is None and target_type in {None, "prototype"}:
                target = next(
                    (item for item in self._active_manifest.prototype_slides if item.key == target_key),
                    None,
                )
        return runtime_profile_key_for_target(
            target,
            fallback_layout_key=slide_spec.runtime_profile_key or target_key,
            slide_kind=slide_spec.kind.value,
        )

    def _target_key_for_slide(self, slide_spec: SlideSpec) -> str | None:
        if slide_spec.render_target is not None and slide_spec.render_target.key:
            return slide_spec.render_target.key
        return slide_spec.preferred_layout_key

    def _target_type_for_slide(self, slide_spec: SlideSpec) -> str | None:
        if slide_spec.render_target is None:
            return None
        return slide_spec.render_target.type.value

    def _replace_tokens_in_slide(self, slide, prototype: PrototypeSlideSpec, slide_spec: SlideSpec, presentation_title: str) -> None:
        token_values = self._build_token_value_map(slide_spec, presentation_title)
        used_shapes: set[str] = set()
        runtime_profile_key = self._runtime_profile_key_for_slide(slide_spec, prototype=prototype)
        layout_profile = profile_for_layout(runtime_profile_key)

        if getattr(slide_spec, "background_only", False):
            for shape in list(slide.shapes):
                if getattr(shape, "has_text_frame", False) or getattr(shape, "is_placeholder", False):
                    self._clear_placeholder(shape)
            return

        # Preferred path for real templates: bind by explicit shape name from manifest.
        for token_spec in prototype.tokens:
            if not token_spec.shape_name:
                continue
            target_shape = next((shape for shape in slide.shapes if shape.name == token_spec.shape_name), None)
            if target_shape is None:
                continue
            self._apply_shape_spec_metadata(target_shape, token_spec)
            shape_profile = self._capacity_profile_for_shape(runtime_profile_key, target_shape, layout_profile)
            self._fill_shape_by_binding(
                target_shape,
                token_spec.binding,
                slide_spec,
                presentation_title,
                shape_profile,
                placeholder_spec=token_spec,
            )
            self._apply_shape_spec_metadata(target_shape, token_spec, apply_text_style=False)
            if slide_spec.kind == SlideKind.CHART and token_spec.binding == "title":
                self._configure_title_text_frame(target_shape)
                self._apply_font_size(
                    target_shape,
                    self._component_font_size("chart", "title", fallback=35.0),
                )
            elif slide_spec.kind == SlideKind.CHART and token_spec.binding == "subtitle":
                self._configure_subtitle_text_frame(target_shape)
                self._apply_font_size(
                    target_shape,
                    self._component_font_size("chart", "subtitle", fallback=20.0),
                )
            used_shapes.add(token_spec.shape_name)

        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            if shape.name in used_shapes:
                continue

            original_text = shape.text or ""
            matches = self.TOKEN_PATTERN.findall(original_text)
            if not matches:
                continue

            normalized = original_text.strip()
            single_token_match = self.TOKEN_PATTERN.fullmatch(normalized)
            if single_token_match:
                token_name = single_token_match.group(1)
                token_value = token_values.get(token_name, "")
                shape_profile = self._capacity_profile_for_shape(runtime_profile_key, shape, layout_profile)
                if isinstance(token_value, list):
                    self._set_bullets(shape, token_value, shape_profile)
                else:
                    self._set_text(shape, str(token_value), shape_profile)
                continue

            replaced_text = original_text
            for token_name in matches:
                token_value = token_values.get(token_name, "")
                if isinstance(token_value, list):
                    token_value = "\n".join(token_value)
                replaced_text = re.sub(r"{{\s*" + re.escape(token_name) + r"\s*}}", str(token_value), replaced_text)
            shape_profile = self._capacity_profile_for_shape(runtime_profile_key, shape, layout_profile)
            self._set_text(shape, replaced_text, shape_profile)

    def _build_token_value_map(self, slide_spec: SlideSpec, presentation_title: str) -> dict[str, str | list[str]]:
        token_map: dict[str, str | list[str]] = {
            "title": slide_spec.title or "",
            "subtitle": slide_spec.subtitle or "",
            "text": slide_spec.text or "",
            "body": slide_spec.text or "",
            "summary": slide_spec.text or "",
            "notes": slide_spec.notes or "",
            "footer": slide_spec.notes or presentation_title,
            "presentation_title": presentation_title,
            "presentation_name": presentation_title,
            "cover_title": presentation_title if not slide_spec.title else slide_spec.title,
            "cover_meta": slide_spec.notes or presentation_title,
            "left_note": slide_spec.left_bullets[0] if slide_spec.left_bullets else "",
            "right_note": slide_spec.right_bullets[0] if slide_spec.right_bullets else "",
            "main_text": slide_spec.text or "",
            "secondary_text": slide_spec.notes or "",
            "left_text": "\n".join(slide_spec.left_bullets) if slide_spec.left_bullets else slide_spec.text or "",
            "right_list": slide_spec.right_bullets or slide_spec.bullets,
            "contact_title": slide_spec.title or "",
            "contact_name_or_title": slide_spec.title or "",
            "contact_role": slide_spec.subtitle or "",
            "contact_phone": slide_spec.left_bullets[0] if slide_spec.left_bullets else "",
            "contact_email": slide_spec.right_bullets[0] if slide_spec.right_bullets else "",
            "address": slide_spec.text or "",
            "phone": slide_spec.left_bullets[0] if slide_spec.left_bullets else "",
            "email": slide_spec.right_bullets[0] if slide_spec.right_bullets else "",
            "website": slide_spec.notes or "",
            "bullets": slide_spec.bullets,
            "left_bullets": slide_spec.left_bullets,
            "right_bullets": slide_spec.right_bullets,
        }

        for index in range(1, 11):
            token_map[f"bullet_{index}"] = slide_spec.bullets[index - 1] if len(slide_spec.bullets) >= index else ""
            token_map[f"left_bullet_{index}"] = slide_spec.left_bullets[index - 1] if len(slide_spec.left_bullets) >= index else ""
            token_map[f"right_bullet_{index}"] = slide_spec.right_bullets[index - 1] if len(slide_spec.right_bullets) >= index else ""
            token_map[f"card_{index}"] = slide_spec.bullets[index - 1] if len(slide_spec.bullets) >= index else ""
            token_map[f"icon_{index}"] = ""

        if slide_spec.table is not None:
            token_map["table"] = [" | ".join(row) for row in slide_spec.table.rows]
            token_map["table_headers"] = slide_spec.table.headers
            token_map["table_rows"] = [" | ".join(row) for row in slide_spec.table.rows]
        if slide_spec.image_base64:
            token_map["image"] = slide_spec.image_base64

        return token_map

    def _resolve_layout(self, manifest: TemplateManifest, slide_spec: SlideSpec) -> LayoutSpec:
        target_key = self._target_key_for_slide(slide_spec)
        target_type = self._target_type_for_slide(slide_spec)
        if target_key and target_type in {None, "layout", "auto_layout", "direct_shape_binding"}:
            for layout in manifest.layouts:
                if layout.key == target_key:
                    return layout

        for layout in manifest.layouts:
            if slide_spec.kind.value in layout.supported_slide_kinds:
                return layout

        if manifest.default_layout_key:
            for layout in manifest.layouts:
                if layout.key == manifest.default_layout_key:
                    return layout

        if not manifest.layouts:
            raise ValueError(f"Template '{manifest.template_id}' does not contain layouts")
        return manifest.layouts[0]

    def _effective_placeholder_kind(
        self,
        placeholder_spec: PlaceholderSpec,
        *,
        logical_layout_key: str,
        slide_kind: SlideKind,
    ) -> PlaceholderKind:
        if placeholder_spec.binding:
            return placeholder_spec.kind

        idx = placeholder_spec.idx
        if idx is None:
            return placeholder_spec.kind

        if logical_layout_key == "cover":
            role_map = {
                0: PlaceholderKind.TITLE,
                15: PlaceholderKind.FOOTER,
            }
            return role_map.get(idx, PlaceholderKind.UNKNOWN)

        if logical_layout_key in {"text_full_width", "dense_text_full_width", "list_full_width"}:
            role_map = {
                0: PlaceholderKind.TITLE,
                13: PlaceholderKind.SUBTITLE,
                14: PlaceholderKind.BODY,
                17: PlaceholderKind.FOOTER,
            }
            return role_map.get(idx, PlaceholderKind.UNKNOWN)

        if logical_layout_key == "table" or slide_kind == SlideKind.CHART:
            role_map = {
                0: PlaceholderKind.TITLE,
                13: PlaceholderKind.SUBTITLE,
                14: PlaceholderKind.CHART if slide_kind == SlideKind.CHART else PlaceholderKind.TABLE,
                15: PlaceholderKind.FOOTER,
            }
            return role_map.get(idx, PlaceholderKind.UNKNOWN)

        if logical_layout_key == "image_text":
            role_map = {
                0: PlaceholderKind.TITLE,
                13: PlaceholderKind.SUBTITLE,
                14: PlaceholderKind.BODY,
                16: PlaceholderKind.IMAGE,
                17: PlaceholderKind.FOOTER,
            }
            return role_map.get(idx, PlaceholderKind.UNKNOWN)

        if logical_layout_key in {"cards_3", "cards_kpi"}:
            role_map = {
                0: PlaceholderKind.TITLE,
                11: PlaceholderKind.BODY,
                12: PlaceholderKind.BODY,
                13: PlaceholderKind.BODY,
            }
            if logical_layout_key == "cards_3":
                role_map[15] = PlaceholderKind.FOOTER
            return role_map.get(idx, PlaceholderKind.UNKNOWN)

        if logical_layout_key == "list_with_icons":
            role_map = {
                0: PlaceholderKind.TITLE,
                12: PlaceholderKind.BODY,
                14: PlaceholderKind.BODY,
                21: PlaceholderKind.FOOTER,
            }
            return role_map.get(idx, PlaceholderKind.UNKNOWN)

        if logical_layout_key == "contacts":
            return PlaceholderKind.BODY if idx in {10, 11, 12, 13} else PlaceholderKind.UNKNOWN

        return placeholder_spec.kind

    def _fill_slide_from_layout(self, slide, slide_spec: SlideSpec, layout: LayoutSpec, presentation_title: str) -> None:
        target_layout_key = self._target_key_for_slide(slide_spec) or layout.key
        runtime_profile_key = self._runtime_profile_key_for_slide(slide_spec, layout=layout)
        use_builtin_flow = target_layout_key in self.BUILTIN_LAYOUT_KEYS
        explicit_background_xml = getattr(slide_spec, "background_xml", None)
        if explicit_background_xml:
            self._apply_background_xml(slide, explicit_background_xml)
        elif layout.background_xml and not self._background_xml_requires_relationships(layout.background_xml):
            self._apply_background_xml(slide, layout.background_xml)
            self._apply_background_style(slide, layout.background_style)
        elif self._background_xml_requires_relationships(layout.background_xml) or layout.background_image_base64:
            if self._active_manifest is not None and self._active_manifest.generation_mode == GenerationMode.LAYOUT:
                self._apply_background_style(slide, layout.background_style)
            else:
                self._apply_layout_background_image(slide, layout)
        else:
            self._apply_background_style(slide, layout.background_style)
        if getattr(slide_spec, "background_only", False):
            for placeholder in slide.placeholders:
                self._clear_placeholder(placeholder)
            return
        if use_builtin_flow and target_layout_key == "cover":
            self._populate_cover_slide(slide, slide_spec)
            return
        if use_builtin_flow and runtime_profile_key == "cards_kpi":
            for placeholder in slide.placeholders:
                self._clear_placeholder(placeholder)
            self._populate_kpi_cards_slide(slide, slide_spec, presentation_title)
            return
        layout_profile = profile_for_layout(runtime_profile_key)
        placeholders = {placeholder.placeholder_format.idx: placeholder for placeholder in slide.placeholders}
        used_placeholder_indices: set[int] = set()
        materialized_roles: set[str] = set()

        for placeholder_spec in layout.placeholders:
            shape = None
            if placeholder_spec.idx is not None and placeholder_spec.idx in placeholders:
                shape = placeholders[placeholder_spec.idx]
                used_placeholder_indices.add(placeholder_spec.idx)
            elif not use_builtin_flow:
                shape = self._materialize_shape_from_layout_spec(
                    slide,
                    placeholder_spec,
                    slide_spec,
                    layout_profile,
                    runtime_profile_key,
                    materialized_roles,
                )
            if shape is None:
                continue
            self._apply_shape_spec_metadata(shape, placeholder_spec)
            if placeholder_spec.binding:
                self._fill_shape_by_binding(
                    shape,
                    placeholder_spec.binding,
                    slide_spec,
                    presentation_title,
                    layout_profile,
                    placeholder_spec=placeholder_spec,
                )
                if placeholder_spec.binding != "table":
                    self._apply_shape_spec_metadata(shape, placeholder_spec)
                continue
            effective_kind = self._effective_placeholder_kind(
                placeholder_spec,
                logical_layout_key=runtime_profile_key,
                slide_kind=slide_spec.kind,
            )
            if effective_kind == PlaceholderKind.UNKNOWN:
                self._clear_placeholder(shape)
            elif effective_kind == PlaceholderKind.TITLE:
                self._set_text(shape, slide_spec.title or "", layout_profile)
            elif effective_kind == PlaceholderKind.SUBTITLE:
                if (slide_spec.subtitle or "").strip():
                    self._set_text(shape, slide_spec.subtitle or "", layout_profile)
                else:
                    self._clear_placeholder(shape)
            elif effective_kind == PlaceholderKind.BODY:
                if runtime_profile_key in {"cards_3", "cards_kpi"}:
                    self._fill_card_body(slide, shape, slide_spec, layout_profile)
                else:
                    self._fill_body(shape, slide_spec, layout_profile)
            elif effective_kind == PlaceholderKind.FOOTER:
                self._set_text(shape, presentation_title, layout_profile)
            elif effective_kind == PlaceholderKind.TABLE:
                self._fill_table_or_chart(shape, slide_spec, placeholder_spec)
            elif effective_kind == PlaceholderKind.CHART:
                self._fill_chart(shape, slide_spec, placeholder_spec)
            self._apply_shape_spec_metadata(
                shape,
                placeholder_spec,
                apply_text_style=not (runtime_profile_key in {"cards_3", "cards_kpi"} and effective_kind == PlaceholderKind.BODY),
                preserve_font_size=effective_kind in {PlaceholderKind.BODY, PlaceholderKind.FOOTER},
            )
            if effective_kind == PlaceholderKind.BODY:
                self._clamp_text_frame_font_size(shape, layout_profile)
            if effective_kind == PlaceholderKind.SUBTITLE and slide_spec.kind == SlideKind.CHART:
                self._configure_subtitle_text_frame(shape)
                body_style = self._active_manifest.theme.master_text_styles.get("body") if self._active_manifest is not None else None
                self._apply_font_size(
                    shape,
                    body_style.font_size_pt if body_style is not None and body_style.font_size_pt else self._table_subtitle_font_size_points(slide_spec.subtitle or ""),
                )

        for placeholder in slide.placeholders:
            placeholder_idx = placeholder.placeholder_format.idx
            if placeholder_idx in used_placeholder_indices:
                continue
            self._clear_placeholder(placeholder)

        if getattr(slide_spec, "background_only", False):
            for placeholder in slide.placeholders:
                self._clear_placeholder(placeholder)
            return

        if use_builtin_flow:
            self._apply_layout_expansion_and_flow(slide, runtime_profile_key, slide_spec)
        if runtime_profile_key == "cards_kpi":
            self._add_kpi_fourth_card(slide, slide_spec, layout_profile)

    def _materialize_shape_from_layout_spec(
        self,
        slide,
        placeholder_spec: PlaceholderSpec,
        slide_spec: SlideSpec,
        layout_profile: LayoutCapacityProfile,
        runtime_profile_key: str,
        materialized_roles: set[str],
    ):
        role = self._materialized_role_key(placeholder_spec, slide_spec, runtime_profile_key)
        if role is None or role in materialized_roles:
            return None
        geometry_values = (
            placeholder_spec.left_emu,
            placeholder_spec.top_emu,
            placeholder_spec.width_emu,
            placeholder_spec.height_emu,
        )
        if not all(isinstance(value, int) and value > 0 for value in geometry_values):
            return None
        if role not in {"title", "subtitle", "body", "footer"}:
            return None

        shape = slide.shapes.add_textbox(
            placeholder_spec.left_emu,
            placeholder_spec.top_emu,
            placeholder_spec.width_emu,
            placeholder_spec.height_emu,
        )
        if placeholder_spec.shape_name:
            shape.name = placeholder_spec.shape_name
        self._apply_shape_spec_metadata(shape, placeholder_spec)
        materialized_roles.add(role)
        return shape

    def _materialized_role_key(
        self,
        placeholder_spec: PlaceholderSpec,
        slide_spec: SlideSpec,
        runtime_profile_key: str,
    ) -> str | None:
        binding = (placeholder_spec.binding or "").strip().lower()
        if binding in {"title", "cover_title"}:
            return "title"
        if binding in {"subtitle", "cover_meta", "presentation_name"}:
            return "subtitle"
        if binding in {
            "body",
            "text",
            "main_text",
            "secondary_text",
            "summary",
            "notes",
            "bullets",
            "left_bullets",
            "right_bullets",
        }:
            return "body"
        if binding == "footer":
            return "footer"
        if binding:
            return None

        effective_kind = self._effective_placeholder_kind(
            placeholder_spec,
            logical_layout_key=runtime_profile_key,
            slide_kind=slide_spec.kind,
        )
        if effective_kind == PlaceholderKind.TITLE:
            return "title"
        if effective_kind == PlaceholderKind.SUBTITLE:
            return "subtitle"
        if effective_kind == PlaceholderKind.BODY:
            return "body"
        if effective_kind == PlaceholderKind.FOOTER:
            return "footer"
        return None

    def _apply_background_xml(self, slide, background_xml: str | None) -> None:
        if not background_xml:
            return
        try:
            background = parse_xml(background_xml)
            existing = slide._element.cSld.bg
            if existing is not None:
                slide._element.cSld.remove(existing)
            slide._element.cSld.insert(0, background)
        except Exception:
            pass

    def _background_xml_requires_relationships(self, background_xml: str | None) -> bool:
        if not background_xml:
            return False
        return "embed=" in background_xml or "link=" in background_xml

    def _apply_layout_background_image(self, slide, layout: LayoutSpec) -> None:
        blob = None
        if layout.background_image_base64:
            try:
                blob = base64.b64decode(layout.background_image_base64)
            except Exception:
                blob = None
        if blob is None:
            blob = self._background_image_blob_for_layout(layout)
        if blob is None or self._active_presentation is None:
            self._apply_background_style(slide, layout.background_style)
            return
        picture = slide.shapes.add_picture(
            BytesIO(blob),
            0,
            0,
            width=self._active_presentation.slide_width,
            height=self._active_presentation.slide_height,
        )
        sp_tree = slide.shapes._spTree
        pic_element = picture._element
        sp_tree.remove(pic_element)
        insert_at = 2
        for index, child in enumerate(sp_tree):
            if not (child.tag.endswith("nvGrpSpPr") or child.tag.endswith("grpSpPr")):
                insert_at = index
                break
        sp_tree.insert(insert_at, pic_element)

    def _background_image_blob_for_layout(self, layout: LayoutSpec) -> bytes | None:
        if self._active_manifest is None or self._active_presentation is None:
            return None
        donor_layout = self._active_presentation.slide_masters[layout.slide_master_index].slide_layouts[
            layout.slide_layout_index
        ]
        background = donor_layout._element.cSld.bg
        if background is not None:
            blob = self._background_image_blob_from_part(donor_layout.part, background)
            if blob is not None:
                return blob
        donor_master = self._active_presentation.slide_masters[layout.slide_master_index]
        master_background = donor_master._element.cSld.bg
        if master_background is not None:
            return self._background_image_blob_from_part(donor_master.part, master_background)
        return None

    def _background_image_blob_from_part(self, part, background) -> bytes | None:
        relationship_id = None
        relationship_attr = f"{{{self.RELATIONSHIP_NAMESPACE}}}embed"
        for element in background.iter():
            relationship_id = element.get(relationship_attr)
            if relationship_id:
                break
        if not relationship_id:
            return None
        rel = part.rels.get(relationship_id)
        target_part = getattr(rel, "target_part", None) if rel is not None else None
        return getattr(target_part, "blob", None)

    def _apply_background_style(self, slide, background_style: TemplateShapeStyleSpec | None) -> None:
        if background_style is None:
            return
        if background_style.fill_type == "solid" and background_style.fill_color:
            try:
                color = background_style.fill_color.lstrip("#").upper()
                background = parse_xml(
                    '<p:bg xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
                    'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
                    "<p:bgPr>"
                    f'<a:solidFill><a:srgbClr val="{color}"/></a:solidFill>'
                    "<a:effectLst/>"
                    "</p:bgPr>"
                    "</p:bg>"
                )
                existing = slide._element.cSld.bg
                if existing is not None:
                    slide._element.cSld.remove(existing)
                slide._element.cSld.insert(0, background)
            except Exception:
                pass

    def _capacity_profile_for_shape(
        self,
        layout_key: str,
        shape,
        base_profile: LayoutCapacityProfile,
    ) -> LayoutCapacityProfile:
        width = getattr(shape, "width", None)
        height = getattr(shape, "height", None)
        if width is None or height is None:
            return base_profile
        reference_width = None
        reference_height = None
        placeholder_idx = None
        if getattr(shape, "is_placeholder", False):
            try:
                placeholder_idx = shape.placeholder_format.idx
            except Exception:
                placeholder_idx = None
        base_geometry = geometry_policy_for_layout(layout_key)
        reference_placeholder = base_geometry.placeholders.get(placeholder_idx or -1)
        if reference_placeholder is None and placeholder_idx in {None, 14}:
            reference_placeholder = base_geometry.placeholders.get(14)
        if reference_placeholder is not None:
            reference_width = reference_placeholder.width_emu
            reference_height = reference_placeholder.height_emu
            if (
                abs(width - reference_width) <= self.GEOMETRY_PROFILE_TOLERANCE_EMU
                and abs(height - reference_height) <= self.GEOMETRY_PROFILE_TOLERANCE_EMU
            ):
                return base_profile
        return derive_capacity_profile_for_geometry(
            layout_key,
            width_emu=width,
            height_emu=height,
            reference_width_emu=reference_width,
            reference_height_emu=reference_height,
        )

    def _apply_layout_expansion_and_flow(self, slide, layout_key: str, slide_spec: SlideSpec | None = None) -> None:
        geometry_layout_key = "text_full_width" if layout_key == "dense_text_full_width" else layout_key
        if geometry_layout_key in {"text_full_width", "list_full_width"}:
            self._expand_text_full_width_layout(slide)
        elif geometry_layout_key == "image_text":
            self._expand_image_text_layout(slide)
        elif geometry_layout_key == "table":
            self._expand_table_layout(slide)
        elif geometry_layout_key in {"cards_3", "cards_kpi"}:
            self._expand_cards_layout(slide, geometry_layout_key)
        elif geometry_layout_key == "list_with_icons":
            self._expand_list_with_icons_layout(slide)
        elif geometry_layout_key == "contacts":
            self._expand_contacts_layout(slide)
        else:
            return

        self._adjust_title_and_flow(slide, layout_key, slide_spec)

    def _apply_shape_spec_metadata(self, shape, spec, *, apply_text_style: bool = True, preserve_font_size: bool = False) -> None:
        geometry_values = (spec.left_emu, spec.top_emu, spec.width_emu, spec.height_emu)
        if all(isinstance(value, int) and value > 0 for value in geometry_values):
            try:
                shape.left = spec.left_emu
                shape.top = spec.top_emu
                shape.width = spec.width_emu
                shape.height = spec.height_emu
            except Exception:
                pass

        if not getattr(shape, "has_text_frame", False):
            return

        text_frame = shape.text_frame
        if isinstance(spec.margin_left_emu, int) and spec.margin_left_emu >= 0:
            text_frame.margin_left = spec.margin_left_emu
        if isinstance(spec.margin_right_emu, int) and spec.margin_right_emu >= 0:
            text_frame.margin_right = spec.margin_right_emu
        if isinstance(spec.margin_top_emu, int) and spec.margin_top_emu >= 0:
            text_frame.margin_top = spec.margin_top_emu
        if isinstance(spec.margin_bottom_emu, int) and spec.margin_bottom_emu >= 0:
            text_frame.margin_bottom = spec.margin_bottom_emu
        shape_style = getattr(spec, "shape_style", None)
        text_style = getattr(spec, "text_style", None) if apply_text_style else None
        paragraph_styles = getattr(spec, "paragraph_styles", None) if apply_text_style else None
        if shape_style is not None:
            self._apply_shape_style(shape, shape_style)
            if isinstance(shape_style.inset_left_emu, int) and shape_style.inset_left_emu >= 0:
                text_frame.margin_left = shape_style.inset_left_emu
            if isinstance(shape_style.inset_right_emu, int) and shape_style.inset_right_emu >= 0:
                text_frame.margin_right = shape_style.inset_right_emu
            if isinstance(shape_style.inset_top_emu, int) and shape_style.inset_top_emu >= 0:
                text_frame.margin_top = shape_style.inset_top_emu
            if isinstance(shape_style.inset_bottom_emu, int) and shape_style.inset_bottom_emu >= 0:
                text_frame.margin_bottom = shape_style.inset_bottom_emu
            self._apply_vertical_anchor(text_frame, shape_style.vertical_anchor)
        if text_style is not None:
            self._apply_vertical_anchor(text_frame, text_style.vertical_anchor)
            self._apply_text_style(text_frame, text_style, preserve_font_size=preserve_font_size)
        if paragraph_styles is not None and paragraph_styles.level_styles:
            self._apply_paragraph_style_catalog(text_frame, paragraph_styles.level_styles)

    def _apply_vertical_anchor(self, text_frame, anchor: str | None) -> None:
        if not anchor:
            return
        try:
            body_pr = text_frame._txBody.bodyPr
            body_pr.set("anchor", anchor)
        except Exception:
            pass

    def _apply_text_style(
        self,
        text_frame,
        style: TemplateTextStyleSpec,
        *,
        apply_font_family: bool = False,
        preserve_font_size: bool = False,
    ) -> None:
        for paragraph in text_frame.paragraphs:
            if style.line_spacing is not None:
                paragraph.line_spacing = style.line_spacing
            if style.space_after_pt is not None:
                paragraph.space_after = Pt(style.space_after_pt)
            if style.space_before_pt is not None:
                paragraph.space_before = Pt(style.space_before_pt)
                if not paragraph.runs and paragraph.text:
                    run = paragraph.add_run()
                    run.text = paragraph.text
                    paragraph.text = ""
            for run in paragraph.runs:
                if style.font_size_pt is not None and not preserve_font_size:
                    run.font.size = Pt(style.font_size_pt)
                if apply_font_family and style.font_family:
                    self._apply_run_font_family(run, style.font_family)
                if style.bold is not None:
                    run.font.bold = style.bold
                if style.italic is not None:
                    run.font.italic = style.italic
                if style.underline is not None:
                    run.font.underline = style.underline
                if style.color:
                    try:
                        run.font.color.rgb = RGBColor.from_string(style.color.lstrip("#"))
                    except Exception:
                        pass

    def _apply_paragraph_style_catalog(self, text_frame, level_styles: dict[str, TemplateTextStyleSpec]) -> None:
        for paragraph in text_frame.paragraphs:
            level_key = str(getattr(paragraph, "level", 0))
            style = level_styles.get(level_key) or level_styles.get("0")
            if style is None:
                continue
            self._apply_paragraph_style(paragraph, style)

    def _apply_paragraph_style(self, paragraph, style: TemplateTextStyleSpec) -> None:
        if style.line_spacing is not None:
            paragraph.line_spacing = style.line_spacing
        if style.space_after_pt is not None:
            paragraph.space_after = Pt(style.space_after_pt)
        if style.space_before_pt is not None:
            paragraph.space_before = Pt(style.space_before_pt)
        ppr = paragraph._p.get_or_add_pPr()
        if style.margin_left_emu is not None:
            ppr.set("marL", str(style.margin_left_emu))
        if style.margin_right_emu is not None:
            ppr.set("marR", str(style.margin_right_emu))
        if style.default_tab_size_emu is not None:
            ppr.set("defTabSz", str(style.default_tab_size_emu))
        if style.rtl:
            ppr.set("rtl", "1")
        if style.hanging_emu is not None:
            ppr.set("indent", str(-abs(style.hanging_emu)))
        elif style.indent_emu is not None:
            ppr.set("indent", str(style.indent_emu))
        if style.bullet_type == "char" and style.bullet_char:
            for child in list(ppr):
                if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
                    ppr.remove(child)
            bullet = OxmlElement("a:buChar")
            bullet.set("char", style.bullet_char)
            if style.bullet_font:
                bullet.set("typeface", style.bullet_font)
            ppr.insert(0, bullet)
        self._apply_text_style_to_paragraph(paragraph, style)

    def _apply_text_style_to_paragraph(
        self,
        paragraph,
        style: TemplateTextStyleSpec,
        *,
        apply_font_family: bool = False,
    ) -> None:
        if not paragraph.runs and paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""
        for run in paragraph.runs:
            if style.font_size_pt is not None:
                run.font.size = Pt(style.font_size_pt)
            if apply_font_family and style.font_family:
                self._apply_run_font_family(run, style.font_family)
            if style.bold is not None:
                run.font.bold = style.bold
            if style.color:
                try:
                    run.font.color.rgb = RGBColor.from_string(style.color.lstrip("#"))
                except Exception:
                    pass

    def _clamp_text_frame_font_size(self, shape, layout_profile: LayoutCapacityProfile) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is None:
                    continue
                points = run.font.size.pt
                clamped = min(max(points, layout_profile.min_font_pt), layout_profile.max_font_pt)
                if clamped != points:
                    run.font.size = Pt(clamped)

    def _apply_shape_style(self, shape, style: TemplateShapeStyleSpec) -> None:
        try:
            sp_pr = shape._element.spPr
        except Exception:
            return
        if sp_pr is None:
            return
        if style.fill_type == "solid" and style.fill_color:
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor.from_string(style.fill_color.lstrip("#"))
            except Exception:
                pass
        line = None
        try:
            line = shape.line
        except Exception:
            line = None
        if line is not None and style.line_color:
            try:
                line.color.rgb = RGBColor.from_string(style.line_color.lstrip("#"))
            except Exception:
                pass
        ln = sp_pr.get_or_add_ln()
        if style.line_compound:
            ln.set("cmpd", style.line_compound)
        if style.line_cap:
            ln.set("cap", style.line_cap)
        if style.line_join == "bevel" and not ln.xpath("./a:bevel"):
            ln.append(OxmlElement("a:bevel"))
        if style.theme_fill_ref and not sp_pr.xpath("./a:style/a:fillRef"):
            style_el = OxmlElement("a:style")
            fill_ref = OxmlElement("a:fillRef")
            fill_ref.set("idx", style.theme_fill_ref)
            fill_scheme = OxmlElement("a:schemeClr")
            fill_scheme.set("val", "accent1")
            fill_ref.append(fill_scheme)
            style_el.append(fill_ref)
            if style.theme_line_ref:
                ln_ref = OxmlElement("a:lnRef")
                ln_ref.set("idx", style.theme_line_ref)
                line_scheme = OxmlElement("a:schemeClr")
                line_scheme.set("val", "accent1")
                ln_ref.append(line_scheme)
                style_el.append(ln_ref)
            sp_pr.append(style_el)

    def _populate_cover_slide(self, slide, slide_spec: SlideSpec) -> None:
        cover_title_left = self._component_spacing_emu("cover", "title_left_emu", self.COVER_TITLE_LEFT_EMU)
        cover_title_top = self._component_spacing_emu("cover", "title_top_emu", self.COVER_TITLE_TOP_EMU)
        cover_title_width = self._component_spacing_emu("cover", "title_width_emu", self.COVER_TITLE_WIDTH_EMU)
        cover_title_min_height = self._component_spacing_emu("cover", "title_min_height_emu", self.COVER_TITLE_MIN_HEIGHT_EMU)
        cover_meta_left = self._component_spacing_emu("cover", "meta_left_emu", self.COVER_META_LEFT_EMU)
        cover_meta_top = self._component_spacing_emu("cover", "meta_top_emu", self.COVER_META_DEFAULT_TOP_EMU)
        cover_meta_width = self._component_spacing_emu("cover", "meta_width_emu", self.COVER_META_WIDTH_EMU)
        cover_title_font_pt = self._component_font_size("cover", "title", fallback=46.0)
        cover_meta_font_pt = self._component_font_size("cover", "meta", fallback=22.0)
        cover_text_color = self._component_font_color("cover", "title", fallback=RGBColor(0xF5, 0xF9, 0xFE))
        cover_meta_color = self._component_font_color("cover", "meta", fallback=RGBColor(0xF5, 0xF9, 0xFE))
        title_shape = self._find_cover_title_shape(slide)
        meta_shape = self._find_cover_meta_shape(slide)
        keep_shape_ids = {
            shape.shape_id
            for shape in (title_shape, meta_shape)
            if shape is not None
        }

        if title_shape is not None:
            title_shape.left = cover_title_left
            title_shape.top = cover_title_top
            title_shape.width = cover_title_width
            title_shape.height = cover_title_min_height
            self._set_cover_text(
                title_shape,
                slide_spec.title or "",
                font_size=Pt(cover_title_font_pt),
                bold=True,
                color=cover_text_color,
                align=PP_ALIGN.LEFT,
            )

        if meta_shape is not None:
            meta_text = (slide_spec.notes or "").strip()
            if meta_text:
                meta_shape.left = cover_meta_left
                meta_shape.top = cover_meta_top
                meta_shape.width = cover_meta_width
                meta_shape.height = 1400000
                self._set_cover_text(
                    meta_shape,
                    meta_text,
                    font_size=Pt(cover_meta_font_pt),
                    bold=False,
                    color=cover_meta_color,
                    align=PP_ALIGN.LEFT,
                )
            else:
                self._clear_placeholder(meta_shape)

        for shape in slide.shapes:
            if shape.shape_id in keep_shape_ids:
                continue
            if getattr(shape, "is_placeholder", False):
                self._clear_placeholder(shape)

        if title_shape is None:
            title_shape = slide.shapes.add_textbox(
                cover_title_left,
                cover_title_top,
                cover_title_width,
                cover_title_min_height,
            )
            self._set_cover_text(
                title_shape,
                slide_spec.title or "",
                font_size=Pt(cover_title_font_pt),
                bold=True,
                color=cover_text_color,
                align=PP_ALIGN.LEFT,
            )

        if meta_shape is None and (slide_spec.notes or "").strip():
            meta_shape = slide.shapes.add_textbox(442913, 6120605, 3371850, 277813)
            meta_shape.left = cover_meta_left
            meta_shape.top = cover_meta_top
            meta_shape.width = cover_meta_width
            meta_shape.height = 1400000
            self._set_cover_text(
                meta_shape,
                (slide_spec.notes or "").strip(),
                font_size=Pt(cover_meta_font_pt),
                bold=False,
                color=cover_meta_color,
                align=PP_ALIGN.LEFT,
            )

        self._adjust_cover_layout(title_shape, meta_shape)

    def _find_cover_title_shape(self, slide):
        candidates = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.top < 2500000
            and shape.width > 5000000
        ]
        if candidates:
            return min(candidates, key=lambda shape: (shape.top, shape.left))
        return None

    def _find_cover_meta_shape(self, slide):
        candidates = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and shape.top > 5000000
            and shape.width < 5000000
        ]
        if candidates:
            return min(candidates, key=lambda shape: (shape.top, shape.left))
        return None

    def _set_cover_text(self, shape, text: str, *, font_size, bold: bool, color: RGBColor, align) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        self._apply_text_frame_margins(text_frame)
        lines = text.splitlines() or [text]
        for index, line in enumerate(lines):
            paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
            paragraph.alignment = align
            self._apply_paragraph_spacing(paragraph, "cover", "text_full_width")
            run = paragraph.add_run()
            run.text = line
            run.font.size = font_size
            run.font.bold = bold
            run.font.color.rgb = color
            theme_role = "title" if bold else "other"
            theme_style = self._fallback_theme_text_style(theme_role)
            if theme_style.font_family:
                self._apply_run_font_family(run, theme_style.font_family)

    def _apply_run_font_family(self, run, font_family: str) -> None:
        run.font.name = font_family
        try:
            r_pr = run._r.get_or_add_rPr()
            for tag in ("a:latin", "a:ea", "a:cs"):
                existing = r_pr.find(f"{{http://schemas.openxmlformats.org/drawingml/2006/main}}{tag.split(':', 1)[1]}")
                if existing is None:
                    existing = OxmlElement(tag)
                    r_pr.append(existing)
                existing.set("typeface", font_family)
        except Exception:
            pass

    def _adjust_cover_layout(self, title_shape, meta_shape) -> None:
        if title_shape is None or not getattr(title_shape, "has_text_frame", False):
            return

        cover_title_min_height = self._component_spacing_emu("cover", "title_min_height_emu", self.COVER_TITLE_MIN_HEIGHT_EMU)
        cover_meta_left = self._component_spacing_emu("cover", "meta_left_emu", self.COVER_META_LEFT_EMU)
        cover_meta_width = self._component_spacing_emu("cover", "meta_width_emu", self.COVER_META_WIDTH_EMU)
        cover_meta_gap = self._component_spacing_emu("cover", "meta_gap_emu", self.COVER_META_GAP_EMU)
        cover_meta_top = self._component_spacing_emu("cover", "meta_top_emu", self.COVER_META_DEFAULT_TOP_EMU)
        cover_meta_min_height = self._component_spacing_emu("cover", "meta_min_height_emu", self.COVER_META_MIN_HEIGHT_EMU)
        cover_bottom_limit = self._component_spacing_emu("cover", "bottom_limit_emu", self.COVER_BOTTOM_LIMIT_EMU)
        cover_meta_font_pt = self._component_font_size("cover", "meta", fallback=22.0)
        title_text = (getattr(title_shape, "text", "") or "").strip()
        if not title_text:
            return

        title_font_size_pt = self._fit_cover_title_font_size_points(title_text, title_shape.width)
        self._apply_font_size(title_shape, title_font_size_pt)
        self._configure_title_text_frame(title_shape)
        required_title_height = self._estimate_title_height_emu(title_shape, title_text, title_font_size_pt)
        title_shape.height = max(cover_title_min_height, required_title_height)

        if meta_shape is None or not getattr(meta_shape, "has_text_frame", False):
            return

        meta_text = (getattr(meta_shape, "text", "") or "").strip()
        if not meta_text:
            return

        meta_shape.left = cover_meta_left
        meta_shape.width = cover_meta_width
        desired_meta_top = title_shape.top + title_shape.height + cover_meta_gap
        meta_shape.top = max(cover_meta_top, desired_meta_top)
        meta_required_height = self._estimate_text_height_emu(meta_text, meta_shape.width, cover_meta_font_pt)
        meta_shape.height = max(cover_meta_min_height, meta_required_height)

        available_meta_height = cover_bottom_limit - meta_shape.top
        if available_meta_height < cover_meta_min_height:
            # If the title becomes too tall, tighten the title first before collapsing the meta block.
            title_font_size_pt = self._fit_cover_title_font_size_points(title_text, title_shape.width, max_height_emu=1900000)
            self._apply_font_size(title_shape, title_font_size_pt)
            required_title_height = self._estimate_title_height_emu(title_shape, title_text, title_font_size_pt)
            title_shape.height = max(cover_title_min_height, required_title_height)
            meta_shape.top = max(cover_meta_top, title_shape.top + title_shape.height + cover_meta_gap)
            available_meta_height = cover_bottom_limit - meta_shape.top

        meta_shape.height = max(cover_meta_min_height, min(meta_shape.height, available_meta_height))

    def _fit_cover_title_font_size_points(self, text: str, width_emu: int, max_height_emu: int = 2200000) -> float:
        base_size = self._component_font_size("cover", "title", fallback=46.0)
        candidates: list[float] = [base_size]
        for delta in (4.0, 8.0, 12.0, 14.0, 16.0, 18.0):
            candidate = max(28.0, base_size - delta)
            if candidate not in candidates:
                candidates.append(candidate)
        if 28.0 not in candidates:
            candidates.append(28.0)

        for candidate in candidates:
            estimated_height = self._estimate_text_height_emu(text, width_emu, candidate)
            if estimated_height <= max_height_emu:
                return candidate
        return 28.0

    def _populate_kpi_cards_slide(self, slide, slide_spec: SlideSpec, presentation_title: str) -> None:
        if self._active_presentation is not None:
            background = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                0,
                0,
                self._active_presentation.slide_width,
                self._active_presentation.slide_height,
            )
            background.fill.solid()
            background.fill.fore_color.rgb = RGBColor(0x34, 0x89, 0xF3)
            background.line.fill.background()

            glow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                int(self._active_presentation.slide_width * 0.48),
                -300000,
                int(self._active_presentation.slide_width * 0.58),
                int(self._active_presentation.slide_width * 0.58),
            )
            glow.fill.solid()
            glow.fill.fore_color.rgb = RGBColor(0x18, 0xC2, 0xFF)
            glow.fill.transparency = 0.28
            glow.line.fill.background()

        title_shape = slide.shapes.add_textbox(
            self.KPI_TITLE_LEFT_EMU,
            self.KPI_TITLE_TOP_EMU,
            self.KPI_TITLE_WIDTH_EMU,
            self.KPI_TITLE_HEIGHT_EMU,
        )
        self._set_cover_text(
            title_shape,
            slide_spec.title or "",
            font_size=Pt(38),
            bold=True,
            color=RGBColor(0xF5, 0xF9, 0xFE),
            align=PP_ALIGN.LEFT,
        )
        self._configure_title_text_frame(title_shape)
        title_shape.text_frame.margin_left = 0
        title_shape.text_frame.margin_right = 0
        title_shape.text_frame.margin_top = 0
        title_shape.text_frame.margin_bottom = 0

        description = (slide_spec.text or "").strip()
        if description:
            description_shape = slide.shapes.add_textbox(
                self.KPI_DESCRIPTION_LEFT_EMU,
                self.KPI_DESCRIPTION_TOP_EMU,
                self.KPI_DESCRIPTION_WIDTH_EMU,
                self.KPI_DESCRIPTION_HEIGHT_EMU,
            )
            self._set_cover_text(
                description_shape,
                description,
                font_size=Pt(20),
                bold=False,
                color=RGBColor(0xF5, 0xF9, 0xFE),
                align=PP_ALIGN.LEFT,
            )
            description_shape.text_frame.margin_left = 0
            description_shape.text_frame.margin_right = 0
            description_shape.text_frame.margin_top = 0
            description_shape.text_frame.margin_bottom = 0

        card_items = [item.strip() for item in slide_spec.bullets if item and item.strip()][:4]
        metric_positions = [
            (self.KPI_METRIC_LEFT_EMU, self.KPI_METRIC_START_TOP_EMU),
            (self.KPI_METRIC_RIGHT_EMU, self.KPI_METRIC_START_TOP_EMU),
            (self.KPI_METRIC_LEFT_EMU, self.KPI_METRIC_START_TOP_EMU + self.KPI_METRIC_HEIGHT_EMU + self.KPI_METRIC_ROW_GAP_EMU),
            (self.KPI_METRIC_RIGHT_EMU, self.KPI_METRIC_START_TOP_EMU + self.KPI_METRIC_HEIGHT_EMU + self.KPI_METRIC_ROW_GAP_EMU),
        ]
        for index, item in enumerate(card_items):
            left, top = metric_positions[index]
            metric_shape = slide.shapes.add_textbox(left, top, self.KPI_METRIC_WIDTH_EMU, self.KPI_METRIC_HEIGHT_EMU)
            self._set_card_text(metric_shape, item, profile_for_layout("cards_kpi"))

        footer_shape = slide.shapes.add_textbox(442913, 6384626, 3371850, 277813)
        self._set_cover_text(
            footer_shape,
            presentation_title,
            font_size=Pt(14),
            bold=False,
            color=RGBColor(0xF5, 0xF9, 0xFE),
            align=PP_ALIGN.LEFT,
        )
        footer_shape.text_frame.margin_left = 0
        footer_shape.text_frame.margin_right = 0
        footer_shape.text_frame.margin_top = 0
        footer_shape.text_frame.margin_bottom = 0

    def _fill_body(self, shape, slide_spec: SlideSpec, layout_profile: LayoutCapacityProfile) -> None:
        if slide_spec.content_blocks:
            self._set_content_blocks(shape, slide_spec.content_blocks, layout_profile)
            return
        if slide_spec.kind == SlideKind.BULLETS:
            if not slide_spec.bullets:
                self._clear_placeholder(shape)
                return
            self._set_bullets(shape, slide_spec.bullets, layout_profile)
            return
        if slide_spec.kind == SlideKind.TWO_COLUMN:
            merged = [*slide_spec.left_bullets, "", *slide_spec.right_bullets]
            if not any(item.strip() for item in merged):
                self._clear_placeholder(shape)
                return
            self._set_bullets(shape, merged, layout_profile)
            return
        if slide_spec.kind == SlideKind.TEXT:
            if not (slide_spec.text or "").strip():
                self._clear_placeholder(shape)
                return
            self._set_text(shape, slide_spec.text or "", layout_profile)
            return
        if slide_spec.kind == SlideKind.TITLE:
            if not (slide_spec.text or "").strip():
                self._clear_placeholder(shape)
                return
            self._set_text(shape, slide_spec.text or "", layout_profile)
            return
        if slide_spec.table is not None:
            rows = [" | ".join(row) for row in slide_spec.table.rows]
            if not rows:
                self._clear_placeholder(shape)
                return
            self._set_bullets(shape, rows, layout_profile)
            return
        if not (slide_spec.text or "").strip():
            self._clear_placeholder(shape)
            return
        self._set_text(shape, slide_spec.text or "", layout_profile)

    def _fill_card_body(self, slide, shape, slide_spec: SlideSpec, layout_profile: LayoutCapacityProfile) -> None:
        placeholder_idx = None
        if getattr(shape, "is_placeholder", False):
            try:
                placeholder_idx = shape.placeholder_format.idx
            except Exception:
                placeholder_idx = None

        card_index_by_placeholder = {11: 0, 12: 1, 13: 2}
        card_index = card_index_by_placeholder.get(placeholder_idx)
        if card_index is None:
            self._clear_placeholder(shape)
            return

        card_items = [item.strip() for item in slide_spec.bullets if item and item.strip()]
        if card_index >= len(card_items):
            self._clear_placeholder(shape)
            return

        common_font_pt = self._card_common_font_size(card_items, layout_profile)
        rich_card = self._parse_rich_card_text(card_items[card_index])
        if layout_profile.layout_key == "cards_3" and rich_card is not None:
            self._set_rich_numeric_card_text(slide, shape, rich_card, layout_profile)
            return
        self._set_card_text(shape, card_items[card_index], layout_profile, common_font_pt=common_font_pt)

    def _card_common_font_size(self, card_items: list[str], layout_profile: LayoutCapacityProfile) -> int:
        return int(min(max(20, layout_profile.min_font_pt), layout_profile.max_font_pt))

    def _card_body_font_size(self, layout_profile: LayoutCapacityProfile) -> int:
        token_value = self._design_token("cards_body_font_size_pt")
        if isinstance(token_value, (int, float)):
            return int(min(max(token_value, layout_profile.min_font_pt), layout_profile.max_font_pt))
        return int(min(max(16, layout_profile.min_font_pt), layout_profile.max_font_pt))

    def _split_card_text(self, text: str) -> tuple[str, str]:
        normalized = "\n".join(line.strip() for line in (text or "").splitlines() if line.strip())
        if not normalized:
            return "", ""
        if "\n" in normalized:
            title, description = normalized.split("\n", 1)
            return title.strip(), " ".join(description.split())

        colon_match = re.match(r"^(.{4,54}?):\s+(.{12,})$", normalized)
        if colon_match:
            return colon_match.group(1).strip(), colon_match.group(2).strip()

        dash_match = re.match(r"^(.{4,54}?)\s+[—-]\s+(.{12,})$", normalized)
        if dash_match:
            return dash_match.group(1).strip(), dash_match.group(2).strip()

        return normalized, ""

    def _split_numeric_card_text(self, text: str) -> tuple[str, str] | None:
        title, description = self._split_card_text(text)
        if title and description and re.search(r"\d", title):
            return title, description
        normalized = " ".join((text or "").split())
        metric_match = re.match(
            r"^([<>~≈]?\s*\d+(?:[.,]\d+)?(?:\s*(?:%|‰|млн|млрд|тыс|трлн|сек(?:унд[аы]?)?|с|мин|ч|дн(?:ей|я)?|₽|руб(?:\.|лей|ля|ль)?))*)\s+(.+)$",
            normalized,
            re.IGNORECASE,
        )
        if metric_match:
            value_text = re.sub(r"\s+([%‰₽])", r"\1", metric_match.group(1).strip())
            return value_text, metric_match.group(2).strip()
        return None

    def _parse_rich_card_text(self, text: str) -> tuple[str, str, list[tuple[str, str]]] | None:
        lines = [line.strip() for line in (text or "").splitlines() if line.strip()]
        if len(lines) < 2:
            return None
        title = lines[0]
        description_lines: list[str] = []
        metrics: list[tuple[str, str]] = []
        for line in lines[1:]:
            metric_parts = self._split_numeric_card_text(line)
            if metric_parts is not None:
                metrics.append(metric_parts)
            else:
                description_lines.append(line)
        if not metrics:
            return None
        return title, " ".join(description_lines).strip(), metrics[:4]

    def _set_card_text(self, shape, text: str, layout_profile: LayoutCapacityProfile, *, common_font_pt: int | None = None) -> None:
        normalized = (text or "").strip()
        if not normalized:
            self._clear_placeholder(shape)
            return

        numeric_parts = self._split_numeric_card_text(normalized)
        if layout_profile.layout_key == "cards_kpi" and numeric_parts is not None:
            self._set_kpi_card_text(shape, numeric_parts[0], numeric_parts[1], layout_profile)
            return

        title, description = self._split_card_text(normalized)
        max_font_pt = common_font_pt or self._card_common_font_size([normalized], layout_profile)
        body_font_pt = self._card_body_font_size(layout_profile)
        card_title_font_pt = self._component_font_size("cards", "title", fallback=min(max_font_pt, 20))
        card_body_color = self._component_font_color("cards", "body", fallback=RGBColor(0xFF, 0xFF, 0xFF))
        card_title_color = self._component_font_color("cards", "title", fallback=RGBColor(0xFF, 0xFF, 0xFF))
        card_profile = replace(layout_profile, max_font_pt=min(layout_profile.max_font_pt, max_font_pt))
        self._configure_card_text_frame(shape.text_frame)
        text_frame = shape.text_frame
        text_frame.clear()
        title_paragraph = text_frame.paragraphs[0]
        title_run = title_paragraph.add_run()
        title_run.text = title
        if description:
            description_paragraph = text_frame.add_paragraph()
            description_run = description_paragraph.add_run()
            description_run.text = description
        self._configure_card_text_frame(text_frame)
        self._set_text_frame_font_size(text_frame, card_profile.max_font_pt, card_profile.layout_key)
        self._set_text_frame_regular(text_frame)
        self._set_text_frame_color(text_frame, card_body_color)
        title_run.font.bold = True
        title_run.font.size = Pt(card_title_font_pt)
        title_run.font.color.rgb = card_title_color
        if description:
            title_paragraph.space_after = Pt(8)
            for run in text_frame.paragraphs[1].runs:
                run.font.size = Pt(body_font_pt)
                run.font.bold = False
                run.font.color.rgb = card_body_color

    def _set_rich_numeric_card_text(
        self,
        slide,
        shape,
        card_content: tuple[str, str, list[tuple[str, str]]],
        layout_profile: LayoutCapacityProfile,
    ) -> None:
        title_text, description_text, metrics = card_content
        self._configure_card_text_frame(shape.text_frame)
        shape.text_frame.clear()

        margin_x = self._component_spacing_emu("cards", "content_margin_x_emu", self.DEFAULT_TEXT_MARGIN_X_EMU)
        margin_y = self._component_spacing_emu("cards", "content_margin_y_emu", self.DEFAULT_TEXT_MARGIN_Y_EMU)
        metric_gap_x = self._component_spacing_emu("cards", "metrics_gap_x_emu", 180000)
        metric_gap_y = self._component_spacing_emu("cards", "metrics_gap_y_emu", 160000)
        title_font_pt = self._component_font_size("cards", "title", fallback=min(layout_profile.max_font_pt, 20))
        body_font_pt = self._numeric_card_body_font_size(description_text, metrics, layout_profile)
        metric_value_font_pt = self._component_behavior_float("cards", "kpi_value_compact_font_pt" if len(metrics) >= 4 else "kpi_value_regular_font_pt", 20 if len(metrics) >= 4 else 22)
        metric_label_font_pt = self._component_font_size("cards", "kpi_label", fallback=12)
        title_body_gap = self._component_spacing_emu("cards", "title_body_gap_emu", 100000)
        body_metrics_gap = self._component_spacing_emu("cards", "body_metrics_gap_emu", 180000)
        card_title_color = self._component_font_color("cards", "title", fallback=RGBColor(0xFF, 0xFF, 0xFF))
        card_body_color = self._component_font_color("cards", "body", fallback=RGBColor(0xF3, 0xF8, 0xFF))
        card_kpi_value_color = self._component_font_color("cards", "kpi_value", fallback=RGBColor(0xFF, 0xFF, 0xFF))
        card_kpi_label_color = self._component_font_color("cards", "kpi_label", fallback=RGBColor(0xE4, 0xF1, 0xFF))
        inner_left = shape.left + margin_x
        inner_top = shape.top + margin_y
        inner_width = max(shape.width - margin_x * 2, 600000)
        title_height = max(300000, self._estimate_text_height_emu(title_text, inner_width, title_font_pt))
        description_height = self._estimate_text_height_emu(description_text, inner_width, body_font_pt) if description_text else 0
        description_bottom = inner_top + title_height
        if description_text:
            description_bottom += title_body_gap + description_height
        metric_top = description_bottom + (body_metrics_gap if description_text else 120000)
        metric_bottom = shape.top + shape.height - margin_y
        metric_area_height = max(metric_bottom - metric_top, 520000)
        metric_boxes = self._numeric_card_metric_boxes(
            count=len(metrics),
            left=inner_left,
            top=metric_top,
            width=inner_width,
            height=metric_area_height,
            gap_x=metric_gap_x,
            gap_y=metric_gap_y,
        )

        title_shape = slide.shapes.add_textbox(inner_left, inner_top, inner_width, title_height)
        title_shape.name = f"A3_CARD_OVERLAY_{shape.placeholder_format.idx}_TITLE"
        self._configure_card_text_frame(title_shape.text_frame)
        title_shape.text_frame.clear()
        title_run = title_shape.text_frame.paragraphs[0].add_run()
        title_run.text = title_text
        title_run.font.bold = True
        title_run.font.size = Pt(title_font_pt)
        title_run.font.color.rgb = card_title_color
        body_style = self._fallback_theme_text_style("body")
        if body_style.font_family:
            self._apply_run_font_family(title_run, body_style.font_family)

        if description_text:
            description_shape = slide.shapes.add_textbox(inner_left, inner_top + title_height + title_body_gap, inner_width, description_height)
            description_shape.name = f"A3_CARD_OVERLAY_{shape.placeholder_format.idx}_DESCRIPTION"
            self._configure_card_text_frame(description_shape.text_frame)
            description_shape.text_frame.clear()
            description_run = description_shape.text_frame.paragraphs[0].add_run()
            description_run.text = description_text
            description_run.font.bold = False
            description_run.font.size = Pt(body_font_pt)
            description_run.font.color.rgb = card_body_color
            if body_style.font_family:
                self._apply_run_font_family(description_run, body_style.font_family)

        for index, (value_text, label_text) in enumerate(metrics):
            metric_left, metric_top_current, metric_width, metric_height = metric_boxes[index]
            metric_shape = slide.shapes.add_textbox(metric_left, metric_top_current, metric_width, metric_height)
            metric_shape.name = f"A3_CARD_OVERLAY_{shape.placeholder_format.idx}_METRIC_{index}"
            self._configure_card_text_frame(metric_shape.text_frame)
            metric_shape.text_frame.clear()
            value_paragraph = metric_shape.text_frame.paragraphs[0]
            value_run = value_paragraph.add_run()
            value_run.text = value_text
            value_run.font.bold = True
            value_run.font.size = Pt(metric_value_font_pt)
            value_run.font.color.rgb = card_kpi_value_color
            if body_style.font_family:
                self._apply_run_font_family(value_run, body_style.font_family)
            if label_text:
                label_paragraph = metric_shape.text_frame.add_paragraph()
                label_paragraph.space_before = Pt(3)
                label_run = label_paragraph.add_run()
                label_run.text = label_text
                label_run.font.bold = False
                label_run.font.size = Pt(metric_label_font_pt)
                label_run.font.color.rgb = card_kpi_label_color
                if body_style.font_family:
                    self._apply_run_font_family(label_run, body_style.font_family)

    def _numeric_card_body_font_size(
        self,
        description_text: str,
        metrics: list[tuple[str, str]],
        layout_profile: LayoutCapacityProfile,
    ) -> int:
        points = self._card_body_font_size(layout_profile)
        text_length = len(description_text or "")
        if text_length >= 90 or len(metrics) >= 3:
            points = min(points, 14)
        if text_length >= 150:
            points = min(points, 13)
        return max(points, 12)

    def _numeric_card_metric_boxes(
        self,
        *,
        count: int,
        left: int,
        top: int,
        width: int,
        height: int,
        gap_x: int,
        gap_y: int,
    ) -> list[tuple[int, int, int, int]]:
        metric_count = max(1, min(count, 4))
        half_width = max((width - gap_x) // 2, 900000)
        if metric_count == 1:
            return [(left, top, width, height)]
        if metric_count == 2:
            return [
                (left, top, half_width, height),
                (left + half_width + gap_x, top, half_width, height),
            ]

        row_height = max((height - gap_y) // 2, 420000)
        boxes = [
            (left, top, half_width, row_height),
            (left + half_width + gap_x, top, half_width, row_height),
            (left, top + row_height + gap_y, half_width, row_height),
        ]
        if metric_count == 4:
            boxes.append((left + half_width + gap_x, top + row_height + gap_y, half_width, row_height))
        return boxes

    def _set_kpi_card_text(self, shape, value_text: str, label_text: str, layout_profile: LayoutCapacityProfile) -> None:
        self._configure_card_text_frame(shape.text_frame)
        text_frame = shape.text_frame
        text_frame.clear()
        value_paragraph = text_frame.paragraphs[0]
        value_paragraph.alignment = PP_ALIGN.LEFT
        value_run = value_paragraph.add_run()
        value_run.text = value_text
        value_run.font.bold = True
        value_run.font.size = Pt(min(layout_profile.max_font_pt, 36))
        value_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
        body_style = self._fallback_theme_text_style("body")
        if body_style.font_family:
            self._apply_run_font_family(value_run, body_style.font_family)

        if label_text:
            label_paragraph = text_frame.add_paragraph()
            label_paragraph.alignment = PP_ALIGN.LEFT
            label_paragraph.space_before = Pt(8)
            label_run = label_paragraph.add_run()
            label_run.text = label_text
            label_run.font.bold = False
            label_run.font.size = Pt(16)
            label_run.font.color.rgb = RGBColor(0xE4, 0xF1, 0xFF)
            if body_style.font_family:
                self._apply_run_font_family(label_run, body_style.font_family)

    def _add_kpi_fourth_card(self, slide, slide_spec: SlideSpec, layout_profile: LayoutCapacityProfile) -> None:
        card_items = [item.strip() for item in slide_spec.bullets if item and item.strip()]
        if len(card_items) < 4:
            return
        metric_shape = slide.shapes.add_textbox(
            self.KPI_FOURTH_CARD_LEFT_EMU,
            self.KPI_FOURTH_CARD_TOP_EMU,
            self.KPI_FOURTH_CARD_WIDTH_EMU,
            self.KPI_FOURTH_CARD_HEIGHT_EMU,
        )
        self._set_card_text(metric_shape, card_items[3], layout_profile)

    def _set_content_blocks(
        self,
        shape,
        content_blocks: list[SlideContentBlock],
        layout_profile: LayoutCapacityProfile,
    ) -> None:
        text_frame = shape.text_frame
        text_frame.clear()

        entries: list[tuple[str, bool, SlideContentBlockKind]] = []
        for block in content_blocks:
            if block.kind in {SlideContentBlockKind.PARAGRAPH, SlideContentBlockKind.CALLOUT, SlideContentBlockKind.QA_ITEM}:
                text = (block.text or "").strip()
                if text:
                    entries.append((text, False, block.kind))
                continue
            if block.kind == SlideContentBlockKind.BULLET_LIST:
                entries.extend((item.strip(), True, block.kind) for item in block.items if item and item.strip())

        if not entries:
            self._clear_placeholder(shape)
            return

        first = True
        for text, is_bullet, block_kind in entries:
            paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            paragraph.text = text
            if is_bullet:
                paragraph.level = 0
                self._apply_bullet_format(paragraph, layout_profile.layout_key)
                self._apply_paragraph_spacing(paragraph, "body", layout_profile.layout_key)
            else:
                self._apply_paragraph_spacing(paragraph, "body", layout_profile.layout_key)
            self._style_content_block_paragraph(paragraph, block_kind)
            first = False

        self._configure_body_text_frame(text_frame)
        self._apply_body_font_size(text_frame, [text for text, _, _ in entries], shape, layout_profile)

    def _style_content_block_paragraph(self, paragraph, block_kind: SlideContentBlockKind) -> None:
        if not paragraph.runs and paragraph.text:
            run = paragraph.add_run()
            run.text = paragraph.text
            paragraph.text = ""

        if block_kind == SlideContentBlockKind.QA_ITEM:
            paragraph.level = 0
            for run in paragraph.runs:
                run.font.bold = True
                run.font.italic = False
            return

        if block_kind == SlideContentBlockKind.CALLOUT:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.italic = True
            return

    def _fill_shape_by_binding(
        self,
        shape,
        binding: str,
        slide_spec: SlideSpec,
        presentation_title: str,
        layout_profile: LayoutCapacityProfile,
        *,
        placeholder_spec: PlaceholderSpec | PrototypeTokenSpec | None = None,
    ) -> None:
        placeholder_idx = None
        if getattr(shape, "is_placeholder", False):
            placeholder_format = getattr(shape, "placeholder_format", None)
            placeholder_idx = getattr(placeholder_format, "idx", None)

        first_block_text = ""
        if slide_spec.content_blocks:
            for block in slide_spec.content_blocks:
                candidate = (block.text or "").strip()
                if candidate:
                    first_block_text = candidate
                    break

        if slide_spec.content_blocks and placeholder_idx == 14 and binding in {"body", "main_text", "bullets"}:
            self._set_content_blocks(shape, slide_spec.content_blocks, layout_profile)
            return
        if slide_spec.content_blocks and placeholder_idx == 15 and binding in {"secondary_text", "notes"}:
            self._clear_placeholder(shape)
            return
        if slide_spec.content_blocks and placeholder_idx == 13 and binding == "subtitle":
            subtitle_text = (slide_spec.subtitle or "").strip()
            if not subtitle_text:
                self._clear_placeholder(shape)
                return
            if first_block_text and first_block_text.startswith(subtitle_text):
                self._clear_placeholder(shape)
                return

        binding_value = self._build_token_value_map(slide_spec, presentation_title).get(binding, "")
        if binding == "table":
            self._fill_table_or_chart(shape, slide_spec, placeholder_spec)
            return
        if binding in {"chart", "chart_image"}:
            self._fill_chart(shape, slide_spec, placeholder_spec)
            return
        if binding == "image":
            self._fill_image(shape, slide_spec)
            return
        if binding == "icon_grid":
            self._clear_placeholder(shape)
            return
        if self._is_empty_binding_value(binding_value) and binding not in {"presentation_name", "cover_title", "title"}:
            self._clear_placeholder(shape)
            return
        if not getattr(shape, "has_text_frame", False):
            return
        if isinstance(binding_value, list):
            self._set_bullets(shape, [str(item) for item in binding_value], layout_profile)
            return
        if slide_spec.kind == SlideKind.CHART and binding == "title":
            self._set_text(shape, str(binding_value), layout_profile)
            self._configure_title_text_frame(shape)
            self._apply_font_size(
                shape,
                self._component_font_size("chart", "title", fallback=35.0),
            )
            return
        if slide_spec.kind == SlideKind.CHART and binding == "subtitle":
            self._set_text(shape, str(binding_value), layout_profile)
            self._configure_subtitle_text_frame(shape)
            self._apply_font_size(
                shape,
                self._component_font_size("chart", "subtitle", fallback=20.0),
            )
            return
        self._set_text(shape, str(binding_value), layout_profile)

    def _expand_text_full_width_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("text_full_width")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }

        title = placeholders.get(0)
        subtitle = placeholders.get(13)
        main_text = placeholders.get(14)
        secondary_text = placeholders.get(15)

        if title is not None:
            policy = geometry.placeholders[0]
            title.left = policy.left_emu
            title.top = policy.top_emu
            title.width = policy.width_emu
            title.height = max(title.height or 0, policy.height_emu)

        if subtitle is not None:
            policy = geometry.placeholders[13]
            subtitle.left = policy.left_emu
            subtitle.top = policy.top_emu
            subtitle.width = policy.width_emu
            subtitle.height = policy.height_emu

        if main_text is not None:
            policy = geometry.placeholders[14]
            main_text.left = policy.left_emu
            main_text.top = policy.top_emu
            main_text.width = policy.width_emu
            main_text.height = 1700000 if secondary_text is not None and getattr(secondary_text, "text", "").strip() else policy.height_emu

        if secondary_text is not None and getattr(secondary_text, "text", "").strip():
            policy = geometry.placeholders[15]
            secondary_text.left = policy.left_emu
            secondary_text.top = policy.top_emu
            secondary_text.width = policy.width_emu
            secondary_text.height = policy.height_emu

        footer = placeholders.get(17)
        if footer is not None:
            policy = geometry.placeholders[17]
            footer.left = policy.left_emu
            footer.top = policy.top_emu
            footer.width = policy.width_emu
            footer.height = policy.height_emu

    def _expand_table_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("table")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }

        title = placeholders.get(0)
        subtitle = placeholders.get(13)
        left_note = placeholders.get(11)
        right_note = placeholders.get(12)
        footer = placeholders.get(15)

        if title is not None:
            policy = geometry.placeholders[0]
            title.left = policy.left_emu
            title.top = policy.top_emu
            title.width = policy.width_emu
            title.height = max(title.height or 0, policy.height_emu)

        if footer is not None:
            policy = geometry.placeholders[15]
            footer.left = policy.left_emu
            footer.top = policy.top_emu
            footer.width = policy.width_emu
            footer.height = policy.height_emu

        if subtitle is None or not getattr(subtitle, "text", "").strip():
            return

        has_side_notes = any(
            shape is not None and getattr(shape, "text", "").strip()
            for shape in (left_note, right_note)
        )
        if has_side_notes:
            return

        policy = geometry.placeholders[13]
        subtitle.left = policy.left_emu
        subtitle.top = policy.top_emu
        subtitle.width = policy.width_emu
        subtitle.height = policy.height_emu

    def _expand_image_text_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("image_text")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _expand_cards_layout(self, slide, layout_key: str = "cards_3") -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _expand_list_with_icons_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("list_with_icons")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _expand_contacts_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("contacts")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _expand_cards_layout(self, slide, layout_key: str = "cards_3") -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _expand_list_with_icons_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("list_with_icons")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _expand_contacts_layout(self, slide) -> None:
        geometry = geometry_policy_for_layout("contacts")
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        for idx, policy in geometry.placeholders.items():
            shape = placeholders.get(idx)
            if shape is None:
                continue
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _adjust_title_and_flow(self, slide, layout_key: str, slide_spec: SlideSpec | None = None) -> None:
        if layout_key in {"text_full_width", "dense_text_full_width", "list_full_width"}:
            self._stack_text_content(slide, layout_key)
            return
        if layout_key == "image_text":
            self._stack_image_text_content(slide, layout_key)
            return
        if layout_key == "table":
            self._stack_table_content(slide, layout_key, slide_spec)
            return
        if layout_key in {"cards_3", "cards_kpi"}:
            self._stack_cards_content(slide, layout_key)
            return
        if layout_key == "list_with_icons":
            self._stack_two_column_content(slide, layout_key)
            return
        if layout_key == "contacts":
            self._stack_contacts_content(slide, layout_key)
            return
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title_shape = placeholders.get(0)
        if title_shape is None or not getattr(title_shape, "has_text_frame", False):
            return

        title_text = (getattr(title_shape, "text", "") or "").strip()
        if not title_text:
            return

        font_size_pt = self._fit_title_font_size_points(title_text, title_shape.width, layout_key)
        self._apply_font_size(title_shape, font_size_pt)
        required_height = self._estimate_text_height_emu(title_text, title_shape.width, font_size_pt)
        base_height = title_shape.height or 0
        final_title_height = max(base_height, required_height)
        title_shape.height = final_title_height

        protected_indices = {0, 17 if layout_key in {"text_full_width", "list_full_width"} else -1}
        flow_shapes = [
            shape
            for placeholder_idx, shape in placeholders.items()
            if placeholder_idx not in protected_indices and shape.top > title_shape.top
        ]
        if not flow_shapes:
            return

        current_flow_top = min(shape.top for shape in flow_shapes)
        desired_flow_top = title_shape.top + final_title_height + self.TITLE_CONTENT_GAP_EMU
        delta = max(0, desired_flow_top - current_flow_top)

        for shape in flow_shapes:
            shape.top += delta

    def _stack_text_content(self, slide, layout_key: str) -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title = placeholders.get(0)
        body = placeholders.get(14)
        footer = placeholders.get(17)
        subtitle = placeholders.get(13)
        secondary = placeholders.get(15)

        if title is None or body is None or footer is None:
            return

        title_text = (getattr(title, "text", "") or "").strip()
        if title_text:
            font_size_pt = self._fit_title_font_size_points(title_text, title.width, layout_key)
            self._apply_font_size(title, font_size_pt)
            self._configure_title_text_frame(title)
            required_height = self._estimate_title_height_emu(title, title_text, font_size_pt)
            title.height = max(self._minimum_title_height_emu(layout_key), required_height)

        if subtitle is not None and self._subtitle_duplicates_body(subtitle, body):
            self._clear_placeholder(subtitle)
            subtitle = None

        has_subtitle = subtitle is not None and getattr(subtitle, "text", "").strip()
        subtitle_body_gap = geometry.title_content_gap_emu
        compact_subtitle = False
        subtitle_text = ""
        body_text = (getattr(body, "text", "") or "").strip()
        if has_subtitle and subtitle is not None:
            subtitle_text = subtitle.text.strip()
            compact_subtitle = self._should_compact_subtitle_stack(title_text, subtitle_text, body_text, layout_key=layout_key)
        title_gap = geometry.title_content_gap_emu if has_subtitle else geometry.title_body_gap_no_subtitle_emu
        if compact_subtitle:
            title_gap = min(title_gap, 120000)
            subtitle_body_gap = 100000
        cursor = title.top + title.height + title_gap

        if has_subtitle:
            subtitle_font_pt = 20.0
            self._configure_subtitle_text_frame(subtitle)
            self._apply_font_size(subtitle, subtitle_font_pt)
            subtitle.height = max(280000 if compact_subtitle else 360000, self._estimate_text_height_emu(subtitle_text, subtitle.width, subtitle_font_pt))
            subtitle.top = cursor
            cursor = subtitle.top + subtitle.height + subtitle_body_gap

        secondary_has_text = secondary is not None and getattr(secondary, "text", "").strip()
        available_bottom = footer.top - geometry.content_footer_gap_emu

        if secondary_has_text and secondary is not None:
            secondary_text = secondary.text.strip()
            secondary.height = max(secondary.height or 0, self._estimate_text_height_emu(secondary_text, secondary.width, 16.0))
            secondary.height = min(secondary.height, max(700000, available_bottom - cursor - 900000))
            body.top = cursor
            body.height = max(900000, secondary.top - geometry.title_content_gap_emu - body.top)
            secondary.top = body.top + body.height + geometry.title_content_gap_emu
            secondary.height = max(700000, min(secondary.height, available_bottom - secondary.top))
            return

        body.top = cursor
        body.height = max(900000, available_bottom - body.top)

    def _should_compact_subtitle_stack(self, title_text: str, subtitle_text: str, body_text: str, *, layout_key: str) -> bool:
        if layout_key not in {"text_full_width", "dense_text_full_width"}:
            return False
        if not subtitle_text or not body_text:
            return False
        if len(subtitle_text) > 72 or len(subtitle_text.split()) > 10:
            return False
        if len(body_text) > 900:
            return False
        if body_text.count("\n") >= 5:
            return False
        return True

    def _subtitle_duplicates_body(self, subtitle, body) -> bool:
        subtitle_text = (getattr(subtitle, "text", "") or "").strip()
        body_text = (getattr(body, "text", "") or "").strip()
        if not subtitle_text or not body_text:
            return False
        if body_text.startswith(subtitle_text):
            return True
        if len(subtitle_text) >= 24 and body_text.startswith(subtitle_text[:-1]):
            return True
        return False

    def _stack_table_content(self, slide, layout_key: str, slide_spec: SlideSpec | None = None) -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title = placeholders.get(0)
        subtitle = placeholders.get(13)
        table = placeholders.get(14)
        chart = next((shape for shape in slide.shapes if getattr(shape, "has_chart", False)), None)
        content = table or chart
        footer = placeholders.get(15)

        if title is None or content is None or footer is None:
            return

        title_text = (getattr(title, "text", "") or "").strip()
        if title_text:
            font_size_pt = self._fit_title_font_size_points(title_text, title.width, layout_key)
            self._apply_font_size(title, font_size_pt)
            self._configure_title_text_frame(title)
            required_height = self._estimate_title_height_emu(title, title_text, font_size_pt)
            title.height = max(self._minimum_title_height_emu(layout_key), required_height)

        cursor = title.top + title.height + geometry.title_content_gap_emu
        if subtitle is not None and getattr(subtitle, "text", "").strip():
            subtitle_text = subtitle.text.strip()
            body_style = self._active_manifest.theme.master_text_styles.get("body") if self._active_manifest is not None else None
            subtitle_font_size = body_style.font_size_pt if body_style is not None and body_style.font_size_pt else self._table_subtitle_font_size_points(subtitle_text)
            self._apply_font_size(subtitle, subtitle_font_size)
            self._configure_subtitle_text_frame(subtitle)
            subtitle.height = max(360000, self._estimate_text_height_emu(subtitle_text, subtitle.width, subtitle_font_size))
            subtitle.top = cursor
            cursor = subtitle.top + subtitle.height + geometry.title_content_gap_emu

        content.top = cursor
        content.height = max(900000, footer.top - geometry.content_footer_gap_emu - content.top)

    def _stack_image_text_content(self, slide, layout_key: str) -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title = placeholders.get(0)
        body = placeholders.get(14)
        footer = placeholders.get(17)
        subtitle = placeholders.get(13)
        secondary = placeholders.get(15)
        image = placeholders.get(16)

        if title is None or body is None or footer is None:
            return

        title_text = (getattr(title, "text", "") or "").strip()
        if title_text:
            font_size_pt = self._fit_title_font_size_points(title_text, title.width, layout_key)
            self._apply_font_size(title, font_size_pt)
            self._configure_title_text_frame(title)
            required_height = self._estimate_title_height_emu(title, title_text, font_size_pt)
            title.height = max(self._minimum_title_height_emu(layout_key), required_height)

        has_subtitle = subtitle is not None and getattr(subtitle, "text", "").strip()
        title_content_gap = self._component_spacing_emu("image", "title_content_gap_emu", geometry.title_content_gap_emu)
        title_body_gap_no_subtitle = self._component_spacing_emu(
            "image",
            "title_body_gap_no_subtitle_emu",
            geometry.title_body_gap_no_subtitle_emu,
        )
        content_footer_gap = self._component_spacing_emu("image", "content_footer_gap_emu", geometry.content_footer_gap_emu)
        min_image_height = self._component_spacing_emu("image", "min_image_height_emu", 1200000)
        secondary_min_height = self._component_spacing_emu("image", "secondary_min_height_emu", 700000)
        secondary_reserved_gap = self._component_spacing_emu("image", "secondary_reserved_image_gap_emu", 900000)
        title_gap = title_content_gap if has_subtitle else title_body_gap_no_subtitle
        cursor = title.top + title.height + title_gap
        if has_subtitle:
            subtitle_text = subtitle.text.strip()
            self._configure_subtitle_text_frame(subtitle)
            subtitle_font_pt = self._component_font_size("image", "subtitle", fallback=18.0)
            self._apply_font_size(subtitle, subtitle_font_pt)
            subtitle.height = max(360000, self._estimate_text_height_emu(subtitle_text, subtitle.width, subtitle_font_pt))
            subtitle.top = cursor
            cursor = subtitle.top + subtitle.height + title_content_gap

        available_bottom = footer.top - content_footer_gap
        if image is not None:
            image.top = cursor
            image.height = max(min_image_height, available_bottom - image.top)

        secondary_has_text = secondary is not None and getattr(secondary, "text", "").strip()
        if secondary_has_text and secondary is not None:
            secondary_text = secondary.text.strip()
            secondary_font_pt = self._component_font_size("image", "body", fallback=16.0)
            secondary.height = max(secondary.height or 0, self._estimate_text_height_emu(secondary_text, secondary.width, secondary_font_pt))
            secondary.height = min(secondary.height, max(secondary_min_height, available_bottom - cursor - secondary_reserved_gap))
            body.top = cursor
            body.height = max(900000, secondary.top - title_content_gap - body.top)
            secondary.top = body.top + body.height + title_content_gap
            secondary.height = max(secondary_min_height, min(secondary.height, available_bottom - secondary.top))
            return

        body.top = cursor
        body.height = max(900000, available_bottom - body.top)

    def _stack_cards_content(self, slide, layout_key: str) -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title = placeholders.get(0)
        cards = [placeholders.get(idx) for idx in (11, 12, 13)]
        cards = [card for card in cards if card is not None]
        if title is None or not cards:
            return

        title_text = (getattr(title, "text", "") or "").strip()
        first_card_top = min(
            geometry.placeholders[card.placeholder_format.idx].top_emu
            for card in cards
            if card.placeholder_format.idx in geometry.placeholders
        )
        max_title_height = max(520000, first_card_top - geometry.title_body_gap_no_subtitle_emu - title.top)
        if title_text:
            font_size_pt = self._fit_title_font_size_for_height(title, title_text, layout_key, max_title_height)
            self._apply_font_size(title, font_size_pt)
            self._configure_title_text_frame(title)
            if layout_key == "cards_kpi":
                self._set_text_frame_color(title.text_frame, RGBColor(0xFF, 0xFF, 0xFF))
            title.text_frame.auto_size = MSO_AUTO_SIZE.NONE
            required_height = self._estimate_title_height_emu(title, title_text, font_size_pt)
            title.height = min(max(520000, required_height), max_title_height)

        for card in cards:
            policy = geometry.placeholders.get(card.placeholder_format.idx)
            if policy is None:
                continue
            desired_top = policy.top_emu
            bottom = policy.top_emu + policy.height_emu
            card.left = policy.left_emu
            card.top = desired_top
            card.width = policy.width_emu
            card.height = max(900000, bottom - desired_top)
            if getattr(card, "has_text_frame", False):
                self._configure_card_text_frame(card.text_frame)

    def _stack_two_column_content(self, slide, layout_key: str) -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title = placeholders.get(0)
        subtitle = placeholders.get(13)
        left = placeholders.get(12)
        right = placeholders.get(14)
        footer = placeholders.get(21)
        if title is None or left is None or right is None or footer is None:
            return

        title_text = (getattr(title, "text", "") or "").strip()
        if title_text:
            font_size_pt = self._fit_title_font_size_points(title_text, title.width, layout_key)
            self._apply_font_size(title, font_size_pt)
            self._configure_title_text_frame(title)
            required_height = self._estimate_title_height_emu(title, title_text, font_size_pt)
            title.height = max(self._minimum_title_height_emu(layout_key), required_height)

        has_subtitle = subtitle is not None and getattr(subtitle, "text", "").strip()
        title_content_gap = self._component_spacing_emu("list_with_icons", "title_content_gap_emu", geometry.title_content_gap_emu)
        title_no_subtitle_gap = self._component_spacing_emu(
            "list_with_icons",
            "title_body_gap_no_subtitle_emu",
            geometry.title_body_gap_no_subtitle_emu,
        )
        content_footer_gap = self._component_spacing_emu("list_with_icons", "content_footer_gap_emu", geometry.content_footer_gap_emu)
        title_gap = title_content_gap if has_subtitle else title_no_subtitle_gap
        cursor = title.top + title.height + title_gap
        if has_subtitle:
            subtitle_text = subtitle.text.strip()
            self._configure_subtitle_text_frame(subtitle)
            subtitle_font_pt = self._component_font_size("list_with_icons", "subtitle", fallback=18.0)
            self._apply_font_size(subtitle, subtitle_font_pt)
            subtitle.height = max(360000, self._estimate_text_height_emu(subtitle_text, subtitle.width, subtitle_font_pt))
            subtitle.top = cursor
            cursor = subtitle.top + subtitle.height + title_content_gap

        content_indices = [12, 14, 15, 16, 17, 18, 19, 20]
        content_shapes = [placeholders[idx] for idx in content_indices if idx in placeholders]
        if not content_shapes:
            return
        base_top = min(shape.top for shape in content_shapes)
        delta = max(0, cursor - base_top)
        max_height = max(900000, footer.top - content_footer_gap - (base_top + delta))
        for shape in content_shapes:
            shape.top += delta
            if shape.placeholder_format.idx in {12, 14}:
                shape.height = min(shape.height, max_height)

    def _stack_contacts_content(self, slide, layout_key: str) -> None:
        geometry = geometry_policy_for_layout(layout_key)
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        primary_font = self._component_font_size("contacts", "primary", fallback=18.0)
        secondary_font = self._component_font_size("contacts", "secondary", fallback=14.0)
        primary_threshold = int(self._component_behavior_float("contacts", "primary_threshold_chars", 60.0))
        secondary_threshold = int(self._component_behavior_float("contacts", "secondary_threshold_chars", 40.0))
        decrement_pt = self._component_behavior_float("contacts", "font_decrement_pt", 2.0)
        for idx in (10, 11, 12, 13):
            shape = placeholders.get(idx)
            policy = geometry.placeholders.get(idx)
            if shape is None or policy is None:
                continue
            text = (getattr(shape, "text", "") or "").strip()
            if not text:
                continue
            font_size_pt = primary_font if idx == 10 else secondary_font
            if len(text) >= (primary_threshold if idx == 10 else secondary_threshold):
                font_size_pt -= decrement_pt
            self._apply_font_size(shape, font_size_pt)
            self._configure_subtitle_text_frame(shape)
            shape.left = policy.left_emu
            shape.top = policy.top_emu
            shape.width = policy.width_emu
            shape.height = policy.height_emu

    def _title_font_size_points(self, layout_key: str) -> float:
        theme = self._active_manifest.theme if self._active_manifest is not None else None
        title_style = theme.master_text_styles.get("title") if theme is not None else None
        if title_style is not None and title_style.font_size_pt:
            return title_style.font_size_pt
        return 28.0

    def _table_subtitle_font_size_points(self, text: str) -> float:
        return 18.0

    def _configure_subtitle_text_frame(self, shape) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        self._apply_text_frame_margins(text_frame)
        for paragraph in text_frame.paragraphs:
            self._apply_paragraph_spacing(paragraph, "subtitle", "text_full_width")

    def _minimum_title_height_emu(self, layout_key: str) -> int:
        if layout_key == "table":
            return 500000
        return 520000

    def _fit_title_font_size_points(self, text: str, width_emu: int, layout_key: str) -> float:
        base_size = self._title_font_size_points(layout_key)
        for candidate in (base_size, base_size - 2, base_size - 4, base_size - 6, base_size - 8):
            font_size = max(candidate, 22.0)
            estimated_height = self._estimate_text_height_emu(text, width_emu, font_size)
            if estimated_height <= 1650000:
                return font_size
        return 22.0

    def _fit_title_font_size_for_height(self, shape, text: str, layout_key: str, max_height_emu: int) -> float:
        base_size = self._title_font_size_points(layout_key)
        for font_size in range(int(base_size), 15, -1):
            if self._estimate_title_height_emu(shape, text, float(font_size)) <= max_height_emu:
                return float(font_size)
        return 16.0

    def _apply_font_size(self, shape, font_size_pt: float) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        for paragraph in text_frame.paragraphs:
            if not paragraph.runs and paragraph.text:
                run = paragraph.add_run()
                run.text = paragraph.text
                paragraph.text = ""
            for run in paragraph.runs:
                run.font.size = Pt(font_size_pt)

    def _configure_title_text_frame(self, shape) -> None:
        if not getattr(shape, "has_text_frame", False):
            return
        text_frame = shape.text_frame
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.SHAPE_TO_FIT_TEXT
        self._apply_text_frame_margins(text_frame)
        for paragraph in text_frame.paragraphs:
            self._apply_paragraph_spacing(paragraph, "title", "text_full_width")

    def _estimate_title_height_emu(self, shape, text: str, font_size_pt: float) -> int:
        text_frame = shape.text_frame
        effective_width = max(shape.width - text_frame.margin_left - text_frame.margin_right, shape.width // 2)
        width_pt = effective_width / self.EMU_PER_PT
        average_char_width_pt = max(font_size_pt * 0.52, 1.0)
        wrapped_lines = self._estimate_wrapped_line_count(text, width_pt, average_char_width_pt, min_chars_per_line=6)

        if len(text) >= 90:
            wrapped_lines = max(wrapped_lines, 2)
        if len(text) >= 135:
            wrapped_lines = max(wrapped_lines, 3)

        line_height_pt = font_size_pt * 1.24
        vertical_padding_pt = font_size_pt * 1.15
        return int((wrapped_lines * line_height_pt + vertical_padding_pt) * self.EMU_PER_PT)

    def _estimate_text_height_emu(self, text: str, width_emu: int, font_size_pt: float) -> int:
        if not text.strip() or width_emu <= 0:
            return 0

        width_pt = width_emu / self.EMU_PER_PT
        average_char_width_pt = max(font_size_pt * 0.52, 1.0)
        wrapped_lines = self._estimate_wrapped_line_count(text, width_pt, average_char_width_pt, min_chars_per_line=8)

        line_height_pt = font_size_pt * 1.18
        vertical_padding_pt = font_size_pt * 0.7
        return int((wrapped_lines * line_height_pt + vertical_padding_pt) * self.EMU_PER_PT)

    def _estimate_wrapped_line_count(
        self,
        text: str,
        width_pt: float,
        average_char_width_pt: float,
        *,
        min_chars_per_line: int,
    ) -> int:
        chars_per_line = max(int(width_pt / average_char_width_pt), min_chars_per_line)
        wrapped_lines = 0

        for paragraph in text.splitlines() or [text]:
            normalized = paragraph.strip()
            if not normalized:
                wrapped_lines += 1
                continue

            words = normalized.split()
            if len(words) <= 1:
                wrapped_lines += max(1, math.ceil(len(normalized) / chars_per_line))
                continue

            current_line_len = 0
            paragraph_lines = 1
            for word in words:
                word_len = len(word)
                projected = word_len if current_line_len == 0 else current_line_len + 1 + word_len
                if projected <= chars_per_line:
                    current_line_len = projected
                    continue
                if current_line_len == 0:
                    paragraph_lines += max(math.ceil(word_len / chars_per_line) - 1, 0)
                    current_line_len = word_len % chars_per_line or chars_per_line
                    continue
                paragraph_lines += 1
                current_line_len = word_len

            wrapped_lines += paragraph_lines

        return wrapped_lines

    def _fill_table(self, shape, slide_spec: SlideSpec, placeholder_spec: PlaceholderSpec | None = None) -> None:
        if slide_spec.table is None:
            if getattr(shape, "has_text_frame", False):
                self._set_text(shape, "", profile_for_layout("text_full_width"))
            return

        headers = slide_spec.table.headers
        rows = slide_spec.table.rows
        row_count = len(rows) + (1 if headers else 0)
        col_count = len(headers) if headers else max((len(row) for row in rows), default=0)
        render_as_shapes = self._component_behavior_bool(
            "table",
            "render_as_shapes",
            self._design_token_bool("table_render_as_shapes", default=False),
        )
        if row_count == 0 or col_count == 0:
            if getattr(shape, "has_text_frame", False):
                self._set_text(shape, "", profile_for_layout("text_full_width"))
            return

        if hasattr(shape, "insert_table"):
            try:
                target_width = shape.width
                target_height = self._table_target_height(shape)
                slide_shapes = shape.part.slide.shapes
                graphic_frame = shape.insert_table(row_count, col_count)
                graphic_frame.width = target_width
                graphic_frame.height = target_height
                table = graphic_frame.table
                current_row = 0
                if headers:
                    for col_index, value in enumerate(headers):
                        table.cell(0, col_index).text = value
                    current_row = 1
                for row_index, row in enumerate(rows, start=current_row):
                    for col_index, value in enumerate(row):
                        if col_index < col_count:
                            table.cell(row_index, col_index).text = value
                final_height = self._format_table(
                    table,
                    slide_spec.table,
                    graphic_frame.width,
                    graphic_frame.height,
                    placeholder_spec=placeholder_spec,
                )
                graphic_frame.height = final_height
                if render_as_shapes:
                    original_left = graphic_frame.left
                    original_top = graphic_frame.top
                    original_width = graphic_frame.width
                    original_height = graphic_frame.height
                    graphic_frame.left = -graphic_frame.width - 1000000
                    self._render_visible_table_grid(
                        slide_shapes,
                        slide_spec.table,
                        table,
                        left=original_left,
                        top=original_top,
                        width=original_width,
                        height=original_height,
                    )
                else:
                    self._overlay_table_grid_lines(
                        slide_shapes,
                        table,
                        left=graphic_frame.left,
                        top=graphic_frame.top,
                        width=graphic_frame.width,
                        height=graphic_frame.height,
                    )
                return
            except (AttributeError, TypeError, ValueError):
                pass

        if getattr(shape, "has_table", False):
            table = shape.table
            max_rows = len(table.rows)
            max_cols = len(table.columns)
            current_row = 0
            if headers and max_rows > 0:
                for col_index, value in enumerate(headers[:max_cols]):
                    table.cell(0, col_index).text = value
                current_row = 1
            for row_index, row in enumerate(rows, start=current_row):
                if row_index >= max_rows:
                    break
                for col_index, value in enumerate(row[:max_cols]):
                    table.cell(row_index, col_index).text = value
            final_height = self._format_table(
                table,
                slide_spec.table,
                shape.width,
                self._table_target_height(shape),
                placeholder_spec=placeholder_spec,
            )
            shape.height = final_height
            if render_as_shapes:
                original_left = shape.left
                original_top = shape.top
                original_width = shape.width
                original_height = shape.height
                shape.left = -shape.width - 1000000
                self._render_visible_table_grid(
                    shape.part.slide.shapes,
                    slide_spec.table,
                    table,
                    left=original_left,
                    top=original_top,
                    width=original_width,
                    height=original_height,
                )
            else:
                self._overlay_table_grid_lines(
                    shape.part.slide.shapes,
                    table,
                    left=shape.left,
                    top=shape.top,
                    width=shape.width,
                    height=shape.height,
                )
            return

        as_lines = []
        if headers:
            as_lines.append(" | ".join(headers))
        as_lines.extend(" | ".join(row) for row in rows)
        if getattr(shape, "has_text_frame", False):
            self._set_bullets(shape, as_lines, profile_for_layout("list_full_width"))

    def _table_target_height(self, shape) -> int:
        shape_height = int(getattr(shape, "height", 0) or 0)
        shape_top = int(getattr(shape, "top", 0) or 0)
        available_height = max(0, self.FOOTER_TOP_EMU - self.CONTENT_FOOTER_GAP_EMU - shape_top)
        return max(shape_height, available_height, 900000)

    def _fill_table_or_chart(self, shape, slide_spec: SlideSpec, placeholder_spec: PlaceholderSpec | None = None) -> None:
        if slide_spec.chart is not None:
            self._fill_chart(shape, slide_spec, placeholder_spec)
            return
        self._fill_table(shape, slide_spec, placeholder_spec)

    def _fill_chart(self, shape, slide_spec: SlideSpec, placeholder_spec: PlaceholderSpec | None = None) -> None:
        if slide_spec.chart is None:
            self._clear_placeholder(shape)
            return

        chart_spec = slide_spec.chart
        if not chart_spec.categories or not chart_spec.series:
            self._clear_placeholder(shape)
            return
        resolved_chart_spec = render_chart_spec(chart_spec)
        if resolved_chart_spec is None or not resolved_chart_spec.series:
            self._clear_placeholder(shape)
            return

        try:
            chart_type = self._resolve_chart_type(resolved_chart_spec)
            chart_data = CategoryChartData()
            chart_data.categories = resolved_chart_spec.categories
            for series in resolved_chart_spec.series:
                chart_data.add_series(series.name or "Ряд", series.values)

            slide_shapes = shape.part.slide.shapes
            left = shape.left
            top = shape.top
            width = shape.width
            height = shape.height
            shape_style = self._merged_component_shape_style(
                "chart",
                placeholder_spec.shape_style if placeholder_spec is not None else None,
            )
            if shape_style is not None:
                plot_left = shape_style.chart_plot_left_factor or 0.0
                plot_top = shape_style.chart_plot_top_factor or 0.0
                plot_width = shape_style.chart_plot_width_factor or 1.0
                plot_height = shape_style.chart_plot_height_factor or 1.0
                left = left + int(width * plot_left)
                top = top + int(height * plot_top)
                width = int(width * plot_width)
                height = int(height * plot_height)
            if getattr(shape, "is_placeholder", False):
                self._clear_placeholder(shape)

            graphic_frame = slide_shapes.add_chart(chart_type, left, top, width, height, chart_data)
            chart = graphic_frame.chart
            if resolved_chart_spec.chart_type == ChartType.COMBO:
                self._convert_chart_to_combo(chart, resolved_chart_spec)
            self._style_chart(chart, resolved_chart_spec)
            if shape_style is not None:
                self._apply_chart_shape_style(chart, shape_style)
        except Exception:
            if getattr(shape, "has_text_frame", False):
                self._set_text(shape, "Не удалось построить график", profile_for_layout("text_full_width"))
            else:
                self._clear_placeholder(shape)

    def _convert_chart_to_combo(self, chart, chart_spec: ChartSpec) -> None:
        chart_space = chart._chartSpace
        plot_area = chart_space.chart.plotArea
        bar_charts = plot_area.xpath("./c:barChart")
        if not bar_charts:
            return

        bar_chart = bar_charts[0]
        series = bar_chart.xpath("./c:ser")
        if len(series) < 2:
            return

        line_series = series[-1]
        bar_chart.remove(line_series)

        line_chart = OxmlElement("c:lineChart")
        grouping = OxmlElement("c:grouping")
        grouping.set("val", "standard")
        line_chart.append(grouping)

        self._ensure_line_series_shape(line_series)
        line_chart.append(line_series)

        bar_axis_ids = [ax_id.get("val") for ax_id in bar_chart.xpath("./c:axId")]
        if uses_secondary_value_axis(chart_spec) and len(bar_axis_ids) >= 1:
            secondary_axis_id = self._ensure_secondary_value_axis(chart)
            for axis_id in [bar_axis_ids[0], secondary_axis_id]:
                cloned_ax_id = OxmlElement("c:axId")
                cloned_ax_id.set("val", axis_id)
                line_chart.append(cloned_ax_id)
        else:
            for ax_id in bar_chart.xpath("./c:axId"):
                cloned_ax_id = OxmlElement("c:axId")
                cloned_ax_id.set("val", ax_id.get("val"))
                line_chart.append(cloned_ax_id)

        insert_at = list(plot_area).index(bar_chart) + 1
        plot_area.insert(insert_at, line_chart)

    def _ensure_secondary_value_axis(self, chart) -> str:
        chart_space = chart._chartSpace
        plot_area = chart_space.chart.plotArea
        value_axes = plot_area.xpath("./c:valAx")
        if len(value_axes) > 1:
            return value_axes[1].xpath("./c:axId")[0].get("val")

        category_axis = (plot_area.xpath("./c:catAx") or plot_area.xpath("./c:dateAx"))[0]
        primary_value_axis = value_axes[0]
        category_axis_id = category_axis.xpath("./c:axId")[0].get("val")
        secondary_axis_id = self._next_chart_axis_id(chart)
        secondary_axis = deepcopy(primary_value_axis)
        secondary_axis.xpath("./c:axId")[0].set("val", secondary_axis_id)
        secondary_axis.xpath("./c:axPos")[0].set("val", "r")
        secondary_axis.xpath("./c:crossAx")[0].set("val", category_axis_id)
        for gridlines in secondary_axis.xpath("./c:majorGridlines"):
            secondary_axis.remove(gridlines)
        plot_area.insert(list(plot_area).index(primary_value_axis) + 1, secondary_axis)
        return secondary_axis_id

    def _next_chart_axis_id(self, chart) -> str:
        axis_ids = [
            int(axis_id.get("val"))
            for axis_id in chart._chartSpace.xpath(".//c:axId")
            if axis_id.get("val") is not None
        ]
        next_id = max((abs(value) for value in axis_ids), default=1) + 1
        return str(next_id)

    def _ensure_line_series_shape(self, series_element) -> None:
        if not series_element.xpath("./c:marker"):
            marker = OxmlElement("c:marker")
            symbol = OxmlElement("c:symbol")
            symbol.set("val", "none")
            marker.append(symbol)
            insert_at = 0
            for idx, child in enumerate(series_element):
                if child.tag.endswith("}spPr"):
                    insert_at = idx + 1
                    break
            series_element.insert(insert_at, marker)

        if not series_element.xpath("./c:smooth"):
            smooth = OxmlElement("c:smooth")
            smooth.set("val", "0")
            series_element.append(smooth)

    def _resolve_chart_type(self, chart_spec: ChartSpec) -> XL_CHART_TYPE:
        chart_type_map = {
            ChartType.BAR: XL_CHART_TYPE.BAR_CLUSTERED,
            ChartType.COLUMN: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartType.LINE: XL_CHART_TYPE.LINE,
            ChartType.STACKED_BAR: XL_CHART_TYPE.BAR_STACKED,
            ChartType.STACKED_COLUMN: XL_CHART_TYPE.COLUMN_STACKED,
            ChartType.COMBO: XL_CHART_TYPE.COLUMN_CLUSTERED,
            ChartType.PIE: XL_CHART_TYPE.PIE,
        }
        return chart_type_map.get(chart_spec.chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    def _style_chart(self, chart, chart_spec: ChartSpec) -> None:
        self._style_chart_plot(chart, chart_spec)
        chart.has_legend = chart_spec.legend_visible and len(chart.series) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False
            self._style_chart_legend(chart)

        if chart_spec.title:
            chart.has_title = True
            chart.chart_title.text_frame.text = chart_spec.title
            self._style_chart_title(chart)
        else:
            chart.has_title = False

        for index, series in enumerate(chart.series):
            self._style_chart_series(series, chart_spec, index)
            if chart_spec.data_labels_visible:
                series.has_data_labels = True
                self._style_data_labels(series, chart_spec)

        if hasattr(chart, "category_axis"):
            try:
                chart.category_axis.has_title = bool(chart_spec.x_axis_title)
                if chart_spec.x_axis_title:
                    chart.category_axis.axis_title.text_frame.text = chart_spec.x_axis_title
                    self._style_axis_title(chart.category_axis)
                self._style_category_axis(chart.category_axis, chart_spec)
            except Exception:
                pass

        for axis_index, axis in enumerate(self._value_axes(chart)):
            axis_role = PRIMARY_AXIS if axis_index == 0 else SECONDARY_AXIS
            try:
                if axis_role == PRIMARY_AXIS:
                    axis.has_title = bool(chart_spec.y_axis_title)
                    if chart_spec.y_axis_title:
                        axis.axis_title.text_frame.text = chart_spec.y_axis_title
                        self._style_axis_title(axis)
                else:
                    axis.has_title = False
                axis.has_major_gridlines = axis_role == PRIMARY_AXIS
                self._style_value_axis(axis, chart_spec, axis_role=axis_role)
            except Exception:
                pass

    def _style_chart_plot(self, chart, chart_spec: ChartSpec) -> None:
        if not chart.plots:
            return
        plot = chart.plots[0]
        dense_threshold = self._chart_style_config()["denseCategoryThreshold"]
        very_dense_threshold = self._chart_style_config()["veryDenseCategoryThreshold"]
        try:
            if chart_spec.chart_type in {ChartType.COLUMN, ChartType.BAR, ChartType.STACKED_BAR, ChartType.STACKED_COLUMN}:
                if len(chart_spec.categories) >= very_dense_threshold:
                    plot.gap_width = 124
                elif len(chart_spec.categories) >= dense_threshold:
                    plot.gap_width = 96
                else:
                    plot.gap_width = 72
                plot.overlap = 100 if chart_spec.chart_type in {ChartType.STACKED_BAR, ChartType.STACKED_COLUMN} else 0
        except Exception:
            pass

    def _style_chart_title(self, chart) -> None:
        try:
            text_frame = chart.chart_title.text_frame
            font_size = Pt(self._component_font_size("chart", "title", fallback=self._title_font_size_points("table")))
            color = self._component_font_color("chart", "title", fallback=self._style_rgb("textColor"))
            self._style_text_frame_runs(text_frame, font_size=font_size, bold=True, color=color)
            self._apply_theme_text_style(text_frame, "title")
        except Exception:
            pass

    def _style_chart_legend(self, chart) -> None:
        try:
            chart.legend.font.size = Pt(11)
            chart.legend.font.color.rgb = self._style_rgb("mutedTextColor")
        except Exception:
            pass

    def _style_chart_series(self, series, chart_spec: ChartSpec, index: int) -> None:
        color = self._chart_series_color(chart_spec, index)
        if chart_spec.chart_type == ChartType.PIE:
            self._style_pie_points(series)
            return

        is_combo_line_series = chart_spec.chart_type == ChartType.COMBO and index == len(chart_spec.series) - 1

        try:
            if not is_combo_line_series:
                series.format.fill.solid()
                series.format.fill.fore_color.rgb = color
        except Exception:
            pass

        try:
            series.format.line.color.rgb = color
            if chart_spec.chart_type == ChartType.LINE or is_combo_line_series:
                series.format.line.width = Pt(2.5)
                self._style_line_marker(series, color)
        except Exception:
            pass

        if chart_spec.chart_type in {ChartType.BAR, ChartType.COLUMN} and len(chart_spec.series) == 1:
            self._style_ranked_points(series, chart_spec.series[0].values)

    def _style_pie_points(self, series) -> None:
        ranked_colors = self._ranked_point_colors(
            [point for point in getattr(series, "points", [])],
            getattr(series, "values", []),
        )
        for index, point in enumerate(getattr(series, "points", [])):
            color = ranked_colors.get(index, self._series_color(index))
            try:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
            except Exception:
                pass
            try:
                point.format.line.color.rgb = self._style_rgb("surfaceColor")
                point.format.line.width = Pt(1)
            except Exception:
                pass

    def _style_line_marker(self, series, color: RGBColor) -> None:
        try:
            series.marker.style = XL_MARKER_STYLE.CIRCLE
            series.marker.size = 7
            series.marker.format.fill.solid()
            series.marker.format.fill.fore_color.rgb = color
            series.marker.format.line.color.rgb = self._style_rgb("surfaceColor")
            series.marker.format.line.width = Pt(1)
        except Exception:
            pass

    def _style_ranked_points(self, series, values: list[float]) -> None:
        ranked_colors = self._ranked_point_colors(
            [point for point in getattr(series, "points", [])],
            values,
        )
        for index, point in enumerate(getattr(series, "points", [])):
            color = ranked_colors.get(index)
            if color is None:
                continue
            try:
                point.format.fill.solid()
                point.format.fill.fore_color.rgb = color
            except Exception:
                pass
            try:
                point.format.line.color.rgb = self._style_rgb("surfaceColor")
                point.format.line.width = Pt(1)
            except Exception:
                pass

    def _style_data_labels(self, series, chart_spec: ChartSpec) -> None:
        try:
            data_labels = series.data_labels
            data_labels.font.size = Pt(10)
            data_labels.font.color.rgb = self._style_rgb("mutedTextColor")
            data_labels.number_format = self._chart_number_format(chart_spec)
        except Exception:
            pass

    def _style_category_axis(self, axis, chart_spec: ChartSpec) -> None:
        self._style_axis_line(axis)
        self._style_tick_labels(axis, chart_spec)

    def _style_value_axis(self, axis, chart_spec: ChartSpec, *, axis_role: str = PRIMARY_AXIS) -> None:
        self._style_axis_line(axis)
        self._style_tick_labels(axis, chart_spec)
        try:
            if axis_role == PRIMARY_AXIS and axis.has_major_gridlines:
                axis.major_gridlines.format.line.color.rgb = self._style_rgb("gridColor")
                axis.major_gridlines.format.line.width = Pt(0.8)
        except Exception:
            pass
        try:
            axis.tick_labels.number_format = chart_axis_number_format_for_axis(chart_spec, axis_role) or self._chart_axis_number_format(chart_spec)
        except Exception:
            pass

    def _value_axes(self, chart) -> list[ValueAxis]:
        return [ValueAxis(element) for element in chart._chartSpace.valAx_lst]

    def _style_axis_line(self, axis) -> None:
        try:
            axis.format.line.color.rgb = self._style_rgb("axisColor")
            axis.format.line.width = Pt(0.8)
        except Exception:
            pass

    def _style_tick_labels(self, axis, chart_spec: ChartSpec) -> None:
        tick_size = self._chart_tick_font_size(chart_spec)
        try:
            axis.tick_labels.font.size = Pt(tick_size)
            axis.tick_labels.font.color.rgb = self._style_rgb("mutedTextColor")
        except Exception:
            pass

    def _chart_number_format(self, chart_spec: ChartSpec) -> str:
        if chart_spec.value_format == "currency":
            return '#,##0" ₽"'
        if chart_spec.value_format == "percent":
            return '0"%"'
        return "General"

    def _chart_axis_number_format(self, chart_spec: ChartSpec) -> str:
        return chart_axis_number_format(chart_spec) or "#,##0"

    def _chart_tick_font_size(self, chart_spec: ChartSpec) -> int:
        config = self._chart_style_config()
        if len(chart_spec.categories) >= config["veryDenseCategoryThreshold"]:
            return config["tickFontSizes"]["veryDense"]
        if len(chart_spec.categories) >= config["denseCategoryThreshold"]:
            return config["tickFontSizes"]["dense"]
        return config["tickFontSizes"]["default"]

    def _style_axis_title(self, axis) -> None:
        try:
            self._style_text_frame_runs(axis.axis_title.text_frame, font_size=Pt(11), bold=False, color=self._style_rgb("textColor"))
        except Exception:
            pass

    def _apply_chart_shape_style(self, chart, style: TemplateShapeStyleSpec) -> None:
        try:
            plot_area = chart._chartSpace.chart.plotArea
            if any(value is not None for value in (
                style.chart_plot_left_factor,
                style.chart_plot_top_factor,
                style.chart_plot_width_factor,
                style.chart_plot_height_factor,
            )):
                manual_layout = OxmlElement("c:manualLayout")
                for tag, value in (
                    ("c:xMode", "factor"),
                    ("c:yMode", "factor"),
                ):
                    el = OxmlElement(tag)
                    el.set("val", value)
                    manual_layout.append(el)
                for tag, value in (
                    ("c:x", style.chart_plot_left_factor),
                    ("c:y", style.chart_plot_top_factor),
                    ("c:w", style.chart_plot_width_factor),
                    ("c:h", style.chart_plot_height_factor),
                ):
                    if value is None:
                        continue
                    el = OxmlElement(tag)
                    el.set("val", str(value))
                    manual_layout.append(el)
                layout_nodes = plot_area.xpath("./c:layout")
                layout = layout_nodes[0] if layout_nodes else OxmlElement("c:layout")
                for child in list(layout):
                    layout.remove(child)
                layout.append(manual_layout)
                if not layout_nodes:
                    plot_area.insert(0, layout)
            if style.chart_category_axis_label_offset is not None:
                for axis in chart._chartSpace.xpath(".//c:catAx"):
                    self._set_axis_label_offset(axis, style.chart_category_axis_label_offset)
            if style.chart_value_axis_label_offset is not None:
                for axis in chart._chartSpace.xpath(".//c:valAx"):
                    self._set_axis_label_offset(axis, style.chart_value_axis_label_offset)
            if chart.has_legend and chart.legend is not None:
                legend = chart._chartSpace.chart.legend
                if style.chart_legend_offset_x_emu is not None or style.chart_legend_offset_y_emu is not None:
                    manual = OxmlElement("c:layout")
                    manual_layout = OxmlElement("c:manualLayout")
                    if style.chart_legend_offset_x_emu is not None:
                        x = OxmlElement("c:x")
                        x.set("val", str(style.chart_legend_offset_x_emu))
                        manual_layout.append(x)
                    if style.chart_legend_offset_y_emu is not None:
                        y = OxmlElement("c:y")
                        y.set("val", str(style.chart_legend_offset_y_emu))
                        manual_layout.append(y)
                    manual.append(manual_layout)
                    existing = legend.xpath("./c:layout")
                    if existing:
                        legend.remove(existing[0])
                    legend.insert(0, manual)
        except Exception:
            pass

    def _set_axis_label_offset(self, axis_element, value: int) -> None:
        existing = axis_element.xpath("./c:lblOffset")
        if existing:
            existing[0].set("val", str(value))
            return
        lbl = OxmlElement("c:lblOffset")
        lbl.set("val", str(value))
        axis_element.append(lbl)

    def _style_text_frame_runs(self, text_frame, *, font_size, bold: bool, color: RGBColor) -> None:
        for paragraph in text_frame.paragraphs:
            if not paragraph.runs and paragraph.text:
                run = paragraph.add_run()
                run.text = paragraph.text
                paragraph.text = ""
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.bold = bold
                theme = self._active_manifest.theme if self._active_manifest is not None else None
                title_style = theme.master_text_styles.get("title") if theme is not None else None
                if title_style is not None and title_style.font_family:
                    run.font.name = title_style.font_family
                run.font.color.rgb = color

    def _set_text_frame_color(self, text_frame, color: RGBColor) -> None:
        for paragraph in text_frame.paragraphs:
            if not paragraph.runs and paragraph.text:
                run = paragraph.add_run()
                run.text = paragraph.text
                paragraph.text = ""
            for run in paragraph.runs:
                run.font.color.rgb = color

    def _fill_image(self, shape, slide_spec: SlideSpec) -> None:
        if not slide_spec.image_base64:
            self._clear_placeholder(shape)
            return

        try:
            image_stream = BytesIO(base64.b64decode(slide_spec.image_base64))
        except Exception:
            self._clear_placeholder(shape)
            return

        if getattr(shape, "is_placeholder", False) and hasattr(shape, "insert_picture"):
            try:
                shape.insert_picture(image_stream)
                return
            except Exception:
                image_stream.seek(0)

        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height
        if getattr(shape, "is_placeholder", False):
            self._clear_placeholder(shape)

        try:
            slide_shapes = shape.part.slide.shapes
            slide_shapes.add_picture(image_stream, left, top, width=width, height=height)
        except Exception:
            if getattr(shape, "has_text_frame", False):
                self._set_text(shape, "Изображение из документа", profile_for_layout("text_full_width"))

    def _set_text(self, shape, text: str, layout_profile: LayoutCapacityProfile) -> None:
        text_frame = shape.text_frame
        text_frame.clear()
        text_frame.text = text
        if getattr(shape, "is_placeholder", False):
            placeholder_format = getattr(shape, "placeholder_format", None)
            if placeholder_format is not None and placeholder_format.idx in {15, 17, 21}:
                self._apply_theme_text_style(text_frame, self._text_role_for_shape(shape))
                self._apply_footer_font_size(text_frame, text)
                return
        self._configure_body_text_frame(text_frame)
        self._apply_body_font_size(text_frame, [text], shape, layout_profile)
        if getattr(shape, "is_placeholder", False) and (
            self._active_manifest is None or self._active_manifest.generation_mode != GenerationMode.PROTOTYPE
        ):
            self._apply_theme_text_style(text_frame, self._text_role_for_shape(shape), preserve_font_size=True)

    def _set_bullets(self, shape, items: list[str], layout_profile: LayoutCapacityProfile) -> None:
        text_frame = shape.text_frame
        text_frame.clear()
        if not items:
            return

        first = True
        for item in items:
            paragraph = text_frame.paragraphs[0] if first else text_frame.add_paragraph()
            paragraph.text = item
            if item:
                paragraph.level = 0
                self._apply_bullet_format(paragraph, layout_profile.layout_key)
                self._apply_paragraph_spacing(paragraph, "body", layout_profile.layout_key)
            first = False
        self._configure_body_text_frame(text_frame)
        self._apply_body_font_size(text_frame, items, shape, layout_profile)
        if getattr(shape, "is_placeholder", False) and (
            self._active_manifest is None or self._active_manifest.generation_mode != GenerationMode.PROTOTYPE
        ):
            self._apply_theme_text_style(text_frame, self._text_role_for_shape(shape), preserve_font_size=True)

    def _text_role_for_shape(self, shape) -> str:
        idx = None
        if getattr(shape, "is_placeholder", False):
            try:
                idx = shape.placeholder_format.idx
            except Exception:
                idx = None
        if idx == 0:
            return "title"
        if idx in {10, 11, 12, 14}:
            return "body"
        if idx in {13, 15, 17, 21}:
            return "other"
        return "body"

    def _fallback_theme_text_style(self, role: str) -> TemplateTextStyleSpec:
        fallback_styles = {
            "title": TemplateTextStyleSpec(font_family="Mont SemiBold", font_size_pt=28.0, color="#081C4F"),
            "body": TemplateTextStyleSpec(font_family="Mont Regular", font_size_pt=18.0, color="#081C4F", line_spacing=0.9, space_before_pt=10.0),
            "other": TemplateTextStyleSpec(font_family="Mont Regular", font_size_pt=18.0, color="#3489F3"),
        }
        return fallback_styles.get(role, fallback_styles["body"])

    def _apply_theme_text_style(self, text_frame, role: str, *, preserve_font_size: bool = False) -> None:
        theme = self._active_manifest.theme if self._active_manifest is not None else None
        style = None
        if theme is not None:
            style = theme.master_text_styles.get(role) or theme.master_text_styles.get("body")
        if style is None:
            style = self._fallback_theme_text_style(role)
        self._apply_text_style(text_frame, style, apply_font_family=True, preserve_font_size=preserve_font_size)

    def _body_theme_font_family(self) -> str:
        theme = self._active_manifest.theme if self._active_manifest is not None else None
        body_style = theme.master_text_styles.get("body") if theme is not None else None
        return (body_style.font_family if body_style is not None else None) or self._body_regular_font_family()

    def _body_regular_font_family(self) -> str:
        theme = self._active_manifest.theme if self._active_manifest is not None else None
        return (
            getattr(theme.font_scheme, "minor_latin", None) if theme is not None and getattr(theme, "font_scheme", None) is not None else None
        ) or (
            theme.master_text_styles.get("other").font_family
            if theme is not None and theme.master_text_styles.get("other") is not None
            else None
        ) or "Mont Regular"

    def _configure_body_text_frame(self, text_frame) -> None:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        self._apply_text_frame_margins(text_frame)

    def _configure_card_text_frame(self, text_frame) -> None:
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
        text_frame.margin_left = self.DEFAULT_TEXT_MARGIN_X_EMU
        text_frame.margin_right = self.DEFAULT_TEXT_MARGIN_X_EMU
        text_frame.margin_top = self.DEFAULT_TEXT_MARGIN_Y_EMU
        text_frame.margin_bottom = self.DEFAULT_TEXT_MARGIN_Y_EMU

    def _apply_paragraph_spacing(self, paragraph, role: str, layout_key: str) -> None:
        spacing = spacing_policy_for_layout(layout_key)
        role_policy = getattr(spacing, role)
        paragraph.line_spacing = role_policy.line_spacing
        paragraph.space_after = Pt(role_policy.space_after_pt)

    def _apply_text_frame_margins(self, text_frame) -> None:
        text_frame.margin_left = self.DEFAULT_TEXT_MARGIN_X_EMU
        text_frame.margin_right = self.DEFAULT_TEXT_MARGIN_X_EMU
        text_frame.margin_top = self.DEFAULT_TEXT_MARGIN_Y_EMU
        text_frame.margin_bottom = self.DEFAULT_TEXT_MARGIN_Y_EMU

    def _apply_body_font_size(self, text_frame, items: list[str], shape, layout_profile: LayoutCapacityProfile) -> None:
        non_empty_items = [item.strip() for item in items if item and item.strip()]
        if not non_empty_items:
            return

        total_chars = sum(len(item) for item in non_empty_items)
        max_item_len = max(len(item) for item in non_empty_items)
        item_count = len(non_empty_items)
        effective_max_item_len = 0 if item_count == 1 else max_item_len
        shape_height = getattr(shape, "height", 0) if shape is not None else 0
        shape_width = getattr(shape, "width", 0) if shape is not None else 0
        shape_height = shape_height or 0
        shape_width = shape_width or 0

        points = self._body_base_font_size(layout_profile)
        if layout_profile.layout_key == "list_full_width":
            points = self._body_base_font_size(layout_profile)
            if item_count == 1 and total_chars >= 260:
                points = max(points - 2, layout_profile.min_font_pt)
            elif item_count == 1 and total_chars >= 180:
                points = max(points - 1, layout_profile.min_font_pt)
            elif total_chars >= 700 or item_count >= 7 or effective_max_item_len >= 260:
                points = max(points - 4, layout_profile.min_font_pt)
            elif total_chars >= 520 or item_count >= 6 or effective_max_item_len >= 220:
                points = max(points - 3, layout_profile.min_font_pt)
            elif total_chars >= 420 or item_count >= 5 or effective_max_item_len >= 180:
                points = max(points - 2, layout_profile.min_font_pt)
            elif total_chars >= 320:
                points = max(points - 1, layout_profile.min_font_pt)
        elif shape_height >= 4_000_000 and shape_width >= 8_000_000 and item_count <= max(layout_profile.max_items + 2, 8):
            points = self._body_base_font_size(layout_profile)
        elif total_chars >= layout_profile.max_chars * 2.4 or effective_max_item_len >= 420:
            points = max(layout_profile.min_font_pt, points - 4)
        elif total_chars >= int(layout_profile.max_chars * 2.0) or effective_max_item_len >= 320 or item_count >= layout_profile.max_items + 2:
            points = max(points - 3, layout_profile.min_font_pt)
        elif total_chars >= int(layout_profile.max_chars * 1.75) or effective_max_item_len >= 260 or item_count >= layout_profile.max_items:
            points = max(points - 2, layout_profile.min_font_pt)
        elif (
            (total_chars >= int(layout_profile.max_chars * 1.6) and item_count > 5)
            or effective_max_item_len >= 200
            or item_count >= layout_profile.max_items
        ):
            points = max(points - 1, layout_profile.min_font_pt)

        if not getattr(shape, "is_placeholder", False) and shape_width < 5_500_000 and shape_height < 2_000_000:
            points = min(points, max(layout_profile.min_font_pt, 18))

        # Tight containers need one extra step down to avoid overflow on dense appendix-like slides.
        if shape_height and shape_height < 4000000 and (total_chars >= 900 or item_count >= 7):
            points = max(points - 1, layout_profile.min_font_pt)
        if shape_width and shape_width < 8000000 and total_chars >= 600:
            points = max(points - 1, layout_profile.min_font_pt)

        points = self._fit_body_font_size_to_shape(text_frame, shape, layout_profile, points)
        self._set_text_frame_font_size(text_frame, points, layout_profile.layout_key)

    def _fit_body_font_size_to_shape(self, text_frame, shape, layout_profile: LayoutCapacityProfile, points: int) -> int:
        shape_height = getattr(shape, "height", 0) if shape is not None else 0
        shape_width = getattr(shape, "width", 0) if shape is not None else 0
        if not shape_height or not shape_width:
            return points

        margin_top = getattr(text_frame, "margin_top", self.DEFAULT_TEXT_MARGIN_Y_EMU) or 0
        margin_bottom = getattr(text_frame, "margin_bottom", self.DEFAULT_TEXT_MARGIN_Y_EMU) or 0
        available_height = max(shape_height - margin_top - margin_bottom, 200000)
        for candidate in range(points, layout_profile.min_font_pt - 1, -1):
            if self._text_frame_height_fits_shape(text_frame, shape, layout_profile.layout_key, candidate, available_height):
                return candidate
        return layout_profile.min_font_pt

    def _text_frame_height_fits_shape(self, text_frame, shape, layout_key: str, font_size_pt: float, available_height_emu: int) -> bool:
        margin_left = getattr(text_frame, "margin_left", self.DEFAULT_TEXT_MARGIN_X_EMU) or 0
        margin_right = getattr(text_frame, "margin_right", self.DEFAULT_TEXT_MARGIN_X_EMU) or 0
        effective_width = max(int((shape.width - margin_left - margin_right) * 0.82), shape.width // 2)
        spacing = spacing_policy_for_layout(layout_key).body
        paragraph_gap_emu = int(Pt(spacing.space_after_pt).emu)
        total_height = 0
        non_empty_count = 0
        for paragraph in text_frame.paragraphs:
            paragraph_text = paragraph.text.strip()
            if not paragraph_text:
                continue
            non_empty_count += 1
            total_height += self._estimate_text_height_emu(paragraph_text, effective_width, font_size_pt)
        if non_empty_count > 1:
            total_height += paragraph_gap_emu * (non_empty_count - 1)
        return int(total_height * 1.04) <= available_height_emu

    def _set_text_frame_font_size(self, text_frame, points: float, layout_key: str) -> None:
        font_size = Pt(points)
        body_family = self._body_theme_font_family()
        for paragraph in text_frame.paragraphs:
            if not paragraph.runs and paragraph.text:
                run = paragraph.add_run()
                run.text = paragraph.text
                paragraph.text = ""
            self._apply_paragraph_spacing(paragraph, "body", layout_key)
            paragraph.font.size = font_size
            paragraph.font.name = body_family
            paragraph.font.bold = False
            for run in paragraph.runs:
                run.font.size = font_size
                self._apply_run_font_family(run, body_family)
                if run.font.bold is not True:
                    run.font.bold = False

    def _set_text_frame_regular(self, text_frame) -> None:
        body_family = self._body_regular_font_family()
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = body_family
            paragraph.font.bold = False
            for run in paragraph.runs:
                self._apply_run_font_family(run, body_family)
                run.font.bold = False

    def _body_base_font_size(self, layout_profile: LayoutCapacityProfile) -> int:
        if layout_profile.layout_key in {"text_full_width", "dense_text_full_width", "list_full_width", "image_text"}:
            return min(max(18, layout_profile.min_font_pt), layout_profile.max_font_pt)
        theme = self._active_manifest.theme if self._active_manifest is not None else None
        body_style = theme.master_text_styles.get("body") if theme is not None else None
        if body_style is not None and body_style.font_size_pt:
            return min(max(int(round(body_style.font_size_pt)), layout_profile.min_font_pt), layout_profile.max_font_pt)
        return layout_profile.max_font_pt

    def _apply_footer_font_size(self, text_frame, text: str) -> None:
        normalized = (text or "").strip()
        if not normalized:
            return

        points = self.FOOTER_FONT_PT
        if len(normalized) >= 160:
            points = 9.0
        elif len(normalized) >= 120:
            points = 10.0
        elif len(normalized) >= 80:
            points = 11.0

        font_size = Pt(points)

        for paragraph in text_frame.paragraphs:
            paragraph.font.size = font_size
            for run in paragraph.runs:
                run.font.size = font_size

    def _apply_bullet_format(self, paragraph, layout_key: str) -> None:
        spacing = spacing_policy_for_layout(layout_key).bullet
        paragraph_properties = paragraph._p.get_or_add_pPr()
        for child in list(paragraph_properties):
            if child.tag.endswith("}buNone") or child.tag.endswith("}buChar") or child.tag.endswith("}buAutoNum"):
                paragraph_properties.remove(child)

        bullet = OxmlElement("a:buChar")
        bullet.set("char", "•")
        paragraph_properties.insert(0, bullet)
        paragraph_properties.set("marL", str(spacing.margin_left_emu))
        paragraph_properties.set("indent", str(spacing.indent_emu))

    def _format_table(
        self,
        table,
        table_block: TableBlock | None,
        target_width: int,
        target_height: int,
        *,
        placeholder_spec: PlaceholderSpec | None = None,
    ) -> int:
        if table_block is None:
            headers: list[str] = []
            rows: list[list[str]] = []
            header_fill_colors: list[str | None] = []
            row_fill_colors: list[list[str | None]] = []
        else:
            headers = table_block.headers
            rows = table_block.rows
            header_fill_colors = table_block.header_fill_colors
            row_fill_colors = table_block.row_fill_colors
        all_rows = [headers, *rows] if headers else rows
        column_stats = self._column_stats(headers, rows)
        max_lengths = [item["max_len"] for item in column_stats]
        row_count = len(all_rows)
        col_count = len(max_lengths)
        max_cell_length = max(max_lengths, default=0)
        self._neutralize_table_style(table, has_header=bool(headers))
        avg_cell_length = (
            sum(len((headers[col] if headers and col < len(headers) else "")) for col in range(col_count))
            + sum(len(value or "") for row in rows for value in row[:col_count])
        ) / max(1, row_count * max(1, col_count))
        self._apply_table_geometry(
            table,
            column_stats,
            target_width,
            target_height,
            row_count,
            avg_cell_length=avg_cell_length,
            rows=all_rows,
            has_header=bool(headers),
        )
        font_size = self._estimate_table_font_size(
            row_count=row_count,
            col_count=col_count,
            max_cell_length=max_cell_length,
            avg_cell_length=avg_cell_length,
        )
        margins = self._estimate_table_margins(
            row_count=row_count,
            col_count=col_count,
            max_cell_length=max_cell_length,
            avg_cell_length=avg_cell_length,
        )
        margins = (
            self._component_spacing_emu("table", "cell_margin_left_emu", margins[0]),
            self._component_spacing_emu("table", "cell_margin_right_emu", margins[1]),
            self._component_spacing_emu("table", "cell_margin_top_emu", margins[2]),
            self._component_spacing_emu("table", "cell_margin_bottom_emu", margins[3]),
        )
        if placeholder_spec is not None and placeholder_spec.shape_style is not None:
            shape_style = placeholder_spec.shape_style
            margins = (
                shape_style.table_cell_margin_left_emu or margins[0],
                shape_style.table_cell_margin_right_emu or margins[1],
                shape_style.table_cell_margin_top_emu or margins[2],
                shape_style.table_cell_margin_bottom_emu or margins[3],
            )

        for row_index, row in enumerate(table.rows):
            is_header = headers and row_index == 0
            for col_index, cell in enumerate(row.cells):
                fill_color = (
                    header_fill_colors[col_index]
                    if is_header and col_index < len(header_fill_colors)
                    else (
                        row_fill_colors[row_index - 1][col_index]
                        if not is_header
                        and row_index - 1 < len(row_fill_colors)
                        and col_index < len(row_fill_colors[row_index - 1])
                        else None
                    )
                )
                self._style_table_cell(
                    cell,
                    is_header=is_header,
                    font_size=font_size,
                    margins=margins,
                    fill_color=fill_color,
                    placeholder_spec=placeholder_spec,
                )
        return max(sum(row.height for row in table.rows), 600000)

    def _column_stats(self, headers: list[str], rows: list[list[str]]) -> list[dict[str, float]]:
        col_count = len(headers) if headers else max((len(row) for row in rows), default=0)
        stats: list[dict[str, float]] = []
        for col_index in range(col_count):
            values: list[str] = []
            if col_index < len(headers):
                values.append(headers[col_index])
            for row in rows:
                if col_index < len(row):
                    values.append(row[col_index])
            lengths = [len(value or "") for value in values]
            stats.append(
                {
                    "max_len": max(lengths, default=8),
                    "avg_len": (sum(lengths) / len(lengths)) if lengths else 8.0,
                    "header_len": len(headers[col_index]) if col_index < len(headers) else 0,
                }
            )
        return stats

    def _apply_table_geometry(
        self,
        table,
        column_stats: list[dict[str, float]],
        target_width: int,
        target_height: int,
        row_count: int,
        *,
        avg_cell_length: float,
        rows: list[list[str]] | None = None,
        has_header: bool = False,
    ) -> None:
        if column_stats and target_width > 0:
            weights = self._column_width_weights(column_stats)
            weight_sum = sum(weights) or len(weights)
            assigned = 0
            for index, column in enumerate(table.columns):
                if index == len(weights) - 1:
                    width = max(target_width - assigned, int(target_width * 0.08))
                else:
                    min_share = 0.14 if len(weights) >= 3 else 0.1
                    width = max(int(target_width * weights[index] / weight_sum), int(target_width * min_share))
                    assigned += width
                column.width = width

        if row_count > 0 and target_height > 0:
            row_heights = self._table_row_heights(table, target_height, row_count, avg_cell_length, rows or [], has_header=has_header)
            for row, row_height in zip(table.rows, row_heights, strict=False):
                row.height = row_height

    def _table_row_height(self, target_height: int, row_count: int, avg_cell_length: float) -> int:
        computed = max(int(target_height / row_count), 200000)
        return computed

    def _table_row_heights(
        self,
        table,
        target_height: int,
        row_count: int,
        avg_cell_length: float,
        rows: list[list[str]],
        *,
        has_header: bool = False,
    ) -> list[int]:
        if row_count <= 0:
            return []
        base_height = self._table_row_height(target_height, row_count, avg_cell_length)
        if not rows:
            return [base_height] * row_count

        column_widths = [int(column.width) for column in table.columns]
        weights: list[float] = []
        for row_index in range(row_count):
            row = rows[row_index] if row_index < len(rows) else []
            if has_header and row_index == 0:
                weights.append(0.85)
                continue
            max_lines = 1.0
            for col_index, value in enumerate(row):
                width = column_widths[min(col_index, len(column_widths) - 1)] if column_widths else 1
                # Approximate wrapped text demand; enough for proportional row-height allocation.
                chars_per_line = max(8, int(width / 95000))
                explicit_lines = str(value or "").splitlines() or [str(value or "")]
                line_count = sum(max(1, (len(line) + chars_per_line - 1) // chars_per_line) for line in explicit_lines)
                max_lines = max(max_lines, float(line_count))
            weights.append(min(max(max_lines, 1.0), 3.5))

        total_weight = sum(weights) or float(row_count)
        min_height = 240000
        heights = [max(min_height, int(target_height * weight / total_weight)) for weight in weights]
        if sum(heights) > target_height:
            return [max(int(target_height / row_count), 180000)] * row_count
        delta = target_height - sum(heights)
        if heights:
            heights[-1] = max(min_height, heights[-1] + delta)
        return heights

    def _column_width_weights(self, column_stats: list[dict[str, float]]) -> list[float]:
        col_count = len(column_stats)
        weights: list[float] = []
        for stat in column_stats:
            # Bias toward typical cell size, not one outlier.
            weight = (stat["header_len"] * 0.9) + (stat["avg_len"] * 1.35) + (min(stat["max_len"], stat["avg_len"] * 1.8) * 0.55)
            weights.append(max(weight, 8.0))

        if col_count == 2:
            first_weight = weights[0]
            second_weight = weights[1]
            total = first_weight + second_weight
            if total > 0:
                first_share = first_weight / total
                first_share = min(max(first_share, 0.22), 0.42)
                return [first_share, 1.0 - first_share]

        if col_count == 3:
            total = sum(weights) or 1.0
            shares = [weight / total for weight in weights]
            normalized = []
            for share in shares:
                normalized.append(min(max(share, 0.14), 0.58))
            scale = sum(normalized) or 1.0
            return [share / scale for share in normalized]

        return weights

    def _estimate_table_font_size(
        self,
        *,
        row_count: int,
        col_count: int,
        max_cell_length: int,
        avg_cell_length: float,
    ) -> Pt:
        points = 11
        if row_count >= 8:
            points = 10
        if row_count >= 12:
            points = 9
        if row_count >= 16:
            points = 8

        if col_count >= 3:
            points -= 1
        if max_cell_length >= 90 or avg_cell_length >= 45:
            points -= 1
        if max_cell_length >= 180 or avg_cell_length >= 75:
            points -= 1
        if max_cell_length >= 240:
            points -= 2
        if row_count >= 10 and (max_cell_length >= 140 or avg_cell_length >= 60):
            points -= 1

        if row_count <= 4:
            points = max(points, 7 if max_cell_length >= 240 or avg_cell_length >= 100 else (8 if max_cell_length >= 180 or avg_cell_length >= 75 else 9))
        elif row_count <= 7:
            points = max(points, 7 if max_cell_length >= 220 else 8)
        else:
            points = max(points, 7 if max_cell_length >= 220 else 8)

        points = min(points, 11)
        return Pt(points)

    def _estimate_table_margins(
        self,
        *,
        row_count: int,
        col_count: int,
        max_cell_length: int,
        avg_cell_length: float,
    ) -> tuple[int, int, int, int]:
        if row_count >= 8 or col_count >= 3 or max_cell_length >= 90 or avg_cell_length >= 40:
            return (40000, 40000, 20000, 20000)
        if row_count >= 5 or max_cell_length >= 60:
            return (60000, 60000, 30000, 30000)
        return (80000, 80000, 40000, 40000)

    def _style_table_cell(
        self,
        cell,
        *,
        is_header: bool,
        font_size: Pt,
        margins: tuple[int, int, int, int],
        fill_color: str | None = None,
        placeholder_spec: PlaceholderSpec | None = None,
    ) -> None:
        fill = cell.fill
        header_fill_rgb = self._component_behavior_color_rgb("table", "header_fill_color") or self._design_token_color_rgb("table_header_fill_color")
        header_text_rgb = self._component_behavior_color_rgb("table", "header_text_color") or self._design_token_color_rgb("table_header_text_color")
        body_transparent = self._component_behavior_bool("table", "body_fill_transparent", self._design_token_bool("table_body_fill_transparent", default=False))
        preserve_source_fill = self._component_behavior_bool(
            "table",
            "preserve_source_fill_colors",
            self._design_token_bool("table_preserve_source_fill_colors", default=True),
        )

        effective_fill_color = fill_color if preserve_source_fill else None
        resolved_fill = None
        if is_header and header_fill_rgb is not None:
            fill.solid()
            fill.fore_color.rgb = header_fill_rgb
            resolved_fill = header_fill_rgb
        elif not is_header and body_transparent:
            try:
                fill.background()
            except Exception:
                fill.solid()
                fill.transparency = 1.0
        else:
            fill.solid()
            resolved_fill = self._table_fill_rgb(effective_fill_color, is_header=is_header)
            fill.fore_color.rgb = resolved_fill

        text_frame = cell.text_frame
        text_frame.word_wrap = True
        cell.margin_left, cell.margin_right, cell.margin_top, cell.margin_bottom = margins
        self._set_table_cell_borders(cell, self._table_border_rgb())
        text_style = self._component_text_style("table", "header" if is_header else "body")

        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.size = font_size
                run.font.bold = bool(text_style.bold) if text_style is not None and text_style.bold is not None else bool(is_header)
                if is_header and header_text_rgb is not None:
                    run.font.color.rgb = header_text_rgb
                else:
                    run.font.color.rgb = self._table_text_rgb_for_fill(resolved_fill, is_header=is_header)
                if text_style is not None and text_style.font_family:
                    run.font.name = text_style.font_family
                else:
                    theme = self._active_manifest.theme if self._active_manifest is not None else None
                    body_style = theme.master_text_styles.get("body") if theme is not None else None
                    if body_style is not None and body_style.font_family:
                        run.font.name = body_style.font_family

    def _table_fill_rgb(self, fill_color: str | None, *, is_header: bool) -> RGBColor:
        if fill_color:
            normalized = fill_color.strip().lstrip("#").upper()
            if re.fullmatch(r"[0-9A-F]{6}", normalized):
                return RGBColor.from_string(normalized)
        return RGBColor(0x08, 0x1C, 0x4F) if is_header else RGBColor(0xEB, 0xF3, 0xFE)

    def _table_text_rgb_for_fill(self, fill_rgb: RGBColor | None, *, is_header: bool) -> RGBColor:
        if fill_rgb is None:
            return RGBColor(0x08, 0x1C, 0x4F)
        red, green, blue = fill_rgb[0], fill_rgb[1], fill_rgb[2]
        luminance = (0.2126 * red + 0.7152 * green + 0.0722 * blue) / 255
        if luminance < 0.45:
            return RGBColor(0xFF, 0xFF, 0xFF)
        if is_header:
            return RGBColor(0x08, 0x1C, 0x4F)
        return RGBColor(0x08, 0x1C, 0x4F)

    def _design_token(self, key: str):
        if self._active_manifest is None:
            return None
        return self._active_manifest.design_tokens.get(key)

    def _component_style(self, component_key: str):
        if self._active_manifest is None:
            return None
        return self._active_manifest.component_styles.get(component_key)

    def _component_text_style(self, component_key: str, role: str):
        component_style = self._component_style(component_key)
        if component_style is None:
            return None
        return component_style.text_styles.get(role)

    def _component_shape_style(self, component_key: str):
        component_style = self._component_style(component_key)
        if component_style is None:
            return None
        return component_style.shape_style

    def _merged_component_shape_style(
        self,
        component_key: str,
        placeholder_style: TemplateShapeStyleSpec | None,
    ) -> TemplateShapeStyleSpec | None:
        component_style = self._component_shape_style(component_key)
        if component_style is None:
            return placeholder_style
        if placeholder_style is None:
            return component_style

        merged = component_style.model_dump()
        for key, value in placeholder_style.model_dump().items():
            if value is None:
                continue
            if isinstance(value, list) and not value:
                continue
            merged[key] = value
        return TemplateShapeStyleSpec.model_validate(merged)

    def _component_spacing_emu(self, component_key: str, token: str, fallback: int) -> int:
        component_style = self._component_style(component_key)
        value = component_style.spacing_tokens.get(token) if component_style is not None else None
        return int(value) if isinstance(value, (int, float)) else fallback

    def _component_behavior_float(self, component_key: str, token: str, fallback: float) -> float:
        component_style = self._component_style(component_key)
        value = component_style.behavior_tokens.get(token) if component_style is not None else None
        return float(value) if isinstance(value, (int, float)) else fallback

    def _component_behavior_bool(self, component_key: str, token: str, fallback: bool) -> bool:
        component_style = self._component_style(component_key)
        value = component_style.behavior_tokens.get(token) if component_style is not None else None
        return value if isinstance(value, bool) else fallback

    def _component_behavior_color_rgb(self, component_key: str, token: str) -> RGBColor | None:
        component_style = self._component_style(component_key)
        value = component_style.behavior_tokens.get(token) if component_style is not None else None
        if not isinstance(value, str):
            return None
        normalized = value.strip().lstrip("#")
        if re.fullmatch(r"[0-9A-Fa-f]{6}", normalized):
            return RGBColor.from_string(normalized.upper())
        return None

    def _component_font_size(self, component_key: str, role: str, *, fallback: float) -> float:
        style = self._component_text_style(component_key, role)
        if style is not None and style.font_size_pt:
            return float(style.font_size_pt)
        return fallback

    def _component_font_color(self, component_key: str, role: str, *, fallback: RGBColor) -> RGBColor:
        style = self._component_text_style(component_key, role)
        if style is not None and style.color:
            normalized = style.color.strip().lstrip("#")
            if len(normalized) == 6:
                return RGBColor.from_string(normalized)
        return fallback

    def _design_token_bool(self, key: str, *, default: bool) -> bool:
        value = self._design_token(key)
        if isinstance(value, bool):
            return value
        return default

    def _design_token_color_rgb(self, key: str) -> RGBColor | None:
        value = self._design_token(key)
        if not isinstance(value, str):
            return None
        normalized = value.strip()
        if not normalized:
            return None
        hex_color = normalized.lstrip("#")
        if re.fullmatch(r"[0-9A-Fa-f]{6}", hex_color):
            return RGBColor.from_string(hex_color.upper())
        rgba_match = re.fullmatch(
            r"rgba\(\s*(\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(\d{1,3})\s*,\s*(0|0?\.\d+|1(?:\.0+)?)\s*\)",
            normalized,
            re.IGNORECASE,
        )
        if rgba_match:
            red, green, blue = (int(rgba_match.group(index)) for index in range(1, 4))
            if all(0 <= channel <= 255 for channel in (red, green, blue)):
                return RGBColor(red, green, blue)
        return None

    def _table_border_rgb(self) -> RGBColor:
        return self._component_behavior_color_rgb("table", "border_color") or self._design_token_color_rgb("table_border_color") or RGBColor(0xC6, 0xDF, 0xFF)

    def _set_table_cell_borders(self, cell, color: RGBColor) -> None:
        tc_pr = cell._tc.get_or_add_tcPr()
        for edge_tag in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
            for existing in tc_pr.xpath(f"./{edge_tag}"):
                tc_pr.remove(existing)
            line = OxmlElement(edge_tag)
            line.set("w", "12700")
            solid_fill = OxmlElement("a:solidFill")
            rgb = OxmlElement("a:srgbClr")
            rgb.set("val", f"{color[0]:02X}{color[1]:02X}{color[2]:02X}")
            solid_fill.append(rgb)
            line.append(solid_fill)
            dash = OxmlElement("a:prstDash")
            dash.set("val", "solid")
            line.append(dash)
            round_join = OxmlElement("a:round")
            line.append(round_join)
            head_end = OxmlElement("a:headEnd")
            head_end.set("type", "none")
            head_end.set("w", "med")
            head_end.set("len", "med")
            line.append(head_end)
            tail_end = OxmlElement("a:tailEnd")
            tail_end.set("type", "none")
            tail_end.set("w", "med")
            tail_end.set("len", "med")
            line.append(tail_end)
            tc_pr.append(line)

    def _neutralize_table_style(self, table, *, has_header: bool) -> None:
        try:
            tbl_pr = table._tbl.tblPr
        except Exception:
            return
        if tbl_pr is None:
            return
        table_style_ids = list(tbl_pr.xpath("./a:tableStyleId"))
        for style_id in table_style_ids:
            tbl_pr.remove(style_id)
        if has_header:
            tbl_pr.set("firstRow", "1")
        else:
            tbl_pr.attrib.pop("firstRow", None)
        tbl_pr.set("bandRow", "0")

    def _render_visible_table_grid(self, slide_shapes, table_block: TableBlock | None, native_table, *, left: int, top: int, width: int, height: int) -> None:
        if table_block is None:
            return
        headers = table_block.headers
        rows = table_block.rows
        all_rows = [headers, *rows] if headers else rows
        if not all_rows:
            return

        row_count = len(all_rows)
        col_count = max(len(headers), max((len(row) for row in rows), default=0))
        max_cell_length = max((len(str(value or "")) for row in all_rows for value in row), default=0)
        total_cells = max(1, sum(len(row) for row in all_rows))
        avg_cell_length = sum(len(str(value or "")) for row in all_rows for value in row) / total_cells
        font_size = self._estimate_table_font_size(
            row_count=row_count,
            col_count=col_count,
            max_cell_length=max_cell_length,
            avg_cell_length=avg_cell_length,
        )
        margins = self._estimate_table_margins(
            row_count=row_count,
            col_count=col_count,
            max_cell_length=max_cell_length,
            avg_cell_length=avg_cell_length,
        )
        margins = (
            self._component_spacing_emu("table", "cell_margin_left_emu", margins[0]),
            self._component_spacing_emu("table", "cell_margin_right_emu", margins[1]),
            self._component_spacing_emu("table", "cell_margin_top_emu", margins[2]),
            self._component_spacing_emu("table", "cell_margin_bottom_emu", margins[3]),
        )
        header_fill_rgb = self._component_behavior_color_rgb("table", "header_fill_color") or self._design_token_color_rgb("table_header_fill_color") or RGBColor(0xC6, 0xDF, 0xFF)
        header_text_rgb = self._component_behavior_color_rgb("table", "header_text_color") or self._design_token_color_rgb("table_header_text_color") or RGBColor(0x09, 0x1E, 0x38)
        body_text_rgb = self._component_font_color("table", "body", fallback=RGBColor(0x08, 0x1C, 0x4F))
        border_rgb = self._table_border_rgb()
        header_style = self._component_text_style("table", "header")
        body_style = self._component_text_style("table", "body")

        column_widths = [int(column.width) for column in native_table.columns]
        row_heights = [int(row.height) for row in native_table.rows]

        current_top = top
        for row_index, row in enumerate(all_rows):
            current_left = left
            row_height = row_heights[min(row_index, len(row_heights) - 1)]
            for col_index in range(col_count):
                cell_width = column_widths[min(col_index, len(column_widths) - 1)]
                value = row[col_index] if col_index < len(row) else ""
                if row_index == 0 and headers:
                    background = slide_shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, current_left, current_top, cell_width, row_height)
                    background.fill.solid()
                    background.fill.fore_color.rgb = header_fill_rgb
                    background.line.fill.background()
                text_box = slide_shapes.add_textbox(current_left, current_top, cell_width, row_height)
                text_box.name = f"A3_TABLE_CELL_{row_index}_{col_index}"
                text_box.fill.background()
                text_box.line.fill.background()
                text_frame = text_box.text_frame
                text_frame.clear()
                text_frame.word_wrap = True
                text_frame.margin_left, text_frame.margin_right, text_frame.margin_top, text_frame.margin_bottom = margins
                role_style = header_style if row_index == 0 and headers else body_style
                role_fallback_size = font_size.pt if hasattr(font_size, "pt") else float(font_size)
                configured_size = self._component_font_size("table", "header" if row_index == 0 and headers else "body", fallback=role_fallback_size)
                effective_font_size = min(configured_size, role_fallback_size)
                lines = str(value or "").splitlines() or [str(value or "")]
                for line_index, line in enumerate(lines):
                    paragraph = text_frame.paragraphs[0] if line_index == 0 else text_frame.add_paragraph()
                    paragraph.line_spacing = 1.0
                    paragraph.space_after = Pt(0)
                    paragraph.space_before = Pt(0)
                    run = paragraph.add_run()
                    run.text = line
                    run.font.size = Pt(effective_font_size)
                    run.font.bold = bool(role_style.bold) if role_style is not None and role_style.bold is not None else bool(row_index == 0 and headers)
                    run.font.color.rgb = header_text_rgb if row_index == 0 and headers else body_text_rgb
                    if role_style is not None and role_style.font_family:
                        run.font.name = role_style.font_family
                current_left += cell_width
            current_top += row_height

        self._draw_visible_table_borders(
            slide_shapes,
            column_widths,
            row_heights,
            left=left,
            top=top,
            border_rgb=border_rgb,
        )

    def _draw_visible_table_borders(self, slide_shapes, column_widths: list[int], row_heights: list[int], *, left: int, top: int, border_rgb: RGBColor) -> None:
        thickness = 28000
        total_width = sum(column_widths)
        total_height = sum(row_heights)

        x_positions = [left]
        current_x = left
        for value in column_widths:
            current_x += value
            x_positions.append(current_x)

        y_positions = [top]
        current_y = top
        for value in row_heights:
            current_y += value
            y_positions.append(current_y)

        for x in x_positions:
            rect = slide_shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, x - thickness // 2, top, thickness, total_height)
            rect.fill.solid()
            rect.fill.fore_color.rgb = border_rgb
            rect.line.fill.background()

        for y in y_positions:
            rect = slide_shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, y - thickness // 2, total_width, thickness)
            rect.fill.solid()
            rect.fill.fore_color.rgb = border_rgb
            rect.line.fill.background()

    def _overlay_table_grid_lines(self, slide_shapes, table, *, left: int, top: int, width: int, height: int) -> None:
        border_rgb = self._table_border_rgb()
        thickness = 42000
        cover = thickness * 2

        x_positions = [left]
        current_x = left
        for column in table.columns:
            current_x += int(column.width)
            x_positions.append(current_x)

        y_positions = [top]
        current_y = top
        for row in table.rows:
            current_y += int(row.height)
            y_positions.append(current_y)

        # Clamp last boundaries to the actual rendered frame.
        x_positions[-1] = left + width
        y_positions[-1] = top + height

        for x in x_positions:
            rect = slide_shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                x - thickness,
                top - thickness,
                cover,
                height + cover,
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = border_rgb
            rect.line.fill.background()

        for y in y_positions:
            rect = slide_shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                left - thickness,
                y - thickness,
                width + cover,
                cover,
            )
            rect.fill.solid()
            rect.fill.fore_color.rgb = border_rgb
            rect.line.fill.background()

    def _chart_palette_rgb(self) -> list[RGBColor]:
        colors = [
            self._component_behavior_color_rgb("chart", "rank_color_1"),
            self._component_behavior_color_rgb("chart", "rank_color_2"),
            self._component_behavior_color_rgb("chart", "rank_color_3"),
            self._component_behavior_color_rgb("chart", "rank_color_4"),
        ]
        if not any(color is not None for color in colors):
            token_keys = (
                "chart_rank_color_1",
                "chart_rank_color_2",
                "chart_rank_color_3",
                "chart_rank_color_4",
            )
            colors = [self._design_token_color_rgb(key) for key in token_keys]
        filtered = [color for color in colors if color is not None]
        if filtered:
            return filtered
        return [
            RGBColor(0x09, 0x1E, 0x38),
            RGBColor(0x34, 0x89, 0xF3),
            RGBColor(0x26, 0x45, 0x95),
            RGBColor(0xBF, 0xCE, 0xF5),
        ]

    def _chart_series_color(self, chart_spec: ChartSpec, index: int) -> RGBColor:
        return self._series_color(index)

    def _ranked_point_colors(self, points, values) -> dict[int, RGBColor]:
        palette = self._chart_palette_rgb()
        ranked = sorted(
            [(index, float(values[index])) for index in range(min(len(points), len(values)))],
            key=lambda item: item[1],
            reverse=True,
        )
        return {
            point_index: palette[min(rank, len(palette) - 1)]
            for rank, (point_index, _) in enumerate(ranked)
        }

    def _clear_placeholder(self, shape) -> None:
        if getattr(shape, "has_text_frame", False):
            self._remove_shape(shape)
            return
        if getattr(shape, "has_table", False):
            return
        self._remove_shape(shape)

    def _remove_shape(self, shape) -> None:
        parent = shape._element.getparent()
        if parent is not None:
            parent.remove(shape._element)

    def _is_empty_binding_value(self, value) -> bool:
        if isinstance(value, list):
            return not any(str(item).strip() for item in value)
        return not str(value or "").strip()

    def _validate_output_file(self, output_path: Path, expected_slide_count: int) -> None:
        try:
            with zipfile.ZipFile(output_path) as archive:
                bad_entry = archive.testzip()
                if bad_entry is not None:
                    raise ValueError(f"Generated PPTX contains a corrupted archive entry: {bad_entry}")

                archive_names = archive.namelist()
                duplicates = sorted({name for name in archive_names if archive_names.count(name) > 1})
                if duplicates:
                    raise ValueError(f"Generated PPTX contains duplicate package entries: {', '.join(duplicates[:5])}")

            presentation = Presentation(str(output_path))
            if len(presentation.slides) != expected_slide_count:
                raise ValueError(
                    f"Generated PPTX slide count mismatch: expected {expected_slide_count}, got {len(presentation.slides)}"
                )
        except Exception:
            output_path.unlink(missing_ok=True)
            raise
