from __future__ import annotations

import json
import re
from pathlib import Path

from pptx import Presentation

from a3presentation.domain.template import (
    ComponentConfidence,
    ComponentEditability,
    ComponentGeometry,
    ComponentStyle,
    ExtractedComponent,
    ExtractedComponentRole,
    ExtractedComponentType,
    ExtractedPresentationInventory,
    ExtractedSlideInventory,
    GenerationMode,
    InventorySourceKind,
    LayoutSpec,
    PlaceholderKind,
    PlaceholderSpec,
    TemplateParagraphStyleCatalog,
    PrototypeSlideSpec,
    PrototypeTokenSpec,
    TemplateComponentStyleSpec,
    TemplateShapeStyleSpec,
    TemplateManifest,
    TemplateTextStyleSpec,
)


class TemplateAnalyzer:
    TOKEN_PATTERN = re.compile(r"{{\s*([a-zA-Z0-9_]+)\s*}}")

    def analyze(self, template_id: str, template_path: Path, display_name: str | None = None) -> TemplateManifest:
        manifest_path = template_path.with_name("manifest.json")
        if manifest_path.exists():
            payload = json.loads(manifest_path.read_text(encoding="utf-8"))
            payload["template_id"] = template_id
            if display_name:
                payload["display_name"] = display_name
            payload["source_pptx"] = template_path.name
            manifest = TemplateManifest.model_validate(payload)
            analyzed = self._analyze_presentation(template_path, template_id, display_name)
            self._backfill_geometry(manifest, analyzed)
            if not manifest.component_styles:
                manifest.component_styles = analyzed.component_styles
            return manifest

        return self._analyze_presentation(template_path, template_id, display_name)

    def _analyze_presentation(self, template_path: Path, template_id: str, display_name: str | None = None) -> TemplateManifest:
        presentation = Presentation(str(template_path))
        layouts: list[LayoutSpec] = []
        prototype_slides: list[PrototypeSlideSpec] = []
        inventory_components: list[ExtractedComponent] = []
        inventory_slides: list[ExtractedSlideInventory] = []
        used_layout_keys: set[str] = set()
        has_explicit_token_prototypes = False
        has_synthesized_prototypes = False

        for master_index, slide_master in enumerate(presentation.slide_masters):
            for layout_index, slide_layout in enumerate(slide_master.slide_layouts):
                placeholders: list[PlaceholderSpec] = []
                layout_component_ids: list[str] = []
                for shape in slide_layout.placeholders:
                    placeholder_format = shape.placeholder_format
                    placeholder_kind = self._map_placeholder_kind(placeholder_format.type)
                    placeholder = PlaceholderSpec(
                        name=shape.name,
                        kind=placeholder_kind,
                        idx=placeholder_format.idx,
                        shape_name=shape.name,
                        binding=self._infer_placeholder_binding(placeholder_kind),
                        editable_role=self._editable_role_for_placeholder_kind(placeholder_kind),
                        editable_capabilities=self._editable_capabilities_for_binding(
                            self._infer_placeholder_binding(placeholder_kind),
                            placeholder_kind.value,
                        ),
                        left_emu=int(shape.left),
                        top_emu=int(shape.top),
                        width_emu=int(shape.width),
                        height_emu=int(shape.height),
                        margin_left_emu=self._text_frame_margin(shape, "left"),
                        margin_right_emu=self._text_frame_margin(shape, "right"),
                        margin_top_emu=self._text_frame_margin(shape, "top"),
                        margin_bottom_emu=self._text_frame_margin(shape, "bottom"),
                    )
                    placeholders.append(placeholder)
                    component = self._inventory_component_from_placeholder(
                        placeholder=placeholder,
                        source_kind=InventorySourceKind.LAYOUT,
                        source_index=layout_index,
                        source_name=slide_layout.name,
                    )
                    inventory_components.append(component)
                    layout_component_ids.append(component.component_id)

                supported_slide_kinds = self._infer_slide_kinds(placeholders, slide_layout.name)
                layout_key = self._make_unique_layout_key(
                    base_key=self._slugify(slide_layout.name, fallback=f"layout_{master_index}_{layout_index}"),
                    master_index=master_index,
                    layout_index=layout_index,
                    used_keys=used_layout_keys,
                )
                layouts.append(
                    LayoutSpec(
                        key=layout_key,
                        name=slide_layout.name,
                        slide_master_index=master_index,
                        slide_layout_index=layout_index,
                        supported_slide_kinds=supported_slide_kinds,
                        representation_hints=self._infer_representation_hints_from_placeholders(placeholders),
                        placeholders=placeholders,
                    )
                )
                inventory_slides.append(
                    ExtractedSlideInventory(
                        source_kind=InventorySourceKind.LAYOUT,
                        source_index=layout_index,
                        name=slide_layout.name,
                        component_ids=layout_component_ids,
                        supported_slide_kinds=supported_slide_kinds,
                        representation_hints=self._infer_representation_hints_from_placeholders(placeholders),
                    )
                )

        for index, slide in enumerate(presentation.slides):
            tokens: list[PrototypeTokenSpec] = []
            seen_tokens: set[tuple[str, str | None]] = set()
            slide_components = self._extract_slide_inventory_components(slide, index)
            inventory_components.extend(slide_components)
            inventory_slides.append(
                ExtractedSlideInventory(
                    source_kind=InventorySourceKind.SLIDE,
                    source_index=index,
                    name=slide.name or f"Slide {index + 1}",
                    component_ids=[component.component_id for component in slide_components],
                )
            )

            for shape in slide.shapes:
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = shape.text or ""
                for match in self.TOKEN_PATTERN.finditer(text):
                    token = match.group(1)
                    key = (token, shape.name)
                    if key in seen_tokens:
                        continue
                    seen_tokens.add(key)
                    binding = self._infer_binding(token)
                    tokens.append(
                        PrototypeTokenSpec(
                            token=token,
                            binding=binding,
                            shape_name=shape.name,
                            editable_role=self._editable_role_for_binding(binding),
                            editable_capabilities=self._editable_capabilities_for_binding(binding),
                            slot_group=self._infer_slot_group(token),
                            slot_group_order=self._infer_slot_group_order(token),
                            left_emu=int(shape.left),
                            top_emu=int(shape.top),
                            width_emu=int(shape.width),
                            height_emu=int(shape.height),
                            margin_left_emu=self._text_frame_margin(shape, "left"),
                            margin_right_emu=self._text_frame_margin(shape, "right"),
                            margin_top_emu=self._text_frame_margin(shape, "top"),
                            margin_bottom_emu=self._text_frame_margin(shape, "bottom"),
                        )
                    )

            if tokens:
                has_explicit_token_prototypes = True
            else:
                tokens = self._synthesize_tokens_from_slide(slide, presentation)
                if not tokens:
                    continue
                has_synthesized_prototypes = True

            prototype_slides.append(
                PrototypeSlideSpec(
                    key=f"slide_{index + 1}",
                    name=slide.name or f"Slide {index + 1}",
                    source_slide_index=index,
                    supported_slide_kinds=self._infer_slide_kinds_from_tokens(tokens),
                    representation_hints=self._infer_representation_hints_from_tokens(tokens),
                    tokens=tokens,
                )
            )

        default_layout_key = layouts[1].key if len(layouts) > 1 else (layouts[0].key if layouts else None)
        generation_mode = GenerationMode.LAYOUT
        if (
            has_explicit_token_prototypes
            or has_synthesized_prototypes
            or (prototype_slides and not self._has_usable_layout_inventory(layouts))
        ):
            generation_mode = GenerationMode.PROTOTYPE
        inventory = ExtractedPresentationInventory(
            components=inventory_components,
            slides=inventory_slides,
            warnings=self._inventory_warnings(layouts, prototype_slides, inventory_components),
            degradation_mode=self._inventory_degradation_mode(layouts, prototype_slides, inventory_components),
            has_usable_layout_inventory=self._has_usable_layout_inventory(layouts),
            has_prototype_inventory=bool(prototype_slides),
        )
        return TemplateManifest(
            template_id=template_id,
            display_name=display_name or template_id,
            source_pptx=template_path.name,
            generation_mode=generation_mode,
            default_layout_key=default_layout_key,
            component_styles=self._infer_component_styles(layouts),
            layouts=layouts,
            prototype_slides=prototype_slides,
            inventory=inventory,
        )

    def _has_usable_layout_inventory(self, layouts: list[LayoutSpec]) -> bool:
        return any(
            layout.placeholders
            and (
                layout.supported_slide_kinds
                or layout.representation_hints
                or any(
                    placeholder.binding or placeholder.editable_role or placeholder.editable_capabilities
                    for placeholder in layout.placeholders
                )
            )
            for layout in layouts
        )

    def _synthesize_tokens_from_slide(self, slide, presentation: Presentation) -> list[PrototypeTokenSpec]:
        tokens: list[PrototypeTokenSpec] = []
        seen_shapes: set[str] = set()
        slide_width = int(presentation.slide_width or 0)
        slide_height = int(presentation.slide_height or 0)

        text_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and self._shape_text_content(shape)
        ]
        text_shapes.sort(key=lambda shape: (int(shape.top), int(shape.left), -int(shape.width)))
        title_shape, subtitle_shape = self._select_title_and_subtitle_shapes(text_shapes, slide_width, slide_height)
        two_column_shapes = self._detect_two_column_text_shapes(
            [shape for shape in text_shapes if shape is not title_shape and shape is not subtitle_shape],
            slide_width,
        )

        for shape in text_shapes:
            binding = None
            if shape == title_shape:
                binding = "title"
            elif shape == subtitle_shape:
                binding = "subtitle"
            elif two_column_shapes and shape in two_column_shapes["left"]:
                binding = "left_bullets" if self._is_bullet_like_shape(shape) else "left_text"
            elif two_column_shapes and shape in two_column_shapes["right"]:
                binding = "right_bullets" if self._is_bullet_like_shape(shape) else "right_list"
            elif self._is_bullet_like_shape(shape):
                binding = "bullets"
            elif not any(token.binding == "body" for token in tokens):
                binding = "body"
            elif not any(token.binding == "secondary_text" for token in tokens):
                binding = "secondary_text"
            else:
                binding = "notes"

            tokens.append(self._prototype_token_from_shape(shape, binding))
            seen_shapes.add(shape.name)

        for shape in slide.shapes:
            if shape.name in seen_shapes:
                continue
            binding = None
            if getattr(shape, "has_table", False):
                binding = "table"
            elif getattr(shape, "has_chart", False):
                binding = "chart"
            elif self._shape_looks_like_picture(shape):
                binding = "image"
            if not binding:
                continue
            tokens.append(self._prototype_token_from_shape(shape, binding))
            seen_shapes.add(shape.name)

        return tokens

    def _prototype_token_from_shape(self, shape, binding: str) -> PrototypeTokenSpec:
        return PrototypeTokenSpec(
            token=binding,
            binding=binding,
            shape_name=shape.name,
            editable_role=self._editable_role_for_binding(binding),
            editable_capabilities=self._editable_capabilities_for_binding(binding),
            slot_group=self._infer_slot_group(binding),
            slot_group_order=self._infer_slot_group_order(binding),
            left_emu=int(shape.left),
            top_emu=int(shape.top),
            width_emu=int(shape.width),
            height_emu=int(shape.height),
            margin_left_emu=self._text_frame_margin(shape, "left"),
            margin_right_emu=self._text_frame_margin(shape, "right"),
            margin_top_emu=self._text_frame_margin(shape, "top"),
            margin_bottom_emu=self._text_frame_margin(shape, "bottom"),
            text_style=self._extract_text_style_from_shape(shape),
            paragraph_styles=self._extract_paragraph_styles_from_shape(shape),
            shape_style=self._extract_shape_style_from_shape(shape),
        )

    def _inventory_component_from_placeholder(
        self,
        placeholder: PlaceholderSpec,
        source_kind: InventorySourceKind,
        source_index: int,
        source_name: str | None,
    ) -> ExtractedComponent:
        role = self._inventory_role_from_binding_or_kind(placeholder.binding, placeholder.kind)
        return ExtractedComponent(
            component_id=f"{source_kind.value}_{source_index}_{placeholder.shape_name or placeholder.name}",
            source_kind=source_kind,
            source_index=source_index,
            source_name=source_name,
            shape_name=placeholder.shape_name or placeholder.name,
            component_type=self._inventory_type_from_placeholder_kind(placeholder.kind),
            role=role,
            binding=placeholder.binding,
            confidence=ComponentConfidence.HIGH if placeholder.binding or placeholder.kind != PlaceholderKind.UNKNOWN else ComponentConfidence.MEDIUM,
            editability=self._inventory_editability_for_role(role),
            capabilities=list(placeholder.editable_capabilities),
            geometry=ComponentGeometry(
                left_emu=placeholder.left_emu,
                top_emu=placeholder.top_emu,
                width_emu=placeholder.width_emu,
                height_emu=placeholder.height_emu,
                margin_left_emu=placeholder.margin_left_emu,
                margin_right_emu=placeholder.margin_right_emu,
                margin_top_emu=placeholder.margin_top_emu,
                margin_bottom_emu=placeholder.margin_bottom_emu,
            ),
        )

    def _extract_slide_inventory_components(self, slide, slide_index: int) -> list[ExtractedComponent]:
        components: list[ExtractedComponent] = []
        for shape_index, shape in enumerate(slide.shapes):
            component = self._inventory_component_from_shape(shape, slide_index, shape_index, slide.name or f"Slide {slide_index + 1}")
            if component is not None:
                components.append(component)
        return components

    def _inventory_component_from_shape(
        self,
        shape,
        slide_index: int,
        shape_index: int,
        slide_name: str,
    ) -> ExtractedComponent | None:
        component_type = self._inventory_type_from_shape(shape)
        role = self._inventory_role_from_shape(shape)
        confidence = self._inventory_confidence_for_shape(shape, role)
        editability = self._inventory_editability_for_role(role)
        text_excerpt = self._shape_text_content(shape) or None
        return ExtractedComponent(
            component_id=f"slide_{slide_index}_shape_{shape_index}",
            source_kind=InventorySourceKind.SLIDE,
            source_index=slide_index,
            source_name=slide_name,
            shape_name=getattr(shape, "name", None),
            component_type=component_type,
            role=role,
            binding=self._inventory_binding_from_shape(shape, role),
            confidence=confidence,
            editability=editability,
            capabilities=self._inventory_capabilities_from_shape(shape, role),
            geometry=ComponentGeometry(
                left_emu=int(getattr(shape, "left", 0)) if getattr(shape, "left", None) is not None else None,
                top_emu=int(getattr(shape, "top", 0)) if getattr(shape, "top", None) is not None else None,
                width_emu=int(getattr(shape, "width", 0)) if getattr(shape, "width", None) is not None else None,
                height_emu=int(getattr(shape, "height", 0)) if getattr(shape, "height", None) is not None else None,
                margin_left_emu=self._text_frame_margin(shape, "left"),
                margin_right_emu=self._text_frame_margin(shape, "right"),
                margin_top_emu=self._text_frame_margin(shape, "top"),
                margin_bottom_emu=self._text_frame_margin(shape, "bottom"),
            ),
            style=ComponentStyle(
                text_style=self._extract_text_style_from_shape(shape),
                paragraph_styles=self._extract_paragraph_styles_from_shape(shape),
                shape_style=self._extract_shape_style_from_shape(shape),
            ),
            text_excerpt=text_excerpt[:160] if text_excerpt else None,
        )

    def _inventory_type_from_placeholder_kind(self, kind: PlaceholderKind) -> ExtractedComponentType:
        if kind in {PlaceholderKind.TITLE, PlaceholderKind.SUBTITLE, PlaceholderKind.BODY}:
            return ExtractedComponentType.PLACEHOLDER
        if kind == PlaceholderKind.IMAGE:
            return ExtractedComponentType.IMAGE
        if kind == PlaceholderKind.TABLE:
            return ExtractedComponentType.TABLE
        if kind == PlaceholderKind.CHART:
            return ExtractedComponentType.CHART
        if kind == PlaceholderKind.FOOTER:
            return ExtractedComponentType.FOOTER
        return ExtractedComponentType.PLACEHOLDER

    def _inventory_type_from_shape(self, shape) -> ExtractedComponentType:
        if getattr(shape, "has_table", False):
            return ExtractedComponentType.TABLE
        if getattr(shape, "has_chart", False):
            return ExtractedComponentType.CHART
        if self._shape_looks_like_picture(shape):
            return ExtractedComponentType.IMAGE
        if "GROUP" in str(getattr(shape, "shape_type", "")):
            return ExtractedComponentType.GROUP
        if getattr(shape, "has_text_frame", False):
            return ExtractedComponentType.TEXT
        return ExtractedComponentType.UNKNOWN

    def _inventory_role_from_binding_or_kind(self, binding: str | None, kind: PlaceholderKind) -> ExtractedComponentRole:
        if binding:
            return self._inventory_role_from_binding(binding)
        if kind == PlaceholderKind.TITLE:
            return ExtractedComponentRole.TITLE
        if kind == PlaceholderKind.SUBTITLE:
            return ExtractedComponentRole.SUBTITLE
        if kind == PlaceholderKind.BODY:
            return ExtractedComponentRole.BODY
        if kind == PlaceholderKind.TABLE:
            return ExtractedComponentRole.TABLE
        if kind == PlaceholderKind.CHART:
            return ExtractedComponentRole.CHART
        if kind == PlaceholderKind.IMAGE:
            return ExtractedComponentRole.IMAGE
        if kind == PlaceholderKind.FOOTER:
            return ExtractedComponentRole.FOOTER
        return ExtractedComponentRole.UNKNOWN

    def _inventory_role_from_binding(self, binding: str) -> ExtractedComponentRole:
        role = self._editable_role_for_binding(binding)
        if role == "title":
            return ExtractedComponentRole.TITLE
        if role == "subtitle":
            return ExtractedComponentRole.SUBTITLE
        if role == "body":
            return ExtractedComponentRole.BODY
        if role == "bullet_list":
            return ExtractedComponentRole.BULLET_LIST
        if role == "bullet_item":
            return ExtractedComponentRole.BULLET_ITEM
        if role == "table":
            return ExtractedComponentRole.TABLE
        if role == "chart":
            return ExtractedComponentRole.CHART
        if role == "image":
            return ExtractedComponentRole.IMAGE
        return ExtractedComponentRole.UNKNOWN

    def _inventory_role_from_shape(self, shape) -> ExtractedComponentRole:
        if getattr(shape, "is_placeholder", False):
            placeholder_kind = self._map_placeholder_kind(shape.placeholder_format.type)
            return self._inventory_role_from_binding_or_kind(self._infer_placeholder_binding(placeholder_kind), placeholder_kind)
        if getattr(shape, "has_table", False):
            return ExtractedComponentRole.TABLE
        if getattr(shape, "has_chart", False):
            return ExtractedComponentRole.CHART
        if self._shape_looks_like_picture(shape):
            return ExtractedComponentRole.IMAGE
        text = self._shape_text_content(shape)
        if text:
            if self._is_footer_like_text_shape(shape, text):
                return ExtractedComponentRole.FOOTER
            if self._is_bullet_like_shape(shape):
                return ExtractedComponentRole.BULLET_LIST
            return ExtractedComponentRole.BODY
        return ExtractedComponentRole.DECORATIVE if self._shape_has_visual_style(shape) else ExtractedComponentRole.UNKNOWN

    def _inventory_confidence_for_shape(self, shape, role: ExtractedComponentRole) -> ComponentConfidence:
        if getattr(shape, "is_placeholder", False):
            return ComponentConfidence.HIGH
        if role in {ExtractedComponentRole.TABLE, ExtractedComponentRole.CHART, ExtractedComponentRole.IMAGE}:
            return ComponentConfidence.HIGH
        if role in {ExtractedComponentRole.FOOTER, ExtractedComponentRole.BULLET_LIST, ExtractedComponentRole.BODY}:
            return ComponentConfidence.MEDIUM
        return ComponentConfidence.LOW

    def _inventory_editability_for_role(self, role: ExtractedComponentRole) -> ComponentEditability:
        if role in {
            ExtractedComponentRole.TITLE,
            ExtractedComponentRole.SUBTITLE,
            ExtractedComponentRole.BODY,
            ExtractedComponentRole.BULLET_LIST,
            ExtractedComponentRole.BULLET_ITEM,
            ExtractedComponentRole.TABLE,
            ExtractedComponentRole.CHART,
            ExtractedComponentRole.IMAGE,
            ExtractedComponentRole.FOOTER,
        }:
            return ComponentEditability.EDITABLE
        if role == ExtractedComponentRole.AUXILIARY:
            return ComponentEditability.SEMI_EDITABLE
        if role == ExtractedComponentRole.DECORATIVE:
            return ComponentEditability.DECORATIVE
        return ComponentEditability.UNSAFE

    def _inventory_binding_from_shape(self, shape, role: ExtractedComponentRole) -> str | None:
        if getattr(shape, "is_placeholder", False):
            placeholder_kind = self._map_placeholder_kind(shape.placeholder_format.type)
            return self._infer_placeholder_binding(placeholder_kind)
        role_map = {
            ExtractedComponentRole.TITLE: "title",
            ExtractedComponentRole.SUBTITLE: "subtitle",
            ExtractedComponentRole.BODY: "body",
            ExtractedComponentRole.BULLET_LIST: "bullets",
            ExtractedComponentRole.TABLE: "table",
            ExtractedComponentRole.CHART: "chart",
            ExtractedComponentRole.IMAGE: "image",
            ExtractedComponentRole.FOOTER: "footer",
        }
        return role_map.get(role)

    def _inventory_capabilities_from_shape(self, shape, role: ExtractedComponentRole) -> list[str]:
        binding = self._inventory_binding_from_shape(shape, role)
        if binding:
            return self._editable_capabilities_for_binding(binding)
        if role == ExtractedComponentRole.DECORATIVE:
            return []
        return ["text"] if getattr(shape, "has_text_frame", False) else []

    def _inventory_warnings(
        self,
        layouts: list[LayoutSpec],
        prototype_slides: list[PrototypeSlideSpec],
        inventory_components: list[ExtractedComponent],
    ) -> list[str]:
        warnings: list[str] = []
        if not self._has_usable_layout_inventory(layouts):
            warnings.append("No usable layout inventory detected; runtime should rely on prototype or direct-shape paths.")
        if not prototype_slides:
            warnings.append("No prototype slide inventory detected.")
        if not any(component.role in {ExtractedComponentRole.TITLE, ExtractedComponentRole.BODY, ExtractedComponentRole.BULLET_LIST} for component in inventory_components):
            warnings.append("No reliable text-capable path detected in presentation inventory.")
        return warnings

    def _inventory_degradation_mode(
        self,
        layouts: list[LayoutSpec],
        prototype_slides: list[PrototypeSlideSpec],
        inventory_components: list[ExtractedComponent],
    ) -> str | None:
        if self._has_usable_layout_inventory(layouts):
            return None
        if prototype_slides:
            return "prototype_only"
        if inventory_components:
            return "direct_shape_binding"
        return "unusable"

    def _is_footer_like_text_shape(self, shape, text: str) -> bool:
        top = int(getattr(shape, "top", 0) or 0)
        height = int(getattr(shape, "height", 0) or 0)
        if top <= 0:
            return False
        if top > 5_000_000:
            return True
        lowered = text.lower()
        return any(marker in lowered for marker in {"@", "тел", "email", "confidential", "copyright"})

    def _shape_has_visual_style(self, shape) -> bool:
        return self._extract_fill_color(shape) is not None or self._extract_line_color(shape) is not None

    def _shape_text_content(self, shape) -> str:
        if not getattr(shape, "has_text_frame", False):
            return ""
        parts = [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text and paragraph.text.strip()]
        return "\n".join(parts).strip()

    def _is_bullet_like_shape(self, shape) -> bool:
        if not getattr(shape, "has_text_frame", False):
            return False
        paragraphs = [paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text and paragraph.text.strip()]
        if len(paragraphs) >= 2:
            return True
        text = paragraphs[0] if paragraphs else ""
        return text.startswith(("•", "-", "–", "—"))

    def _select_title_and_subtitle_shapes(self, text_shapes: list, slide_width: int, slide_height: int):
        if not text_shapes:
            return None, None

        title_shape = None
        for shape in text_shapes:
            text = self._shape_text_content(shape)
            if not text:
                continue
            if int(shape.top) <= max(int(slide_height * 0.28), 1_200_000) and len(text) <= 160:
                title_shape = shape
                break
        if title_shape is None:
            title_shape = text_shapes[0]

        subtitle_shape = None
        for shape in text_shapes:
            if shape == title_shape:
                continue
            text = self._shape_text_content(shape)
            if not text or self._is_bullet_like_shape(shape):
                continue
            if int(shape.top) <= max(int(slide_height * 0.45), 2_200_000) and int(shape.width) >= int(slide_width * 0.25):
                subtitle_shape = shape
                break
        return title_shape, subtitle_shape

    def _detect_two_column_text_shapes(self, text_shapes: list, slide_width: int) -> dict[str, list] | None:
        if len(text_shapes) < 2:
            return None
        left_candidates = sorted(text_shapes, key=lambda shape: int(shape.left))
        distinct_columns: list[int] = []
        for shape in left_candidates:
            current = int(shape.left)
            if not distinct_columns or abs(current - distinct_columns[-1]) > max(int(slide_width * 0.12), 900000):
                distinct_columns.append(current)
        if len(distinct_columns) != 2:
            return None

        threshold = (distinct_columns[0] + distinct_columns[1]) / 2
        left = [shape for shape in text_shapes if int(shape.left) < threshold]
        right = [shape for shape in text_shapes if int(shape.left) >= threshold]
        if not left or not right:
            return None
        return {"left": left, "right": right}

    def _shape_looks_like_picture(self, shape) -> bool:
        return "PICTURE" in str(getattr(shape, "shape_type", ""))

    def _extract_text_style_from_shape(self, shape) -> TemplateTextStyleSpec | None:
        if not getattr(shape, "has_text_frame", False):
            return None
        text_frame = shape.text_frame
        paragraph = next(
            (item for item in text_frame.paragraphs if (item.text or "").strip()),
            text_frame.paragraphs[0] if text_frame.paragraphs else None,
        )
        if paragraph is None:
            return None

        run = next((item for item in paragraph.runs if (item.text or "").strip()), None)
        font = run.font if run is not None else paragraph.font
        color = self._extract_font_color(font)
        font_size = font.size.pt if getattr(font, "size", None) is not None else None
        style = TemplateTextStyleSpec(
            font_family=getattr(font, "name", None),
            font_size_pt=float(font_size) if font_size is not None else None,
            bold=getattr(font, "bold", None),
            italic=getattr(font, "italic", None),
            underline=getattr(font, "underline", None),
            color=color,
            alignment=str(paragraph.alignment).split(".")[-1].lower() if getattr(paragraph, "alignment", None) is not None else None,
            vertical_anchor=str(text_frame.vertical_anchor).split(".")[-1].lower() if getattr(text_frame, "vertical_anchor", None) is not None else None,
            word_wrap=getattr(text_frame, "word_wrap", None),
            margin_left_emu=self._text_frame_margin(shape, "left"),
            margin_right_emu=self._text_frame_margin(shape, "right"),
        )
        return style if any(value is not None for value in style.model_dump().values()) else None

    def _extract_paragraph_styles_from_shape(self, shape) -> TemplateParagraphStyleCatalog | None:
        if not getattr(shape, "has_text_frame", False):
            return None
        level_styles: dict[str, TemplateTextStyleSpec] = {}
        for paragraph in shape.text_frame.paragraphs:
            if not (paragraph.text or "").strip():
                continue
            level_key = str(getattr(paragraph, "level", 0) or 0)
            if level_key in level_styles:
                continue
            style = TemplateTextStyleSpec(
                level=getattr(paragraph, "level", None),
                alignment=str(paragraph.alignment).split(".")[-1].lower() if getattr(paragraph, "alignment", None) is not None else None,
                line_spacing=float(paragraph.line_spacing) if isinstance(getattr(paragraph, "line_spacing", None), (int, float)) else None,
                space_before_pt=float(paragraph.space_before.pt) if getattr(paragraph, "space_before", None) is not None else None,
                space_after_pt=float(paragraph.space_after.pt) if getattr(paragraph, "space_after", None) is not None else None,
            )
            if any(value is not None for value in style.model_dump().values()):
                level_styles[level_key] = style
        return TemplateParagraphStyleCatalog(level_styles=level_styles) if level_styles else None

    def _extract_shape_style_from_shape(self, shape) -> TemplateShapeStyleSpec | None:
        fill_color = self._extract_fill_color(shape)
        line_color = self._extract_line_color(shape)
        style = TemplateShapeStyleSpec(
            fill_color=fill_color,
            line_color=line_color,
        )
        return style if any(value is not None and value != [] for value in style.model_dump().values()) else None

    def _extract_font_color(self, font) -> str | None:
        color_format = getattr(font, "color", None)
        if color_format is None:
            return None
        rgb = getattr(color_format, "rgb", None)
        return self._rgb_to_hex(rgb)

    def _extract_fill_color(self, shape) -> str | None:
        fill = getattr(shape, "fill", None)
        if fill is None:
            return None
        try:
            fore_color = getattr(fill, "fore_color", None)
        except Exception:
            return None
        if fore_color is None:
            return None
        return self._rgb_to_hex(getattr(fore_color, "rgb", None))

    def _extract_line_color(self, shape) -> str | None:
        line = getattr(shape, "line", None)
        if line is None:
            return None
        try:
            color = getattr(line, "color", None)
        except Exception:
            color = None
        if color is None:
            try:
                color = getattr(line, "fore_color", None)
            except Exception:
                color = None
        if color is None:
            return None
        return self._rgb_to_hex(getattr(color, "rgb", None))

    def _rgb_to_hex(self, rgb) -> str | None:
        if rgb is None:
            return None
        value = str(rgb).replace("0x", "").replace("#", "").strip()
        if len(value) != 6:
            return None
        return f"#{value.upper()}"

    def _backfill_geometry(self, manifest: TemplateManifest, analyzed: TemplateManifest) -> None:
        analyzed_layouts = {(layout.name, layout.slide_layout_index): layout for layout in analyzed.layouts}
        for layout in manifest.layouts:
            source_layout = analyzed_layouts.get((layout.name, layout.slide_layout_index))
            if source_layout is None:
                continue
            if not layout.representation_hints:
                layout.representation_hints = list(source_layout.representation_hints)
            for placeholder in layout.placeholders:
                source_placeholder = next(
                    (
                        item
                        for item in source_layout.placeholders
                        if item.shape_name == placeholder.shape_name
                        or (item.idx is not None and item.idx == placeholder.idx)
                    ),
                    None,
                )
                if source_placeholder is None:
                    continue
                if placeholder.binding is None:
                    placeholder.binding = source_placeholder.binding
                if placeholder.editable_role is None:
                    placeholder.editable_role = source_placeholder.editable_role
                if not placeholder.editable_capabilities:
                    placeholder.editable_capabilities = list(source_placeholder.editable_capabilities)
                if placeholder.slot_group is None:
                    placeholder.slot_group = source_placeholder.slot_group
                if placeholder.slot_group_order is None:
                    placeholder.slot_group_order = source_placeholder.slot_group_order
                for field in (
                    "left_emu",
                    "top_emu",
                    "width_emu",
                    "height_emu",
                    "margin_left_emu",
                    "margin_right_emu",
                    "margin_top_emu",
                    "margin_bottom_emu",
                ):
                    if getattr(placeholder, field, None) is None:
                        setattr(placeholder, field, getattr(source_placeholder, field, None))

        analyzed_prototypes = {slide.name: slide for slide in analyzed.prototype_slides}
        for prototype in manifest.prototype_slides:
            source_slide = analyzed_prototypes.get(prototype.name)
            if source_slide is not None:
                if not prototype.representation_hints:
                    prototype.representation_hints = list(source_slide.representation_hints)
                for index, token in enumerate(prototype.tokens):
                    source_token = next(
                        (item for item in source_slide.tokens if item.shape_name == token.shape_name or item.token == token.token),
                        source_slide.tokens[index] if index < len(source_slide.tokens) else None,
                    )
                    if source_token is None:
                        continue
                    if token.editable_role is None:
                        token.editable_role = source_token.editable_role
                    if not token.editable_capabilities:
                        token.editable_capabilities = list(source_token.editable_capabilities)
                    if token.slot_group is None:
                        token.slot_group = source_token.slot_group
                    if token.slot_group_order is None:
                        token.slot_group_order = source_token.slot_group_order
                    for field in (
                        "left_emu",
                        "top_emu",
                        "width_emu",
                        "height_emu",
                        "margin_left_emu",
                        "margin_right_emu",
                        "margin_top_emu",
                        "margin_bottom_emu",
                    ):
                        if getattr(token, field, None) is None:
                            setattr(token, field, getattr(source_token, field, None))
            text_token = next(
                (
                    item
                    for item in prototype.tokens
                    if item.binding in {"title", "subtitle", "text", "body", "main_text", "secondary_text", "cover_title"}
                ),
                prototype.tokens[0] if prototype.tokens else None,
            )
            if text_token is None:
                continue
            if not any(
                all(isinstance(v, int) and v > 0 for v in [token.left_emu, token.top_emu, token.width_emu, token.height_emu])
                for token in prototype.tokens
            ):
                token = text_token
                token.left_emu = token.left_emu or 1
                token.top_emu = token.top_emu or 1
                token.width_emu = token.width_emu or 1
                token.height_emu = token.height_emu or 1
            if not any(
                item.binding in {"title", "subtitle", "text", "body", "main_text", "secondary_text", "cover_title"}
                and any(getattr(item, field, None) is not None for field in (
                    "margin_left_emu",
                    "margin_right_emu",
                    "margin_top_emu",
                    "margin_bottom_emu",
                ))
                for item in prototype.tokens
            ):
                token = text_token
                token.margin_left_emu = token.margin_left_emu if token.margin_left_emu is not None else 0
                token.margin_right_emu = token.margin_right_emu if token.margin_right_emu is not None else 0
                token.margin_top_emu = token.margin_top_emu if token.margin_top_emu is not None else 0
                token.margin_bottom_emu = token.margin_bottom_emu if token.margin_bottom_emu is not None else 0

    def _infer_component_styles(self, layouts: list[LayoutSpec]) -> dict[str, TemplateComponentStyleSpec]:
        card_layout = next((layout for layout in layouts if "карточ" in layout.name.lower()), None)
        table_layout = next((layout for layout in layouts if "table" in layout.supported_slide_kinds or "табл" in layout.name.lower()), None)
        image_layout = next((layout for layout in layouts if "image" in layout.supported_slide_kinds), None)
        cover_layout = next((layout for layout in layouts if "тит" in layout.name.lower() or "title" in layout.name.lower()), None)
        list_icons_layout = next((layout for layout in layouts if any(item.idx == 21 for item in layout.placeholders)), None)
        contacts_layout = next((layout for layout in layouts if any(item.idx == 10 for item in layout.placeholders)), None)
        card_margin_x = 91440
        card_margin_y = 45720
        if card_layout is not None:
            body_placeholder = next((item for item in card_layout.placeholders if item.idx in {11, 12, 13}), None)
            if body_placeholder is not None:
                card_margin_x = body_placeholder.margin_left_emu or card_margin_x
                card_margin_y = body_placeholder.margin_top_emu or card_margin_y
        table_margin_left = 80000
        table_margin_right = 80000
        table_margin_top = 40000
        table_margin_bottom = 40000
        if table_layout is not None:
            table_placeholder = next((item for item in table_layout.placeholders if item.idx == 14), None)
            if table_placeholder is not None and table_placeholder.shape_style is not None:
                table_margin_left = table_placeholder.shape_style.table_cell_margin_left_emu or table_margin_left
                table_margin_right = table_placeholder.shape_style.table_cell_margin_right_emu or table_margin_right
                table_margin_top = table_placeholder.shape_style.table_cell_margin_top_emu or table_margin_top
                table_margin_bottom = table_placeholder.shape_style.table_cell_margin_bottom_emu or table_margin_bottom
        image_margin_x = 91440
        image_margin_y = 45720
        if image_layout is not None:
            body_placeholder = next((item for item in image_layout.placeholders if item.idx == 14), None)
            if body_placeholder is not None:
                image_margin_x = body_placeholder.margin_left_emu or image_margin_x
                image_margin_y = body_placeholder.margin_top_emu or image_margin_y
        cover_title_top = 651176
        cover_title_left = 444249
        cover_title_width = 10693901
        cover_title_min_height = 1422646
        cover_meta_left = 444249
        cover_meta_top = 2438400
        cover_meta_width = 8200000
        if cover_layout is not None:
            title_placeholder = next((item for item in cover_layout.placeholders if item.idx == 0), None)
            if title_placeholder is not None:
                cover_title_top = title_placeholder.top_emu or cover_title_top
                cover_title_left = title_placeholder.left_emu or cover_title_left
                cover_title_width = title_placeholder.width_emu or cover_title_width
                cover_title_min_height = title_placeholder.height_emu or cover_title_min_height
            meta_placeholder = next((item for item in cover_layout.placeholders if item.idx in {13, 15}), None)
            if meta_placeholder is not None:
                cover_meta_left = meta_placeholder.left_emu or cover_meta_left
                cover_meta_top = meta_placeholder.top_emu or cover_meta_top
                cover_meta_width = meta_placeholder.width_emu or cover_meta_width
        list_title_gap = 180000
        list_no_subtitle_gap = 300000
        list_footer_gap = 180000
        if list_icons_layout is not None:
            title_placeholder = next((item for item in list_icons_layout.placeholders if item.idx == 0), None)
            left_placeholder = next((item for item in list_icons_layout.placeholders if item.idx == 12), None)
            footer_placeholder = next((item for item in list_icons_layout.placeholders if item.idx == 21), None)
            if title_placeholder is not None and left_placeholder is not None:
                list_title_gap = max(left_placeholder.top_emu - (title_placeholder.top_emu + title_placeholder.height_emu), 120000)
            if footer_placeholder is not None and left_placeholder is not None:
                list_footer_gap = max(footer_placeholder.top_emu - (left_placeholder.top_emu + left_placeholder.height_emu), 120000)
        return {
            "cards": TemplateComponentStyleSpec(
                text_styles={
                    "title": TemplateTextStyleSpec(font_size_pt=20.0, bold=True, color="#FFFFFF"),
                    "body": TemplateTextStyleSpec(font_size_pt=16.0, color="#FFFFFF"),
                    "kpi_value": TemplateTextStyleSpec(font_size_pt=22.0, bold=True, color="#FFFFFF"),
                    "kpi_label": TemplateTextStyleSpec(font_size_pt=12.0, color="#E4F1FF"),
                },
                spacing_tokens={
                    "content_margin_x_emu": card_margin_x,
                    "content_margin_y_emu": card_margin_y,
                    "title_body_gap_emu": 100000,
                    "body_metrics_gap_emu": 180000,
                    "metrics_gap_x_emu": 180000,
                    "metrics_gap_y_emu": 160000,
                },
                behavior_tokens={
                    "kpi_max_metrics": 4,
                    "kpi_value_compact_font_pt": 20,
                    "kpi_value_regular_font_pt": 22,
                },
            ),
            "table": TemplateComponentStyleSpec(
                text_styles={
                    "header": TemplateTextStyleSpec(font_size_pt=16.0, bold=True, color="#091E38"),
                    "body": TemplateTextStyleSpec(font_size_pt=16.0, color="#081C4F"),
                },
                spacing_tokens={
                    "cell_margin_left_emu": table_margin_left,
                    "cell_margin_right_emu": table_margin_right,
                    "cell_margin_top_emu": table_margin_top,
                    "cell_margin_bottom_emu": table_margin_bottom,
                },
                behavior_tokens={
                    "header_fill_color": "#C6DFFF",
                    "header_text_color": "#091E38",
                    "border_color": "#C6DFFF",
                    "body_fill_transparent": True,
                    "preserve_source_fill_colors": False,
                    "render_as_shapes": True,
                },
            ),
            "chart": TemplateComponentStyleSpec(
                text_styles={
                    "title": TemplateTextStyleSpec(font_size_pt=35.0, bold=True, color="#081C4F"),
                    "subtitle": TemplateTextStyleSpec(font_size_pt=20.0, color="#081C4F"),
                },
                shape_style=TemplateShapeStyleSpec(
                    role="chart",
                    chart_plot_left_factor=0.0,
                    chart_plot_top_factor=0.0,
                    chart_plot_width_factor=1.0,
                    chart_plot_height_factor=1.0,
                ),
                behavior_tokens={
                    "rank_color_1": "#091E38",
                    "rank_color_2": "#3489F3",
                    "rank_color_3": "#264595",
                    "rank_color_4": "#BFCEF5",
                },
            ),
            "image": TemplateComponentStyleSpec(
                text_styles={
                    "title": TemplateTextStyleSpec(font_size_pt=35.0, bold=True, color="#081C4F"),
                    "subtitle": TemplateTextStyleSpec(font_size_pt=18.0, color="#081C4F"),
                    "body": TemplateTextStyleSpec(font_size_pt=20.0, color="#081C4F"),
                },
                spacing_tokens={
                    "content_margin_x_emu": image_margin_x,
                    "content_margin_y_emu": image_margin_y,
                    "title_content_gap_emu": 180000,
                    "title_body_gap_no_subtitle_emu": 300000,
                    "content_footer_gap_emu": 180000,
                    "min_image_height_emu": 1200000,
                    "secondary_min_height_emu": 700000,
                    "secondary_reserved_image_gap_emu": 900000,
                },
            ),
            "cover": TemplateComponentStyleSpec(
                text_styles={
                    "title": TemplateTextStyleSpec(font_size_pt=46.0, bold=True, color="#F5F9FE"),
                    "meta": TemplateTextStyleSpec(font_size_pt=22.0, color="#F5F9FE"),
                    "footer": TemplateTextStyleSpec(font_size_pt=14.0, color="#F5F9FE"),
                },
                spacing_tokens={
                    "title_top_emu": cover_title_top,
                    "title_left_emu": cover_title_left,
                    "title_width_emu": cover_title_width,
                    "title_min_height_emu": cover_title_min_height,
                    "meta_left_emu": cover_meta_left,
                    "meta_top_emu": cover_meta_top,
                    "meta_width_emu": cover_meta_width,
                    "meta_min_height_emu": 700000,
                    "meta_gap_emu": 220000,
                    "bottom_limit_emu": 6200000,
                },
            ),
            "list_with_icons": TemplateComponentStyleSpec(
                text_styles={
                    "subtitle": TemplateTextStyleSpec(font_size_pt=18.0, color="#081C4F"),
                },
                spacing_tokens={
                    "title_content_gap_emu": list_title_gap,
                    "title_body_gap_no_subtitle_emu": list_no_subtitle_gap,
                    "content_footer_gap_emu": list_footer_gap,
                },
            ),
            "contacts": TemplateComponentStyleSpec(
                text_styles={
                    "primary": TemplateTextStyleSpec(font_size_pt=18.0, color="#081C4F"),
                    "secondary": TemplateTextStyleSpec(font_size_pt=14.0, color="#081C4F"),
                },
                behavior_tokens={
                    "primary_threshold_chars": 60,
                    "secondary_threshold_chars": 40,
                    "font_decrement_pt": 2.0,
                },
            ),
        }

    def _map_placeholder_kind(self, placeholder_type) -> PlaceholderKind:
        name = getattr(placeholder_type, "name", str(placeholder_type)).lower()
        if "title" in name and "ctr" not in name:
            return PlaceholderKind.TITLE
        if "sub" in name:
            return PlaceholderKind.SUBTITLE
        if "pic" in name or "media" in name or "obj" in name:
            return PlaceholderKind.IMAGE
        if "tbl" in name:
            return PlaceholderKind.TABLE
        if "chart" in name:
            return PlaceholderKind.CHART
        if "footer" in name or "dt" in name or "sld_num" in name:
            return PlaceholderKind.FOOTER
        if "body" in name or "content" in name or "text" in name:
            return PlaceholderKind.BODY
        return PlaceholderKind.UNKNOWN

    def _infer_slide_kinds(self, placeholders: list[PlaceholderSpec], layout_name: str = "") -> list[str]:
        kinds = {item.kind for item in placeholders}
        supported: list[str] = []
        layout_name_lower = layout_name.lower()

        if PlaceholderKind.TITLE in kinds and PlaceholderKind.SUBTITLE in kinds and len(placeholders) <= 3:
            supported.append("title")
        elif PlaceholderKind.TITLE in kinds and ("тит" in layout_name_lower or "title" in layout_name_lower):
            supported.append("title")
        if PlaceholderKind.BODY in kinds:
            supported.extend(["bullets", "text", "table"])
        if PlaceholderKind.IMAGE in kinds:
            supported.append("image")
        return list(dict.fromkeys(supported))

    def _infer_slide_kinds_from_tokens(self, tokens: list[PrototypeTokenSpec]) -> list[str]:
        token_names = {token.token.lower() for token in tokens}
        supported: list[str] = []

        if "title" in token_names and ("subtitle" in token_names or len(token_names) <= 2):
            supported.append("title")
        if "bullets" in token_names or any(name.startswith("bullet_") for name in token_names):
            supported.append("bullets")
        if "text" in token_names or "body" in token_names or "summary" in token_names:
            supported.append("text")
        if "left_bullets" in token_names or "right_bullets" in token_names:
            supported.append("two_column")
        if "image" in token_names or any(name.startswith("image_") for name in token_names):
            supported.append("image")
        if not supported:
            supported.append("text")
        return list(dict.fromkeys(supported))

    def _infer_representation_hints_from_placeholders(self, placeholders: list[PlaceholderSpec]) -> list[str]:
        hints: list[str] = []
        if self._looks_like_contacts_from_slots(placeholders):
            hints.append("contacts")
        if self._looks_like_cards_from_slots(placeholders):
            hints.append("cards")
        if self._looks_like_two_column_text_from_slots(placeholders):
            hints.append("two_column")
        if any(self._slot_has_capability(slot, "table") for slot in placeholders):
            hints.append("table")
        if any(self._slot_has_capability(slot, "chart") for slot in placeholders):
            hints.append("chart")
        if any(self._slot_has_capability(slot, "image") for slot in placeholders):
            hints.append("image")
        return list(dict.fromkeys(hints))

    def _infer_representation_hints_from_tokens(self, tokens: list[PrototypeTokenSpec]) -> list[str]:
        hints: list[str] = []
        if self._looks_like_contacts_from_slots(tokens):
            hints.append("contacts")
        if self._looks_like_cards_from_slots(tokens):
            hints.append("cards")
        if self._looks_like_two_column_text_from_slots(tokens):
            hints.append("two_column")
        if any(self._slot_has_capability(slot, "table") for slot in tokens):
            hints.append("table")
        if any(self._slot_has_capability(slot, "chart") for slot in tokens):
            hints.append("chart")
        if any(self._slot_has_capability(slot, "image") for slot in tokens):
            hints.append("image")
        return list(dict.fromkeys(hints))

    def _infer_binding(self, token: str) -> str:
        token_name = token.lower()
        if token_name in {"title", "subtitle", "text", "body", "summary", "notes"}:
            return token_name
        if token_name in {"bullets", "left_bullets", "right_bullets"}:
            return token_name
        if token_name.startswith("bullet_"):
            return token_name
        if token_name.startswith("left_bullet_") or token_name.startswith("right_bullet_"):
            return token_name
        return "text"

    def _looks_like_cards_from_slots(self, slots: list[PlaceholderSpec | PrototypeTokenSpec]) -> bool:
        text_slots = [
            slot
            for slot in slots
            if self._is_card_text_slot(slot)
            and isinstance(slot.left_emu, int)
            and isinstance(slot.top_emu, int)
            and isinstance(slot.width_emu, int)
            and isinstance(slot.height_emu, int)
            and slot.width_emu > 0
            and slot.height_emu > 0
        ]
        if len(text_slots) < 2:
            return False

        for base_slot in text_slots:
            same_row = [
                slot
                for slot in text_slots
                if abs((slot.top_emu or 0) - (base_slot.top_emu or 0))
                <= max(slot.height_emu or 0, base_slot.height_emu or 0) * 0.45
            ]
            if len(same_row) < 2 or len(same_row) > 4:
                continue

            widths = [slot.width_emu or 0 for slot in same_row]
            heights = [slot.height_emu or 0 for slot in same_row]
            lefts = sorted(slot.left_emu or 0 for slot in same_row)
            if len({round(value / 10000) for value in lefts}) < len(same_row):
                continue

            width_spread = max(widths) / max(min(widths), 1)
            height_spread = max(heights) / max(min(heights), 1)
            if width_spread <= 1.8 and height_spread <= 1.8:
                return True
        return False

    def _is_card_text_slot(self, slot: PlaceholderSpec | PrototypeTokenSpec) -> bool:
        role = getattr(slot, "editable_role", None)
        capabilities = getattr(slot, "editable_capabilities", [])
        if role in {"body", "bullet_item", "bullet_list"}:
            return True
        if role in {"title", "subtitle", "image", "table", "chart"}:
            return False
        return "text" in capabilities or "list_item" in capabilities

    def _looks_like_two_column_text_from_slots(self, slots: list[PlaceholderSpec | PrototypeTokenSpec]) -> bool:
        text_slots = [
            slot
            for slot in slots
            if self._is_card_text_slot(slot)
            and isinstance(slot.left_emu, int)
            and isinstance(slot.width_emu, int)
            and slot.width_emu > 0
        ]
        if len(text_slots) < 2:
            return False

        lefts = sorted(slot.left_emu or 0 for slot in text_slots)
        widths = [slot.width_emu or 0 for slot in text_slots]
        distinct_columns = sorted({round(value / 100000) for value in lefts})
        if len(distinct_columns) != 2:
            return False
        return max(widths) / max(min(widths), 1) <= 2.2

    def _looks_like_contacts_from_slots(self, slots: list[PlaceholderSpec | PrototypeTokenSpec]) -> bool:
        contact_like = [
            slot
            for slot in slots
            if getattr(slot, "binding", "") in {"contact_name_or_title", "contact_role", "contact_phone", "contact_email"}
        ]
        return len(contact_like) >= 3

    def _slot_has_capability(self, slot: PlaceholderSpec | PrototypeTokenSpec, capability: str) -> bool:
        capabilities = getattr(slot, "editable_capabilities", [])
        role = getattr(slot, "editable_role", None)
        if capability in capabilities:
            return True
        return role == capability

    def _infer_placeholder_binding(self, kind: PlaceholderKind) -> str | None:
        if kind == PlaceholderKind.TITLE:
            return "title"
        if kind == PlaceholderKind.SUBTITLE:
            return "subtitle"
        if kind == PlaceholderKind.BODY:
            return "body"
        if kind == PlaceholderKind.IMAGE:
            return "image"
        if kind == PlaceholderKind.TABLE:
            return "table"
        if kind == PlaceholderKind.CHART:
            return "chart"
        if kind == PlaceholderKind.FOOTER:
            return "footer"
        return None

    def _editable_role_for_placeholder_kind(self, kind: PlaceholderKind) -> str | None:
        binding = self._infer_placeholder_binding(kind)
        return self._editable_role_for_binding(binding)

    def _editable_role_for_binding(self, binding: str | None) -> str | None:
        if not binding:
            return None
        binding_name = binding.lower()
        if binding_name in {"title", "cover_title"}:
            return "title"
        if binding_name in {"subtitle", "cover_meta", "presentation_name"}:
            return "subtitle"
        if binding_name in {"text", "body", "summary", "main_text", "secondary_text", "notes", "footer"}:
            return "body"
        if binding_name in {"bullets", "left_bullets", "right_bullets"}:
            return "bullet_list"
        if binding_name.startswith("bullet_") or binding_name.startswith("left_bullet_") or binding_name.startswith("right_bullet_"):
            return "bullet_item"
        if binding_name in {"image", "chart_image"} or binding_name.startswith("image_"):
            return "image"
        if binding_name == "table":
            return "table"
        if binding_name == "chart":
            return "chart"
        return "body"

    def _editable_capabilities_for_binding(self, binding: str | None, fallback_role: str | None = None) -> list[str]:
        role = self._editable_role_for_binding(binding) or fallback_role
        if role in {"title", "subtitle", "body"}:
            return ["text"]
        if role == "bullet_list":
            return ["bullet_list", "text"]
        if role == "bullet_item":
            return ["text", "list_item"]
        if role == "image":
            return ["image"]
        if role == "table":
            return ["table"]
        if role == "chart":
            return ["chart"]
        return ["text"] if role else []

    def _infer_slot_group(self, value: str | None) -> str | None:
        if not value:
            return None
        normalized = value.strip().lower()
        if not normalized:
            return None
        if normalized.startswith("{{") and normalized.endswith("}}"):
            normalized = normalized[2:-2].strip().lower()
        normalized = normalized.replace("-", "_").replace(" ", "_")
        while "__" in normalized:
            normalized = normalized.replace("__", "_")

        if normalized.startswith("left_"):
            return "left_column"
        if normalized.startswith("right_"):
            return "right_column"

        if "_" in normalized:
            prefix, suffix = normalized.rsplit("_", 1)
            if suffix.isdigit():
                return prefix
        return None

    def _infer_slot_group_order(self, value: str | None) -> int | None:
        if not value:
            return None
        normalized = value.strip().lower()
        if normalized.startswith("{{") and normalized.endswith("}}"):
            normalized = normalized[2:-2].strip().lower()
        normalized = normalized.replace("-", "_").replace(" ", "_")
        if "_" not in normalized:
            return None
        _, suffix = normalized.rsplit("_", 1)
        if suffix.isdigit():
            return int(suffix)
        return None

    def _slugify(self, value: str, fallback: str) -> str:
        cleaned = "".join(char.lower() if char.isalnum() else "_" for char in value).strip("_")
        while "__" in cleaned:
            cleaned = cleaned.replace("__", "_")
        return cleaned or fallback

    def _make_unique_layout_key(self, base_key: str, master_index: int, layout_index: int, used_keys: set[str]) -> str:
        candidate = base_key
        if candidate in used_keys:
            candidate = f"{base_key}_m{master_index}_l{layout_index}"
        used_keys.add(candidate)
        return candidate

    def _text_frame_margin(self, shape, side: str) -> int | None:
        if not getattr(shape, "has_text_frame", False):
            return None
        text_frame = shape.text_frame
        value = getattr(text_frame, f"margin_{side}", None)
        return int(value) if value is not None else None
