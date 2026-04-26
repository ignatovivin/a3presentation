from __future__ import annotations

import re
from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.chart.axis import ValueAxis
from a3presentation.domain.chart import ChartSpec, ChartType
from a3presentation.domain.presentation import PresentationPlan, SlideKind, SlideSpec
from a3presentation.domain.template import PlaceholderKind, TemplateManifest
from a3presentation.services.layout_capacity import (
    LayoutCapacityProfile,
    LayoutGeometryPolicy,
    PlaceholderGeometryPolicy,
    derive_capacity_profile_for_geometry,
    geometry_policy_for_layout,
    runtime_profile_key_for_target,
)
from a3presentation.services.pptx_generator import PptxGenerator
from a3presentation.services.chart_render_contract import (
    PRIMARY_AXIS,
    SECONDARY_AXIS,
    chart_axis_number_format,
    chart_axis_number_format_for_axis,
    render_chart_series_count,
    render_chart_spec,
    render_chart_type,
    uses_secondary_value_axis,
)

CONTINUATION_FONT_DELTA_TOLERANCE_PT = 2.0
CONTINUATION_UNDERFILL_GRACE = 0.09
CONTINUATION_AUDIT_EPSILON = 0.01
GEOMETRY_TOLERANCE_EMU = 90000
BODY_HEIGHT_UNDERFILL_RATIO = 0.45
EXPECTED_CHART_TITLE_FONT_PT = 35.0
EXPECTED_CHART_SUBTITLE_FONT_PT = 20.0
CHART_TITLE_FONT_TOLERANCE_PT = 0.1
MAX_REASONABLE_STACK_GAP_EMU = 450000
RUNTIME_EXPANDED_LAYOUT_KEYS = {
    "text_full_width",
    "dense_text_full_width",
    "list_full_width",
    "table",
    "image_text",
    "cards_3",
    "cards_kpi",
    "list_with_icons",
    "contacts",
}


@dataclass(frozen=True)
class SlideAudit:
    slide_index: int
    title: str
    kind: str
    layout_key: str
    body_char_count: int
    body_font_sizes: tuple[float, ...]
    profile: LayoutCapacityProfile
    title_font_sizes: tuple[float, ...] = ()
    subtitle_font_sizes: tuple[float, ...] = ()
    content_width: int | None = None
    footer_width: int | None = None
    has_table: bool = False
    has_chart: bool = False
    has_image: bool = False
    expected_items: tuple[str, ...] = ()
    rendered_items: tuple[str, ...] = ()
    title_top: int | None = None
    title_height: int | None = None
    title_width: int | None = None
    subtitle_top: int | None = None
    subtitle_height: int | None = None
    subtitle_width: int | None = None
    body_top: int | None = None
    body_height: int | None = None
    body_left: int | None = None
    footer_top: int | None = None
    footer_left: int | None = None
    expected_footer_width: int | None = None
    auxiliary_widths: dict[int, int] | None = None
    auxiliary_lefts: dict[int, int] | None = None
    auxiliary_tops: dict[int, int] | None = None
    auxiliary_char_counts: dict[int, int] | None = None
    expected_auxiliary_char_counts: dict[int, int] | None = None
    placeholder_char_counts: dict[int, int] | None = None
    expected_placeholder_char_counts: dict[int, int] | None = None
    expected_subtitle_char_count: int = 0
    image_width: int | None = None
    image_left: int | None = None
    image_top: int | None = None
    body_margin_left: int | None = None
    body_margin_right: int | None = None
    body_margin_top: int | None = None
    body_margin_bottom: int | None = None
    expected_body_margin_left: int | None = None
    expected_body_margin_right: int | None = None
    expected_body_margin_top: int | None = None
    expected_body_margin_bottom: int | None = None
    expected_body_max_font_pt: float | None = None
    title_placeholder_idx: int | None = None
    subtitle_placeholder_idx: int | None = None
    body_placeholder_idx: int | None = None
    footer_placeholder_idx: int | None = None
    expected_chart_type: str | None = None
    rendered_chart_type: str | None = None
    expected_chart_series_count: int | None = None
    rendered_chart_series_count: int | None = None
    rendered_chart_bar_series_count: int | None = None
    rendered_chart_line_series_count: int | None = None
    expected_chart_value_axis_number_format: str | None = None
    rendered_chart_value_axis_number_format: str | None = None
    expected_chart_secondary_value_axis: bool = False
    rendered_chart_secondary_value_axis: bool = False
    expected_chart_secondary_value_axis_number_format: str | None = None
    rendered_chart_secondary_value_axis_number_format: str | None = None
    expected_chart_content_width: int | None = None
    expected_title_font_pt: float | None = None
    expected_subtitle_font_pt: float | None = None
    table_overlay_overflow_count: int = 0
    table_overlay_overflow_details: tuple[str, ...] = ()
    runtime_profile_key: str = ""
    geometry: LayoutGeometryPolicy | None = None
    target_type: str | None = None
    target_source: str | None = None
    target_degradation_reasons: tuple[str, ...] = ()

    @property
    def fill_ratio(self) -> float:
        if self.profile.max_chars <= 0:
            return 0.0
        return self.body_char_count / self.profile.max_chars

    @property
    def min_font_size(self) -> float | None:
        if not self.body_font_sizes:
            return None
        return min(self.body_font_sizes)

    @property
    def max_font_size(self) -> float | None:
        if not self.body_font_sizes:
            return None
        return max(self.body_font_sizes)

    @property
    def representative_font_size(self) -> float | None:
        if not self.body_font_sizes:
            return None
        return min(self.body_font_sizes)

    @property
    def within_font_bounds(self) -> bool:
        if not self.body_font_sizes:
            return True
        return (
            min(self.body_font_sizes) >= self.profile.min_font_pt
            and max(self.body_font_sizes) <= self.profile.max_font_pt
        )

    @property
    def content_width_ratio(self) -> float:
        if not self.content_width:
            return 0.0
        return self.content_width / PptxGenerator.FULL_CONTENT_WIDTH_EMU

    @property
    def footer_width_ratio(self) -> float:
        if not self.footer_width:
            return 0.0
        reference_width = self.expected_footer_width or PptxGenerator.FULL_CONTENT_WIDTH_EMU
        if reference_width <= 0:
            return 0.0
        return self.footer_width / reference_width

    @property
    def chart_content_width_ratio(self) -> float:
        if not self.content_width:
            return 0.0
        reference_width = self.expected_chart_content_width or PptxGenerator.FULL_CONTENT_WIDTH_EMU
        if reference_width <= 0:
            return 0.0
        return self.content_width / reference_width

    @property
    def title_bottom(self) -> int | None:
        if self.title_top is None or self.title_height is None:
            return None
        return self.title_top + self.title_height

    @property
    def subtitle_bottom(self) -> int | None:
        if self.subtitle_top is None or self.subtitle_height is None:
            return None
        return self.subtitle_top + self.subtitle_height

    @property
    def degraded_but_valid(self) -> bool:
        if self.target_degradation_reasons:
            return True
        return self.target_type in {"auto_layout", "prototype", "direct_shape_binding"}


@dataclass(frozen=True)
class CapacityViolation:
    slide_index: int
    title: str
    rule: str
    details: str


def audit_generated_presentation(
    output_path: Path,
    plan: PresentationPlan,
    manifest: TemplateManifest | None = None,
) -> list[SlideAudit]:
    presentation = Presentation(str(output_path))
    audits: list[SlideAudit] = []

    for slide_index, slide_spec in enumerate(plan.slides, start=1):
        if slide_index > len(presentation.slides):
            break
        if slide_spec.kind not in {
            SlideKind.TEXT,
            SlideKind.BULLETS,
            SlideKind.TWO_COLUMN,
            SlideKind.TABLE,
            SlideKind.CHART,
            SlideKind.IMAGE,
        }:
            continue

        slide = presentation.slides[slide_index - 1]
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        title_idx = _placeholder_idx_for_role(slide_spec, manifest, PlaceholderKind.TITLE, (0,))
        subtitle_idx = _placeholder_idx_for_role(slide_spec, manifest, PlaceholderKind.SUBTITLE, (13,))
        body_idx = _placeholder_idx_for_role(slide_spec, manifest, PlaceholderKind.BODY, (14,))
        footer_idx = _placeholder_idx_for_role(slide_spec, manifest, PlaceholderKind.FOOTER, (15, 17, 21))
        title_geometry_idx = _geometry_idx_for_role(slide_spec, manifest, PlaceholderKind.TITLE, title_idx)
        subtitle_geometry_idx = _geometry_idx_for_role(slide_spec, manifest, PlaceholderKind.SUBTITLE, subtitle_idx)
        body_geometry_idx = _geometry_idx_for_role(slide_spec, manifest, PlaceholderKind.BODY, body_idx)
        footer_geometry_idx = _geometry_idx_for_role(slide_spec, manifest, PlaceholderKind.FOOTER, footer_idx)
        body_candidates = _shapes_for_role(slide, placeholders, slide_spec, manifest, PlaceholderKind.BODY, body_idx)
        body = _preferred_text_shape(body_candidates)
        footer = _shape_for_role(slide, placeholders, slide_spec, manifest, PlaceholderKind.FOOTER, footer_idx)
        title = _shape_for_role(slide, placeholders, slide_spec, manifest, PlaceholderKind.TITLE, title_idx)
        subtitle = _shape_for_role(slide, placeholders, slide_spec, manifest, PlaceholderKind.SUBTITLE, subtitle_idx)
        layout_key = slide_spec.preferred_layout_key or _infer_layout_key(slide_spec.kind.value)
        runtime_profile_key = _runtime_profile_key_for_slide(slide_spec, manifest)
        resolved_geometry = _geometry_policy_for_slide(slide_spec, manifest)
        expected_body_spec = _shape_spec_for_role(slide_spec, manifest, PlaceholderKind.BODY)
        expected_footer_spec = _shape_spec_for_role(slide_spec, manifest, PlaceholderKind.FOOTER)
        expected_body_geometry = resolved_geometry.placeholders.get(body_geometry_idx or 14)
        effective_profile = derive_capacity_profile_for_geometry(
            runtime_profile_key,
            width_emu=expected_body_geometry.width_emu if expected_body_geometry is not None else None,
            height_emu=expected_body_geometry.height_emu if expected_body_geometry is not None else None,
        )
        excluded_auxiliary_indices = {
            value for value in {title_idx, subtitle_idx, body_idx, footer_idx, 17} if value is not None
        }
        if runtime_profile_key in {"cards_3", "cards_kpi"}:
            excluded_auxiliary_indices -= {11, 12, 13}
        auxiliary_widths = {
            idx: getattr(shape, "width", None)
            for idx, shape in placeholders.items()
            if idx not in excluded_auxiliary_indices
        }
        auxiliary_lefts = {
            idx: getattr(shape, "left", None)
            for idx, shape in placeholders.items()
            if idx not in excluded_auxiliary_indices
        }
        auxiliary_tops = {
            idx: getattr(shape, "top", None)
            for idx, shape in placeholders.items()
            if idx not in excluded_auxiliary_indices
        }
        auxiliary_char_counts = {
            idx: _shape_text_char_count(shape)
            for idx, shape in placeholders.items()
            if idx not in excluded_auxiliary_indices
        }
        expected_auxiliary_char_counts = _expected_auxiliary_char_counts_for_slide(slide_spec, runtime_profile_key)
        placeholder_char_counts = {
            idx: _shape_text_char_count(shape)
            for idx, shape in placeholders.items()
        }
        expected_placeholder_char_counts = _expected_placeholder_char_counts_for_slide(slide_spec, runtime_profile_key)
        expected_subtitle_char_count = _expected_subtitle_char_count_for_slide(slide_spec, runtime_profile_key)
        content_width = getattr(body, "width", None) if body is not None else None
        footer_width = getattr(footer, "width", None) if footer is not None else None
        expected_footer_width = None
        if expected_footer_spec is not None:
            explicit_footer_width = getattr(expected_footer_spec, "width_emu", None)
            if explicit_footer_width:
                expected_footer_width = explicit_footer_width
            elif getattr(expected_footer_spec, "shape_name", None) and footer_width is not None:
                expected_footer_width = footer_width

        body_char_count = 0
        body_font_sizes: tuple[float, ...] = ()
        rendered_items: tuple[str, ...] = ()
        if runtime_profile_key in {"cards_3", "cards_kpi"}:
            card_texts: list[str] = []
            body_paragraphs: list[str] = []
            card_font_size_values: set[float] = set()
            weighted_body_char_count = 0
            card_shapes = [placeholders[idx] for idx in (11, 12, 13) if idx in placeholders]
            if runtime_profile_key == "cards_kpi":
                metric_geometry = geometry_policy_for_layout("cards_kpi")
                metric_top = metric_geometry.placeholders[11].top_emu
                metric_bottom_limit = (
                    metric_geometry.placeholders[13].top_emu
                    + metric_geometry.placeholders[13].height_emu
                    - GEOMETRY_TOLERANCE_EMU
                )
                card_shapes.extend(
                    shape
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                    and not getattr(shape, "is_placeholder", False)
                    and getattr(shape, "text", "").strip()
                    and getattr(shape, "top", 0) >= metric_top
                    and getattr(shape, "top", 0) <= metric_bottom_limit
                )
                card_shapes = sorted(card_shapes, key=lambda shape: (getattr(shape, "top", 0), getattr(shape, "left", 0)))
                card_texts = [
                    " ".join(
                        paragraph.text.strip()
                        for paragraph in shape.text_frame.paragraphs
                        if paragraph.text.strip()
                    )
                    for shape in card_shapes
                    if getattr(shape, "has_text_frame", False)
                ]
                body_paragraphs = [
                    paragraph.text.strip()
                    for shape in card_shapes
                    if getattr(shape, "has_text_frame", False)
                    for paragraph in shape.text_frame.paragraphs
                    if paragraph.text.strip()
                ]
                weighted_body_char_count = sum(len(paragraph) for paragraph in body_paragraphs)
            else:
                overlay_shapes = [
                    shape
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False)
                    and not getattr(shape, "is_placeholder", False)
                    and getattr(shape, "name", "").startswith("A3_CARD_OVERLAY_")
                    and getattr(shape, "text", "").strip()
                ]
                overlay_shapes_by_idx: dict[int, list] = {}
                for shape in overlay_shapes:
                    match = re.match(r"^A3_CARD_OVERLAY_(\d+)_", getattr(shape, "name", ""))
                    if not match:
                        continue
                    overlay_shapes_by_idx.setdefault(int(match.group(1)), []).append(shape)
                for idx in (11, 12, 13):
                    placeholder = placeholders.get(idx)
                    if placeholder is None:
                        continue
                    current_shapes = overlay_shapes_by_idx.get(idx)
                    if current_shapes:
                        current_shapes = sorted(current_shapes, key=lambda shape: (getattr(shape, "top", 0), getattr(shape, "left", 0)))
                        card_texts.append(
                            " ".join(
                                paragraph.text.strip()
                                for shape in current_shapes
                                for paragraph in shape.text_frame.paragraphs
                                if paragraph.text.strip()
                            )
                        )
                        body_paragraphs.extend(
                            paragraph.text.strip()
                            for shape in current_shapes
                            for paragraph in shape.text_frame.paragraphs
                            if paragraph.text.strip()
                        )
                        weighted_body_char_count += sum(_weighted_card_overlay_char_count(shape) for shape in current_shapes)
                        for shape in current_shapes:
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.size is not None:
                                        card_font_size_values.add(run.font.size.pt)
                        continue
                    if getattr(placeholder, "has_text_frame", False):
                        paragraphs = [paragraph.text.strip() for paragraph in placeholder.text_frame.paragraphs if paragraph.text.strip()]
                        if paragraphs:
                            card_texts.append(" ".join(paragraphs))
                            body_paragraphs.extend(paragraphs)
                            weighted_body_char_count += sum(len(paragraph) for paragraph in paragraphs)
                            for paragraph in placeholder.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    if run.font.size is not None:
                                        card_font_size_values.add(run.font.size.pt)
            body_char_count = weighted_body_char_count or sum(len(paragraph) for paragraph in body_paragraphs)
            rendered_items = tuple(text for text in card_texts if text)
            body_font_sizes = tuple(sorted(card_font_size_values))
        elif body_candidates:
            body_paragraphs: list[str] = []
            body_font_size_values: set[float] = set()
            for candidate in body_candidates:
                if not getattr(candidate, "has_text_frame", False):
                    continue
                body_paragraphs.extend(
                    paragraph.text.strip()
                    for paragraph in candidate.text_frame.paragraphs
                    if paragraph.text.strip()
                )
                body_font_size_values.update(_text_frame_font_sizes(candidate))
            body_char_count = sum(len(paragraph) for paragraph in body_paragraphs)
            rendered_items = tuple(body_paragraphs)
            expected_body_font_size = getattr(getattr(expected_body_spec, "text_style", None), "font_size_pt", None)
            if expected_body_font_size is None and manifest is not None and manifest.theme.master_text_styles.get("body") is not None:
                expected_body_font_size = manifest.theme.master_text_styles.get("body").font_size_pt
            if expected_body_font_size is None:
                expected_body_font_size = effective_profile.max_font_pt
            if not body_font_size_values and expected_body_font_size is not None:
                body_font_size_values.add(expected_body_font_size)
            body_font_sizes = tuple(sorted(body_font_size_values))
        title_font_sizes = _text_frame_font_sizes(title)
        subtitle_font_sizes = _text_frame_font_sizes(subtitle)
        subtitle_has_text = bool((slide_spec.subtitle or "").strip())

        has_table = any(getattr(shape, "has_table", False) for shape in slide.shapes)
        has_chart = any(getattr(shape, "has_chart", False) for shape in slide.shapes)
        table_overlay_overflows = _table_overlay_overflows(slide) if slide_spec.kind == SlideKind.TABLE else ()
        image_shape = _content_image_shape(slide, placeholders)
        has_image = image_shape is not None
        chart_shape = next((shape for shape in slide.shapes if getattr(shape, "has_chart", False)), None)
        expected_chart_type = _expected_render_chart_type(slide_spec.chart) if slide_spec.kind == SlideKind.CHART else None
        expected_chart_series_count = _expected_render_chart_series_count(slide_spec.chart) if slide_spec.kind == SlideKind.CHART else None
        expected_chart_axis_number_format = (
            _expected_chart_axis_number_format(slide_spec.chart) if slide_spec.kind == SlideKind.CHART else None
        )
        expected_chart_secondary_axis_number_format = (
            chart_axis_number_format_for_axis(slide_spec.chart, SECONDARY_AXIS) if slide_spec.kind == SlideKind.CHART else None
        )
        rendered_chart_semantics = _chart_semantics(getattr(chart_shape, "chart", None)) if chart_shape is not None else {}
        if content_width is None:
            if has_table:
                table_shape = next((shape for shape in slide.shapes if getattr(shape, "has_table", False)), None)
                content_width = getattr(table_shape, "width", None) if table_shape is not None else None
            elif has_chart:
                content_width = getattr(chart_shape, "width", None) if chart_shape is not None else None
            elif has_image:
                content_width = getattr(image_shape, "width", None) if image_shape is not None else None
        expected_chart_shape_spec = (
            _shape_spec_for_role(slide_spec, manifest, PlaceholderKind.CHART)
            if slide_spec.kind == SlideKind.CHART
            else None
        )
        expected_chart_content_width = None
        if expected_chart_shape_spec is not None:
            explicit_width = getattr(expected_chart_shape_spec, "width_emu", None)
            if explicit_width:
                expected_chart_content_width = explicit_width
            elif getattr(expected_chart_shape_spec, "shape_name", None) and content_width is not None:
                expected_chart_content_width = content_width
        expected_items = _expected_items_for_slide(slide_spec)
        fallback_body_font_sizes = body_font_sizes
        if not fallback_body_font_sizes and body is not None:
            fallback_font_size = getattr(getattr(expected_body_spec, "text_style", None), "font_size_pt", None)
            if fallback_font_size is None and manifest is not None and manifest.theme.master_text_styles.get("body") is not None:
                fallback_font_size = manifest.theme.master_text_styles.get("body").font_size_pt
            if fallback_font_size is None:
                fallback_font_size = effective_profile.max_font_pt
            if fallback_font_size is not None:
                fallback_body_font_sizes = (float(fallback_font_size),)
        render_target = slide_spec.render_target
        audits.append(
            SlideAudit(
                slide_index=slide_index,
                title=slide_spec.title or "",
                kind=slide_spec.kind.value,
                layout_key=layout_key,
                runtime_profile_key=runtime_profile_key,
                body_char_count=body_char_count,
                body_font_sizes=fallback_body_font_sizes,
                title_font_sizes=title_font_sizes,
                subtitle_font_sizes=subtitle_font_sizes if subtitle_has_text else (),
                profile=effective_profile,
                content_width=content_width,
                footer_width=footer_width,
                has_table=has_table,
                has_chart=has_chart,
                has_image=has_image,
                expected_items=expected_items,
                rendered_items=rendered_items,
                title_top=getattr(title, "top", None) if title is not None else None,
                title_height=getattr(title, "height", None) if title is not None else None,
                title_width=getattr(title, "width", None) if title is not None else None,
                subtitle_top=getattr(subtitle, "top", None) if subtitle_has_text else None,
                subtitle_height=getattr(subtitle, "height", None) if subtitle_has_text else None,
                subtitle_width=getattr(subtitle, "width", None) if subtitle is not None else None,
                body_top=getattr(body, "top", None) if body is not None else None,
                body_height=getattr(body, "height", None) if body is not None else None,
                body_left=getattr(body, "left", None) if body is not None else None,
                footer_top=getattr(footer, "top", None) if footer is not None else None,
                footer_left=getattr(footer, "left", None) if footer is not None else None,
                auxiliary_widths=auxiliary_widths,
                auxiliary_lefts=auxiliary_lefts,
            auxiliary_tops=auxiliary_tops,
            auxiliary_char_counts=auxiliary_char_counts,
            expected_auxiliary_char_counts=expected_auxiliary_char_counts,
            placeholder_char_counts=placeholder_char_counts,
            expected_placeholder_char_counts=expected_placeholder_char_counts,
            expected_subtitle_char_count=expected_subtitle_char_count,
            image_width=getattr(image_shape, "width", None) if image_shape is not None else None,
            image_left=getattr(image_shape, "left", None) if image_shape is not None else None,
            image_top=getattr(image_shape, "top", None) if image_shape is not None else None,
                body_margin_left=getattr(getattr(body, "text_frame", None), "margin_left", None) if body is not None and getattr(body, "has_text_frame", False) else None,
                body_margin_right=getattr(getattr(body, "text_frame", None), "margin_right", None) if body is not None and getattr(body, "has_text_frame", False) else None,
                body_margin_top=getattr(getattr(body, "text_frame", None), "margin_top", None) if body is not None and getattr(body, "has_text_frame", False) else None,
                body_margin_bottom=getattr(getattr(body, "text_frame", None), "margin_bottom", None) if body is not None and getattr(body, "has_text_frame", False) else None,
                expected_body_margin_left=getattr(expected_body_spec, "margin_left_emu", None),
                expected_body_margin_right=getattr(expected_body_spec, "margin_right_emu", None),
                expected_body_margin_top=getattr(expected_body_spec, "margin_top_emu", None),
                expected_body_margin_bottom=getattr(expected_body_spec, "margin_bottom_emu", None),
                expected_body_max_font_pt=getattr(getattr(expected_body_spec, "text_style", None), "font_size_pt", None),
                expected_footer_width=expected_footer_width,
                title_placeholder_idx=title_geometry_idx,
                subtitle_placeholder_idx=subtitle_geometry_idx,
                body_placeholder_idx=body_geometry_idx,
                footer_placeholder_idx=footer_geometry_idx,
                expected_chart_type=expected_chart_type.value if expected_chart_type is not None else None,
                rendered_chart_type=rendered_chart_semantics.get("chart_type"),
                expected_chart_series_count=expected_chart_series_count,
                rendered_chart_series_count=rendered_chart_semantics.get("series_count"),
                rendered_chart_bar_series_count=rendered_chart_semantics.get("bar_series_count"),
                rendered_chart_line_series_count=rendered_chart_semantics.get("line_series_count"),
                expected_chart_value_axis_number_format=expected_chart_axis_number_format,
                rendered_chart_value_axis_number_format=rendered_chart_semantics.get("value_axis_number_format"),
                expected_chart_secondary_value_axis=uses_secondary_value_axis(slide_spec.chart) if slide_spec.kind == SlideKind.CHART else False,
                rendered_chart_secondary_value_axis=bool(rendered_chart_semantics.get("secondary_value_axis_number_format")),
                expected_chart_secondary_value_axis_number_format=expected_chart_secondary_axis_number_format,
                rendered_chart_secondary_value_axis_number_format=rendered_chart_semantics.get("secondary_value_axis_number_format"),
                expected_chart_content_width=expected_chart_content_width,
                expected_title_font_pt=(
                    manifest.theme.master_text_styles.get("title").font_size_pt
                    if manifest is not None and manifest.theme.master_text_styles.get("title") is not None
                    else None
                ),
                expected_subtitle_font_pt=(
                    manifest.theme.master_text_styles.get("body").font_size_pt
                    if manifest is not None and manifest.theme.master_text_styles.get("body") is not None
                    else None
                ),
                table_overlay_overflow_count=len(table_overlay_overflows),
                table_overlay_overflow_details=table_overlay_overflows,
                geometry=resolved_geometry,
                target_type=render_target.type.value if render_target is not None else None,
                target_source=render_target.source if render_target is not None else None,
                target_degradation_reasons=tuple(render_target.degradation_reasons) if render_target is not None else (),
            )
        )

    return audits


def _expected_render_chart_type(chart_spec: ChartSpec | None) -> ChartType | None:
    if chart_spec is None:
        return None
    return render_chart_type(chart_spec)


def _table_overlay_overflows(slide) -> tuple[str, ...]:
    overflows: list[str] = []
    for shape in slide.shapes:
        name = getattr(shape, "name", "") or ""
        if not name.startswith("A3_TABLE_CELL_") or not getattr(shape, "has_text_frame", False):
            continue
        text = "\n".join(paragraph.text.strip() for paragraph in shape.text_frame.paragraphs if paragraph.text.strip())
        if not text:
            continue
        font_sizes = _text_frame_font_sizes(shape)
        font_pt = min(font_sizes) if font_sizes else 8.0
        margin_left = getattr(shape.text_frame, "margin_left", 0) or 0
        margin_right = getattr(shape.text_frame, "margin_right", 0) or 0
        margin_top = getattr(shape.text_frame, "margin_top", 0) or 0
        margin_bottom = getattr(shape.text_frame, "margin_bottom", 0) or 0
        available_width = max(int(getattr(shape, "width", 0) - margin_left - margin_right), 1)
        available_height = max(int(getattr(shape, "height", 0) - margin_top - margin_bottom), 1)
        char_width_emu = max(int(font_pt * 12700 * 0.52), 1)
        chars_per_line = max(4, int(available_width / char_width_emu))
        explicit_lines = text.splitlines() or [text]
        line_count = sum(max(1, (len(line) + chars_per_line - 1) // chars_per_line) for line in explicit_lines)
        required_height = int(line_count * font_pt * 12700 * 1.18)
        if required_height > int(available_height * 1.08):
            overflows.append(
                f"{name}: required={required_height} available={available_height} text_len={len(text)} font={font_pt:.1f}"
            )
    return tuple(overflows)


def _expected_render_chart_series_count(chart_spec: ChartSpec | None) -> int | None:
    return render_chart_series_count(chart_spec)


def _chart_semantics(chart) -> dict[str, int | str | None]:
    if chart is None:
        return {}

    chart_space = chart._chartSpace
    bar_charts = chart_space.xpath(".//c:barChart")
    line_charts = chart_space.xpath(".//c:lineChart")
    pie_charts = chart_space.xpath(".//c:pieChart")
    bar_series_count = sum(len(element.xpath("./c:ser")) for element in bar_charts)
    line_series_count = sum(len(element.xpath("./c:ser")) for element in line_charts)
    pie_series_count = sum(len(element.xpath("./c:ser")) for element in pie_charts)
    series_count = bar_series_count + line_series_count + pie_series_count

    chart_type = "unknown"
    if bar_charts and line_charts:
        chart_type = ChartType.COMBO.value
    elif pie_charts:
        chart_type = ChartType.PIE.value
    elif line_charts:
        chart_type = ChartType.LINE.value
    elif bar_charts:
        bar_chart = bar_charts[0]
        bar_direction = next((element.get("val") for element in bar_chart.xpath("./c:barDir")), "col")
        grouping = next((element.get("val") for element in bar_chart.xpath("./c:grouping")), "clustered")
        is_stacked = grouping in {"stacked", "percentStacked"}
        if bar_direction == "bar":
            chart_type = ChartType.STACKED_BAR.value if is_stacked else ChartType.BAR.value
        else:
            chart_type = ChartType.STACKED_COLUMN.value if is_stacked else ChartType.COLUMN.value

    return {
        "chart_type": chart_type,
        "series_count": series_count,
        "bar_series_count": bar_series_count,
        "line_series_count": line_series_count,
        "value_axis_number_format": _chart_value_axis_number_format(chart, PRIMARY_AXIS),
        "secondary_value_axis_number_format": _chart_value_axis_number_format(chart, SECONDARY_AXIS),
    }


def _chart_value_axis_number_format(chart, axis_role: str = PRIMARY_AXIS) -> str | None:
    try:
        value_axes = chart._chartSpace.valAx_lst
        index = 0 if axis_role == PRIMARY_AXIS else 1
        if len(value_axes) <= index:
            return None
        return ValueAxis(value_axes[index]).tick_labels.number_format
    except Exception:
        return None


def _expected_chart_axis_number_format(chart_spec: ChartSpec | None) -> str | None:
    return chart_axis_number_format(chart_spec)


def _expected_render_chart_spec(chart_spec: ChartSpec | None) -> ChartSpec | None:
    return render_chart_spec(chart_spec)


def _text_frame_font_sizes(shape) -> tuple[float, ...]:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return ()
    font_sizes: set[float] = set()
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size is not None:
                font_sizes.add(run.font.size.pt)
        if not font_sizes and getattr(paragraph, "font", None) is not None and paragraph.font.size is not None:
            font_sizes.add(paragraph.font.size.pt)
    return tuple(sorted(font_sizes))


def _font_sizes_match(font_sizes: tuple[float, ...], expected: float) -> bool:
    return bool(font_sizes) and all(abs(size - expected) <= CHART_TITLE_FONT_TOLERANCE_PT for size in font_sizes)


def continuation_groups(audits: list[SlideAudit]) -> dict[str, list[SlideAudit]]:
    groups: dict[str, list[SlideAudit]] = {}
    current_title = ""
    current_items: list[SlideAudit] = []

    def flush() -> None:
        nonlocal current_title, current_items
        if len(current_items) > 1:
            groups.setdefault(current_title, []).extend(current_items)
        current_title = ""
        current_items = []

    for audit in audits:
        if audit.kind not in {SlideKind.TEXT.value, SlideKind.BULLETS.value, SlideKind.TWO_COLUMN.value}:
            flush()
            continue
        base_title = re.sub(r"\s+\(\d+\)$", "", audit.title.strip())
        if current_items and base_title != current_title:
            flush()
        current_title = base_title
        current_items.append(audit)
    flush()
    return {title: items for title, items in groups.items() if len(items) > 1}


def find_capacity_violations(audits: list[SlideAudit]) -> list[CapacityViolation]:
    violations: list[CapacityViolation] = []

    for audit in audits:
        strict_layout_contracts = audit.target_type not in {"auto_layout", "direct_shape_binding"}
        geometry = audit.geometry or geometry_policy_for_layout(
            audit.runtime_profile_key or audit.profile.layout_key or audit.layout_key
        )
        if not audit.within_font_bounds:
            violations.append(
                CapacityViolation(
                    slide_index=audit.slide_index,
                    title=audit.title,
                    rule="font_bounds",
                    details=(
                        f"fonts={audit.body_font_sizes} profile={audit.profile.min_font_pt}-{audit.profile.max_font_pt}"
                    ),
                )
            )

        if audit.fill_ratio > (audit.profile.max_fill_ratio + 0.02):
            violations.append(
                CapacityViolation(
                    slide_index=audit.slide_index,
                    title=audit.title,
                    rule="overflow_risk",
                    details=f"fill_ratio={audit.fill_ratio:.2f} max={audit.profile.max_fill_ratio:.2f}",
                )
            )

        if audit.kind == SlideKind.TABLE.value:
            if not audit.has_table:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="missing_table_shape",
                        details="table slide does not contain rendered table shape",
                    )
                )
            if audit.content_width_ratio and audit.content_width_ratio < 0.9:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_table_content",
                        details=f"content_width_ratio={audit.content_width_ratio:.2f}",
                    )
                )
            if audit.table_overlay_overflow_count:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="table_overlay_text_overflow",
                        details="; ".join(audit.table_overlay_overflow_details[:4]),
                    )
                )

        if audit.kind == SlideKind.CHART.value:
            expected_title_font_pt = audit.expected_title_font_pt or EXPECTED_CHART_TITLE_FONT_PT
            expected_subtitle_font_pt = audit.expected_subtitle_font_pt or EXPECTED_CHART_SUBTITLE_FONT_PT
            if not audit.has_chart:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="missing_chart_shape",
                        details="chart slide does not contain rendered chart shape",
                    )
                )
            if audit.title_font_sizes and not _font_sizes_match(audit.title_font_sizes, expected_title_font_pt):
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="chart_title_font_mismatch",
                        details=f"title_font_sizes={audit.title_font_sizes} expected={expected_title_font_pt}",
                    )
                )
            if audit.subtitle_font_sizes and not _font_sizes_match(audit.subtitle_font_sizes, expected_subtitle_font_pt):
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="chart_subtitle_font_mismatch",
                        details=f"subtitle_font_sizes={audit.subtitle_font_sizes} expected={expected_subtitle_font_pt}",
                    )
                )
            if audit.has_chart and audit.expected_chart_type and audit.rendered_chart_type != audit.expected_chart_type:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="chart_type_mismatch",
                        details=f"rendered={audit.rendered_chart_type} expected={audit.expected_chart_type}",
                    )
                )
            if (
                audit.has_chart
                and audit.expected_chart_series_count is not None
                and audit.rendered_chart_series_count != audit.expected_chart_series_count
            ):
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="chart_series_count_mismatch",
                        details=f"rendered={audit.rendered_chart_series_count} expected={audit.expected_chart_series_count}",
                    )
                )
            if (
                audit.has_chart
                and audit.expected_chart_value_axis_number_format
                and audit.rendered_chart_value_axis_number_format != audit.expected_chart_value_axis_number_format
            ):
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="chart_value_axis_number_format_mismatch",
                        details=(
                            f"rendered={audit.rendered_chart_value_axis_number_format} "
                            f"expected={audit.expected_chart_value_axis_number_format}"
                        ),
                    )
                )
            if audit.expected_chart_type == ChartType.COMBO.value and audit.has_chart:
                expected_bar_series = max((audit.expected_chart_series_count or 0) - 1, 0)
                if audit.rendered_chart_bar_series_count != expected_bar_series or audit.rendered_chart_line_series_count != 1:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="combo_chart_structure_mismatch",
                            details=(
                                f"bar_series={audit.rendered_chart_bar_series_count} expected={expected_bar_series}; "
                                f"line_series={audit.rendered_chart_line_series_count} expected=1"
                            ),
                        )
                    )
            if audit.expected_chart_secondary_value_axis and not audit.rendered_chart_secondary_value_axis:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="missing_secondary_value_axis",
                        details="chart should render a secondary value axis but it is missing",
                    )
                )
            if (
                audit.has_chart
                and audit.expected_chart_secondary_value_axis_number_format
                and audit.rendered_chart_secondary_value_axis_number_format != audit.expected_chart_secondary_value_axis_number_format
            ):
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="chart_secondary_value_axis_number_format_mismatch",
                        details=(
                            f"rendered={audit.rendered_chart_secondary_value_axis_number_format} "
                            f"expected={audit.expected_chart_secondary_value_axis_number_format}"
                        ),
                    )
                )
            if audit.chart_content_width_ratio and audit.chart_content_width_ratio < 0.9:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_chart_content",
                        details=f"content_width_ratio={audit.chart_content_width_ratio:.2f}",
                    )
                )

        if audit.kind == SlideKind.IMAGE.value:
            if not audit.has_image:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="missing_image_shape",
                        details="image slide does not contain rendered image shape",
                    )
                )
            if audit.content_width_ratio and audit.content_width_ratio < 0.35:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_image_content",
                        details=f"content_width_ratio={audit.content_width_ratio:.2f}",
                    )
                )
            if audit.layout_key == "image_text" and audit.image_left is not None and audit.body_left is not None and audit.content_width:
                separation = audit.image_left - (audit.body_left + audit.content_width)
                if separation < 300000:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="image_text_overlap",
                            details=f"separation={separation} min=300000",
                        )
                    )

        if audit.kind in {SlideKind.TABLE.value, SlideKind.CHART.value}:
            if audit.footer_width_ratio and audit.footer_width_ratio < 0.9:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_table_footer",
                        details=f"footer_width_ratio={audit.footer_width_ratio:.2f}",
                    )
                )

        expected_footer_geometry = None
        if strict_layout_contracts and audit.footer_placeholder_idx is not None:
            for footer_idx in (audit.footer_placeholder_idx, 15, 17, 21):
                if footer_idx is None:
                    continue
                expected_footer_geometry = geometry.placeholders.get(footer_idx)
                if expected_footer_geometry is not None:
                    break
        if strict_layout_contracts and expected_footer_geometry is not None:
            expected_footer_width = audit.expected_footer_width or expected_footer_geometry.width_emu
            if audit.footer_width is not None and audit.footer_width < expected_footer_width - GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_footer",
                        details=f"footer_width={audit.footer_width} expected_min={expected_footer_width - GEOMETRY_TOLERANCE_EMU}",
                    )
                )
            if audit.footer_left is not None and abs(audit.footer_left - expected_footer_geometry.left_emu) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="footer_left_misalignment",
                        details=f"footer_left={audit.footer_left} expected={expected_footer_geometry.left_emu}",
                    )
                )

        if strict_layout_contracts and audit.layout_key in {"text_full_width", "dense_text_full_width", "list_full_width"}:
            if audit.footer_width_ratio and audit.footer_width_ratio < 0.9:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_text_footer",
                        details=f"footer_width_ratio={audit.footer_width_ratio:.2f}",
                    )
                )

            expected_body_geometry = geometry.placeholders.get(audit.body_placeholder_idx or 14)
            if expected_body_geometry is not None and audit.body_left is not None and abs(audit.body_left - expected_body_geometry.left_emu) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="body_left_misalignment",
                        details=(
                            f"body_left={audit.body_left} expected={expected_body_geometry.left_emu} "
                            f"tolerance={GEOMETRY_TOLERANCE_EMU}"
                        ),
                    )
                )
            expected_margin_left = audit.expected_body_margin_left or PptxGenerator.DEFAULT_TEXT_MARGIN_X_EMU
            expected_margin_right = audit.expected_body_margin_right or PptxGenerator.DEFAULT_TEXT_MARGIN_X_EMU
            expected_margin_top = audit.expected_body_margin_top or PptxGenerator.DEFAULT_TEXT_MARGIN_Y_EMU
            expected_margin_bottom = audit.expected_body_margin_bottom or PptxGenerator.DEFAULT_TEXT_MARGIN_Y_EMU
            if audit.body_margin_left is not None and abs(audit.body_margin_left - expected_margin_left) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="body_margin_mismatch",
                        details=f"margin_left={audit.body_margin_left} expected={expected_margin_left}",
                    )
                )
            if audit.body_margin_right is not None and abs(audit.body_margin_right - expected_margin_right) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="body_margin_mismatch",
                        details=f"margin_right={audit.body_margin_right} expected={expected_margin_right}",
                    )
                )
            if audit.body_margin_top is not None and abs(audit.body_margin_top - expected_margin_top) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="body_margin_mismatch",
                        details=f"margin_top={audit.body_margin_top} expected={expected_margin_top}",
                    )
                )
            if audit.body_margin_bottom is not None and abs(audit.body_margin_bottom - expected_margin_bottom) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="body_margin_mismatch",
                        details=f"margin_bottom={audit.body_margin_bottom} expected={expected_margin_bottom}",
                    )
                )

        if strict_layout_contracts and audit.layout_key == "image_text" and audit.kind == SlideKind.IMAGE.value:
            if audit.image_width is not None and audit.image_width < geometry.placeholders[16].width_emu - GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="narrow_image_panel",
                        details=f"image_width={audit.image_width} min={geometry.placeholders[16].width_emu - GEOMETRY_TOLERANCE_EMU}",
                    )
                )
            if audit.body_left is not None and abs(audit.body_left - geometry.placeholders[14].left_emu) > GEOMETRY_TOLERANCE_EMU:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="image_text_body_misalignment",
                        details=f"body_left={audit.body_left} expected={geometry.placeholders[14].left_emu}",
                    )
                )

        if audit.layout_key in {"cards_3", "cards_kpi"}:
            card_lefts = audit.auxiliary_lefts or {}
            card_widths = audit.auxiliary_widths or {}
            card_positions = [(idx, card_lefts.get(idx), card_widths.get(idx)) for idx in (11, 12, 13)]
            if all(left is not None and width is not None for _, left, width in card_positions):
                for idx, left, width in card_positions:
                    if width < geometry.placeholders[idx].width_emu - GEOMETRY_TOLERANCE_EMU:
                        violations.append(
                            CapacityViolation(
                                slide_index=audit.slide_index,
                                title=audit.title,
                                rule="narrow_card_placeholder",
                                details=f"idx={idx} width={width}",
                            )
                        )
                if audit.layout_key == "cards_3":
                    previous_right = None
                    for idx, left, width in card_positions:
                        if previous_right is not None and left < previous_right + 120000:
                            violations.append(
                                CapacityViolation(
                                    slide_index=audit.slide_index,
                                    title=audit.title,
                                    rule="card_overlap",
                                    details=f"idx={idx} left={left} previous_right={previous_right}",
                                )
                            )
                        previous_right = left + width
                elif audit.layout_key == "cards_kpi":
                    row_positions = {
                        "top": [(idx, left, width) for idx, left, width in card_positions if idx in {11, 12}],
                        "bottom": [(idx, left, width) for idx, left, width in card_positions if idx in {13}],
                    }
                    top_cards = row_positions["top"]
                    if len(top_cards) == 2 and top_cards[1][1] < top_cards[0][1] + top_cards[0][2] + 120000:
                        violations.append(
                            CapacityViolation(
                                slide_index=audit.slide_index,
                                title=audit.title,
                                rule="card_overlap",
                                details=f"idx={top_cards[1][0]} left={top_cards[1][1]} previous_right={top_cards[0][1] + top_cards[0][2]}",
                            )
                        )

        if audit.layout_key == "list_with_icons":
            aux_lefts = audit.auxiliary_lefts or {}
            aux_widths = audit.auxiliary_widths or {}
            left_left = aux_lefts.get(12)
            left_width = aux_widths.get(12)
            right_left = aux_lefts.get(14)
            right_width = aux_widths.get(14)
            if None not in {left_left, left_width, right_left, right_width}:
                if left_width < geometry.placeholders[12].width_emu - GEOMETRY_TOLERANCE_EMU:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="narrow_left_column",
                            details=f"width={left_width}",
                        )
                    )
                if right_width < geometry.placeholders[14].width_emu - GEOMETRY_TOLERANCE_EMU:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="narrow_right_column",
                            details=f"width={right_width}",
                        )
                    )
                if right_left < left_left + left_width + 300000:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="two_column_overlap",
                            details=f"left_right_gap={right_left - (left_left + left_width)}",
                        )
                    )

        if audit.layout_key == "contacts":
            aux_lefts = audit.auxiliary_lefts or {}
            aux_widths = audit.auxiliary_widths or {}
            for idx in (10, 11, 12, 13):
                expected = geometry.placeholders[idx]
                left = aux_lefts.get(idx)
                width = aux_widths.get(idx)
                if left is not None and abs(left - expected.left_emu) > GEOMETRY_TOLERANCE_EMU:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="contact_block_misalignment",
                            details=f"idx={idx} left={left} expected={expected.left_emu}",
                        )
                    )
                if width is not None and width < expected.width_emu - GEOMETRY_TOLERANCE_EMU:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="narrow_contact_block",
                            details=f"idx={idx} width={width}",
                        )
                    )

        if strict_layout_contracts and audit.layout_key in {"text_full_width", "dense_text_full_width", "list_full_width", "image_text"}:
            expected_body_geometry = geometry.placeholders.get(audit.body_placeholder_idx or 14)
            expected_body_height = expected_body_geometry.height_emu if expected_body_geometry is not None else None
            minimum_fill_ratio = _minimum_placeholder_body_fill_ratio(audit)
            if (
                expected_body_height is not None
                and audit.body_height is not None
                and audit.body_height < int(expected_body_height * minimum_fill_ratio)
                and audit.fill_ratio < max(audit.profile.target_fill_ratio - 0.2, 0.0)
            ):
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="underfilled_placeholder_fill",
                        details=(
                            f"body_height={audit.body_height} "
                            f"expected_min={int(expected_body_height * minimum_fill_ratio)} "
                            f"height_fill_ratio={audit.body_height / expected_body_height:.2f}"
                        ),
                    )
                )

        if audit.expected_auxiliary_char_counts:
            rendered_auxiliary_chars = audit.auxiliary_char_counts or {}
            for idx, expected_chars in audit.expected_auxiliary_char_counts.items():
                if expected_chars <= 0:
                    continue
                rendered_chars = rendered_auxiliary_chars.get(idx, 0)
                if rendered_chars <= 0:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="underfilled_auxiliary_placeholder_fill",
                            details=f"idx={idx} expected_chars={expected_chars} rendered_chars={rendered_chars}",
                        )
                    )

        if audit.layout_key == "contacts" and audit.expected_placeholder_char_counts:
            rendered_placeholder_chars = audit.placeholder_char_counts or {}
            for idx, expected_chars in audit.expected_placeholder_char_counts.items():
                if expected_chars <= 0:
                    continue
                rendered_chars = rendered_placeholder_chars.get(idx, 0)
                if rendered_chars <= 0:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="underfilled_contact_placeholder_fill",
                            details=f"idx={idx} expected_chars={expected_chars} rendered_chars={rendered_chars}",
                        )
                    )

        if audit.layout_key == "list_with_icons" and audit.expected_placeholder_char_counts:
            rendered_placeholder_chars = audit.placeholder_char_counts or {}
            for idx, expected_chars in audit.expected_placeholder_char_counts.items():
                if expected_chars <= 0:
                    continue
                rendered_chars = rendered_placeholder_chars.get(idx, 0)
                if rendered_chars <= 0:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="underfilled_two_column_placeholder_fill",
                            details=f"idx={idx} expected_chars={expected_chars} rendered_chars={rendered_chars}",
                        )
                    )

        if audit.layout_key == "cards_3" and audit.expected_placeholder_char_counts:
            rendered_placeholder_chars = audit.placeholder_char_counts or {}
            for idx, expected_chars in audit.expected_placeholder_char_counts.items():
                if expected_chars <= 0:
                    continue
                rendered_chars = rendered_placeholder_chars.get(idx, 0)
                if rendered_chars <= 0 and not audit.rendered_items:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="underfilled_card_placeholder_fill",
                            details=f"idx={idx} expected_chars={expected_chars} rendered_chars={rendered_chars}",
                        )
                    )

        if (
            strict_layout_contracts
            and audit.expected_subtitle_char_count > 0
            and audit.placeholder_char_counts
            and audit.subtitle_placeholder_idx is not None
        ):
            rendered_subtitle_chars = audit.placeholder_char_counts.get(audit.subtitle_placeholder_idx, 0)
            if rendered_subtitle_chars <= 0:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="underfilled_subtitle_placeholder_fill",
                        details=(
                            f"idx={audit.subtitle_placeholder_idx} "
                            f"expected_chars={audit.expected_subtitle_char_count} "
                            f"rendered_chars={rendered_subtitle_chars}"
                        ),
                    )
                )

        if audit.body_top is not None and audit.body_height is not None and audit.footer_top is not None:
            minimum_bottom_gap = geometry.content_footer_gap_emu - GEOMETRY_TOLERANCE_EMU
            body_bottom = audit.body_top + audit.body_height
            footer_gap = audit.footer_top - body_bottom
            if footer_gap < minimum_bottom_gap:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="content_footer_overlap",
                        details=f"gap={footer_gap} min={minimum_bottom_gap}",
                    )
                )

        if audit.title_bottom is not None and audit.body_top is not None:
            expected_title_geometry = geometry.placeholders.get(audit.title_placeholder_idx or 0)
            expected_subtitle_geometry = geometry.placeholders.get(audit.subtitle_placeholder_idx or 13)
            expected_body_geometry = geometry.placeholders.get(audit.body_placeholder_idx or 14)
            skip_stack_overlap_check = audit.layout_key == "dense_text_full_width"
            if audit.subtitle_top is not None and audit.subtitle_height is not None:
                title_subtitle_gap = audit.subtitle_top - audit.title_bottom
                minimum_gap = geometry.title_content_gap_emu - GEOMETRY_TOLERANCE_EMU
                if not skip_stack_overlap_check and title_subtitle_gap < minimum_gap:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="title_subtitle_overlap",
                            details=f"gap={title_subtitle_gap} min={minimum_gap}",
                        )
                    )
                expected_title_subtitle_gap = _expected_gap_between(expected_title_geometry, expected_subtitle_geometry)
                title_subtitle_gap_limit = _maximum_reasonable_gap(expected_title_subtitle_gap)
                if title_subtitle_gap_limit is not None and title_subtitle_gap > title_subtitle_gap_limit:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="title_subtitle_gap_drift",
                            details=f"gap={title_subtitle_gap} max={title_subtitle_gap_limit}",
                        )
                    )
                subtitle_body_gap = audit.body_top - audit.subtitle_bottom
                if not skip_stack_overlap_check and subtitle_body_gap < minimum_gap:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="subtitle_body_overlap",
                            details=f"gap={subtitle_body_gap} min={minimum_gap}",
                        )
                    )
                expected_subtitle_body_gap = _expected_gap_between(expected_subtitle_geometry, expected_body_geometry)
                subtitle_body_gap_limit = _maximum_reasonable_gap(expected_subtitle_body_gap)
                if subtitle_body_gap_limit is not None and subtitle_body_gap > subtitle_body_gap_limit:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="subtitle_body_gap_drift",
                            details=f"gap={subtitle_body_gap} max={subtitle_body_gap_limit}",
                        )
                    )
            else:
                title_body_gap = audit.body_top - audit.title_bottom
                if audit.kind in {SlideKind.TABLE.value, SlideKind.CHART.value}:
                    expected_gap = geometry.title_content_gap_emu
                else:
                    expected_gap = geometry.title_body_gap_no_subtitle_emu
                minimum_gap = expected_gap - GEOMETRY_TOLERANCE_EMU
                if title_body_gap < minimum_gap:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="title_body_overlap",
                            details=f"gap={title_body_gap} min={minimum_gap}",
                        )
                    )
                expected_title_body_gap = _expected_gap_between(expected_title_geometry, expected_body_geometry)
                title_body_gap_limit = _maximum_reasonable_gap(expected_title_body_gap)
                if audit.layout_key != "cards_kpi" and title_body_gap_limit is not None and title_body_gap > title_body_gap_limit:
                    violations.append(
                        CapacityViolation(
                            slide_index=audit.slide_index,
                            title=audit.title,
                            rule="title_body_gap_drift",
                            details=f"gap={title_body_gap} max={title_body_gap_limit}",
                        )
                    )

        should_check_body_order = audit.kind == SlideKind.BULLETS.value or (
            audit.kind == SlideKind.TEXT.value and audit.expected_items and any(audit.expected_items)
        )
        if should_check_body_order and audit.expected_items:
            expected_source = audit.expected_items
            if audit.layout_key in {"cards_3", "cards_kpi"}:
                expected_source = tuple(_normalize_card_audit_item(item) for item in audit.expected_items)
            expected = [_normalize_audit_text(item) for item in expected_source if _normalize_audit_text(item)]
            rendered = [_normalize_audit_text(item) for item in audit.rendered_items if _normalize_audit_text(item)]
            if rendered and expected != rendered:
                violations.append(
                    CapacityViolation(
                        slide_index=audit.slide_index,
                        title=audit.title,
                        rule="content_order_mismatch",
                        details=f"expected={expected} rendered={rendered}",
                    )
                )

    for title, items in continuation_groups(audits).items():
        fills = [item.fill_ratio for item in items]
        fill_delta = max(fills) - min(fills)
        tolerance = items[0].profile.continuation_balance_tolerance
        has_group_overflow = any(item.fill_ratio > (item.profile.max_fill_ratio + 0.02) for item in items)
        has_material_underfill = any(
            item.fill_ratio + CONTINUATION_AUDIT_EPSILON
            < max(item.profile.target_fill_ratio - item.profile.continuation_balance_tolerance - CONTINUATION_UNDERFILL_GRACE, 0.0)
            for item in items[1:]
        )
        if fill_delta > tolerance and (has_group_overflow or has_material_underfill):
            violations.append(
                CapacityViolation(
                    slide_index=items[-1].slide_index,
                    title=title,
                    rule="continuation_balance",
                    details=f"delta={fill_delta:.2f} tolerance={tolerance:.2f}",
                )
            )

        for item in items[1:]:
            min_fill = max(
                item.profile.target_fill_ratio
                - item.profile.continuation_balance_tolerance
                - CONTINUATION_UNDERFILL_GRACE,
                0.0,
            )
            if item.fill_ratio + CONTINUATION_AUDIT_EPSILON < min_fill:
                violations.append(
                    CapacityViolation(
                        slide_index=item.slide_index,
                        title=item.title,
                        rule="underfilled_continuation",
                        details=f"fill_ratio={item.fill_ratio:.2f} min={min_fill:.2f}",
                    )
                )
            elif item.fill_ratio > (item.profile.max_fill_ratio + 0.02):
                violations.append(
                    CapacityViolation(
                        slide_index=item.slide_index,
                        title=item.title,
                        rule="overflow_continuation",
                        details=f"fill_ratio={item.fill_ratio:.2f} max={item.profile.max_fill_ratio:.2f}",
                    )
                )

        for previous, current in zip(items, items[1:]):
            previous_font = previous.representative_font_size
            current_font = current.representative_font_size
            if previous_font is None or current_font is None:
                continue
            delta = abs(previous_font - current_font)
            if delta > CONTINUATION_FONT_DELTA_TOLERANCE_PT:
                violations.append(
                    CapacityViolation(
                        slide_index=current.slide_index,
                        title=current.title,
                        rule="continuation_font_delta",
                        details=(
                            f"delta={delta:.2f} tolerance={CONTINUATION_FONT_DELTA_TOLERANCE_PT:.2f} "
                            f"previous={previous_font:.2f} current={current_font:.2f}"
                        ),
                    )
                )

        tracked_items = [item for item in items if item.expected_items]
        expected_group = [
            _normalize_audit_text(entry)
            for item in tracked_items
            for entry in item.expected_items
            if _normalize_audit_text(entry)
        ]
        rendered_group = [
            _normalize_audit_text(entry)
            for item in tracked_items
            for entry in item.rendered_items
            if _normalize_audit_text(entry)
        ]
        if expected_group and rendered_group and expected_group != rendered_group:
            violations.append(
                CapacityViolation(
                    slide_index=items[-1].slide_index,
                    title=title,
                    rule="continuation_order_mismatch",
                    details=f"expected={expected_group} rendered={rendered_group}",
                )
            )

    return violations


def _minimum_placeholder_body_fill_ratio(audit: SlideAudit) -> float:
    if audit.layout_key == "dense_text_full_width":
        return 0.55
    if audit.layout_key == "image_text":
        return 0.4
    if audit.layout_key == "list_full_width":
        return 0.42
    return BODY_HEIGHT_UNDERFILL_RATIO


def _expected_auxiliary_char_counts_for_slide(slide_spec: SlideSpec, layout_key: str) -> dict[int, int]:
    def list_len(values: list[str]) -> int:
        return len("\n".join(item.strip() for item in values if item.strip()))

    if layout_key == "list_with_icons":
        left_column_chars = list_len(slide_spec.left_bullets)
        return {12: left_column_chars} if left_column_chars else {}
    return {}


def _expected_placeholder_char_counts_for_slide(slide_spec: SlideSpec, layout_key: str) -> dict[int, int]:
    def text_len(value: str | None) -> int:
        return len((value or "").strip())

    def list_len(values: list[str]) -> int:
        return len("\n".join(item.strip() for item in values if item.strip()))

    if layout_key == "contacts":
        expected = {
            10: text_len(slide_spec.title),
            11: text_len(slide_spec.subtitle),
            12: list_len(slide_spec.left_bullets),
            13: list_len(slide_spec.right_bullets),
        }
        return {idx: count for idx, count in expected.items() if count > 0}
    if layout_key == "list_with_icons":
        expected = {
            12: list_len(slide_spec.left_bullets),
            14: list_len(slide_spec.right_bullets),
        }
        return {idx: count for idx, count in expected.items() if count > 0}
    if layout_key == "cards_3":
        bullets = [item.strip() for item in slide_spec.bullets if item.strip()]
        return {
            idx: len(text)
            for idx, text in zip((11, 12, 13), bullets, strict=False)
            if len(text) > 0
        }
    return {}


def _expected_subtitle_char_count_for_slide(slide_spec: SlideSpec, layout_key: str) -> int:
    subtitle_text = (slide_spec.subtitle or "").strip()
    if not subtitle_text:
        return 0
    if layout_key in {"contacts", "list_with_icons"}:
        return 0
    if layout_key in {"text_full_width", "dense_text_full_width", "list_full_width"} and _subtitle_duplicates_body_payload(slide_spec, subtitle_text):
        return 0
    return len(subtitle_text)


def _subtitle_duplicates_body_payload(slide_spec: SlideSpec, subtitle_text: str) -> bool:
    body_text = ""
    if slide_spec.content_blocks:
        for block in slide_spec.content_blocks:
            if block.text and block.text.strip():
                body_text = block.text.strip()
                break
    if not body_text:
        body_text = (slide_spec.text or "").strip()
    if not body_text:
        return False
    if body_text.startswith(subtitle_text):
        return True
    if len(subtitle_text) >= 24 and body_text.startswith(subtitle_text[:-1]):
        return True
    return False


def _shape_text_char_count(shape) -> int:
    if shape is None or not getattr(shape, "has_text_frame", False):
        return 0
    return sum(
        len(paragraph.text.strip())
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text and paragraph.text.strip()
    )


def _expected_gap_between(upper_geometry: PlaceholderGeometryPolicy | None, lower_geometry: PlaceholderGeometryPolicy | None) -> int | None:
    if upper_geometry is None or lower_geometry is None:
        return None
    expected_gap = lower_geometry.top_emu - (upper_geometry.top_emu + upper_geometry.height_emu)
    if expected_gap < 0:
        return None
    return expected_gap


def _maximum_reasonable_gap(expected_gap: int | None) -> int | None:
    if expected_gap is None:
        return None
    return max(expected_gap + GEOMETRY_TOLERANCE_EMU, MAX_REASONABLE_STACK_GAP_EMU)


def _placeholder_idx_for_role(
    slide_spec: SlideSpec,
    manifest: TemplateManifest | None,
    kind: PlaceholderKind,
    fallback_indices: tuple[int, ...],
) -> int | None:
    placeholder = _shape_spec_for_role(slide_spec, manifest, kind)
    if placeholder is not None:
        placeholder_idx = getattr(placeholder, "idx", None)
        if placeholder_idx is not None:
            return placeholder_idx
    if manifest is not None and _template_has_explicit_layout_or_prototype(slide_spec, manifest):
        return None
    return fallback_indices[0] if fallback_indices else None


def _geometry_idx_for_role(
    slide_spec: SlideSpec,
    manifest: TemplateManifest | None,
    kind: PlaceholderKind,
    placeholder_idx: int | None,
) -> int | None:
    if placeholder_idx is not None:
        return placeholder_idx
    if manifest is None:
        return None
    prototype = _prototype_slide_for_spec(slide_spec, manifest)
    if prototype is None:
        return None
    token_spec = _prototype_token_spec_for_role(prototype, kind, slide_spec.kind)
    if token_spec is None:
        return None
    synthetic_indices = {
        PlaceholderKind.TITLE: 0,
        PlaceholderKind.SUBTITLE: 13,
        PlaceholderKind.BODY: 14,
        PlaceholderKind.FOOTER: 17,
    }
    return synthetic_indices.get(kind)


def _shape_spec_for_role(
    slide_spec: SlideSpec,
    manifest: TemplateManifest | None,
    kind: PlaceholderKind,
):
    if manifest is None:
        return None
    layout = _layout_target_for_spec(slide_spec, manifest)
    if layout is not None:
        typed = [placeholder for placeholder in layout.placeholders if placeholder.kind == kind and placeholder.idx is not None]
        if manifest.generation_mode.value == "layout" and typed:
            explicit_typed = [
                placeholder
                for placeholder in typed
                if None not in {
                    getattr(placeholder, "left_emu", None),
                    getattr(placeholder, "top_emu", None),
                    getattr(placeholder, "width_emu", None),
                    getattr(placeholder, "height_emu", None),
                }
            ]
            if explicit_typed:
                return explicit_typed[0]
        preferred_indices = _preferred_placeholder_indices_for_role(slide_spec, kind, manifest)
        for preferred_idx in preferred_indices:
            preferred = next((placeholder for placeholder in layout.placeholders if placeholder.idx == preferred_idx), None)
            if preferred is not None:
                return preferred
        if typed:
            return typed[0]
    prototype = _prototype_slide_for_spec(slide_spec, manifest)
    if prototype is not None:
        bindings = _prototype_bindings_for_role(kind, slide_spec.kind)
        for binding in bindings:
            token = next((item for item in prototype.tokens if item.binding == binding and item.shape_name), None)
            if token is not None:
                return token
    return None


def _preferred_placeholder_indices_for_role(
    slide_spec: SlideSpec,
    kind: PlaceholderKind,
    manifest: TemplateManifest | None,
) -> tuple[int, ...]:
    layout_key = _runtime_profile_key_for_slide(slide_spec, manifest)
    if kind == PlaceholderKind.TITLE:
        return (0,)
    if kind == PlaceholderKind.SUBTITLE:
        return (13,)
    if kind == PlaceholderKind.BODY:
        if layout_key in {"cards_3", "cards_kpi"}:
            return (11, 12, 13)
        if layout_key in {"text_full_width", "dense_text_full_width", "list_full_width", "image_text"}:
            return (14,)
        if layout_key == "table" or slide_spec.kind == SlideKind.CHART:
            return (14, 11, 12)
        return ()
    if kind == PlaceholderKind.FOOTER:
        if layout_key in {"text_full_width", "dense_text_full_width", "list_full_width", "image_text"}:
            return (17,)
        if layout_key == "table":
            return (15,)
        return ()
    return ()


def _template_has_explicit_layout_or_prototype(slide_spec: SlideSpec, manifest: TemplateManifest | None) -> bool:
    if manifest is None:
        return False
    if _layout_target_for_spec(slide_spec, manifest) is not None:
        return True
    return _prototype_slide_for_spec(slide_spec, manifest) is not None


def _geometry_policy_for_slide(slide_spec: SlideSpec, manifest: TemplateManifest | None) -> LayoutGeometryPolicy:
    runtime_profile_key = _runtime_profile_key_for_slide(slide_spec, manifest)
    base_policy = geometry_policy_for_layout(runtime_profile_key)
    if manifest is None:
        return base_policy
    layout = _layout_target_for_spec(slide_spec, manifest)
    if layout is None:
        return base_policy
    placeholders: dict[int, PlaceholderGeometryPolicy] = {}
    for placeholder in layout.placeholders:
        if placeholder.idx is None:
            continue
        if None in {placeholder.left_emu, placeholder.top_emu, placeholder.width_emu, placeholder.height_emu}:
            continue
        placeholders[placeholder.idx] = PlaceholderGeometryPolicy(
            placeholder_idx=placeholder.idx,
            left_emu=placeholder.left_emu,
            top_emu=placeholder.top_emu,
            width_emu=placeholder.width_emu,
            height_emu=placeholder.height_emu,
        )
    if not placeholders:
        prototype = _prototype_slide_for_spec(slide_spec, manifest)
        if prototype is None:
            return base_policy
        role_specs = {
            0: _prototype_token_spec_for_role(prototype, PlaceholderKind.TITLE, slide_spec.kind),
            13: _prototype_token_spec_for_role(prototype, PlaceholderKind.SUBTITLE, slide_spec.kind),
            14: _prototype_token_spec_for_role(prototype, PlaceholderKind.BODY, slide_spec.kind),
            17: _prototype_token_spec_for_role(prototype, PlaceholderKind.FOOTER, slide_spec.kind),
        }
        for synthetic_idx, token_spec in role_specs.items():
            if token_spec is None:
                continue
            if None in {token_spec.left_emu, token_spec.top_emu, token_spec.width_emu, token_spec.height_emu}:
                continue
            placeholders[synthetic_idx] = PlaceholderGeometryPolicy(
                placeholder_idx=synthetic_idx,
                left_emu=token_spec.left_emu,
                top_emu=token_spec.top_emu,
                width_emu=token_spec.width_emu,
                height_emu=token_spec.height_emu,
            )
        if not placeholders:
            return base_policy
    return LayoutGeometryPolicy(
        layout_key=layout.key,
        placeholders=placeholders,
        title_content_gap_emu=base_policy.title_content_gap_emu,
        title_body_gap_no_subtitle_emu=base_policy.title_body_gap_no_subtitle_emu,
        content_footer_gap_emu=base_policy.content_footer_gap_emu,
    )


def _shape_for_role(slide, placeholders: dict[int, object], slide_spec: SlideSpec, manifest: TemplateManifest | None, kind: PlaceholderKind, placeholder_idx: int | None):
    shapes = _shapes_for_role(slide, placeholders, slide_spec, manifest, kind, placeholder_idx)
    return _preferred_text_shape(shapes)


def _shapes_for_role(
    slide,
    placeholders: dict[int, object],
    slide_spec: SlideSpec,
    manifest: TemplateManifest | None,
    kind: PlaceholderKind,
    placeholder_idx: int | None,
) -> list[object]:
    if placeholder_idx is not None and placeholder_idx in placeholders:
        return [placeholders[placeholder_idx]]
    shape_spec = _shape_spec_for_role(slide_spec, manifest, kind)
    if shape_spec is None:
        return []
    shape_name = getattr(shape_spec, "shape_name", None)
    named_matches = []
    if shape_name:
        named_matches = [shape for shape in slide.shapes if getattr(shape, "name", None) == shape_name]
        if named_matches:
            return named_matches
    geometry_matches = _shapes_matching_geometry(slide, shape_spec)
    if geometry_matches:
        return geometry_matches
    return named_matches


def _shapes_matching_geometry(slide, shape_spec) -> list[object]:
    expected_left = getattr(shape_spec, "left_emu", None)
    expected_top = getattr(shape_spec, "top_emu", None)
    expected_width = getattr(shape_spec, "width_emu", None)
    expected_height = getattr(shape_spec, "height_emu", None)
    if None in {expected_left, expected_top, expected_width, expected_height}:
        return []
    matches = [
        shape
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False)
        and abs(getattr(shape, "left", 0) - expected_left) <= GEOMETRY_TOLERANCE_EMU
        and abs(getattr(shape, "top", 0) - expected_top) <= GEOMETRY_TOLERANCE_EMU
        and abs(getattr(shape, "width", 0) - expected_width) <= GEOMETRY_TOLERANCE_EMU
        and abs(getattr(shape, "height", 0) - expected_height) <= GEOMETRY_TOLERANCE_EMU
    ]
    return sorted(matches, key=lambda shape: (getattr(shape, "top", 0), getattr(shape, "left", 0)))


def _preferred_text_shape(shapes: list[object]) -> object | None:
    if not shapes:
        return None
    with_text = [shape for shape in shapes if getattr(shape, "has_text_frame", False) and getattr(shape, "text", "").strip()]
    if with_text:
        return with_text[0]
    return shapes[0]


def _prototype_slide_for_spec(slide_spec: SlideSpec, manifest: TemplateManifest | None):
    if manifest is None or not manifest.prototype_slides:
        return None
    if slide_spec.kind == SlideKind.CHART:
        for prototype in manifest.prototype_slides:
            if _prototype_supports_chart(prototype):
                return prototype
    if slide_spec.preferred_layout_key:
        preferred = next((item for item in manifest.prototype_slides if item.key == slide_spec.preferred_layout_key), None)
        if preferred is not None:
            return preferred
    return next((item for item in manifest.prototype_slides if slide_spec.kind.value in item.supported_slide_kinds), None)


def _layout_target_for_spec(slide_spec: SlideSpec, manifest: TemplateManifest | None):
    if manifest is None:
        return None
    preferred_key = slide_spec.preferred_layout_key or ""
    if not preferred_key:
        return None
    return next((item for item in manifest.layouts if item.key == preferred_key), None)


def _runtime_profile_key_for_slide(slide_spec: SlideSpec, manifest: TemplateManifest | None) -> str:
    target = None
    if manifest is not None:
        target = _layout_target_for_spec(slide_spec, manifest)
        if target is None:
            target = _prototype_slide_for_spec(slide_spec, manifest)
    return runtime_profile_key_for_target(
        target,
        fallback_layout_key=slide_spec.runtime_profile_key or slide_spec.preferred_layout_key or _infer_layout_key(slide_spec.kind.value),
        slide_kind=slide_spec.kind.value,
    )


def _prototype_supports_chart(prototype) -> bool:
    if getattr(prototype, "key", None) == "chart" or SlideKind.CHART.value in getattr(prototype, "supported_slide_kinds", []):
        return True
    return any(getattr(token, "binding", None) in {"chart", "chart_image"} for token in getattr(prototype, "tokens", []))


def _prototype_token_spec_for_role(prototype, kind: PlaceholderKind, slide_kind: SlideKind):
    bindings = _prototype_bindings_for_role(kind, slide_kind)
    for binding in bindings:
        token = next((item for item in prototype.tokens if item.binding == binding and item.shape_name), None)
        if token is not None:
            return token
    return None


def _prototype_bindings_for_role(kind: PlaceholderKind, slide_kind: SlideKind) -> tuple[str, ...]:
    if kind == PlaceholderKind.TITLE:
        return ("cover_title", "title", "contact_title")
    if kind == PlaceholderKind.SUBTITLE:
        return ("subtitle", "cover_meta", "contact_role")
    if kind == PlaceholderKind.CHART:
        return ("chart", "chart_image")
    if kind == PlaceholderKind.TABLE:
        return ("table",)
    if kind == PlaceholderKind.IMAGE:
        return ("image",)
    if kind == PlaceholderKind.BODY:
        if slide_kind == SlideKind.BULLETS:
            return (
                "bullets",
                "right_list",
                "left_bullets",
                "right_bullets",
                "main_text",
                "text",
                "body",
                "summary",
                "secondary_text",
                "left_text",
                "right_text",
                "left_note",
                "right_note",
                "contact_phone",
                "contact_email",
            )
        if slide_kind == SlideKind.TWO_COLUMN:
            return ("left_text", "main_text", "text", "body", "summary", "secondary_text")
        if slide_kind == SlideKind.TITLE:
            return ("cover_meta", "subtitle", "main_text", "text", "body", "summary")
        return (
            "main_text",
            "text",
            "body",
            "summary",
            "secondary_text",
            "left_text",
            "right_text",
            "left_note",
            "right_note",
            "contact_phone",
            "contact_email",
            "bullets",
            "right_list",
        )
    if kind == PlaceholderKind.FOOTER:
        return ("presentation_name", "notes", "cover_meta")
    return ()


def _infer_layout_key(kind: str) -> str:
    if kind == SlideKind.BULLETS.value:
        return "list_full_width"
    if kind == SlideKind.IMAGE.value:
        return "image_text"
    return "text_full_width"


def self_has_image(shape) -> bool:
    return hasattr(shape, "image") or "Picture Placeholder" in (getattr(shape, "name", "") or "")


def _content_image_shape(slide, placeholders: dict[int, object]):
    for preferred_idx in (16,):
        shape = placeholders.get(preferred_idx)
        if shape is not None:
            return shape
    for shape in slide.shapes:
        if "Picture Placeholder" in (getattr(shape, "name", "") or ""):
            return shape
    for shape in slide.shapes:
        if hasattr(shape, "image"):
            return shape
    return None


def _expected_items_for_slide(slide_spec: SlideSpec) -> tuple[str, ...]:
    if slide_spec.content_blocks:
        has_list_like_block = any(block.items for block in slide_spec.content_blocks)
        if slide_spec.kind == SlideKind.TEXT and not has_list_like_block:
            return ()
        expected: list[str] = []
        for block in slide_spec.content_blocks:
            if block.text and block.text.strip():
                expected.append(block.text.strip())
            expected.extend(item.strip() for item in block.items if item.strip())
        return tuple(expected)
    if slide_spec.kind == SlideKind.BULLETS:
        return tuple(item.strip() for item in slide_spec.bullets if item.strip())
    return ()


def _normalize_audit_text(text: str) -> str:
    normalized = re.sub(r"\s+", " ", (text or "").strip())
    return re.sub(r"\s+([%‰₽])", r"\1", normalized)


def _normalize_card_audit_item(text: str) -> str:
    normalized = "\n".join(line.strip() for line in (text or "").splitlines() if line.strip())
    if not normalized:
        return ""
    if "\n" in normalized:
        title, description = normalized.split("\n", 1)
        return f"{title.strip()} {' '.join(description.split())}".strip()
    colon_match = re.match(r"^(.{4,54}?):\s+(.{12,})$", normalized)
    if colon_match:
        return f"{colon_match.group(1).strip()} {colon_match.group(2).strip()}"
    dash_match = re.match(r"^(.{4,54}?)\s+[—-]\s+(.{12,})$", normalized)
    if dash_match:
        return f"{dash_match.group(1).strip()} {dash_match.group(2).strip()}"
    return normalized


def _weighted_card_overlay_char_count(shape) -> int:
    paragraphs = [
        paragraph.text.strip()
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text.strip()
    ]
    if not paragraphs:
        return 0
    shape_name = getattr(shape, "name", "")
    if "_METRIC_" in shape_name:
        value_text = paragraphs[0]
        label_text = " ".join(paragraphs[1:])
        return min(len(value_text), 8) + int(len(label_text) * 0.35)
    if "_DESCRIPTION" in shape_name:
        return int(sum(len(paragraph) for paragraph in paragraphs) * 0.7)
    if "_TITLE" in shape_name:
        return int(sum(len(paragraph) for paragraph in paragraphs) * 0.8)
    return sum(len(paragraph) for paragraph in paragraphs)
