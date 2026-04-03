from __future__ import annotations

import base64
import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from docx import Document
from pptx import Presentation

from a3presentation.services.document_text_extractor import DocumentTextExtractor
from a3presentation.services.planner import TextToPlanService
from a3presentation.services.pptx_generator import PptxGenerator
from a3presentation.services.template_registry import TemplateRegistry
from a3presentation.settings import get_settings


SMALL_PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9W6i8AAAAASUVORK5CYII="
)


class EndToEndDocxPipelineTests(unittest.TestCase):
    def test_docx_pipeline_generates_branded_presentation(self) -> None:
        document = Document()
        document.add_paragraph("Продуктовый отчет Q4")
        document.add_heading("Основные выводы", level=1)
        document.add_paragraph(
            "Рост выручки ускорился за счет новых сегментов, улучшения конверсии и снижения операционных потерь."
        )
        table = document.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Показатель"
        table.cell(0, 1).text = "Значение"
        table.cell(1, 0).text = "GMV"
        table.cell(1, 1).text = "125"
        document.add_picture(BytesIO(SMALL_PNG_BYTES))

        buffer = BytesIO()
        document.save(buffer)
        content = buffer.getvalue()

        extractor = DocumentTextExtractor()
        text, tables, blocks = extractor.extract("report.docx", content)
        plan = TextToPlanService().build_plan("corp_light_v1", text, None, tables, blocks)

        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            self.assertGreaterEqual(len(presentation.slides), 3)
            self.assertTrue(any(slide.kind.value == "image" for slide in plan.slides))


if __name__ == "__main__":
    unittest.main()
