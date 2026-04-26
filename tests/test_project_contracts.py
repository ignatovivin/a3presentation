from __future__ import annotations

import asyncio
import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt
from starlette.datastructures import UploadFile

from a3presentation.api import routes as routes_module
from a3presentation.domain.api import TextPlanRequest
from a3presentation.domain.chart import ChartConfidence, ChartSeries, ChartSpec, ChartType
from a3presentation.domain.presentation import (
    PresentationPlan,
    RenderTargetType,
    SlideContentBlock,
    SlideContentBlockKind,
    SlideKind,
    SlideRenderTarget,
    SlideSpec,
    TableBlock,
)
from a3presentation.domain.template import (
    ComponentConfidence,
    ComponentEditability,
    ComponentGeometry,
    ExtractedComponent,
    ExtractedComponentRole,
    ExtractedPresentationInventory,
    ExtractedSlideInventory,
    ExtractedComponentType,
    GenerationMode,
    InventorySourceKind,
    LayoutSpec,
    PlaceholderSpec,
    PlaceholderKind,
    PrototypeSlideSpec,
    PrototypeTokenSpec,
    TemplateManifest,
    TemplateShapeStyleSpec,
    TemplateTextStyleSpec,
)
from a3presentation.services.deck_audit import (
    SlideAudit,
    _geometry_policy_for_slide,
    audit_generated_presentation,
    continuation_groups,
    find_capacity_violations,
)
from a3presentation.services.document_text_extractor import DocumentTextExtractor
from a3presentation.services.layout_capacity import (
    LIST_FULL_WIDTH_PROFILE,
    TEXT_FULL_WIDTH_PROFILE,
    derive_capacity_profile_for_geometry,
    geometry_policy_for_layout,
    profile_for_layout,
    spacing_policy_for_layout,
)
from a3presentation.services.planner import TextToPlanService
from a3presentation.services.pptx_generator import PptxGenerator
from a3presentation.services.template_analyzer import TemplateAnalyzer
from a3presentation.services.template_registry import TemplateRegistry
from a3presentation.settings import get_settings


class ProjectContractTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.settings = get_settings()
        cls.registry = TemplateRegistry(cls.settings.templates_dir)
        cls.analyzer = TemplateAnalyzer()
        cls.extractor = DocumentTextExtractor()
        cls.planner = TextToPlanService()
        cls.generator = PptxGenerator()

    def test_template_registry_entries_are_internally_consistent(self) -> None:
        manifests = self.registry.list_templates()

        self.assertGreaterEqual(len(manifests), 1)
        self.assertEqual(len({manifest.template_id for manifest in manifests}), len(manifests))

        for manifest in manifests:
            with self.subTest(template_id=manifest.template_id):
                template_dir = self.settings.templates_dir / manifest.template_id
                self.assertTrue(template_dir.exists())
                self.assertTrue((template_dir / "manifest.json").exists())
                self.assertEqual(manifest.template_id, template_dir.name)

                details = routes_module.get_template(manifest.template_id)
                self.assertEqual(details.manifest.template_id, manifest.template_id)

                template_path = template_dir / manifest.source_pptx
                self.assertEqual(details.has_template_file, template_path.exists())
                if template_path.exists():
                    self.assertEqual(self.registry.get_template_pptx_path(manifest.template_id), template_path)

    def test_template_manifests_have_valid_internal_structure(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            with self.subTest(template_id=manifest.template_id):
                self.assertTrue(manifest.template_id.strip())
                self.assertTrue(manifest.display_name.strip())
                self.assertTrue(manifest.source_pptx.lower().endswith(".pptx"))

                layout_keys = [layout.key for layout in manifest.layouts]
                self.assertEqual(len(set(layout_keys)), len(layout_keys))

                if manifest.layouts:
                    self.assertTrue(manifest.default_layout_key)
                    self.assertIn(manifest.default_layout_key, layout_keys)

                for layout in manifest.layouts:
                    self.assertTrue(layout.key.strip())
                    self.assertTrue(layout.name.strip())
                    self.assertGreaterEqual(layout.slide_layout_index, 0)
                    self.assertEqual(
                        len(set(layout.supported_slide_kinds)),
                        len(layout.supported_slide_kinds),
                    )

                for prototype in manifest.prototype_slides:
                    self.assertTrue(prototype.key.strip())
                    self.assertGreaterEqual(prototype.source_slide_index, 0)
                    self.assertEqual(
                        len(set(prototype.supported_slide_kinds)),
                        len(prototype.supported_slide_kinds),
                    )

    def test_template_manifests_expose_component_style_layer(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")

        self.assertIn("cards", manifest.component_styles)
        self.assertIn("text", manifest.component_styles)
        self.assertIn("table", manifest.component_styles)
        self.assertIn("chart", manifest.component_styles)
        self.assertIn("image", manifest.component_styles)
        self.assertIn("cover", manifest.component_styles)
        self.assertIn("list_with_icons", manifest.component_styles)
        self.assertIn("contacts", manifest.component_styles)
        self.assertEqual(manifest.component_styles["cards"].text_styles["title"].font_size_pt, 20.0)
        self.assertEqual(manifest.component_styles["cards"].spacing_tokens["content_margin_x_emu"], 91440)
        self.assertEqual(manifest.component_styles["cards"].behavior_tokens["kpi_max_metrics"], 4)
        self.assertEqual(manifest.component_styles["table"].spacing_tokens["cell_margin_left_emu"], 80000)
        self.assertTrue(manifest.component_styles["table"].behavior_tokens["render_as_shapes"])
        self.assertEqual(manifest.component_styles["chart"].behavior_tokens["rank_color_1"], "#091E38")
        self.assertEqual(manifest.component_styles["chart"].shape_style.chart_plot_width_factor, 1.0)
        self.assertEqual(manifest.component_styles["image"].text_styles["subtitle"].font_size_pt, 18.0)
        self.assertEqual(manifest.component_styles["image"].spacing_tokens["min_image_height_emu"], 1200000)
        self.assertEqual(manifest.component_styles["cover"].text_styles["meta"].font_size_pt, 22.0)
        self.assertEqual(manifest.component_styles["cover"].spacing_tokens["title_top_emu"], 651176)
        self.assertEqual(manifest.component_styles["list_with_icons"].text_styles["subtitle"].font_size_pt, 18.0)
        self.assertEqual(manifest.component_styles["contacts"].behavior_tokens["primary_threshold_chars"], 60)

    def test_every_template_with_pptx_supports_smoke_generation(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists():
                continue

            with self.subTest(template_id=manifest.template_id):
                plan = PresentationPlan(
                    template_id=manifest.template_id,
                    title=f"{manifest.display_name} Smoke",
                    slides=self._smoke_slides_for_manifest(manifest),
                )

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = self.generator.generate(
                        template_path=template_path,
                        manifest=manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))
                    self.assertEqual(len(presentation.slides), len(plan.slides))

    def test_template_analyzer_output_stays_compatible_with_generator(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists():
                continue

            with self.subTest(template_id=manifest.template_id):
                analyzed_manifest = self.analyzer.analyze(
                    template_id=manifest.template_id,
                    template_path=template_path,
                    display_name=manifest.display_name,
                )
                self.assertEqual(analyzed_manifest.template_id, manifest.template_id)
                self.assertEqual(analyzed_manifest.source_pptx, template_path.name)

                plan = PresentationPlan(
                    template_id=manifest.template_id,
                    title=f"{manifest.display_name} Analyzed Smoke",
                    slides=self._smoke_slides_for_manifest(analyzed_manifest),
                )

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = self.generator.generate(
                        template_path=template_path,
                        manifest=manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))
                    self.assertEqual(len(presentation.slides), len(plan.slides))

    def test_template_analyzer_extracts_placeholder_geometry_for_uploaded_layouts(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists():
                continue

            with self.subTest(template_id=manifest.template_id):
                analyzed_manifest = self.analyzer.analyze(
                    template_id=manifest.template_id,
                    template_path=template_path,
                    display_name=manifest.display_name,
                )

                analyzed_placeholders = [
                    placeholder
                    for layout in analyzed_manifest.layouts
                    for placeholder in layout.placeholders
                ]
                self.assertTrue(analyzed_placeholders)
                self.assertTrue(
                    any(self._has_geometry_metadata(placeholder) for placeholder in analyzed_placeholders)
                )
                self.assertTrue(
                    any(
                        placeholder.kind in {PlaceholderKind.TITLE, PlaceholderKind.SUBTITLE, PlaceholderKind.BODY}
                        and self._has_text_margin_metadata(placeholder)
                        for placeholder in analyzed_placeholders
                    )
                )

    def test_template_analyzer_extracts_editable_slot_metadata_for_uploaded_layouts(self) -> None:
        analyzed_manifest = self.analyzer.analyze(
            template_id="corp_light_v1",
            template_path=self.settings.templates_dir / "corp_light_v1" / "template.pptx",
            display_name="Light Theme",
        )

        editable_placeholders = [
            placeholder
            for layout in analyzed_manifest.layouts
            for placeholder in layout.placeholders
            if placeholder.editable_role is not None
        ]
        self.assertTrue(editable_placeholders)
        self.assertTrue(
            any(
                placeholder.kind in {PlaceholderKind.TITLE, PlaceholderKind.SUBTITLE, PlaceholderKind.BODY}
                and placeholder.binding in {"title", "subtitle", "body"}
                and placeholder.editable_capabilities == ["text"]
                for placeholder in editable_placeholders
            )
        )
        self.assertTrue(
            any(
                placeholder.kind == PlaceholderKind.IMAGE
                and "image" in placeholder.editable_capabilities
                for placeholder in editable_placeholders
            )
        )

    def test_template_analyzer_synthesizes_prototype_slots_for_arbitrary_user_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "arbitrary-user-template.pptx"
            presentation = Presentation()
            blank_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(blank_layout)

            title_shape = slide.shapes.add_textbox(600000, 400000, 8000000, 900000)
            title_shape.text_frame.text = "Произвольный заголовок"
            title_paragraph = title_shape.text_frame.paragraphs[0]
            title_run = title_paragraph.runs[0]
            title_run.font.size = Pt(28)
            title_run.font.bold = True
            title_run.font.color.rgb = RGBColor(0x11, 0x22, 0x33)
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = RGBColor(0xEE, 0xDD, 0xCC)

            subtitle_shape = slide.shapes.add_textbox(600000, 1400000, 8000000, 700000)
            subtitle_shape.text_frame.text = "Краткий подзаголовок"

            body_shape = slide.shapes.add_textbox(600000, 2500000, 4200000, 2200000)
            body_shape.text_frame.text = "Первый тезис"
            body_shape.text_frame.add_paragraph().text = "Второй тезис"
            body_shape.text_frame.add_paragraph().text = "Третий тезис"

            side_shape = slide.shapes.add_textbox(5200000, 2500000, 3200000, 1600000)
            side_shape.text_frame.text = "Поясняющий текст без placeholder-слоя"

            presentation.save(str(template_path))

            manifest = self.analyzer.analyze(
                template_id="uploaded_arbitrary_template",
                template_path=template_path,
                display_name="Uploaded Arbitrary Template",
            )

        self.assertEqual(manifest.generation_mode, GenerationMode.PROTOTYPE)
        self.assertTrue(manifest.prototype_slides)
        prototype = manifest.prototype_slides[0]
        bindings = {token.binding for token in prototype.tokens}
        self.assertIn("title", bindings)
        self.assertIn("subtitle", bindings)
        self.assertTrue({"bullets", "body", "left_bullets", "left_text"} & bindings)
        self.assertTrue({"text", "bullets", "two_column"} & set(prototype.supported_slide_kinds))
        self.assertTrue(all(token.shape_name for token in prototype.tokens))
        title_token = next(token for token in prototype.tokens if token.binding == "title")
        self.assertIsNotNone(title_token.text_style)
        self.assertEqual(title_token.text_style.font_size_pt, 28.0)
        self.assertEqual(title_token.text_style.color, "#112233")
        self.assertTrue(title_token.text_style.bold)
        self.assertIsNotNone(title_token.shape_style)
        self.assertEqual(title_token.shape_style.fill_color, "#EEDDCC")
        body_token = next(token for token in prototype.tokens if token.binding in {"bullets", "body", "left_bullets", "left_text"})
        self.assertIsNotNone(body_token.paragraph_styles)
        self.assertIn("0", body_token.paragraph_styles.level_styles)

    def test_template_analyzer_builds_generic_inventory_for_arbitrary_user_pptx(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "inventory-user-template.pptx"
            presentation = Presentation()
            blank_layout = presentation.slide_layouts[6]
            slide = presentation.slides.add_slide(blank_layout)

            title_shape = slide.shapes.add_textbox(600000, 400000, 8000000, 900000)
            title_shape.text_frame.text = "Произвольный заголовок"

            body_shape = slide.shapes.add_textbox(600000, 1900000, 7600000, 1800000)
            body_shape.text_frame.text = "Первый тезис"
            body_shape.text_frame.add_paragraph().text = "Второй тезис"

            image_path = Path(temp_dir) / "inventory-image.png"
            image_path.write_bytes(
                bytes.fromhex("89504E470D0A1A0A0000000D49484452000000010000000108060000001F15C4890000000D49444154789C6360F8CFC0000003010100C9FE92EF0000000049454E44AE426082")
            )
            slide.shapes.add_picture(str(image_path), 7200000, 420000, 900000, 900000)

            presentation.save(str(template_path))

            manifest = self.analyzer.analyze(
                template_id="uploaded_inventory_template",
                template_path=template_path,
                display_name="Uploaded Inventory Template",
            )

        self.assertTrue(manifest.inventory.components)
        self.assertTrue(manifest.inventory.slides)
        self.assertTrue(manifest.inventory.has_prototype_inventory)
        self.assertIn(manifest.inventory.degradation_mode, {None, "prototype_only"})
        slide_inventory = next(item for item in manifest.inventory.slides if item.source_kind == InventorySourceKind.SLIDE)
        slide_components = [item for item in manifest.inventory.components if item.component_id in slide_inventory.component_ids]
        self.assertTrue(any(item.role == ExtractedComponentRole.BODY for item in slide_components))
        self.assertTrue(any(item.component_type == ExtractedComponentType.IMAGE for item in slide_components))
        title_component = next(item for item in slide_components if item.text_excerpt == "Произвольный заголовок")
        self.assertEqual(title_component.component_type, ExtractedComponentType.TEXT)
        self.assertEqual(title_component.role, ExtractedComponentRole.BODY)
        self.assertEqual(title_component.source_kind, InventorySourceKind.SLIDE)

    def test_template_registry_normalize_manifest_syncs_editable_metadata_for_bound_placeholders(self) -> None:
        analyzed_manifest = self.analyzer.analyze(
            template_id="corp_light_v1",
            template_path=self.settings.templates_dir / "corp_light_v1" / "template.pptx",
            display_name="Light Theme",
        )
        normalized_manifest = self.registry.normalize_manifest(analyzed_manifest)

        table_layout = next(layout for layout in normalized_manifest.layouts if layout.key == "table")
        table_placeholder = next(placeholder for placeholder in table_layout.placeholders if placeholder.idx == 14)
        self.assertEqual(table_placeholder.binding, "table")
        self.assertEqual(table_placeholder.editable_role, "table")
        self.assertEqual(table_placeholder.editable_capabilities, ["table"])

    def test_template_registry_normalize_manifest_does_not_add_synthetic_legacy_layout_aliases(self) -> None:
        manifest = TemplateManifest(
            template_id="custom_inventory_only",
            display_name="Custom Inventory Only",
            source_pptx="custom.pptx",
            generation_mode=GenerationMode.LAYOUT,
            default_layout_key="custom_text_target",
            layouts=[
                LayoutSpec(
                    key="custom_text_target",
                    name="Custom Text Target",
                    slide_layout_index=0,
                    supported_slide_kinds=["text"],
                    placeholders=[
                        PlaceholderSpec(name="Title", kind=PlaceholderKind.TITLE, idx=0),
                        PlaceholderSpec(name="Body", kind=PlaceholderKind.BODY, idx=14),
                        PlaceholderSpec(name="Footer", kind=PlaceholderKind.UNKNOWN, idx=17),
                    ],
                )
            ],
        )

        normalized = self.registry.normalize_manifest(manifest)

        self.assertEqual([layout.key for layout in normalized.layouts], ["custom_text_target"])
        self.assertEqual(normalized.default_layout_key, "custom_text_target")

    def test_template_registry_builds_inventory_summary_over_layouts_and_prototypes(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")

        summary = self.registry.build_inventory_summary(manifest)

        self.assertEqual(summary.generation_mode, manifest.generation_mode.value)
        self.assertEqual(summary.has_usable_layout_inventory, manifest.inventory.has_usable_layout_inventory)
        self.assertEqual(summary.has_prototype_inventory, manifest.inventory.has_prototype_inventory)
        self.assertEqual(summary.layout_target_count, len([target for target in summary.targets if target.source == "layout"]))
        self.assertEqual(summary.prototype_target_count, len([target for target in summary.targets if target.source == "prototype"]))
        self.assertTrue(summary.targets)
        self.assertTrue(all(target.key.strip() and target.name.strip() for target in summary.targets))

    def test_template_analyzer_extracts_representation_hints_for_card_like_layouts(self) -> None:
        analyzed_manifest = self.analyzer.analyze(
            template_id="corp_light_v1",
            template_path=self.settings.templates_dir / "corp_light_v1" / "template.pptx",
            display_name="Light Theme",
        )
        normalized_manifest = self.registry.normalize_manifest(analyzed_manifest)

        cards_layout = next(layout for layout in normalized_manifest.layouts if layout.key == "cards_3")
        text_layout = next(layout for layout in normalized_manifest.layouts if layout.key == "text_full_width")
        self.assertIn("cards", cards_layout.representation_hints)
        self.assertNotIn("cards", text_layout.representation_hints)

    def test_template_analyzer_extracts_representation_hints_for_data_and_contacts_layouts(self) -> None:
        analyzed_manifest = self.analyzer.analyze(
            template_id="corp_light_v1",
            template_path=self.settings.templates_dir / "corp_light_v1" / "template.pptx",
            display_name="Light Theme",
        )
        normalized_manifest = self.registry.normalize_manifest(analyzed_manifest)

        table_layout = next(layout for layout in normalized_manifest.layouts if layout.key == "table")
        image_layout = next(layout for layout in normalized_manifest.layouts if layout.key == "image_text")
        contacts_layout = next(layout for layout in normalized_manifest.layouts if layout.key == "contacts")

        self.assertIn("table", table_layout.representation_hints)
        self.assertIn("image", image_layout.representation_hints)
        self.assertIn("contacts", contacts_layout.representation_hints)

    def test_template_registry_maps_plan_slides_to_detected_template_layouts(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Layout Inventory Mapping",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Cover", preferred_layout_key="cover"),
                SlideSpec(kind=SlideKind.TEXT, title="Text", text="Narrative body", preferred_layout_key="text_full_width"),
                SlideSpec(kind=SlideKind.BULLETS, title="List", bullets=["A", "B"], preferred_layout_key="list_full_width"),
                SlideSpec(kind=SlideKind.TABLE, title="Table", table=TableBlock(headers=["A"], rows=[["1"]]), preferred_layout_key="table"),
                SlideSpec(kind=SlideKind.IMAGE, title="Image", text="Image note", preferred_layout_key="image_text"),
                SlideSpec(kind=SlideKind.TEXT, title="Contacts", text="test@example.com", preferred_layout_key="contacts"),
            ],
        )

        adapted = self.registry.apply_layout_inventory_to_plan(manifest, plan)

        self.assertEqual(adapted.slides[0].render_target.type.value, "layout")
        self.assertEqual(adapted.slides[1].render_target.type.value, "layout")
        self.assertIn("table", adapted.slides[3].render_target.binding_keys)
        self.assertTrue(all(slide.render_target is not None for slide in adapted.slides))
        self.assertTrue(all(slide.preferred_layout_key == slide.render_target.key for slide in adapted.slides))
        self.assertTrue(all((slide.runtime_profile_key or "").strip() for slide in adapted.slides))

    def test_template_registry_falls_back_from_logical_aliases_to_real_inventory_targets(self) -> None:
        manifest = TemplateManifest(
            template_id="inventory_fallback_demo",
            display_name="Inventory Fallback Demo",
            source_pptx="demo.pptx",
            generation_mode=GenerationMode.PROTOTYPE,
            default_layout_key="prototype_text",
            prototype_slides=[
                PrototypeSlideSpec(
                    key="prototype_text",
                    name="Prototype Text",
                    source_slide_index=0,
                    supported_slide_kinds=["text"],
                    tokens=[
                        PrototypeTokenSpec(
                            token="body",
                            binding="body",
                            editable_role="body",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=1200000,
                            width_emu=7000000,
                            height_emu=2600000,
                        )
                    ],
                ),
                PrototypeSlideSpec(
                    key="prototype_bullets",
                    name="Prototype Bullets",
                    source_slide_index=1,
                    supported_slide_kinds=["bullets"],
                    tokens=[
                        PrototypeTokenSpec(
                            token="bullets",
                            binding="bullets",
                            editable_role="bullet_list",
                            editable_capabilities=["bullet_list", "text"],
                            left_emu=600000,
                            top_emu=1200000,
                            width_emu=7000000,
                            height_emu=2600000,
                        )
                    ],
                ),
            ],
        )
        plan = PresentationPlan(
            template_id="inventory_fallback_demo",
            title="Inventory Alias Fallback",
            slides=[
                SlideSpec(kind=SlideKind.BULLETS, title="Cards", bullets=["One", "Two"], preferred_layout_key="cards_3"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Contacts",
                    subtitle="CEO",
                    text="HQ address",
                    left_bullets=["+7 999 123-45-67"],
                    right_bullets=["ivan@example.com"],
                    preferred_layout_key="contacts",
                ),
            ],
        )

        adapted = self.registry.apply_layout_inventory_to_plan(manifest, plan)

        self.assertEqual(adapted.slides[0].preferred_layout_key, "prototype_bullets")
        self.assertEqual(adapted.slides[1].preferred_layout_key, "prototype_text")
        self.assertEqual(adapted.slides[1].kind, SlideKind.TEXT)
        self.assertEqual(adapted.slides[1].left_bullets, [])
        self.assertEqual(adapted.slides[1].right_bullets, [])
        self.assertIn("HQ address", adapted.slides[1].text or "")
        self.assertIn("+7 999 123-45-67", adapted.slides[1].text or "")
        self.assertIn("ivan@example.com", adapted.slides[1].text or "")
        self.assertEqual(adapted.slides[0].render_target.type.value, "prototype")
        self.assertEqual(adapted.slides[0].render_target.key, "prototype_bullets")
        self.assertEqual(adapted.slides[0].render_target.degradation_reasons, ["inventory_fallback"])
        self.assertEqual(adapted.slides[1].render_target.type.value, "prototype")
        self.assertEqual(adapted.slides[1].render_target.key, "prototype_text")
        self.assertTrue(all((slide.runtime_profile_key or "").strip() for slide in adapted.slides))

    def test_template_registry_builds_slide_level_layout_review_options(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Layout Review",
            slides=[
                SlideSpec(kind=SlideKind.TEXT, title="Overview", text="Краткий обзор и выводы"),
                SlideSpec(kind=SlideKind.TABLE, title="Metrics", table=TableBlock(headers=["A"], rows=[["1"]])),
            ],
        )

        adapted = self.registry.apply_layout_inventory_to_plan(manifest, plan)
        reviews = self.registry.build_slide_layout_reviews(manifest, adapted)

        self.assertEqual(len(reviews), 2)
        self.assertEqual(reviews[0].slide_index, 0)
        self.assertEqual(reviews[0].current_layout_key, adapted.slides[0].preferred_layout_key)
        self.assertEqual(reviews[0].current_target_key, adapted.slides[0].render_target.key)
        self.assertEqual(reviews[0].current_target_type, adapted.slides[0].render_target.type.value)
        self.assertEqual(reviews[0].current_target_source, adapted.slides[0].render_target.source)
        self.assertEqual(reviews[0].current_target_explanation, adapted.slides[0].render_target.source and "Слайд будет заполнен через layout target из извлеченного inventory.")
        self.assertEqual(reviews[0].current_target_confidence, adapted.slides[0].render_target.confidence)
        self.assertEqual(reviews[0].current_target_degradation_reasons, adapted.slides[0].render_target.degradation_reasons)
        self.assertEqual(reviews[0].current_runtime_profile_key, adapted.slides[0].runtime_profile_key)
        self.assertTrue(reviews[0].available_layouts)
        self.assertEqual(reviews[0].available_layouts[0].key, adapted.slides[0].preferred_layout_key)
        self.assertEqual(reviews[0].available_layouts[0].runtime_profile_key, adapted.slides[0].runtime_profile_key)
        self.assertTrue(reviews[0].available_layouts[0].supports_current_slide_kind)
        self.assertEqual(reviews[0].available_layouts[0].source, "layout")
        self.assertTrue(reviews[0].available_layouts[0].source_label)
        self.assertTrue(reviews[0].available_layouts[0].match_summary)

        self.assertEqual(reviews[1].available_layouts[0].key, adapted.slides[1].preferred_layout_key)
        self.assertEqual(reviews[1].available_layouts[0].runtime_profile_key, adapted.slides[1].runtime_profile_key)
        self.assertIn("table", reviews[1].available_layouts[0].representation_hints)

    def test_template_registry_ranking_prefers_render_target_when_layout_key_is_missing(self) -> None:
        manifest = TemplateManifest(
            template_id="ranking_target_demo",
            display_name="Ranking Target Demo",
            source_pptx="demo.pptx",
            generation_mode=GenerationMode.PROTOTYPE,
            prototype_slides=[
                PrototypeSlideSpec(
                    key="prototype_text",
                    name="Prototype Text",
                    source_slide_index=0,
                    supported_slide_kinds=["text"],
                    tokens=[
                        PrototypeTokenSpec(
                            token="body",
                            binding="body",
                            editable_role="body",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=1200000,
                            width_emu=7000000,
                            height_emu=2600000,
                        )
                    ],
                )
            ],
        )
        slide = SlideSpec(
            kind=SlideKind.TEXT,
            title="Targeted slide",
            text="Narrative body",
            preferred_layout_key=None,
            render_target=SlideRenderTarget(type=RenderTargetType.PROTOTYPE, key="prototype_text"),
        )

        reviews = self.registry.build_slide_layout_reviews(
            manifest,
            PresentationPlan(template_id="ranking_target_demo", title="Ranking", slides=[slide]),
        )

        self.assertEqual(reviews[0].current_target_key, "prototype_text")
        self.assertEqual(reviews[0].current_target_type, "prototype")
        self.assertEqual(reviews[0].available_layouts[0].key, "prototype_text")

    def test_template_registry_preserves_planner_degradation_metadata_when_resolving_target(self) -> None:
        manifest = TemplateManifest(
            template_id="degradation_merge_demo",
            display_name="Degradation Merge Demo",
            source_pptx="demo.pptx",
            generation_mode=GenerationMode.LAYOUT,
            layouts=[
                LayoutSpec(
                    key="text_layout",
                    name="Text Layout",
                    slide_layout_index=0,
                    supported_slide_kinds=["text"],
                    placeholders=[
                        PlaceholderSpec(
                            name="Body",
                            kind=PlaceholderKind.BODY,
                            idx=14,
                            editable_role="body",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=1200000,
                            width_emu=7000000,
                            height_emu=2600000,
                        )
                    ],
                )
            ],
        )
        plan = PresentationPlan(
            template_id="degradation_merge_demo",
            title="Merge Demo",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Fallback text",
                    text="Narrative body",
                    runtime_profile_key="text_layout",
                    render_target=SlideRenderTarget(
                        type=RenderTargetType.LAYOUT,
                        key="text_layout",
                        source="planner fallback",
                        degradation_reasons=["document_fallback"],
                        confidence="medium",
                    ),
                )
            ],
        )

        adapted = self.registry.apply_layout_inventory_to_plan(manifest, plan)
        review = self.registry.build_slide_layout_reviews(manifest, adapted)[0]

        self.assertEqual(adapted.slides[0].render_target.key, "text_layout")
        self.assertIn("document_fallback", adapted.slides[0].render_target.degradation_reasons)
        self.assertEqual(adapted.slides[0].render_target.confidence, "medium")
        self.assertIn("planner fallback", adapted.slides[0].render_target.source or "")
        self.assertEqual(review.current_target_source, adapted.slides[0].render_target.source)
        self.assertEqual(review.current_target_confidence, "medium")
        self.assertIn("document_fallback", review.current_target_degradation_reasons)

    def test_template_registry_marks_unresolved_target_as_auto_layout(self) -> None:
        manifest = TemplateManifest(
            template_id="empty_inventory_demo",
            display_name="Empty Inventory Demo",
            source_pptx="demo.pptx",
            generation_mode=GenerationMode.LAYOUT,
        )
        plan = PresentationPlan(
            template_id="empty_inventory_demo",
            title="Auto Layout Fallback",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Fallback",
                    text="Narrative body",
                    runtime_profile_key="text_full_width",
                )
            ],
        )

        adapted = self.registry.apply_layout_inventory_to_plan(manifest, plan)

        self.assertEqual(adapted.slides[0].render_target.type.value, "auto_layout")
        self.assertEqual(adapted.slides[0].render_target.key, "text_full_width")
        self.assertEqual(adapted.slides[0].render_target.degradation_reasons, ["inventory_unresolved"])

    def test_template_registry_builds_direct_shape_binding_target_from_inventory(self) -> None:
        manifest = TemplateManifest(
            template_id="direct_binding_demo",
            display_name="Direct Binding Demo",
            source_pptx="demo.pptx",
            inventory=ExtractedPresentationInventory(
                degradation_mode="direct_shape_binding",
                slides=[
                    ExtractedSlideInventory(
                        source_kind=InventorySourceKind.SLIDE,
                        source_index=0,
                        name="Direct Slide",
                        component_ids=["slide_0_title", "slide_0_body"],
                        supported_slide_kinds=["text"],
                        representation_hints=["text"],
                    )
                ],
                components=[
                    ExtractedComponent(
                        component_id="slide_0_title",
                        source_kind=InventorySourceKind.SLIDE,
                        source_index=0,
                        source_name="Direct Slide",
                        shape_name="Title Box",
                        component_type=ExtractedComponentType.TEXT,
                        role=ExtractedComponentRole.TITLE,
                        binding="title",
                        confidence=ComponentConfidence.MEDIUM,
                        editability=ComponentEditability.EDITABLE,
                        capabilities=["text"],
                        geometry=ComponentGeometry(left_emu=600000, top_emu=400000, width_emu=8000000, height_emu=900000),
                    ),
                    ExtractedComponent(
                        component_id="slide_0_body",
                        source_kind=InventorySourceKind.SLIDE,
                        source_index=0,
                        source_name="Direct Slide",
                        shape_name="Body Box",
                        component_type=ExtractedComponentType.TEXT,
                        role=ExtractedComponentRole.BODY,
                        binding="body",
                        confidence=ComponentConfidence.MEDIUM,
                        editability=ComponentEditability.EDITABLE,
                        capabilities=["text"],
                        geometry=ComponentGeometry(left_emu=600000, top_emu=1800000, width_emu=7600000, height_emu=2400000),
                    ),
                ],
            ),
        )
        plan = PresentationPlan(
            template_id="direct_binding_demo",
            title="Direct Binding",
            slides=[SlideSpec(kind=SlideKind.TEXT, title="Title", text="Body text")],
        )

        adapted = self.registry.apply_layout_inventory_to_plan(manifest, plan)
        summary = self.registry.build_inventory_summary(manifest)
        editable_targets = self.registry.build_editable_targets(manifest)
        reviews = self.registry.build_slide_layout_reviews(manifest, adapted)

        self.assertEqual(summary.direct_target_count, 1)
        self.assertTrue(any(target.source == "direct_shape_binding" for target in summary.targets))
        self.assertTrue(any(target.source == "direct_shape_binding" for target in editable_targets))
        self.assertEqual(adapted.slides[0].render_target.type.value, "direct_shape_binding")
        self.assertEqual(adapted.slides[0].render_target.key, "direct_slide_0")
        self.assertIn("direct_shape_binding", adapted.slides[0].render_target.degradation_reasons)
        self.assertEqual(reviews[0].current_target_type, "direct_shape_binding")
        self.assertEqual(reviews[0].current_target_key, "direct_slide_0")
        self.assertEqual(
            reviews[0].current_target_explanation,
            "Слайд будет заполнен через извлеченные именованные shapes исходной презентации.",
        )

    def test_template_registry_prefers_direct_shape_binding_target_when_manifest_declares_direct_binding_degradation(self) -> None:
        manifest = TemplateManifest(
            template_id="direct_binding_priority_demo",
            display_name="Direct Binding Priority Demo",
            source_pptx="demo.pptx",
            layouts=[
                LayoutSpec(
                    key="text_layout",
                    name="Text Layout",
                    slide_layout_index=0,
                    supported_slide_kinds=["text"],
                    placeholders=[
                        PlaceholderSpec(
                            name="Body",
                            kind=PlaceholderKind.BODY,
                            idx=14,
                            editable_role="body",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=1200000,
                            width_emu=7000000,
                            height_emu=2600000,
                        )
                    ],
                )
            ],
            inventory=ExtractedPresentationInventory(
                degradation_mode="direct_shape_binding",
                slides=[
                    ExtractedSlideInventory(
                        source_kind=InventorySourceKind.SLIDE,
                        source_index=0,
                        name="Direct Slide",
                        component_ids=["slide_0_body"],
                        supported_slide_kinds=["text"],
                        representation_hints=["text"],
                    )
                ],
                components=[
                    ExtractedComponent(
                        component_id="slide_0_body",
                        source_kind=InventorySourceKind.SLIDE,
                        source_index=0,
                        source_name="Direct Slide",
                        shape_name="Body Box",
                        component_type=ExtractedComponentType.TEXT,
                        role=ExtractedComponentRole.BODY,
                        binding="body",
                        confidence=ComponentConfidence.MEDIUM,
                        editability=ComponentEditability.EDITABLE,
                        capabilities=["text"],
                        geometry=ComponentGeometry(left_emu=600000, top_emu=1800000, width_emu=7600000, height_emu=2400000),
                    )
                ],
            ),
        )
        slide = SlideSpec(kind=SlideKind.TEXT, title="Targeted slide", text="Narrative body")

        reviews = self.registry.build_slide_layout_reviews(
            manifest,
            PresentationPlan(template_id="direct_binding_priority_demo", title="Ranking", slides=[slide]),
        )
        resolved = self.registry.resolve_layout_key_for_slide(manifest, slide)

        self.assertEqual(reviews[0].available_layouts[0].source, "direct_shape_binding")
        self.assertEqual(reviews[0].available_layouts[0].key, "direct_slide_0")
        self.assertIn("direct binding", reviews[0].available_layouts[0].match_summary or "")
        self.assertIn("source shapes", reviews[0].available_layouts[0].match_summary or "")
        self.assertTrue(
            any("shape bindings" in reason for reason in reviews[0].available_layouts[0].recommendation_reasons)
        )
        self.assertEqual(resolved, "direct_slide_0")

    def test_template_registry_ranks_text_layouts_by_semantics_and_capacity(self) -> None:
        manifest = TemplateManifest(
            template_id="ranking_demo",
            display_name="Ranking Demo",
            source_pptx="demo.pptx",
            default_layout_key="wide_text",
            layouts=[
                LayoutSpec(
                    key="wide_text",
                    name="Wide Text",
                    slide_layout_index=0,
                    supported_slide_kinds=["text"],
                    placeholders=[
                        PlaceholderSpec(
                            name="Title",
                            kind=PlaceholderKind.TITLE,
                            idx=0,
                            editable_role="title",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=300000,
                            width_emu=7000000,
                            height_emu=700000,
                        ),
                        PlaceholderSpec(
                            name="Body",
                            kind=PlaceholderKind.BODY,
                            idx=14,
                            editable_role="body",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=1300000,
                            width_emu=8600000,
                            height_emu=3400000,
                        ),
                    ],
                ),
                LayoutSpec(
                    key="narrow_text",
                    name="Narrow Text",
                    slide_layout_index=1,
                    supported_slide_kinds=["text"],
                    placeholders=[
                        PlaceholderSpec(
                            name="Title",
                            kind=PlaceholderKind.TITLE,
                            idx=0,
                            editable_role="title",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=300000,
                            width_emu=5000000,
                            height_emu=700000,
                        ),
                        PlaceholderSpec(
                            name="Body",
                            kind=PlaceholderKind.BODY,
                            idx=14,
                            editable_role="body",
                            editable_capabilities=["text"],
                            left_emu=600000,
                            top_emu=1300000,
                            width_emu=2200000,
                            height_emu=1000000,
                        ),
                    ],
                ),
            ],
        )
        plan = PresentationPlan(
            template_id="ranking_demo",
            title="Ranking Demo",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Overview",
                    text=" ".join(["Подробный текстовый блок для проверки вместимости макета."] * 12),
                ),
            ],
        )

        reviews = self.registry.build_slide_layout_reviews(manifest, plan)

        self.assertEqual(reviews[0].available_layouts[0].key, "wide_text")
        self.assertEqual(reviews[0].available_layouts[1].key, "narrow_text")
        self.assertGreater(
            reviews[0].available_layouts[0].estimated_text_capacity_chars or 0,
            reviews[0].available_layouts[1].estimated_text_capacity_chars or 0,
        )

    def test_template_analyzer_infers_generic_editable_slot_group_metadata_for_tokens(self) -> None:
        self.assertEqual(self.analyzer._infer_slot_group("{{bullet_2}}"), "bullet")
        self.assertEqual(self.analyzer._infer_slot_group_order("{{bullet_2}}"), 2)
        self.assertEqual(self.analyzer._editable_role_for_binding("bullet_2"), "bullet_item")
        self.assertEqual(self.analyzer._editable_capabilities_for_binding("bullet_2"), ["text", "list_item"])

        self.assertEqual(self.analyzer._infer_slot_group("{{left_bullet_3}}"), "left_column")
        self.assertEqual(self.analyzer._infer_slot_group_order("{{left_bullet_3}}"), 3)
        self.assertEqual(self.analyzer._editable_role_for_binding("left_bullets"), "bullet_list")
        self.assertEqual(
            self.analyzer._editable_capabilities_for_binding("left_bullets"),
            ["bullet_list", "text"],
        )

    def test_template_analyzer_extracts_theme_and_text_style_catalog_from_xml(self) -> None:
        template_id = "corp_light_v1"
        template_path = self.settings.templates_dir / template_id / "template.pptx"

        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Light Theme",
        )

        self.assertTrue(analyzed_manifest.theme.color_scheme)
        self.assertEqual(analyzed_manifest.theme.color_scheme.get("dk1"), "#081C4F")
        self.assertEqual(analyzed_manifest.theme.font_scheme.get("major_latin"), "Mont SemiBold")
        self.assertEqual(analyzed_manifest.theme.font_scheme.get("minor_latin"), "Mont Regular")
        self.assertIn("title", analyzed_manifest.theme.master_text_styles)
        self.assertEqual(analyzed_manifest.theme.master_text_styles["title"].font_size_pt, 35.0)
        self.assertEqual(analyzed_manifest.theme.master_text_styles["body"].font_size_pt, 20.0)

        layout_placeholders = [
            placeholder
            for layout in analyzed_manifest.layouts
            for placeholder in layout.placeholders
            if placeholder.kind in {PlaceholderKind.TITLE, PlaceholderKind.BODY, PlaceholderKind.SUBTITLE}
        ]
        self.assertTrue(layout_placeholders)
        self.assertTrue(any(placeholder.text_style is not None for placeholder in layout_placeholders))
        self.assertTrue(
            any(
                placeholder.text_style is not None
                and placeholder.text_style.font_size_pt is not None
                and placeholder.text_style.color is not None
                for placeholder in layout_placeholders
            )
        )

    def test_template_analyzer_extracts_component_shape_styles_from_xml(self) -> None:
        template_id = "corp_light_v1"
        template_path = self.settings.templates_dir / template_id / "template.pptx"

        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Light Theme",
        )

        self.assertIn("background", analyzed_manifest.theme.master_shape_styles)
        background_style = analyzed_manifest.theme.master_shape_styles["background"]
        self.assertEqual(background_style.role, "background")
        self.assertIsNotNone(background_style.fill_type)

        styled_layouts = [layout for layout in analyzed_manifest.layouts if layout.background_style is not None]
        self.assertTrue(styled_layouts)

        component_placeholders = [
            placeholder
            for layout in analyzed_manifest.layouts
            for placeholder in layout.placeholders
            if placeholder.kind in {PlaceholderKind.TABLE, PlaceholderKind.CHART, PlaceholderKind.IMAGE, PlaceholderKind.BODY}
        ]
        self.assertTrue(component_placeholders)
        self.assertTrue(any(placeholder.shape_style is not None for placeholder in component_placeholders))
        self.assertTrue(
            any(
                placeholder.shape_style is not None
                and (
                    placeholder.shape_style.fill_type is not None
                    or placeholder.shape_style.line_color is not None
                    or placeholder.shape_style.geometry_preset is not None
                )
                for placeholder in component_placeholders
            )
        )

    def test_template_analyzer_extracts_background_xml_for_uploaded_layouts(self) -> None:
        layout_manifest = self.analyzer.analyze(
            template_id="corp_light_v1",
            template_path=self.settings.templates_dir / "corp_light_v1" / "template.pptx",
            display_name="Light Theme",
        )
        self.assertTrue(any(layout.background_xml for layout in layout_manifest.layouts))
        self.assertTrue(any(layout.background_image_base64 for layout in layout_manifest.layouts))

    def test_template_analyzer_extracts_advanced_paragraph_and_component_xml_metadata(self) -> None:
        template_id = "corp_light_v1"
        template_path = self.settings.templates_dir / template_id / "template.pptx"

        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Light Theme",
        )

        self.assertIn("body", analyzed_manifest.theme.master_paragraph_styles)
        body_levels = analyzed_manifest.theme.master_paragraph_styles["body"].level_styles
        self.assertIn("1", body_levels)
        self.assertIsNotNone(body_levels["1"].font_size_pt)

        placeholders_with_levels = [
            placeholder
            for layout in analyzed_manifest.layouts
            for placeholder in layout.placeholders
            if placeholder.paragraph_styles is not None and placeholder.paragraph_styles.level_styles
        ]
        self.assertTrue(placeholders_with_levels)
        self.assertTrue(
            any(
                any(
                    style.bullet_type is not None
                    or style.hanging_emu is not None
                    or style.indent_emu is not None
                    or style.margin_right_emu is not None
                    for style in placeholder.paragraph_styles.level_styles.values()
                )
                for placeholder in placeholders_with_levels
            )
        )

        component_styles = [
            placeholder.shape_style
            for layout in analyzed_manifest.layouts
            for placeholder in layout.placeholders
            if placeholder.shape_style is not None
        ]
        self.assertTrue(component_styles)
        self.assertTrue(
            any(
                style.line_compound is not None
                or style.line_cap is not None
                or style.line_join is not None
                or style.theme_fill_ref is not None
                or style.theme_line_ref is not None
                or style.inset_left_emu is not None
                or len(style.effect_list) > 0
                for style in component_styles
            )
        )

    def test_template_analyzer_extracts_shape_geometry_for_uploaded_prototype_tokens(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists():
                continue

            analyzed_manifest = self.analyzer.analyze(
                template_id=manifest.template_id,
                template_path=template_path,
                display_name=manifest.display_name,
            )
            if analyzed_manifest.generation_mode.value != "prototype":
                continue

            with self.subTest(template_id=manifest.template_id):
                tokens = [
                    token
                    for prototype in analyzed_manifest.prototype_slides
                    for token in prototype.tokens
                ]
                self.assertTrue(tokens)
                self.assertTrue(any(self._has_geometry_metadata(token) for token in tokens))
                self.assertTrue(
                    any(
                        token.binding in {"title", "subtitle", "text", "body", "main_text", "secondary_text", "cover_title"}
                        and self._has_text_margin_metadata(token)
                        for token in tokens
                    )
                )

    def test_generator_applies_analyzer_geometry_metadata_for_uploaded_layout_templates(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Размещение согласий",
        )
        analyzed_manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in analyzed_manifest.layouts if "text" in layout.supported_slide_kinds)
        body_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.kind == PlaceholderKind.BODY)
        body_placeholder.left_emu = 777777
        body_placeholder.top_emu = 1888888
        body_placeholder.width_emu = 5555555
        body_placeholder.height_emu = 2222222
        body_placeholder.margin_left_emu = 11111
        body_placeholder.margin_right_emu = 22222
        body_placeholder.margin_top_emu = 33333
        body_placeholder.margin_bottom_emu = 44444

        plan = PresentationPlan(
            template_id=template_id,
            title="Uploaded Layout Geometry",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="",
                    text="Текст должен следовать analyzer-derived placeholder metadata.",
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=analyzed_manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == body_placeholder.idx)
        self.assertEqual(shape.left, body_placeholder.left_emu)
        self.assertEqual(shape.top, body_placeholder.top_emu)
        self.assertEqual(shape.width, body_placeholder.width_emu)
        self.assertEqual(shape.height, body_placeholder.height_emu)
        self.assertEqual(shape.text_frame.margin_left, body_placeholder.margin_left_emu)
        self.assertEqual(shape.text_frame.margin_right, body_placeholder.margin_right_emu)
        self.assertEqual(shape.text_frame.margin_top, body_placeholder.margin_top_emu)
        self.assertEqual(shape.text_frame.margin_bottom, body_placeholder.margin_bottom_emu)

    def test_deck_audit_uses_analyzer_geometry_metadata_for_uploaded_layout_templates(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        if not template_path.exists():
            self.skipTest(f"optional uploaded template is not installed: {template_id}")
        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Размещение согласий",
        )
        analyzed_manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in analyzed_manifest.layouts if "text" in layout.supported_slide_kinds)
        body_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.kind == PlaceholderKind.BODY)
        body_placeholder.left_emu = 777777
        body_placeholder.top_emu = 1888888
        body_placeholder.width_emu = 5555555
        body_placeholder.height_emu = 2222222
        body_placeholder.margin_left_emu = 11111
        body_placeholder.margin_right_emu = 22222
        body_placeholder.margin_top_emu = 33333
        body_placeholder.margin_bottom_emu = 44444

        plan = PresentationPlan(
            template_id=template_id,
            title="Uploaded Layout Geometry",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="",
                    text="Текст должен проходить template-aware deck audit на analyzer-derived layout metadata.",
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=analyzed_manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, analyzed_manifest)
            violations = find_capacity_violations(audits)

        self.assertTrue(audits)
        self.assertFalse(any(item.rule == "body_left_misalignment" for item in violations))
        self.assertFalse(any(item.rule == "body_margin_mismatch" for item in violations))
        self.assertLess(audits[0].profile.max_chars, TEXT_FULL_WIDTH_PROFILE.max_chars)
        self.assertTrue(audits[0].body_font_sizes)
        self.assertLessEqual(
            max(audits[0].body_font_sizes),
            max(audits[0].profile.max_font_pt, audits[0].expected_body_max_font_pt or 0),
        )

    def test_generator_applies_xml_derived_text_and_background_styles_from_manifest(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Размещение согласий",
        )
        analyzed_manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in analyzed_manifest.layouts if "text" in layout.supported_slide_kinds)
        if layout.background_style is None:
            layout.background_style = TemplateShapeStyleSpec()
        layout.background_style.fill_type = "solid"
        layout.background_style.fill_color = "#FFF4CC"

        body_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.kind == PlaceholderKind.BODY)
        if body_placeholder.text_style is None:
            body_placeholder.text_style = TemplateTextStyleSpec()
        if body_placeholder.shape_style is None:
            body_placeholder.shape_style = TemplateShapeStyleSpec()
        body_placeholder.text_style.font_size_pt = 19.0
        body_placeholder.text_style.font_family = "Arial"
        body_placeholder.text_style.color = "#CC3300"
        body_placeholder.text_style.bold = True
        body_placeholder.text_style.vertical_anchor = "ctr"
        body_placeholder.shape_style.fill_type = "solid"
        body_placeholder.shape_style.fill_color = "#DDEEFF"
        body_placeholder.shape_style.line_color = "#224466"
        body_placeholder.shape_style.line_width_pt = 1.5
        body_placeholder.shape_style.inset_left_emu = 54321
        body_placeholder.shape_style.inset_right_emu = 65432
        body_placeholder.shape_style.inset_top_emu = 76543
        body_placeholder.shape_style.inset_bottom_emu = 87654

        plan = PresentationPlan(
            template_id=template_id,
            title="XML Generator Styles",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="",
                    text="Generator должен применить XML-derived style metadata.",
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=analyzed_manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == body_placeholder.idx)
        self.assertEqual(str(slide.background.fill.fore_color.rgb), "FFF4CC")
        self.assertEqual(shape.text_frame.margin_left, 54321)
        self.assertEqual(shape.text_frame.margin_right, 65432)
        self.assertEqual(shape.text_frame.margin_top, 76543)
        self.assertEqual(shape.text_frame.margin_bottom, 87654)
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.runs[0]
        self.assertEqual(paragraph.text, "Generator должен применить XML-derived style metadata.")
        self.assertTrue(run.font.bold)
        xml = shape._element.xml
        self.assertIn('anchor="ctr"', xml)
        self.assertIn('a:latin typeface="Mont Regular"', xml)
        self.assertIn('a:ea typeface="Mont Regular"', xml)
        self.assertIn('a:cs typeface="Mont Regular"', xml)

    def test_generator_applies_manifest_background_xml_for_layout_templates(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        layout = next(layout for layout in manifest.layouts if layout.key == "text_full_width")
        layout.background_xml = (
            '<p:bg xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:bgPr>'
            '<a:solidFill><a:srgbClr val="445566"/></a:solidFill>'
            '<a:effectLst/>'
            '</p:bgPr>'
            '</p:bg>'
        )

        plan = PresentationPlan(
            template_id=template_id,
            title="Manifest Background XML",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    preferred_layout_key=layout.key,
                    text="Проверка фона из manifest background_xml.",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        self.assertIn('val="445566"', presentation.slides[0]._element.cSld.bg.xml)

    def test_generator_can_render_background_only_layout_slide_with_xml_override(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        layout = next(layout for layout in manifest.layouts if layout.key == "text_full_width")
        background_xml = (
            '<p:bg xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main" '
            'xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">'
            '<p:bgPr>'
            '<a:solidFill><a:srgbClr val="112233"/></a:solidFill>'
            '<a:effectLst/>'
            '</p:bgPr>'
            '</p:bg>'
        )

        plan = PresentationPlan(
            template_id=template_id,
            title="Background Only Layout",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    preferred_layout_key=layout.key,
                    background_only=True,
                    background_xml=background_xml,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        self.assertIn('val="112233"', slide._element.cSld.bg.xml)
        self.assertTrue(all(not (placeholder.text or "").strip() for placeholder in slide.placeholders))

    def test_generator_can_render_background_only_prototype_slide_without_bound_token_text(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        self.assertEqual(manifest.generation_mode, GenerationMode.PROTOTYPE)
        prototype = next(item for item in manifest.prototype_slides if item.tokens)

        plan = PresentationPlan(
            template_id=template_id,
            title="Background Only Prototype",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    preferred_layout_key=prototype.key,
                    background_only=True,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        token_shape_names = {token.shape_name for token in prototype.tokens if token.shape_name}
        for shape in slide.shapes:
            if shape.name not in token_shape_names or not getattr(shape, "has_text_frame", False):
                continue
            self.assertFalse((shape.text or "").strip())

    def test_generator_applies_xml_bullet_and_paragraph_spacing_from_manifest(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Размещение согласий",
        )
        analyzed_manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in analyzed_manifest.layouts if "bullets" in layout.supported_slide_kinds or "text" in layout.supported_slide_kinds)
        body_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.kind == PlaceholderKind.BODY)
        if body_placeholder.paragraph_styles is None:
            from a3presentation.domain.template import TemplateParagraphStyleCatalog
            body_placeholder.paragraph_styles = TemplateParagraphStyleCatalog(level_styles={})
        if body_placeholder.text_style is None:
            body_placeholder.text_style = TemplateTextStyleSpec()
        body_placeholder.paragraph_styles.level_styles["0"] = TemplateTextStyleSpec(
            bullet_type="char",
            bullet_char="■",
            bullet_font="Arial",
            hanging_emu=222222,
            margin_left_emu=333333,
            margin_right_emu=444444,
            default_tab_size_emu=555555,
            rtl=True,
            line_spacing=1.4,
            space_after_pt=7.0,
            font_size_pt=18.0,
            color="#116622",
        )

        plan = PresentationPlan(
            template_id=template_id,
            title="XML Bullet Styles",
            slides=[
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Пункты",
                    bullets=["Первый пункт", "Второй пункт"],
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=analyzed_manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == body_placeholder.idx)
        paragraph = next(item for item in shape.text_frame.paragraphs if item.text.strip())
        self.assertAlmostEqual(paragraph.line_spacing, 1.4, places=1)
        self.assertEqual(round(paragraph.space_after.pt), 7)
        run = paragraph.runs[0]
        self.assertEqual(round(run.font.size.pt), 18)
        self.assertEqual(str(run.font.color.rgb), "116622")
        xml = paragraph._p.xml
        self.assertIn('char="■"', xml)
        self.assertIn('typeface="Arial"', xml)
        self.assertIn('marL="333333"', xml)
        self.assertIn('marR="444444"', xml)
        self.assertIn('defTabSz="555555"', xml)
        self.assertIn('rtl="1"', xml)
        self.assertIn('indent="-222222"', xml)

    def test_generator_applies_xml_table_cell_margins_from_manifest(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in manifest.layouts if "table" in layout.supported_slide_kinds)
        table_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.binding == "table")
        if table_placeholder.shape_style is None:
            table_placeholder.shape_style = TemplateShapeStyleSpec()
        table_placeholder.shape_style.table_cell_margin_left_emu = 11111
        table_placeholder.shape_style.table_cell_margin_right_emu = 22222
        table_placeholder.shape_style.table_cell_margin_top_emu = 33333
        table_placeholder.shape_style.table_cell_margin_bottom_emu = 44444

        plan = PresentationPlan(
            template_id=template_id,
            title="XML Table Margins",
            slides=[
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Таблица",
                    table=TableBlock(
                        headers=["A", "B"],
                        header_fill_colors=[None, None],
                        rows=[["1", "2"]],
                        row_fill_colors=[[None, None]],
                    ),
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
        cell = table_shape.table.cell(0, 0)
        self.assertEqual(cell.margin_left, 11111)
        self.assertEqual(cell.margin_right, 22222)
        self.assertEqual(cell.margin_top, 33333)
        self.assertEqual(cell.margin_bottom, 44444)

    def test_generator_applies_component_style_table_contract(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        table_style = manifest.component_styles["table"]
        table_style.spacing_tokens["cell_margin_left_emu"] = 12345
        table_style.spacing_tokens["cell_margin_right_emu"] = 23456
        table_style.spacing_tokens["cell_margin_top_emu"] = 34567
        table_style.spacing_tokens["cell_margin_bottom_emu"] = 45678
        table_style.behavior_tokens["header_fill_color"] = "#ABCDEF"
        table_style.behavior_tokens["header_text_color"] = "#010203"
        table_style.behavior_tokens["border_color"] = "#020304"

        plan = PresentationPlan(
            template_id=template_id,
            title="Component Table Contract",
            slides=[
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Таблица",
                    table=TableBlock(
                        headers=["A", "B"],
                        header_fill_colors=[None, None],
                        rows=[["1", "2"]],
                        row_fill_colors=[[None, None]],
                    ),
                    preferred_layout_key="table",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
        cell = table_shape.table.cell(0, 0)
        self.assertEqual(cell.margin_left, 12345)
        self.assertEqual(cell.margin_right, 23456)
        self.assertEqual(cell.margin_top, 34567)
        self.assertEqual(cell.margin_bottom, 45678)
        cell_xml = cell._tc.xml
        self.assertIn('ABCDEF', cell_xml)
        self.assertIn('010203', cell_xml)
        self.assertIn('020304', cell_xml)

    def test_generator_applies_theme_fonts_to_non_title_text_layers(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx

        plan = PresentationPlan(
            template_id=template_id,
            title="Theme Font Coverage",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Заголовок",
                    subtitle="Подзаголовок",
                    text="Основной текст",
                    notes="Футер",
                    preferred_layout_key="text_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Список",
                    bullets=["Первый пункт", "Второй пункт"],
                    preferred_layout_key="list_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Титульный",
                    notes="Подпись титула",
                    preferred_layout_key="cover",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        text_slide = presentation.slides[0]
        title_shape = next(placeholder for placeholder in text_slide.placeholders if placeholder.placeholder_format.idx == 0)
        subtitle_shape = next(placeholder for placeholder in text_slide.placeholders if placeholder.placeholder_format.idx == 13)
        body_shape = next(placeholder for placeholder in text_slide.placeholders if placeholder.placeholder_format.idx == 14)
        footer_shape = next(placeholder for placeholder in text_slide.placeholders if placeholder.placeholder_format.idx == 17)
        self.assertIn('typeface="Mont SemiBold"', title_shape._element.xml)
        self.assertIn('typeface="Mont Regular"', subtitle_shape._element.xml)
        self.assertIn('typeface="Mont SemiBold"', body_shape._element.xml)
        self.assertIn('typeface="Mont Regular"', footer_shape._element.xml)

        list_slide = presentation.slides[1]
        list_body_shape = next(placeholder for placeholder in list_slide.placeholders if placeholder.placeholder_format.idx == 14)
        self.assertIn('typeface="Mont SemiBold"', list_body_shape._element.xml)

        cover_slide = presentation.slides[2]
        cover_xml = "".join(shape._element.xml for shape in cover_slide.shapes if getattr(shape, "has_text_frame", False))
        self.assertIn('typeface="Mont SemiBold"', cover_xml)
        self.assertIn('typeface="Mont Regular"', cover_xml)

    def test_generator_applies_theme_font_to_table_cell_text(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx

        plan = PresentationPlan(
            template_id=template_id,
            title="Theme Table Font",
            slides=[
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Таблица",
                    table=TableBlock(
                        headers=["A", "B"],
                        header_fill_colors=[None, None],
                        rows=[["1", "2"]],
                        row_fill_colors=[[None, None]],
                    ),
                    preferred_layout_key="table",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
        header_xml = table_shape.table.cell(0, 0)._tc.xml
        body_xml = table_shape.table.cell(1, 0)._tc.xml
        self.assertIn('typeface="Mont SemiBold"', header_xml)
        self.assertIn('typeface="Mont SemiBold"', body_xml)

    def test_generator_applies_xml_chart_offsets_from_manifest(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in manifest.layouts if "table" in layout.supported_slide_kinds)
        chart_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.binding == "table")
        if chart_placeholder.shape_style is None:
            chart_placeholder.shape_style = TemplateShapeStyleSpec()
        chart_placeholder.left_emu = 900000
        chart_placeholder.top_emu = 1500000
        chart_placeholder.width_emu = 6000000
        chart_placeholder.height_emu = 3200000
        chart_placeholder.shape_style.chart_plot_left_factor = 0.1
        chart_placeholder.shape_style.chart_plot_top_factor = 0.05
        chart_placeholder.shape_style.chart_plot_width_factor = 0.8
        chart_placeholder.shape_style.chart_plot_height_factor = 0.75
        chart_placeholder.shape_style.chart_legend_offset_x_emu = 250000
        chart_placeholder.shape_style.chart_legend_offset_y_emu = 120000
        chart_placeholder.shape_style.chart_category_axis_label_offset = 250
        chart_placeholder.shape_style.chart_value_axis_label_offset = 180

        base_left = chart_placeholder.left_emu
        base_top = chart_placeholder.top_emu
        base_width = chart_placeholder.width_emu
        base_height = chart_placeholder.height_emu

        plan = PresentationPlan(
            template_id=template_id,
            title="XML Chart Offsets",
            slides=[
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="График",
                    chart=ChartSpec(
                        chart_id="xml_offsets_chart",
                        source_table_id="xml_offsets_table",
                        title="Выручка",
                        chart_type=ChartType.COLUMN,
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(name="Выручка", values=[120.0, 200.0, 90.0]),
                            ChartSeries(name="EBITDA", values=[60.0, 80.0, 55.0]),
                        ],
                        legend_visible=True,
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        chart_shape = next(shape for shape in slide.shapes if getattr(shape, "has_chart", False))
        chart = chart_shape.chart
        self.assertEqual(chart_shape.left, base_left + int(base_width * 0.1))
        self.assertEqual(chart_shape.width, int(base_width * 0.8))
        plot_xml = chart._chartSpace.chart.plotArea.xml
        self.assertIn('c:xMode val="factor"', plot_xml)
        self.assertIn('c:yMode val="factor"', plot_xml)
        self.assertIn('c:w val="0.8"', plot_xml)
        self.assertIn('c:h val="0.75"', plot_xml)
        self.assertIn('c:lblOffset val="250"', chart._chartSpace.xml)
        self.assertIn('c:lblOffset val="180"', chart._chartSpace.xml)
        self.assertIn('c:y val="', chart._chartSpace.chart.legend.xml)

    def test_generator_applies_component_style_chart_palette(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        chart_style = manifest.component_styles["chart"]
        chart_style.behavior_tokens["rank_color_1"] = "#112233"
        chart_style.behavior_tokens["rank_color_2"] = "#223344"
        chart_style.behavior_tokens["rank_color_3"] = "#334455"
        chart_style.behavior_tokens["rank_color_4"] = "#445566"

        plan = PresentationPlan(
            template_id=template_id,
            title="Component Chart Contract",
            slides=[
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="График",
                    chart=ChartSpec(
                        chart_id="component_palette_chart",
                        source_table_id="component_palette_table",
                        title="Выручка",
                        chart_type=ChartType.COLUMN,
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(name="Выручка", values=[120.0, 200.0, 90.0]),
                            ChartSeries(name="EBITDA", values=[60.0, 80.0, 55.0]),
                        ],
                        legend_visible=True,
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        chart_shape = next(shape for shape in slide.shapes if getattr(shape, "has_chart", False))
        chart_xml = chart_shape.chart._chartSpace.xml
        self.assertIn("112233", chart_xml)
        self.assertIn("223344", chart_xml)

    def test_generator_applies_component_style_chart_geometry(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in manifest.layouts if "table" in layout.supported_slide_kinds)
        chart_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.binding == "table")
        chart_placeholder.shape_style = None
        chart_style = manifest.component_styles["chart"]
        chart_style.shape_style = TemplateShapeStyleSpec(
            role="chart",
            chart_plot_left_factor=0.12,
            chart_plot_top_factor=0.04,
            chart_plot_width_factor=0.72,
            chart_plot_height_factor=0.68,
        )
        base_left = chart_placeholder.left_emu
        base_top = chart_placeholder.top_emu
        base_width = chart_placeholder.width_emu
        base_height = chart_placeholder.height_emu

        plan = PresentationPlan(
            template_id=template_id,
            title="Component Chart Geometry",
            slides=[
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="График",
                    chart=ChartSpec(
                        chart_id="component_geometry_chart",
                        source_table_id="component_geometry_table",
                        title="Выручка",
                        chart_type=ChartType.COLUMN,
                        categories=["Q1", "Q2", "Q3"],
                        series=[ChartSeries(name="Выручка", values=[120.0, 200.0, 90.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        chart_shape = next(shape for shape in slide.shapes if getattr(shape, "has_chart", False))
        self.assertEqual(chart_shape.left, base_left + int(base_width * 0.12))
        self.assertEqual(chart_shape.top, base_top + int(base_height * 0.04))
        self.assertEqual(chart_shape.width, int(base_width * 0.72))
        self.assertEqual(chart_shape.height, int(base_height * 0.68))
        plot_xml = chart_shape.chart._chartSpace.chart.plotArea.xml
        self.assertIn('c:w val="0.72"', plot_xml)
        self.assertIn('c:h val="0.68"', plot_xml)

    def test_generator_applies_component_style_image_subtitle_size(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        image_style = manifest.component_styles["image"]
        image_style.text_styles["subtitle"].font_size_pt = 13.0

        plan = PresentationPlan(
            template_id=template_id,
            title="Component Image Contract",
            slides=[
                SlideSpec(
                    kind=SlideKind.IMAGE,
                    title="Схема процесса",
                    subtitle="Подпись к изображению",
                    text="Основной текст рядом с изображением.",
                    preferred_layout_key="image_text",
                    image_base64="iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9W6i8AAAAASUVORK5CYII=",
                    image_content_type="image/png",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        subtitle_shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == 13)
        self.assertIn('sz="1300"', subtitle_shape._element.xml)

    def test_generator_applies_component_style_image_geometry_tokens(self) -> None:
        template_id = "corp_light_v1"
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx
        image_style = manifest.component_styles["image"]
        image_style.spacing_tokens["min_image_height_emu"] = 1800000
        image_style.spacing_tokens["content_footer_gap_emu"] = 240000

        plan = PresentationPlan(
            template_id=template_id,
            title="Component Image Geometry",
            slides=[
                SlideSpec(
                    kind=SlideKind.IMAGE,
                    title="Схема процесса",
                    subtitle="Подпись к изображению",
                    text="Основной текст рядом с изображением.",
                    preferred_layout_key="image_text",
                    image_base64="iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9W6i8AAAAASUVORK5CYII=",
                    image_content_type="image/png",
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        image_shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == 16)
        self.assertGreaterEqual(image_shape.height, 1800000)

    def test_generator_applies_xml_line_style_and_theme_refs_from_manifest(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Размещение согласий",
        )
        analyzed_manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in analyzed_manifest.layouts if "text" in layout.supported_slide_kinds)
        body_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.kind == PlaceholderKind.BODY)
        if body_placeholder.shape_style is None:
            body_placeholder.shape_style = TemplateShapeStyleSpec()
        body_placeholder.shape_style.line_color = "#224466"
        body_placeholder.shape_style.line_compound = "thickThin"
        body_placeholder.shape_style.line_cap = "rnd"
        body_placeholder.shape_style.line_join = "bevel"
        body_placeholder.shape_style.theme_fill_ref = "3"
        body_placeholder.shape_style.theme_line_ref = "2"

        plan = PresentationPlan(
            template_id=template_id,
            title="XML Shape Line Style",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="",
                    text="Generator должен сохранять XML refs и advanced line style.",
                    preferred_layout_key=layout.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=analyzed_manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == body_placeholder.idx)
        xml = shape._element.xml
        self.assertIn('cmpd="thickThin"', xml)
        self.assertIn('cap="rnd"', xml)
        self.assertIn("<a:bevel/>", xml)
        self.assertIn('<a:fillRef idx="3">', xml)
        self.assertIn('<a:lnRef idx="2">', xml)

    def test_generator_applies_manifest_geometry_metadata_for_uploaded_prototype_templates(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        manifest = self.registry.get_template(template_id)
        prototype = manifest.prototype_slides[0]
        bound_text_token = next(
            token
            for token in prototype.tokens
            if token.binding in {"cover_meta", "subtitle", "text", "body", "main_text", "secondary_text", "presentation_name"}
        )
        bound_text_token.left_emu = 999999
        bound_text_token.top_emu = 2111111
        bound_text_token.width_emu = 4444444
        bound_text_token.height_emu = 1234567
        bound_text_token.margin_left_emu = 21000
        bound_text_token.margin_right_emu = 22000
        bound_text_token.margin_top_emu = 23000
        bound_text_token.margin_bottom_emu = 24000

        plan = PresentationPlan(
            template_id=template_id,
            title="Uploaded Prototype Geometry",
            slides=[
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Заголовок",
                    notes="Метаданные prototype token shape должны примениться к cloned slide.",
                    preferred_layout_key=prototype.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(shape for shape in slide.shapes if shape.name == bound_text_token.shape_name)
        self.assertEqual(shape.left, bound_text_token.left_emu)
        self.assertEqual(shape.top, bound_text_token.top_emu)
        self.assertEqual(shape.width, bound_text_token.width_emu)
        self.assertEqual(shape.height, bound_text_token.height_emu)
        self.assertEqual(shape.text_frame.margin_left, bound_text_token.margin_left_emu)
        self.assertEqual(shape.text_frame.margin_right, bound_text_token.margin_right_emu)
        self.assertEqual(shape.text_frame.margin_top, bound_text_token.margin_top_emu)
        self.assertEqual(shape.text_frame.margin_bottom, bound_text_token.margin_bottom_emu)

    def test_generator_resolves_uploaded_layout_from_render_target_without_preferred_layout_key(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        analyzed_manifest = self.analyzer.analyze(
            template_id=template_id,
            template_path=template_path,
            display_name="Размещение согласий",
        )
        analyzed_manifest.generation_mode = GenerationMode.LAYOUT
        layout = next(layout for layout in analyzed_manifest.layouts if "text" in layout.supported_slide_kinds)
        body_placeholder = next(placeholder for placeholder in layout.placeholders if placeholder.kind == PlaceholderKind.BODY)
        body_placeholder.left_emu = 777777
        body_placeholder.top_emu = 1888888
        body_placeholder.width_emu = 5555555
        body_placeholder.height_emu = 2222222

        plan = PresentationPlan(
            template_id=template_id,
            title="Uploaded Layout Render Target",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="",
                    text="Generator должен использовать render_target для uploaded layout path.",
                    preferred_layout_key=None,
                    render_target=SlideRenderTarget(
                        type=RenderTargetType.LAYOUT,
                        key=layout.key,
                    ),
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=analyzed_manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(placeholder for placeholder in slide.placeholders if placeholder.placeholder_format.idx == body_placeholder.idx)
        self.assertEqual(shape.left, body_placeholder.left_emu)
        self.assertEqual(shape.top, body_placeholder.top_emu)
        self.assertEqual(shape.width, body_placeholder.width_emu)
        self.assertEqual(shape.height, body_placeholder.height_emu)

    def test_generator_resolves_uploaded_prototype_from_render_target_without_preferred_layout_key(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        manifest = self.registry.get_template(template_id)
        prototype = manifest.prototype_slides[0]
        bound_text_token = next(
            token
            for token in prototype.tokens
            if token.binding in {"cover_meta", "subtitle", "text", "body", "main_text", "secondary_text", "presentation_name"}
        )
        bound_text_token.left_emu = 999999
        bound_text_token.top_emu = 2111111
        bound_text_token.width_emu = 4444444
        bound_text_token.height_emu = 1234567

        plan = PresentationPlan(
            template_id=template_id,
            title="Uploaded Prototype Render Target",
            slides=[
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Заголовок",
                    notes="Generator должен использовать render_target для uploaded prototype path.",
                    preferred_layout_key=None,
                    render_target=SlideRenderTarget(
                        type=RenderTargetType.PROTOTYPE,
                        key=prototype.key,
                    ),
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        shape = next(shape for shape in slide.shapes if shape.name == bound_text_token.shape_name)
        self.assertEqual(shape.left, bound_text_token.left_emu)
        self.assertEqual(shape.top, bound_text_token.top_emu)
        self.assertEqual(shape.width, bound_text_token.width_emu)
        self.assertEqual(shape.height, bound_text_token.height_emu)

    def test_generator_renders_direct_shape_binding_target_from_inventory(self) -> None:
        pptx = Presentation()
        slide = pptx.slides.add_slide(pptx.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(600000, 400000, 8000000, 900000)
        title_shape.name = "Title Box"
        body_shape = slide.shapes.add_textbox(600000, 1800000, 7600000, 2400000)
        body_shape.name = "Body Box"

        with tempfile.TemporaryDirectory() as temp_dir:
            template_path = Path(temp_dir) / "direct-binding-template.pptx"
            pptx.save(str(template_path))
            manifest = TemplateManifest(
                template_id="direct_binding_demo",
                display_name="Direct Binding Demo",
                source_pptx=template_path.name,
                inventory=ExtractedPresentationInventory(
                    degradation_mode="direct_shape_binding",
                    slides=[
                        ExtractedSlideInventory(
                            source_kind=InventorySourceKind.SLIDE,
                            source_index=0,
                            name="Direct Slide",
                            component_ids=["slide_0_title", "slide_0_body"],
                            supported_slide_kinds=["text"],
                            representation_hints=["text"],
                        )
                    ],
                    components=[
                        ExtractedComponent(
                            component_id="slide_0_title",
                            source_kind=InventorySourceKind.SLIDE,
                            source_index=0,
                            source_name="Direct Slide",
                            shape_name="Title Box",
                            component_type=ExtractedComponentType.TEXT,
                            role=ExtractedComponentRole.TITLE,
                            binding="title",
                            confidence=ComponentConfidence.MEDIUM,
                            editability=ComponentEditability.EDITABLE,
                            capabilities=["text"],
                            geometry=ComponentGeometry(left_emu=600000, top_emu=400000, width_emu=8000000, height_emu=900000),
                        ),
                        ExtractedComponent(
                            component_id="slide_0_body",
                            source_kind=InventorySourceKind.SLIDE,
                            source_index=0,
                            source_name="Direct Slide",
                            shape_name="Body Box",
                            component_type=ExtractedComponentType.TEXT,
                            role=ExtractedComponentRole.BODY,
                            binding="body",
                            confidence=ComponentConfidence.MEDIUM,
                            editability=ComponentEditability.EDITABLE,
                            capabilities=["text"],
                            geometry=ComponentGeometry(left_emu=600000, top_emu=1800000, width_emu=7600000, height_emu=2400000),
                        ),
                    ],
                ),
            )
            plan = PresentationPlan(
                template_id="direct_binding_demo",
                title="Direct Binding Demo",
                slides=[
                    SlideSpec(
                        kind=SlideKind.TEXT,
                        title="Direct title",
                        text="Direct body text",
                        render_target=SlideRenderTarget(
                            type=RenderTargetType.DIRECT_SHAPE_BINDING,
                            key="direct_slide_0",
                        ),
                    )
                ],
            )

            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            audits = audit_generated_presentation(output_path, plan, manifest)

        output_slide = presentation.slides[0]
        rendered_title = next(shape for shape in output_slide.shapes if shape.name == "Title Box")
        rendered_body = next(shape for shape in output_slide.shapes if shape.name == "Body Box")
        self.assertIn("Direct title", rendered_title.text)
        self.assertIn("Direct body text", rendered_body.text)
        self.assertEqual(audits[0].target_type, "direct_shape_binding")
        self.assertTrue(audits[0].degraded_but_valid)
        self.assertEqual(audits[0].target_degradation_reasons, ())

    def test_deck_audit_uses_manifest_geometry_metadata_for_uploaded_prototype_templates(self) -> None:
        template_id = "razmeshchenie_soglasiy"
        template_path = self.settings.templates_dir / template_id / "template.pptx"
        if not template_path.exists():
            self.skipTest(f"optional uploaded template is not installed: {template_id}")
        manifest = self.registry.get_template(template_id)
        prototype = next(
            item for item in manifest.prototype_slides if any(token.binding == "main_text" for token in item.tokens)
        )
        body_token = next(token for token in prototype.tokens if token.binding == "main_text")
        body_token.left_emu = 999999
        body_token.top_emu = 2111111
        body_token.width_emu = 4444444
        body_token.height_emu = 1234567
        body_token.margin_left_emu = 21000
        body_token.margin_right_emu = 22000
        body_token.margin_top_emu = 23000
        body_token.margin_bottom_emu = 24000

        plan = PresentationPlan(
            template_id=template_id,
            title="Uploaded Prototype Audit Geometry",
            slides=[
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Заголовок",
                    subtitle="Подзаголовок",
                    text="Template-aware deck audit должен читать body geometry и margins из prototype token metadata.",
                    notes="Служебная строка для prototype footer.",
                    preferred_layout_key=prototype.key,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)
            violations = find_capacity_violations(audits)

        self.assertTrue(audits)
        self.assertFalse(any(item.rule == "body_left_misalignment" for item in violations))
        self.assertFalse(any(item.rule == "body_margin_mismatch" for item in violations))
        self.assertEqual(audits[0].footer_placeholder_idx, 17)
        self.assertFalse(any(item.rule == "narrow_footer" for item in violations))
        self.assertFalse(any(item.rule == "footer_left_misalignment" for item in violations))
        self.assertLess(audits[0].profile.max_chars, TEXT_FULL_WIDTH_PROFILE.max_chars)
        self.assertTrue(audits[0].body_font_sizes)
        self.assertLessEqual(
            max(audits[0].body_font_sizes),
            max(audits[0].profile.max_font_pt, audits[0].expected_body_max_font_pt or 0),
        )

    def test_capacity_profile_derivation_reduces_limits_for_narrower_body_geometry(self) -> None:
        reference_body = geometry_policy_for_layout("text_full_width").placeholders[14]
        derived = derive_capacity_profile_for_geometry(
            "text_full_width",
            width_emu=int(reference_body.width_emu * 0.6),
            height_emu=int(reference_body.height_emu * 0.7),
        )

        self.assertLess(derived.max_chars, TEXT_FULL_WIDTH_PROFILE.max_chars)
        self.assertLessEqual(derived.max_font_pt, TEXT_FULL_WIDTH_PROFILE.max_font_pt)

    def test_capacity_profile_derivation_accepts_explicit_reference_geometry(self) -> None:
        derived = derive_capacity_profile_for_geometry(
            "text_full_width",
            width_emu=3600000,
            height_emu=1800000,
            reference_width_emu=6000000,
            reference_height_emu=3000000,
        )

        self.assertLess(derived.max_chars, TEXT_FULL_WIDTH_PROFILE.max_chars)
        self.assertLess(derived.max_items, TEXT_FULL_WIDTH_PROFILE.max_items)
        self.assertLessEqual(derived.max_font_pt, TEXT_FULL_WIDTH_PROFILE.max_font_pt)

    def test_full_pipeline_contract_for_mixed_docx_document(self) -> None:
        document = Document()
        document.add_paragraph("A3")
        document.add_paragraph("Смешанный стратегический документ")
        document.add_heading("1. Основные выводы", level=1)
        document.add_paragraph(
            "Рост выручки обеспечивается за счет новых сегментов, улучшения конверсии и масштабируемой платформы."
        )
        document.add_paragraph(style="List Bullet").add_run("Приоритет 1: рост B2B сегмента")
        document.add_paragraph(style="List Bullet").add_run("Приоритет 2: снижение концентрации выручки")
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Показатель"
        table.cell(0, 1).text = "Значение"
        table.cell(1, 0).text = "GMV"
        table.cell(1, 1).text = "125"
        table.cell(2, 0).text = "NPS"
        table.cell(2, 1).text = "61"
        document.add_heading("2. Следующий раздел", level=1)
        document.add_paragraph("Дополнительный narrative-блок для устойчивой классификации документа.")

        buffer = BytesIO()
        document.save(buffer)
        content = buffer.getvalue()

        text, tables, blocks = self.extractor.extract("mixed-contract.docx", content)
        plan = self.planner.build_plan("corp_light_v1", text, None, tables, blocks)

        self.assertGreaterEqual(len(plan.slides), 3)
        self.assertTrue(any(slide.kind == SlideKind.TABLE for slide in plan.slides))
        self.assertTrue(any(slide.kind in {SlideKind.TEXT, SlideKind.BULLETS} for slide in plan.slides[1:]))

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), len(plan.slides))

    def test_full_pipeline_contract_for_uploaded_prototype_template_keeps_template_aware_audit_green(self) -> None:
        document = Document()
        document.add_paragraph("A3")
        document.add_paragraph("Стратегический документ для пользовательского шаблона")
        document.add_heading("1. Основные выводы", level=1)
        document.add_paragraph(
            "Рост выручки обеспечивается платформенными интеграциями, улучшением конверсии и управляемым контуром сопровождения."
        )
        document.add_paragraph(
            "Следующий абзац нужен для text-flow сценария и проверки того, что template-aware audit читает реальную геометрию пользовательского prototype template."
        )
        document.add_heading("2. Таблица метрик", level=1)
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Метрика"
        table.cell(0, 1).text = "Значение"
        table.cell(1, 0).text = "GMV"
        table.cell(1, 1).text = "125"
        table.cell(2, 0).text = "NPS"
        table.cell(2, 1).text = "61"

        buffer = BytesIO()
        document.save(buffer)
        content = buffer.getvalue()

        template_id = "razmeshchenie_soglasiy"
        if not (self.settings.templates_dir / template_id / "template.pptx").exists():
            self.skipTest(f"optional uploaded template is not installed: {template_id}")
        manifest = self.registry.get_template(template_id)
        template_path = self.settings.templates_dir / template_id / manifest.source_pptx

        text, tables, blocks = self.extractor.extract("uploaded-prototype-contract.docx", content)
        plan = self.planner.build_plan(template_id, text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)
            violations = find_capacity_violations(audits)

        self.assertEqual(len(audits), len([slide for slide in plan.slides if slide.kind in {
            SlideKind.TEXT,
            SlideKind.BULLETS,
            SlideKind.TWO_COLUMN,
            SlideKind.TABLE,
            SlideKind.CHART,
            SlideKind.IMAGE,
        }]))
        self.assertFalse(any(item.rule == "body_margin_mismatch" for item in violations))
        self.assertFalse(any(item.rule == "body_left_misalignment" for item in violations))

    def test_api_roundtrip_contract_keeps_extract_plan_generate_download_in_sync(self) -> None:
        document = Document()
        document.add_paragraph("Стратегический отчёт")
        document.add_heading("Рынок", level=1)
        document.add_paragraph("Компания усиливает позиции за счет платформенных интеграций.")
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Сегмент"
        table.cell(0, 1).text = "Выручка"
        table.cell(1, 0).text = "SMB"
        table.cell(1, 1).text = "120"
        table.cell(2, 0).text = "Enterprise"
        table.cell(2, 1).text = "250"

        buffer = BytesIO()
        document.save(buffer)
        upload = UploadFile(filename="api-contract.docx", file=BytesIO(buffer.getvalue()))

        extracted = asyncio.run(routes_module.extract_document_text(upload))
        self.assertEqual(len(extracted.tables), 1)
        self.assertEqual(len(extracted.chart_assessments), 1)

        plan = routes_module.plan_from_text(
            TextPlanRequest(
                template_id="corp_light_v1",
                title="API Contract",
                raw_text=extracted.text,
                tables=extracted.tables,
                blocks=extracted.blocks,
            )
        )
        self.assertGreaterEqual(len(plan.slides), 2)

        generated = routes_module.generate_presentation(plan)
        downloaded = routes_module.download_presentation(generated.file_name)
        self.assertEqual(Path(downloaded.path).name, generated.file_name)

    def test_layout_capacity_profiles_are_consistent_with_planner_contract(self) -> None:
        self.assertEqual(profile_for_layout("list_full_width"), LIST_FULL_WIDTH_PROFILE)
        self.assertEqual(profile_for_layout("text_full_width"), TEXT_FULL_WIDTH_PROFILE)
        self.assertEqual(self.planner.list_profile.max_items, self.planner.LIST_BATCH_SIZE)
        self.assertEqual(self.planner.list_profile.max_weight, self.planner.LIST_SLIDE_MAX_WEIGHT)
        self.assertEqual(self.planner.text_profile.max_chars, self.planner.TEXT_SLIDE_MAX_CHARS)
        self.assertEqual(self.planner.text_profile.max_primary_chars, self.planner.TEXT_PRIMARY_MAX_CHARS)
        self.assertLessEqual(self.planner.text_profile.min_font_pt, self.planner.text_profile.max_font_pt)
        self.assertLessEqual(self.planner.list_profile.min_font_pt, self.planner.list_profile.max_font_pt)

    def test_layout_geometry_policies_keep_full_width_contracts(self) -> None:
        text_policy = geometry_policy_for_layout("text_full_width")
        list_policy = geometry_policy_for_layout("list_full_width")
        table_policy = geometry_policy_for_layout("table")
        image_policy = geometry_policy_for_layout("image_text")
        cards_policy = geometry_policy_for_layout("cards_3")
        kpi_cards_policy = geometry_policy_for_layout("cards_kpi")
        icons_policy = geometry_policy_for_layout("list_with_icons")
        contacts_policy = geometry_policy_for_layout("contacts")

        self.assertEqual(text_policy.placeholders[0].width_emu, 11198224)
        self.assertEqual(text_policy.placeholders[14].width_emu, 11198224)
        self.assertEqual(text_policy.placeholders[17].top_emu, 6384626)
        self.assertEqual(list_policy.placeholders[14].width_emu, text_policy.placeholders[14].width_emu)
        self.assertEqual(table_policy.placeholders[15].width_emu, 11198224)
        self.assertEqual(image_policy.placeholders[16].width_emu, 4990840)
        self.assertEqual(cards_policy.placeholders[11].width_emu, cards_policy.placeholders[12].width_emu)
        self.assertEqual(kpi_cards_policy.placeholders[11].top_emu, kpi_cards_policy.placeholders[12].top_emu)
        self.assertGreater(kpi_cards_policy.placeholders[13].top_emu, kpi_cards_policy.placeholders[11].top_emu)
        self.assertEqual(icons_policy.placeholders[21].top_emu, 6384626)
        self.assertEqual(contacts_policy.placeholders[10].width_emu, 3724275)

    def test_layout_spacing_policies_keep_bullet_indent_contracts(self) -> None:
        text_spacing = spacing_policy_for_layout("text_full_width")
        list_spacing = spacing_policy_for_layout("list_full_width")

        self.assertEqual(text_spacing.bullet.margin_left_emu, 342900)
        self.assertEqual(text_spacing.bullet.indent_emu, -171450)
        self.assertEqual(text_spacing.body.line_spacing, 1.1)
        self.assertEqual(text_spacing.body.space_after_pt, 6.0)
        self.assertGreater(list_spacing.bullet.margin_left_emu, text_spacing.bullet.margin_left_emu)
        self.assertLess(list_spacing.bullet.indent_emu, text_spacing.bullet.indent_emu)

    def test_deck_audit_reports_body_font_sizes_within_layout_profile_bounds(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Contract",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Contract", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Dense bullets",
                    bullets=[
                        "Первый длинный пункт объясняет стратегическую логику и ограничения.",
                        "Второй длинный пункт добавляет риски, сроки и организационные последствия.",
                        "Третий длинный пункт описывает KPI, unit-экономику и инфраструктурные требования.",
                        "Четвёртый длинный пункт связывает выводы с дорожной картой внедрения.",
                    ],
                    preferred_layout_key="list_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Dense text",
                    text=(
                        "Первый длинный абзац описывает контекст, ограничения, допущения и критерии принятия решения. "
                        "Второй длинный абзац добавляет финансовые ориентиры и риски реализации. "
                        "Третий длинный абзац связывает выводы с KPI и дорожной картой."
                    ),
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)

        self.assertGreaterEqual(len(audits), 2)
        for audit in audits:
            with self.subTest(slide=audit.slide_index):
                self.assertTrue(audit.body_font_sizes)
                self.assertGreaterEqual(min(audit.body_font_sizes), audit.profile.min_font_pt)
                self.assertLessEqual(
                    max(audit.body_font_sizes),
                    max(audit.profile.max_font_pt, audit.expected_body_max_font_pt or 0),
                )

    def test_deck_audit_detects_continuation_groups_for_multislide_sections(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Continuations",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Continuations", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Раздел",
                    bullets=["Пункт 1", "Пункт 2", "Пункт 3", "Пункт 4", "Пункт 5", "Пункт 6"],
                    preferred_layout_key="list_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Раздел (2)",
                    bullets=["Пункт 7", "Пункт 8"],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        groups = continuation_groups(audits)
        self.assertIn("Раздел", groups)
        self.assertEqual(len(groups["Раздел"]), 2)

    def test_deck_audit_flags_underfilled_continuation_pairs(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Violations",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Violations", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Раздел",
                    bullets=[
                        "Первый длинный пункт подробно описывает стратегическую инициативу и контекст принятия решения.",
                        "Второй длинный пункт раскрывает риски, ресурсы, ограничения и организационные зависимости.",
                        "Третий длинный пункт связывает инициативу с метриками, сроками и критериями качества.",
                        "Четвёртый длинный пункт объясняет изменения процесса и требования к операционной модели.",
                        "Пятый длинный пункт добавляет детали по продукту, рынку и коммерческому эффекту.",
                        "Шестой длинный пункт завершает блок ожидаемыми результатами и контрольными точками.",
                    ],
                    preferred_layout_key="list_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Раздел (2)",
                    bullets=["Короткий хвост."],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        violations = find_capacity_violations(audits)
        violation_rules = {violation.rule for violation in violations}
        self.assertIn("continuation_balance", violation_rules)
        self.assertIn("underfilled_continuation", violation_rules)

    def test_table_overlay_text_must_fit_visible_cells(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Dense Table Quality",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Dense Table Quality", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Плотная таблица",
                    table=TableBlock(
                        headers=["Риск", "Описание", "Митигация"],
                        rows=[
                            [
                                "Партнерская зависимость",
                                "Очень длинное описание риска, которое должно переноситься внутри видимой ячейки без наложения на соседние строки и без выхода за границы таблицы.",
                                "Сценарное планирование, договорные SLA, резервные каналы и еженедельный мониторинг статуса с понятной ответственностью.",
                            ],
                            [
                                "Инфраструктура",
                                "Длинное описание инфраструктурного ограничения с несколькими условиями запуска и зависимостями от внешних команд.",
                                "Пошаговая миграция, контрольные точки, fallback-план и замер фактической нагрузки перед масштабированием.",
                            ],
                            [
                                "Коммерческий риск",
                                "Длинный текст должен оставаться читаемым, а генератор обязан уменьшить шрифт и распределить высоту строк до выдачи файла пользователю.",
                                "Единый quality gate проверяет реальные overlay-ячейки после генерации PowerPoint.",
                            ],
                        ],
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            manifest = self.registry.get_template("corp_light_v1")
            output_path = PptxGenerator().generate(
                template_path=self.registry.get_template_pptx_path("corp_light_v1"),
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)

        table_audit = next(audit for audit in audits if audit.kind == SlideKind.TABLE.value)
        self.assertEqual(table_audit.table_overlay_overflow_count, 0)
        self.assertNotIn("table_overlay_text_overflow", {violation.rule for violation in find_capacity_violations(audits)})

    def test_deck_audit_does_not_flag_borderline_underfilled_continuation(self) -> None:
        audits = [
            SlideAudit(
                slide_index=1,
                title="Раздел",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=458,
                body_font_sizes=(14.0,),
                profile=TEXT_FULL_WIDTH_PROFILE,
            ),
            SlideAudit(
                slide_index=2,
                title="Раздел (2)",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=291,
                body_font_sizes=(14.0,),
                profile=TEXT_FULL_WIDTH_PROFILE,
            ),
        ]

        violations = find_capacity_violations(audits)

        violation_rules = {violation.rule for violation in violations}
        self.assertNotIn("underfilled_continuation", violation_rules)
        self.assertNotIn("continuation_balance", violation_rules)

    def test_deck_audit_keeps_balance_violation_for_materially_underfilled_group(self) -> None:
        audits = [
            SlideAudit(
                slide_index=1,
                title="Раздел",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=720,
                body_font_sizes=(14.0,),
                profile=LIST_FULL_WIDTH_PROFILE,
            ),
            SlideAudit(
                slide_index=2,
                title="Раздел (2)",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=270,
                body_font_sizes=(14.0,),
                profile=LIST_FULL_WIDTH_PROFILE,
            ),
        ]

        violations = find_capacity_violations(audits)

        violation_rules = {violation.rule for violation in violations}
        self.assertIn("underfilled_continuation", violation_rules)
        self.assertIn("continuation_balance", violation_rules)

    def test_deck_audit_does_not_flag_borderline_bullet_tail(self) -> None:
        audits = [
            SlideAudit(
                slide_index=1,
                title="Раздел",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=507,
                body_font_sizes=(14.0,),
                profile=LIST_FULL_WIDTH_PROFILE,
            ),
            SlideAudit(
                slide_index=2,
                title="Раздел (2)",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=474,
                body_font_sizes=(14.0,),
                profile=LIST_FULL_WIDTH_PROFILE,
            ),
        ]

        violations = find_capacity_violations(audits)

        self.assertNotIn("underfilled_continuation", {violation.rule for violation in violations})

    def test_deck_audit_keeps_expected_bullet_order_for_mixed_slide(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Order",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Order", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Смешанный блок",
                    bullets=[
                        "Вводный абзац задает контекст.",
                        "Первый тезис фиксирует массовый каталог.",
                        "Второй тезис фиксирует SLA-контур.",
                        "Финальный абзац завершает аргументацию.",
                    ],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        mixed_audit = next(audit for audit in audits if audit.title == "Смешанный блок")
        self.assertEqual(
            tuple(item.strip() for item in mixed_audit.rendered_items if item.strip()),
            tuple(item.strip() for item in mixed_audit.expected_items if item.strip()),
        )
        violations = find_capacity_violations(audits)
        self.assertNotIn("content_order_mismatch", {violation.rule for violation in violations})

    def test_deck_audit_uses_content_blocks_for_mixed_paragraph_and_list_order(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Mixed Blocks",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Mixed Blocks", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Смешанный контейнер",
                    bullets=[
                        "Вводный абзац задает контекст.",
                        "Первый тезис фиксирует массовый каталог.",
                        "Второй тезис фиксирует SLA-контур.",
                        "Финальный абзац завершает аргументацию.",
                    ],
                    content_blocks=[
                        SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text="Вводный абзац задает контекст."),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.BULLET_LIST,
                            items=[
                                "Первый тезис фиксирует массовый каталог.",
                                "Второй тезис фиксирует SLA-контур.",
                            ],
                        ),
                        SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text="Финальный абзац завершает аргументацию."),
                    ],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        mixed_audit = next(audit for audit in audits if audit.title == "Смешанный контейнер")
        self.assertEqual(
            tuple(item.strip() for item in mixed_audit.rendered_items if item.strip()),
            (
                "Вводный абзац задает контекст.",
                "Первый тезис фиксирует массовый каталог.",
                "Второй тезис фиксирует SLA-контур.",
                "Финальный абзац завершает аргументацию.",
            ),
        )
        violations = find_capacity_violations(audits)
        self.assertNotIn("content_order_mismatch", {violation.rule for violation in violations})

    def test_template_binding_keeps_content_blocks_for_text_slide_with_list_layout(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Binding Content Blocks",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Binding Content Blocks", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Смешанный narrative",
                    text="Вводный абзац задает контекст.",
                    content_blocks=[
                        SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text="Вводный абзац задает контекст."),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.BULLET_LIST,
                            items=[
                                "Первый тезис описывает охват каталога.",
                                "Второй тезис описывает SLA-контур.",
                            ],
                        ),
                    ],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        target_audit = next(audit for audit in audits if audit.title == "Смешанный narrative")
        self.assertEqual(
            tuple(item.strip() for item in target_audit.rendered_items if item.strip()),
            (
                "Вводный абзац задает контекст.",
                "Первый тезис описывает охват каталога.",
                "Второй тезис описывает SLA-контур.",
            ),
        )
        violations = find_capacity_violations(audits)
        self.assertNotIn("content_order_mismatch", {violation.rule for violation in violations})

    def test_template_binding_does_not_duplicate_notes_when_content_blocks_fill_body(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="No Duplicate Tail",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="No Duplicate Tail", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Narrative",
                    text="Первый абзац задает контекст.",
                    notes="Хвост не должен дублироваться отдельным placeholder.",
                    content_blocks=[
                        SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text="Первый абзац задает контекст."),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.PARAGRAPH,
                            text="Хвост не должен дублироваться отдельным placeholder.",
                        ),
                    ],
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        self.assertIn(14, placeholders)
        self.assertIn("Хвост не должен дублироваться", placeholders[14].text)
        self.assertNotIn(15, placeholders)

    def test_template_binding_clears_duplicate_subtitle_when_body_already_starts_with_it(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="No Duplicate Subtitle",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="No Duplicate Subtitle", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Narrative",
                    subtitle="Первый абзац задает контекст.",
                    text="Первый абзац задает контекст.",
                    content_blocks=[
                        SlideContentBlock(
                            kind=SlideContentBlockKind.PARAGRAPH,
                            text="Первый абзац задает контекст. Дальше идет основной текст без отдельного subtitle-повтора.",
                        ),
                    ],
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        } 
        self.assertNotIn(13, placeholders)

    def test_generator_clears_duplicate_subtitle_for_plain_text_slide(self) -> None:
        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="No Plain Duplicate Subtitle",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="No Plain Duplicate Subtitle", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Narrative",
                    subtitle="Первый абзац задает контекст и открывает раздел.",
                    text="Первый абзац задает контекст и открывает раздел. Дальше идет основной narrative без отдельного подзаголовка.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        placeholders = {
            shape.placeholder_format.idx: shape
            for shape in slide.placeholders
            if getattr(shape, "is_placeholder", False)
        }
        self.assertNotIn(13, placeholders)
        self.assertIn(14, placeholders)

    def test_deck_audit_ignores_text_only_continuations_when_tracking_group_order(self) -> None:
        audits = [
            SlideAudit(
                slide_index=2,
                title="Группа",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=320,
                body_font_sizes=(16.0,),
                profile=profile_for_layout("text_full_width"),
                expected_items=(),
                rendered_items=("Вводный narrative блок.",),
            ),
            SlideAudit(
                slide_index=3,
                title="Группа (2)",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=220,
                body_font_sizes=(14.0,),
                profile=profile_for_layout("list_full_width"),
                expected_items=("Первый тезис.", "Второй тезис."),
                rendered_items=("Первый тезис.", "Второй тезис."),
            ),
            SlideAudit(
                slide_index=4,
                title="Группа (3)",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=280,
                body_font_sizes=(16.0,),
                profile=profile_for_layout("text_full_width"),
                expected_items=(),
                rendered_items=("Финальный narrative блок.",),
            ),
        ]

        violations = find_capacity_violations(audits)
        self.assertNotIn("continuation_order_mismatch", {violation.rule for violation in violations})

    def test_deck_audit_accepts_question_and_appendix_style_slide_without_order_violations(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Question Appendix",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Question Appendix", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Частые вопросы и выводы",
                    bullets=[
                        "Почему нужен второй контур инфраструктуры?",
                        "Важно: без резервного контура SLA для критичных платежей останется хрупким.",
                        "Итог: резервный контур снижает риск простоя и упрощает масштабирование.",
                    ],
                    content_blocks=[
                        SlideContentBlock(kind=SlideContentBlockKind.QA_ITEM, text="Почему нужен второй контур инфраструктуры?"),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.CALLOUT,
                            text="Важно: без резервного контура SLA для критичных платежей останется хрупким.",
                        ),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.CALLOUT,
                            text="Итог: резервный контур снижает риск простоя и упрощает масштабирование.",
                        ),
                    ],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_deck_audit_accepts_semantic_text_layout_for_question_and_callout_blocks(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Semantic Text Layout",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Semantic Text Layout", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="FAQ и выводы",
                    text="Вопрос: почему нужен второй контур инфраструктуры?",
                    notes=(
                        "Важно: резервный контур удерживает SLA в часы пиковой нагрузки. "
                        "Итог: архитектура снижает риск простоя без усложнения операционной модели."
                    ),
                    content_blocks=[
                        SlideContentBlock(
                            kind=SlideContentBlockKind.QA_ITEM,
                            text="Вопрос: почему нужен второй контур инфраструктуры?",
                        ),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.CALLOUT,
                            text="Важно: резервный контур удерживает SLA в часы пиковой нагрузки.",
                        ),
                        SlideContentBlock(
                            kind=SlideContentBlockKind.CALLOUT,
                            text="Итог: архитектура снижает риск простоя без усложнения операционной модели.",
                        ),
                    ],
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertEqual(len(audits), 1)
        self.assertTrue(audits[0].body_font_sizes)
        violations = find_capacity_violations(audits)
        self.assertNotIn("content_order_mismatch", {violation.rule for violation in violations})

    def test_manifests_keep_generic_text_or_bullet_layout_for_semantic_blocks(self) -> None:
        manifests = self.registry.list_templates()
        self.assertGreaterEqual(len(manifests), 1)

        for manifest in manifests:
            layout_keys = {layout.key for layout in manifest.layouts}
            supported_slide_kinds = {
                supported_kind
                for layout in manifest.layouts
                for supported_kind in layout.supported_slide_kinds
            }
            with self.subTest(template_id=manifest.template_id):
                if {"text", "bullets"} & supported_slide_kinds:
                    self.assertTrue(
                        any(
                            layout_key in layout_keys
                            for layout_key in {"text_full_width", "list_full_width", "content", "слайд_с_перечислением"}
                        )
                    )

    def test_deck_audit_accepts_long_title_and_dense_body_without_capacity_regression(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Long Title Dense Body",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Long Title Dense Body", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title=(
                        "Очень длинный заголовок, который одновременно фиксирует стратегическую развилку, "
                        "операционные ограничения и требования к инфраструктуре платежной платформы"
                    ),
                    text=(
                        "Первый плотный абзац объясняет, почему устойчивость SLA зависит от второго контура, "
                        "какие ограничения есть у текущей архитектуры и как это влияет на масштабирование продукта. "
                        "Второй плотный абзац связывает инфраструктурные решения с коммерческой логикой, "
                        "переговорной позицией с банками и требованиями к государственным интеграциям."
                    ),
                    notes=(
                        "Третий плотный абзац фиксирует риски внедрения, требования к операционной модели и "
                        "критерии, по которым можно считать архитектурное решение успешным после запуска."
                    ),
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_deck_audit_flags_continuation_order_mismatch(self) -> None:
        audits = [
            SlideAudit(
                slide_index=2,
                title="Раздел",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=120,
                body_font_sizes=(14.0,),
                profile=profile_for_layout("list_full_width"),
                expected_items=("Первый", "Второй"),
                rendered_items=("Второй", "Первый"),
            ),
            SlideAudit(
                slide_index=3,
                title="Раздел (2)",
                kind=SlideKind.BULLETS.value,
                layout_key="list_full_width",
                body_char_count=120,
                body_font_sizes=(14.0,),
                profile=profile_for_layout("list_full_width"),
                expected_items=("Третий",),
                rendered_items=("Третий",),
            ),
        ]

        violations = find_capacity_violations(audits)
        rules = {violation.rule for violation in violations}
        self.assertIn("content_order_mismatch", rules)
        self.assertIn("continuation_order_mismatch", rules)

    def test_deck_audit_flags_large_font_jump_between_continuation_slides(self) -> None:
        audits = [
            SlideAudit(
                slide_index=2,
                title="Раздел",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=320,
                body_font_sizes=(16.0,),
                profile=profile_for_layout("text_full_width"),
                expected_items=("Первый блок",),
                rendered_items=("Первый блок",),
            ),
            SlideAudit(
                slide_index=3,
                title="Раздел (2)",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=300,
                body_font_sizes=(12.0,),
                profile=profile_for_layout("text_full_width"),
                expected_items=("Второй блок",),
                rendered_items=("Второй блок",),
            ),
        ]

        violations = find_capacity_violations(audits)
        self.assertIn("continuation_font_delta", {violation.rule for violation in violations})

    def test_deck_audit_accepts_small_font_delta_between_continuation_slides(self) -> None:
        audits = [
            SlideAudit(
                slide_index=2,
                title="Раздел",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=320,
                body_font_sizes=(14.0,),
                profile=profile_for_layout("text_full_width"),
                expected_items=("Первый блок",),
                rendered_items=("Первый блок",),
            ),
            SlideAudit(
                slide_index=3,
                title="Раздел (2)",
                kind=SlideKind.TEXT.value,
                layout_key="text_full_width",
                body_char_count=300,
                body_font_sizes=(13.0,),
                profile=profile_for_layout("text_full_width"),
                expected_items=("Второй блок",),
                rendered_items=("Второй блок",),
            ),
        ]

        violations = find_capacity_violations(audits)
        self.assertNotIn("continuation_font_delta", {violation.rule for violation in violations})

    def test_deck_audit_flags_title_body_overlap_for_text_slide(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=140,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=900000,
            body_top=1700000,
            body_height=2800000,
            footer_top=6384626,
            footer_width=11198224,
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("title_body_overlap", {violation.rule for violation in violations})

    def test_deck_audit_flags_subtitle_body_overlap_for_text_slide(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=140,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            subtitle_top=1600000,
            subtitle_height=400000,
            body_top=2050000,
            body_height=2500000,
            footer_top=6384626,
            footer_width=11198224,
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("subtitle_body_overlap", {violation.rule for violation in violations})

    def test_deck_audit_flags_title_body_gap_drift_for_text_slide(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=320,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=2800000,
            body_height=2400000,
            footer_top=6384626,
            footer_width=11198224,
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("title_body_gap_drift", {violation.rule for violation in violations})

    def test_deck_audit_flags_subtitle_body_gap_drift_for_text_slide(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=320,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            subtitle_top=1550000,
            subtitle_height=350000,
            body_top=3100000,
            body_height=2200000,
            footer_top=6384626,
            footer_width=11198224,
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("subtitle_body_gap_drift", {violation.rule for violation in violations})

    def test_deck_audit_flags_narrow_footer_for_text_layout(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=320,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=3600000,
            footer_top=6384626,
            footer_width=int(PptxGenerator.FULL_CONTENT_WIDTH_EMU * 0.75),
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("narrow_text_footer", {violation.rule for violation in violations})

    def test_deck_audit_flags_narrow_footer_by_placeholder_geometry(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Плотный раздел",
            kind=SlideKind.TEXT.value,
            layout_key="dense_text_full_width",
            body_char_count=320,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("dense_text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=2400000,
            footer_top=6384626,
            footer_left=442913,
            footer_width=3000000,
            body_left=442913,
            footer_placeholder_idx=17,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("narrow_footer", {violation.rule for violation in violations})

    def test_deck_audit_flags_footer_left_misalignment(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=320,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=2400000,
            footer_top=6384626,
            footer_left=900000,
            footer_width=3371850,
            body_left=442913,
            footer_placeholder_idx=17,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("footer_left_misalignment", {violation.rule for violation in violations})

    def test_deck_audit_accepts_balanced_dense_slides_without_capacity_violations(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Healthy",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Healthy", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Сбалансированный список",
                    bullets=[
                        "Пункт 1 подробно описывает цель и ограничения.",
                        "Пункт 2 раскрывает ресурсы и зависимости.",
                        "Пункт 3 фиксирует KPI и ожидаемый эффект.",
                        "Пункт 4 связывает решение с дорожной картой.",
                    ],
                    preferred_layout_key="list_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Сбалансированный текст",
                    text=(
                        "Первый абзац задаёт контекст и основные допущения. "
                        "Второй абзац описывает ожидаемый эффект и критерии контроля. "
                        "Третий абзац связывает решение с финансовыми и операционными метриками."
                    ),
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_deck_audit_validates_table_layout_geometry(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Table Geometry",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Table Geometry", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Ключевые показатели",
                    subtitle="Таблица должна занимать рабочую ширину layout",
                    table=TableBlock(
                        headers=["Показатель", "Значение"],
                        rows=[["Выручка", "120"], ["Маржа", "24%"], ["NPS", "61"]],
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        table_audit = next(audit for audit in audits if audit.kind == SlideKind.TABLE.value)
        self.assertTrue(table_audit.has_table)
        self.assertGreaterEqual(table_audit.content_width_ratio, 0.9)
        if table_audit.footer_width:
            self.assertGreaterEqual(table_audit.footer_width_ratio, 0.9)
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_generator_applies_corp_table_style_tokens(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Table Fill Colors",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Table Fill Colors", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Цветная таблица",
                    table=TableBlock(
                        headers=["Показатель", "Значение"],
                        header_fill_colors=["1F4E79", "1F4E79"],
                        rows=[["Выручка", "120"], ["Маржа", "24%"]],
                        row_fill_colors=[[None, "D9EAF7"], ["FDE7D7", None]],
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
        self.assertEqual(table.cell(0, 0).fill.fore_color.rgb, RGBColor(0xC6, 0xDF, 0xFF))
        self.assertEqual(table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.color.rgb, RGBColor(0x09, 0x1E, 0x38))
        self.assertIn('a:srgbClr val="C6DFFF"', table.cell(0, 0)._tc.tcPr.xml)
        self.assertIn("a:noFill", table.cell(1, 1)._tc.tcPr.xml)
        self.assertIn("a:noFill", table.cell(2, 0)._tc.tcPr.xml)

    def test_deck_audit_validates_chart_layout_geometry(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Chart Geometry",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Chart Geometry", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Выручка по каналам",
                    subtitle="График должен занимать рабочую ширину layout",
                    chart=ChartSpec(
                        chart_id="chart_geometry",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Выручка",
                        categories=["SEO", "Ads", "Referral"],
                        series=[ChartSeries(name="Выручка", values=[120.0, 200.0, 90.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        chart_audit = next(audit for audit in audits if audit.kind == SlideKind.CHART.value)
        self.assertTrue(chart_audit.has_chart)
        self.assertEqual(chart_audit.expected_chart_type, ChartType.COLUMN.value)
        self.assertEqual(chart_audit.rendered_chart_type, ChartType.COLUMN.value)
        self.assertEqual(chart_audit.expected_chart_series_count, 1)
        self.assertEqual(chart_audit.rendered_chart_series_count, 1)
        self.assertEqual(chart_audit.title_font_sizes, (manifest.theme.master_text_styles["title"].font_size_pt,))
        self.assertIn(chart_audit.subtitle_font_sizes, ((), (manifest.theme.master_text_styles["body"].font_size_pt,)))
        self.assertGreaterEqual(chart_audit.content_width_ratio, 0.9)
        if chart_audit.footer_width:
            self.assertGreaterEqual(chart_audit.footer_width_ratio, 0.9)
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_deck_audit_validates_chart_value_axis_number_format(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Chart Axis Format",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Chart Axis Format", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Доход по кварталам",
                    chart=ChartSpec(
                        chart_id="chart_axis_format",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Доход",
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(
                                name="Доход",
                                values=[104_300_000.0, 111_300_000.0, 135_700_000.0],
                            )
                        ],
                        confidence=ChartConfidence.HIGH,
                        value_format="currency",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        chart_audit = next(audit for audit in audits if audit.kind == SlideKind.CHART.value)
        self.assertEqual(chart_audit.expected_chart_value_axis_number_format, '0.0,," млн ₽"')
        self.assertEqual(chart_audit.rendered_chart_value_axis_number_format, '0.0,," млн ₽"')
        self.assertEqual(find_capacity_violations(audits), [])

    def test_deck_audit_validates_secondary_chart_value_axis_number_format(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Chart Secondary Axis Format",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Chart Secondary Axis Format", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Выручка и маржа",
                    chart=ChartSpec(
                        chart_id="chart_combo_secondary_axis_audit",
                        source_table_id="table_1",
                        chart_type=ChartType.COMBO,
                        title="Выручка и маржа",
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(name="Выручка", values=[104_300_000.0, 111_300_000.0, 135_700_000.0], unit="RUB"),
                            ChartSeries(name="Маржа", values=[18.0, 22.0, 27.0], unit="%"),
                        ],
                        confidence=ChartConfidence.HIGH,
                        value_format="number",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        chart_audit = next(audit for audit in audits if audit.kind == SlideKind.CHART.value)
        self.assertTrue(chart_audit.expected_chart_secondary_value_axis)
        self.assertTrue(chart_audit.rendered_chart_secondary_value_axis)
        self.assertEqual(chart_audit.expected_chart_value_axis_number_format, '0.0,," млн ₽"')
        self.assertEqual(chart_audit.rendered_chart_value_axis_number_format, '0.0,," млн ₽"')
        self.assertEqual(chart_audit.expected_chart_secondary_value_axis_number_format, '0"%"')
        self.assertEqual(chart_audit.rendered_chart_secondary_value_axis_number_format, '0"%"')
        self.assertEqual(find_capacity_violations(audits), [])

    def test_deck_audit_uses_trillion_currency_axis_format_for_large_market_combo(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Market Axis Format",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Market Axis Format", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Объем рынков",
                    chart=ChartSpec(
                        chart_id="chart_market_trillion_axis",
                        source_table_id="table_1",
                        chart_type=ChartType.COMBO,
                        title="Объем рынков",
                        categories=["ЖКХ", "Налоги", "Образование"],
                        series=[
                            ChartSeries(
                                name="Объем рынка (2024)",
                                values=[8_496_000_000_000.0, 55_600_000_000_000.0, 851_000_000_000.0],
                                unit="RUB",
                                axis="primary",
                            ),
                            ChartSeries(name="Доля А3 GMV (2025)", values=[2.12, 0.02, 0.06], unit="%", axis="secondary"),
                        ],
                        confidence=ChartConfidence.HIGH,
                        value_format="number",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)

        chart_audit = next(audit for audit in audits if audit.kind == SlideKind.CHART.value)
        self.assertEqual(chart_audit.expected_chart_value_axis_number_format, '0.0,,,," трлн ₽"')
        self.assertEqual(chart_audit.rendered_chart_value_axis_number_format, '0.0,,,," трлн ₽"')
        self.assertEqual(chart_audit.expected_chart_secondary_value_axis_number_format, '0"%"')
        self.assertEqual(chart_audit.rendered_chart_secondary_value_axis_number_format, '0"%"')
        self.assertEqual(find_capacity_violations(audits), [])

    def test_mixed_unit_single_axis_chart_does_not_format_primary_axis_as_percent(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Mixed Axis Safety",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Mixed Axis Safety", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Рынок и доля",
                    chart=ChartSpec(
                        chart_id="chart_mixed_axis_safety",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Рынок и доля",
                        categories=["ЖКХ", "Налоги", "Образование"],
                        series=[
                            ChartSeries(name="Объем рынка", values=[8_496_000_000_000.0, 55_600_000_000_000.0, 851_000_000_000.0]),
                            ChartSeries(name="Доля А3", values=[2.12, 0.02, 0.06], unit="%"),
                        ],
                        confidence=ChartConfidence.HIGH,
                        value_format="percent",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                plan=plan,
                template_path=template_path,
                manifest=manifest,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        chart_audit = next(audit for audit in audits if audit.kind == SlideKind.CHART.value)
        self.assertEqual(chart_audit.expected_chart_value_axis_number_format, '0.0,,,," трлн"')
        self.assertEqual(chart_audit.rendered_chart_value_axis_number_format, '0.0,,,," трлн"')

    def test_deck_audit_flags_chart_semantic_mismatch(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Combo mismatch",
            kind=SlideKind.CHART.value,
            layout_key="table",
            body_char_count=0,
            body_font_sizes=(),
            profile=profile_for_layout("table"),
            has_chart=True,
            content_width=PptxGenerator.FULL_CONTENT_WIDTH_EMU,
            footer_width=PptxGenerator.FULL_CONTENT_WIDTH_EMU,
            expected_chart_type=ChartType.COMBO.value,
            rendered_chart_type=ChartType.COLUMN.value,
            expected_chart_series_count=3,
            rendered_chart_series_count=2,
            rendered_chart_bar_series_count=2,
            rendered_chart_line_series_count=0,
        )

        violations = find_capacity_violations([audit])

        rules = {violation.rule for violation in violations}
        self.assertIn("chart_type_mismatch", rules)
        self.assertIn("chart_series_count_mismatch", rules)
        self.assertIn("combo_chart_structure_mismatch", rules)

    def test_deck_audit_flags_chart_axis_number_format_mismatch(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Axis format mismatch",
            kind=SlideKind.CHART.value,
            layout_key="table",
            body_char_count=0,
            body_font_sizes=(),
            profile=profile_for_layout("table"),
            has_chart=True,
            content_width=PptxGenerator.FULL_CONTENT_WIDTH_EMU,
            footer_width=PptxGenerator.FULL_CONTENT_WIDTH_EMU,
            expected_chart_type=ChartType.COLUMN.value,
            rendered_chart_type=ChartType.COLUMN.value,
            expected_chart_series_count=1,
            rendered_chart_series_count=1,
            expected_chart_value_axis_number_format='0.0,," млн ₽"',
            rendered_chart_value_axis_number_format="#,##0",
        )

        violations = find_capacity_violations([audit])

        rules = {violation.rule for violation in violations}
        self.assertIn("chart_value_axis_number_format_mismatch", rules)

    def test_deck_audit_flags_secondary_chart_axis_number_format_mismatch(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Secondary axis format mismatch",
            kind=SlideKind.CHART.value,
            layout_key="table",
            body_char_count=0,
            body_font_sizes=(),
            profile=profile_for_layout("table"),
            has_chart=True,
            expected_chart_secondary_value_axis=True,
            rendered_chart_secondary_value_axis=True,
            expected_chart_secondary_value_axis_number_format='0"%"',
            rendered_chart_secondary_value_axis_number_format="#,##0",
        )

        violations = find_capacity_violations([audit])

        rules = {violation.rule for violation in violations}
        self.assertIn("chart_secondary_value_axis_number_format_mismatch", rules)

    def test_generated_footer_font_is_visually_secondary_to_subtitle(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Footer Font Contract",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Footer Font Contract", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Квартальный профиль",
                    subtitle="Финансовая модель 2026",
                    chart=ChartSpec(
                        chart_id="footer_font_chart",
                        source_table_id="table_1",
                        chart_type=ChartType.LINE,
                        categories=["Q1", "Q2", "Q3"],
                        series=[ChartSeries(name="Выручка", values=[120.0, 150.0, 190.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        subtitle = next(shape for shape in slide.placeholders if shape.placeholder_format.idx == 13)
        footer = next(shape for shape in slide.placeholders if shape.placeholder_format.idx == 15)
        subtitle_sizes = self._text_frame_font_sizes(subtitle.text_frame)
        footer_sizes = self._text_frame_font_sizes(footer.text_frame)

        self.assertTrue(subtitle_sizes)
        self.assertTrue(footer_sizes)
        self.assertLess(max(footer_sizes), max(subtitle_sizes))
        self.assertEqual(max(footer_sizes), PptxGenerator.FOOTER_FONT_PT)

    def test_deck_audit_flags_chart_title_and_subtitle_font_mismatch(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Chart font mismatch",
            kind=SlideKind.CHART.value,
            layout_key="table",
            body_char_count=0,
            body_font_sizes=(),
            title_font_sizes=(8.0,),
            subtitle_font_sizes=(8.0,),
            profile=profile_for_layout("table"),
            has_chart=True,
            content_width=PptxGenerator.FULL_CONTENT_WIDTH_EMU,
            footer_width=PptxGenerator.FULL_CONTENT_WIDTH_EMU,
            expected_chart_type=ChartType.COLUMN.value,
            rendered_chart_type=ChartType.COLUMN.value,
            expected_chart_series_count=1,
            rendered_chart_series_count=1,
            expected_title_font_pt=35.0,
            expected_subtitle_font_pt=20.0,
        )

        violations = find_capacity_violations([audit])

        rules = {violation.rule for violation in violations}
        self.assertIn("chart_title_font_mismatch", rules)
        self.assertIn("chart_subtitle_font_mismatch", rules)

    def test_deck_audit_validates_image_layout_geometry(self) -> None:
        small_png_base64 = (
            "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9W6i8AAAAASUVORK5CYII="
        )
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Image Geometry",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Image Geometry", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.IMAGE,
                    title="Схема процесса",
                    text="Изображение должно рендериться как picture shape и сохранять рабочую геометрию layout.",
                    preferred_layout_key="image_text",
                    image_base64=small_png_base64,
                    image_content_type="image/png",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        image_audit = next(audit for audit in audits if audit.kind == SlideKind.IMAGE.value)
        self.assertTrue(image_audit.has_image)
        self.assertGreaterEqual(image_audit.content_width_ratio, 0.35)
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_deck_audit_validates_cards_layout_geometry(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Cards Geometry",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Cards Geometry", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Три направления роста",
                    bullets=["Усилить ядро платформы", "Ускорить интеграции", "Повысить retention"],
                    preferred_layout_key="cards_3",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        card_audit = next(audit for audit in audits if audit.layout_key == "cards_3")
        violations = find_capacity_violations(audits)
        self.assertTrue(card_audit.auxiliary_widths)
        self.assertEqual({idx: card_audit.placeholder_char_counts[idx] for idx in (11, 12, 13)}, {11: 22, 12: 19, 13: 18})
        self.assertEqual(card_audit.expected_placeholder_char_counts, {11: 22, 12: 19, 13: 18})
        self.assertEqual([v.rule for v in violations if v.slide_index == card_audit.slide_index], [])

    def test_cards_layout_renders_one_item_per_card_with_fitted_text(self) -> None:
        card_texts = [
            "Первое направление: усилить платформенное ядро и убрать ручные операции в ключевых интеграциях.",
            "Второе направление: расширить партнерскую сеть и ускорить подключение новых каналов продаж.",
            "Третье направление: повысить удержание клиентов за счет аналитики, персонализации и регулярных продуктовых улучшений.",
        ]
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Cards Text Fit",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Cards Text Fit", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Три направления роста",
                    bullets=card_texts,
                    preferred_layout_key="cards_3",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            generated = Presentation(output_path)
            card_slide = generated.slides[1]
            full_slide_pictures = [
                shape
                for shape in card_slide.shapes
                if shape.left == 0
                and shape.top == 0
                and shape.width == generated.slide_width
                and shape.height == generated.slide_height
                and str(shape.shape_type) == "PICTURE (13)"
            ]
            self.assertEqual(full_slide_pictures, [])
            cards = {
                shape.placeholder_format.idx: shape
                for shape in card_slide.placeholders
                if shape.placeholder_format.idx in {11, 12, 13}
            }

            self.assertEqual(
                [cards[idx].text.strip() for idx in (11, 12, 13)],
                [
                    "Первое направление\nусилить платформенное ядро и убрать ручные операции в ключевых интеграциях.",
                    "Второе направление\nрасширить партнерскую сеть и ускорить подключение новых каналов продаж.",
                    "Третье направление\nповысить удержание клиентов за счет аналитики, персонализации и регулярных продуктовых улучшений.",
                ],
            )
            expected_geometry = {
                11: (739775, 1723633, 3259138, 4164013),
                12: (4412456, 1723633, 3259138, 4164013),
                13: (8193087, 1723633, 3259138, 4164013),
            }
            card_font_sizes = []
            for shape in cards.values():
                self.assertEqual(
                    (shape.left, shape.top, shape.width, shape.height),
                    expected_geometry[shape.placeholder_format.idx],
                )
                self.assertEqual(shape.text_frame.margin_left, self.generator.DEFAULT_TEXT_MARGIN_X_EMU)
                self.assertEqual(shape.text_frame.margin_right, self.generator.DEFAULT_TEXT_MARGIN_X_EMU)
                self.assertEqual(shape.text_frame.margin_top, self.generator.DEFAULT_TEXT_MARGIN_Y_EMU)
                self.assertEqual(shape.text_frame.margin_bottom, self.generator.DEFAULT_TEXT_MARGIN_Y_EMU)
                sizes = [
                    run.font.size.pt
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.font.size is not None
                ]
                colors = [
                    str(run.font.color.rgb)
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.font.color.rgb is not None
                ]
                bold_values = [
                    run.font.bold
                    for paragraph in shape.text_frame.paragraphs
                    for run in paragraph.runs
                ]
                self.assertTrue(sizes)
                self.assertEqual(set(round(size, 1) for size in sizes), {16.0, 20.0})
                self.assertEqual(set(colors), {"FFFFFF"})
                self.assertIn(True, set(bold_values))
                self.assertIn(False, set(bold_values))
                card_font_sizes.append(tuple(sorted(set(round(size, 1) for size in sizes))))
            self.assertEqual(len(set(card_font_sizes)), 1)

    def test_kpi_cards_layout_renders_four_numeric_metrics(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="KPI Cards",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="KPI Cards", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="A3 GIS",
                    bullets=[
                        "99,5 %\nуспешных поисков",
                        "53 млн\nактивных начислений",
                        "92 млн\nдокументов для поиска",
                        "0,86 %\nсреднее время поиска в секундах",
                    ],
                    preferred_layout_key="cards_kpi",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)
            generated = Presentation(output_path)

        kpi_audit = next(audit for audit in audits if audit.layout_key == "cards_kpi")
        self.assertEqual([v.rule for v in find_capacity_violations(audits) if v.slide_index == kpi_audit.slide_index], [])

        slide = generated.slides[1]
        rendered_texts = sorted(shape.text.strip() for shape in slide.shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
        self.assertIn("99,5 %\nуспешных поисков", rendered_texts)
        self.assertIn("53 млн\nактивных начислений", rendered_texts)
        self.assertIn("92 млн\nдокументов для поиска", rendered_texts)
        self.assertIn("0,86 %\nсреднее время поиска в секундах", rendered_texts)

        metric_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False) and shape.text.strip() in {
                "99,5 %\nуспешных поисков",
                "53 млн\nактивных начислений",
                "92 млн\nдокументов для поиска",
                "0,86 %\nсреднее время поиска в секундах",
            }
        ]
        self.assertEqual(len(metric_shapes), 4)
        self.assertTrue(any(not getattr(shape, "is_placeholder", False) for shape in metric_shapes))
        max_metric_font = max(
            run.font.size.pt
            for shape in metric_shapes
            for paragraph in shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None
        )
        self.assertGreaterEqual(max_metric_font, 36.0)

    def test_cards_layout_renders_numeric_metrics_inside_each_card(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Cards Numeric Demo",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Cards Numeric Demo", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Проводим платеж",
                    bullets=[
                        "A3 GIS\nЗапрашиваем, обновляем, перепроверяем и кешируем информацию по начислениям из ГИС ГМП и ГИС ЖКХ.\n99,5 % успешных поисков\n53 млн активных начислений\n92 млн документов для поиска\n0,86 % среднее время поиска в секундах",
                        "A3 ETL\nДля поставщиков, с которыми невозможно онлайн-подключение, подстраиваемся под тип и вид файлов.\n100 успешных поисков\n15 млн документов для поиска\n10 млн активных начислений",
                        "A3 ONLINE\nДля поставщиков, которые готовы передавать данные онлайн.\n100500 документов для поиска\n> 100 онлайн-интеграций с крупными поставщиками услуг",
                    ],
                    preferred_layout_key="cards_3",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)
            generated = Presentation(output_path)

        self.assertEqual(find_capacity_violations(audits), [])

        slide = generated.slides[1]
        overlay_texts = {
            shape.name: shape.text.strip()
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and getattr(shape, "name", "").startswith("A3_CARD_OVERLAY_")
            and shape.text.strip()
        }
        self.assertIn("A3_CARD_OVERLAY_11_TITLE", overlay_texts)
        self.assertIn("A3_CARD_OVERLAY_11_METRIC_0", overlay_texts)
        self.assertIn("A3_CARD_OVERLAY_11_METRIC_3", overlay_texts)
        self.assertIn("99,5%\nуспешных поисков", overlay_texts.values())
        self.assertIn("53 млн\nактивных начислений", overlay_texts.values())
        self.assertIn("> 100\nонлайн-интеграций с крупными поставщиками услуг", overlay_texts.values())

    def test_cards_layout_supports_one_two_three_numeric_metrics_per_card(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Cards Metric Counts",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Cards Metric Counts", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Произвольные показатели",
                    bullets=[
                        "Одна метрика\nЛюбой поясняющий текст карточки.\n7 дней средний срок",
                        "Две метрики\nДругой текст без привязки к примеру.\n12,4 сек время ответа\n8500 заявок обработано",
                        "Три метрики\nОписание может быть любым.\n> 100 интеграций подключено\n3 млрд ₽ оборот\n91 % успешных операций",
                    ],
                    preferred_layout_key="cards_3",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan, manifest)
            generated = Presentation(output_path)

        self.assertEqual(find_capacity_violations(audits), [])

        slide = generated.slides[1]
        metric_shapes = [
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and getattr(shape, "name", "").startswith("A3_CARD_OVERLAY_")
            and "_METRIC_" in getattr(shape, "name", "")
        ]
        self.assertEqual(len(metric_shapes), 6)
        self.assertIn("7 дней\nсредний срок", {shape.text.strip() for shape in metric_shapes})
        self.assertIn("12,4 сек\nвремя ответа", {shape.text.strip() for shape in metric_shapes})
        self.assertIn("91%\nуспешных операций", {shape.text.strip() for shape in metric_shapes})

    def test_numeric_cards_skip_empty_description_overlays(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Cards Without Description",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Cards Without Description", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Карточки без заглушек",
                    bullets=[
                        "Одна метрика\n120 заявок обработано",
                        "Две метрики\n45 сек среднее ожидание\n98% успешных ответов",
                        "Три метрики\n> 100 интеграций подключено\n12 млн начислений\n7 дней средний срок",
                    ],
                    preferred_layout_key="cards_3",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            generated = Presentation(output_path)

        slide = generated.slides[1]
        overlay_names = {
            shape.name
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and getattr(shape, "name", "").startswith("A3_CARD_OVERLAY_")
            and shape.text.strip()
        }
        self.assertFalse(any("_DESCRIPTION" in name for name in overlay_names))

    def test_layout_background_image_is_applied_without_template_id_special_case(self) -> None:
        manifest = self.registry.get_template("corp_light_v1").model_copy(
            update={"template_id": "uploaded_demo_template"},
            deep=True,
        )
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        background_layout = next(
            layout
            for layout in manifest.layouts
            if layout.background_image_base64 and layout.slide_master_index == 0
        )
        plan = PresentationPlan(
            template_id=manifest.template_id,
            title="Template Background Contract",
            slides=[
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Контекст",
                    preferred_layout_key=background_layout.key,
                    background_only=True,
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            generated = Presentation(output_path)
        slide = generated.slides[0]
        background = next(
            shape
            for shape in slide.shapes
            if str(shape.shape_type) == "PICTURE (13)"
            and shape.left == 0
            and shape.top == 0
            and shape.width == generated.slide_width
            and shape.height == generated.slide_height
        )
        self.assertEqual(str(background.shape_type), "PICTURE (13)")
        self.assertEqual(background.left, 0)
        self.assertEqual(background.top, 0)
        self.assertEqual(background.width, generated.slide_width)
        self.assertEqual(background.height, generated.slide_height)

    def test_deck_audit_uses_manifest_geometry_without_template_id_special_case(self) -> None:
        manifest = self.registry.get_template("corp_light_v1").model_copy(
            update={"template_id": "uploaded_geometry_contract"},
            deep=True,
        )
        layout = next(item for item in manifest.layouts if item.key == "text_full_width")
        body = next(item for item in layout.placeholders if item.idx == 14)
        body.left_emu = 123456
        body.top_emu = 234567
        body.width_emu = 3456789
        body.height_emu = 456789

        policy = _geometry_policy_for_slide(
            SlideSpec(
                kind=SlideKind.TEXT,
                title="Контекст",
                text="Проверка manifest geometry.",
                preferred_layout_key="text_full_width",
            ),
            manifest,
        )

        self.assertEqual(policy.placeholders[14].left_emu, 123456)
        self.assertEqual(policy.placeholders[14].top_emu, 234567)
        self.assertEqual(policy.placeholders[14].width_emu, 3456789)
        self.assertEqual(policy.placeholders[14].height_emu, 456789)

    def test_deck_audit_uses_inventory_target_geometry_for_custom_target_key(self) -> None:
        manifest = self.registry.get_template("corp_light_v1").model_copy(
            update={"template_id": "uploaded_custom_geometry_contract"},
            deep=True,
        )
        layout = next(item for item in manifest.layouts if item.key == "text_full_width")
        layout.key = "uploaded_text_target"
        if manifest.default_layout_key == "text_full_width":
            manifest.default_layout_key = layout.key
        body = next(item for item in layout.placeholders if item.idx == 14)
        body.left_emu = 654321
        body.top_emu = 765432
        body.width_emu = 2765432
        body.height_emu = 876543

        policy = _geometry_policy_for_slide(
            SlideSpec(
                kind=SlideKind.TEXT,
                title="Контекст",
                text="Проверка extracted target geometry.",
                runtime_profile_key="text_full_width",
                preferred_layout_key="uploaded_text_target",
            ),
            manifest,
        )

        self.assertEqual(policy.layout_key, "uploaded_text_target")
        self.assertEqual(policy.placeholders[14].left_emu, 654321)
        self.assertEqual(policy.placeholders[14].top_emu, 765432)
        self.assertEqual(policy.placeholders[14].width_emu, 2765432)
        self.assertEqual(policy.placeholders[14].height_emu, 876543)

    def test_cards_layout_keeps_clearance_under_wrapped_title(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Cards Wrapped Title",
            slides=[
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Очень длинный заголовок карточного слайда для проверки переноса на две строки",
                    bullets=[
                        "Первое направление: усилить платформенное ядро и убрать ручные операции в ключевых интеграциях.",
                        "Второе направление: расширить партнерскую сеть и ускорить подключение новых каналов продаж.",
                        "Третье направление: повысить удержание клиентов за счет аналитики и персонализации.",
                    ],
                    preferred_layout_key="cards_3",
                )
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            generated = Presentation(output_path)
            slide = generated.slides[0]
            placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}
            title = placeholders[0]
            first_card_top = min(placeholders[idx].top for idx in (11, 12, 13))
            expected_card_bottom = 1723633 + 4164013
            audits = audit_generated_presentation(output_path, plan, manifest)

        self.assertGreaterEqual(first_card_top - (title.top + title.height), 320000)
        for idx in (11, 12, 13):
            self.assertEqual(placeholders[idx].top, 1723633)
            self.assertEqual(placeholders[idx].height, 4164013)
            self.assertEqual(placeholders[idx].top + placeholders[idx].height, expected_card_bottom)
        self.assertEqual(find_capacity_violations(audits), [])

    def test_deck_audit_validates_two_column_layout_geometry(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Two Column Geometry",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Two Column Geometry", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TWO_COLUMN,
                    title="Модель взаимодействия",
                    subtitle="Колонки и иконки должны сохранять рабочую сетку",
                    left_bullets=["Контекст проекта", "Ограничения", "Допущения"],
                    right_bullets=["Шаг 1", "Шаг 2", "Шаг 3", "Шаг 4"],
                    preferred_layout_key="list_with_icons",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        two_column_audit = next(audit for audit in audits if audit.layout_key == "list_with_icons")
        violations = find_capacity_violations(audits)
        self.assertTrue(two_column_audit.auxiliary_widths)
        self.assertEqual(two_column_audit.expected_auxiliary_char_counts, {12: 38})
        self.assertEqual(two_column_audit.expected_placeholder_char_counts, {12: 38, 14: 23})
        self.assertEqual([v.rule for v in violations if v.slide_index == two_column_audit.slide_index], [])

    def test_deck_audit_tracks_expected_subtitle_placeholder_fill_for_text_slide(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Subtitle Fill",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Subtitle Fill", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Раздел",
                    subtitle="Подзаголовок",
                    text="Основной текст для проверки subtitle placeholder.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        text_audit = next(audit for audit in audits if audit.layout_key == "text_full_width")
        self.assertEqual(text_audit.expected_subtitle_char_count, 12)
        self.assertNotIn("underfilled_subtitle_placeholder_fill", {violation.rule for violation in find_capacity_violations(audits)})

    def test_deck_audit_skips_subtitle_placeholder_fill_when_subtitle_duplicates_body(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Subtitle Dedup",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Subtitle Dedup", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Narrative",
                    subtitle="Первый абзац задает контекст.",
                    text="Первый абзац задает контекст. Дальше идет основной narrative без отдельного подзаголовка.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        text_audit = next(audit for audit in audits if audit.layout_key == "text_full_width")
        self.assertEqual(text_audit.expected_subtitle_char_count, 0)
        self.assertNotIn("underfilled_subtitle_placeholder_fill", {violation.rule for violation in find_capacity_violations(audits)})

    def test_deck_audit_validates_contacts_layout_geometry(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Contacts Geometry",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Contacts Geometry", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Иван Иванов",
                    subtitle="CEO",
                    left_bullets=["+7 999 123-45-67"],
                    right_bullets=["ivan@example.com"],
                    preferred_layout_key="contacts",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        contacts_audit = next(audit for audit in audits if audit.layout_key == "contacts")
        violations = find_capacity_violations(audits)
        self.assertTrue(contacts_audit.auxiliary_widths)
        self.assertEqual(contacts_audit.placeholder_char_counts, {10: 11, 11: 3, 12: 16, 13: 16})
        self.assertEqual(contacts_audit.expected_placeholder_char_counts, {10: 11, 11: 3, 12: 16, 13: 16})
        self.assertEqual([v.rule for v in violations if v.slide_index == contacts_audit.slide_index], [])

    def test_list_with_icons_component_style_controls_subtitle_and_spacing(self) -> None:
        class PlaceholderRef:
            def __init__(self, idx: int) -> None:
                self.idx = idx

        class FakeShape:
            def __init__(self, idx: int, text: str, *, top: int, height: int, width: int) -> None:
                self.placeholder_format = PlaceholderRef(idx)
                self.is_placeholder = True
                self.text = text
                self.top = top
                self.height = height
                self.width = width
                self.left = 0

        class FakeSlide:
            def __init__(self, placeholders: list[FakeShape]) -> None:
                self.placeholders = placeholders

        manifest = self.registry.get_template("corp_light_v1").model_copy(deep=True)
        icons_style = manifest.component_styles["list_with_icons"]
        icons_style.text_styles["subtitle"] = icons_style.text_styles["subtitle"].model_copy(update={"font_size_pt": 26.0})
        icons_style.spacing_tokens["title_content_gap_emu"] = 300000

        title = FakeShape(0, "Короткий заголовок", top=100000, height=500000, width=7000000)
        subtitle = FakeShape(13, "Подзаголовок", top=0, height=300000, width=7000000)
        left = FakeShape(12, "Левая колонка", top=1200000, height=1500000, width=3000000)
        right = FakeShape(14, "Правая колонка", top=1200000, height=1500000, width=3000000)
        footer = FakeShape(21, "Footer", top=6200000, height=300000, width=7000000)
        slide = FakeSlide([title, subtitle, left, right, footer])

        original_manifest = self.generator._active_manifest
        original_apply = self.generator._apply_font_size
        original_fit_title = self.generator._fit_title_font_size_points
        original_configure_title = self.generator._configure_title_text_frame
        original_estimate_title = self.generator._estimate_title_height_emu
        original_min_title = self.generator._minimum_title_height_emu
        original_configure_subtitle = self.generator._configure_subtitle_text_frame
        original_estimate_text = self.generator._estimate_text_height_emu

        applied_sizes: dict[int, float] = {}
        self.generator._active_manifest = manifest
        self.generator._apply_font_size = lambda shape, size: applied_sizes.__setitem__(shape.placeholder_format.idx, size)
        self.generator._fit_title_font_size_points = lambda _text, _width, _layout_key: 35.0
        self.generator._configure_title_text_frame = lambda _shape: None
        self.generator._estimate_title_height_emu = lambda _shape, _text, _size: 600000
        self.generator._minimum_title_height_emu = lambda _layout_key: 400000
        self.generator._configure_subtitle_text_frame = lambda _shape: None
        self.generator._estimate_text_height_emu = lambda _text, _width, _size: 200000

        try:
            self.generator._stack_two_column_content(slide, "list_with_icons")
        finally:
            self.generator._active_manifest = original_manifest
            self.generator._apply_font_size = original_apply
            self.generator._fit_title_font_size_points = original_fit_title
            self.generator._configure_title_text_frame = original_configure_title
            self.generator._estimate_title_height_emu = original_estimate_title
            self.generator._minimum_title_height_emu = original_min_title
            self.generator._configure_subtitle_text_frame = original_configure_subtitle
            self.generator._estimate_text_height_emu = original_estimate_text

        self.assertEqual(applied_sizes[13], 26.0)
        self.assertEqual(left.top, subtitle.top + subtitle.height + 300000)
        self.assertEqual(right.top, left.top)

    def test_contacts_component_style_controls_font_decrement_behavior(self) -> None:
        class PlaceholderRef:
            def __init__(self, idx: int) -> None:
                self.idx = idx

        class FakeShape:
            def __init__(self, idx: int, text: str) -> None:
                self.placeholder_format = PlaceholderRef(idx)
                self.is_placeholder = True
                self.text = text
                self.left = 0
                self.top = 0
                self.width = 0
                self.height = 0

        class FakeSlide:
            def __init__(self, placeholders: list[FakeShape]) -> None:
                self.placeholders = placeholders

        manifest = self.registry.get_template("corp_light_v1").model_copy(deep=True)
        contacts_style = manifest.component_styles["contacts"]
        contacts_style.text_styles["primary"] = contacts_style.text_styles["primary"].model_copy(update={"font_size_pt": 20.0})
        contacts_style.text_styles["secondary"] = contacts_style.text_styles["secondary"].model_copy(update={"font_size_pt": 16.0})
        contacts_style.behavior_tokens["primary_threshold_chars"] = 1
        contacts_style.behavior_tokens["secondary_threshold_chars"] = 1
        contacts_style.behavior_tokens["font_decrement_pt"] = 3.0

        slide = FakeSlide(
            [
                FakeShape(10, "Иван Иванов"),
                FakeShape(11, "Chief Executive Officer"),
                FakeShape(12, "+7 999 123-45-67"),
                FakeShape(13, "ivan@example.com"),
            ]
        )

        original_manifest = self.generator._active_manifest
        original_apply = self.generator._apply_font_size
        original_configure_subtitle = self.generator._configure_subtitle_text_frame

        applied_sizes: dict[int, float] = {}
        self.generator._active_manifest = manifest
        self.generator._apply_font_size = lambda shape, size: applied_sizes.__setitem__(shape.placeholder_format.idx, size)
        self.generator._configure_subtitle_text_frame = lambda _shape: None

        try:
            self.generator._stack_contacts_content(slide, "contacts")
        finally:
            self.generator._active_manifest = original_manifest
            self.generator._apply_font_size = original_apply
            self.generator._configure_subtitle_text_frame = original_configure_subtitle

        self.assertEqual(applied_sizes[10], 17.0)
        self.assertEqual(applied_sizes[11], 13.0)
        self.assertEqual(applied_sizes[12], 13.0)
        self.assertEqual(applied_sizes[13], 13.0)

    def test_deck_audit_flags_underfilled_placeholder_fill_for_full_width_text(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=140,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=1200000,
            footer_top=6384626,
            footer_width=11198224,
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_flags_underfilled_placeholder_fill_for_dense_text(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Плотный раздел",
            kind=SlideKind.TEXT.value,
            layout_key="dense_text_full_width",
            body_char_count=180,
            body_font_sizes=(12.0,),
            profile=derive_capacity_profile_for_geometry("dense_text_full_width", width_emu=11198224, height_emu=3550000),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=1600000,
            footer_top=6384626,
            footer_width=11198224,
            body_left=442913,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_skips_strict_safe_layout_rules_for_auto_layout_targets(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Fallback",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=140,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=1200000,
            body_left=900000,
            body_margin_left=0,
            body_margin_right=0,
            body_margin_top=0,
            body_margin_bottom=0,
            footer_top=6384626,
            footer_left=0,
            footer_width=1000000,
            footer_placeholder_idx=17,
            placeholder_char_counts={0: 6, 14: 40, 17: 10},
            subtitle_placeholder_idx=13,
            expected_subtitle_char_count=12,
            target_type="auto_layout",
        )

        violations = {violation.rule for violation in find_capacity_violations([audit])}
        self.assertNotIn("underfilled_placeholder_fill", violations)
        self.assertNotIn("body_left_misalignment", violations)
        self.assertNotIn("body_margin_mismatch", violations)
        self.assertNotIn("narrow_text_footer", violations)
        self.assertNotIn("narrow_footer", violations)
        self.assertNotIn("footer_left_misalignment", violations)
        self.assertNotIn("underfilled_subtitle_placeholder_fill", violations)

    def test_deck_audit_skips_strict_safe_layout_rules_for_direct_shape_binding_targets(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Direct binding",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=140,
            body_font_sizes=(14.0,),
            profile=profile_for_layout("text_full_width"),
            title_top=671247,
            title_height=800000,
            body_top=1900000,
            body_height=1200000,
            body_left=900000,
            body_margin_left=0,
            body_margin_right=0,
            body_margin_top=0,
            body_margin_bottom=0,
            footer_top=6384626,
            footer_left=0,
            footer_width=1000000,
            footer_placeholder_idx=17,
            placeholder_char_counts={0: 6, 14: 40, 17: 10},
            subtitle_placeholder_idx=13,
            expected_subtitle_char_count=12,
            target_type="direct_shape_binding",
        )

        violations = {violation.rule for violation in find_capacity_violations([audit])}
        self.assertNotIn("underfilled_placeholder_fill", violations)
        self.assertNotIn("body_left_misalignment", violations)
        self.assertNotIn("body_margin_mismatch", violations)
        self.assertNotIn("narrow_text_footer", violations)
        self.assertNotIn("narrow_footer", violations)
        self.assertNotIn("footer_left_misalignment", violations)
        self.assertNotIn("underfilled_subtitle_placeholder_fill", violations)

    def test_deck_audit_flags_underfilled_auxiliary_placeholder_fill(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Колонки",
            kind=SlideKind.TWO_COLUMN.value,
            layout_key="list_with_icons",
            body_char_count=0,
            body_font_sizes=(),
            profile=profile_for_layout("text_full_width"),
            auxiliary_char_counts={12: 0},
            expected_auxiliary_char_counts={12: 14},
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_auxiliary_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_flags_underfilled_auxiliary_placeholder_fill_for_right_column(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Колонки",
            kind=SlideKind.TWO_COLUMN.value,
            layout_key="list_with_icons",
            body_char_count=0,
            body_font_sizes=(),
            profile=profile_for_layout("text_full_width"),
            placeholder_char_counts={12: 14, 14: 0},
            expected_placeholder_char_counts={12: 14, 14: 23},
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_two_column_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_flags_underfilled_contact_placeholder_fill(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Иван Иванов",
            kind=SlideKind.TEXT.value,
            layout_key="contacts",
            body_char_count=11,
            body_font_sizes=(18.0,),
            profile=profile_for_layout("text_full_width"),
            placeholder_char_counts={10: 11, 11: 3, 12: 16, 13: 0},
            expected_placeholder_char_counts={10: 11, 11: 3, 12: 16, 13: 16},
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_contact_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_flags_underfilled_card_placeholder_fill(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Три направления роста",
            kind=SlideKind.BULLETS.value,
            layout_key="cards_3",
            body_char_count=60,
            body_font_sizes=(16.0,),
            profile=profile_for_layout("text_full_width"),
            placeholder_char_counts={11: 24, 12: 19, 13: 0},
            expected_placeholder_char_counts={11: 24, 12: 19, 13: 17},
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_card_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_flags_underfilled_subtitle_placeholder_fill(self) -> None:
        audit = SlideAudit(
            slide_index=2,
            title="Раздел",
            kind=SlideKind.TEXT.value,
            layout_key="text_full_width",
            body_char_count=80,
            body_font_sizes=(16.0,),
            profile=profile_for_layout("text_full_width"),
            placeholder_char_counts={0: 6, 14: 40, 17: 10},
            subtitle_placeholder_idx=13,
            expected_subtitle_char_count=12,
        )

        violations = find_capacity_violations([audit])
        self.assertIn("underfilled_subtitle_placeholder_fill", {violation.rule for violation in violations})

    def test_deck_audit_keeps_body_text_frame_margin_contract_for_full_width_text(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Body Margins",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Body Margins", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Раздел",
                    text="Основной текст должен рендериться с явными внутренними margin-параметрами.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        text_audit = next(audit for audit in audits if audit.layout_key == "text_full_width")
        self.assertEqual(text_audit.body_margin_left, PptxGenerator.DEFAULT_TEXT_MARGIN_X_EMU)
        self.assertEqual(text_audit.body_margin_right, PptxGenerator.DEFAULT_TEXT_MARGIN_X_EMU)
        self.assertEqual(text_audit.body_margin_top, PptxGenerator.DEFAULT_TEXT_MARGIN_Y_EMU)
        self.assertEqual(text_audit.body_margin_bottom, PptxGenerator.DEFAULT_TEXT_MARGIN_Y_EMU)
        violations = find_capacity_violations(audits)
        self.assertNotIn("body_margin_mismatch", {violation.rule for violation in violations})

    def test_full_width_text_body_keeps_paragraph_spacing_contract(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Paragraph Spacing",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Paragraph Spacing", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Раздел",
                    text="Первый абзац основного текста.\nВторой абзац основного текста.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        body = next(shape for shape in slide.placeholders if shape.placeholder_format.idx == 14)
        paragraphs = [paragraph for paragraph in body.text_frame.paragraphs if paragraph.text.strip()]
        self.assertTrue(paragraphs)
        expected_body_style = manifest.theme.master_text_styles["body"]
        for paragraph in paragraphs:
            with self.subTest(text=paragraph.text):
                self.assertEqual(paragraph.line_spacing, expected_body_style.line_spacing)
                if expected_body_style.space_after_pt is not None:
                    self.assertEqual(paragraph.space_after.pt, expected_body_style.space_after_pt)

    def test_full_width_bullets_keep_paragraph_spacing_contract(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Audit Bullet Paragraph Spacing",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Audit Bullet Paragraph Spacing", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Раздел",
                    bullets=["Первый пункт", "Второй пункт", "Третий пункт"],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[1]
        body = next(shape for shape in slide.placeholders if shape.placeholder_format.idx == 14)
        paragraphs = [paragraph for paragraph in body.text_frame.paragraphs if paragraph.text.strip()]
        self.assertTrue(paragraphs)
        expected_body_style = manifest.theme.master_text_styles["body"]
        for paragraph in paragraphs:
            with self.subTest(text=paragraph.text):
                self.assertEqual(paragraph.line_spacing, expected_body_style.line_spacing)
                if expected_body_style.space_after_pt is not None:
                    self.assertEqual(paragraph.space_after.pt, expected_body_style.space_after_pt)

    def test_cover_text_frames_keep_margin_contract(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Стратегия А3",
            slides=[
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Стратегия А3",
                    notes="Горизонт планирования: 2026-2030\nМарт 2026",
                    preferred_layout_key="cover",
                )
            ],
        )

        manifest = self.registry.get_template("corp_light_v1")
        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        text_shapes = [shape for shape in slide.shapes if getattr(shape, "has_text_frame", False)]
        self.assertTrue(text_shapes)
        for shape in text_shapes:
            with self.subTest(shape=getattr(shape, "name", "")):
                text_frame = shape.text_frame
                self.assertEqual(text_frame.margin_left, PptxGenerator.DEFAULT_TEXT_MARGIN_X_EMU)
                self.assertEqual(text_frame.margin_right, PptxGenerator.DEFAULT_TEXT_MARGIN_X_EMU)
                self.assertEqual(text_frame.margin_top, PptxGenerator.DEFAULT_TEXT_MARGIN_Y_EMU)
                self.assertEqual(text_frame.margin_bottom, PptxGenerator.DEFAULT_TEXT_MARGIN_Y_EMU)

    def test_cover_component_style_controls_cover_layout_and_fonts(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Cover Style Contract",
            slides=[
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Тестовый заголовок",
                    notes="Линия метаданных",
                    preferred_layout_key="cover",
                )
            ],
        )

        manifest = self.registry.get_template("corp_light_v1").model_copy(deep=True)
        cover_style = manifest.component_styles["cover"]
        cover_style.spacing_tokens["title_top_emu"] = 900000
        cover_style.spacing_tokens["title_left_emu"] = 700000
        cover_style.spacing_tokens["meta_top_emu"] = 3400000
        cover_style.text_styles["title"] = cover_style.text_styles["title"].model_copy(update={"font_size_pt": 40.0})
        cover_style.text_styles["meta"] = cover_style.text_styles["meta"].model_copy(update={"font_size_pt": 19.0})

        template_path = self.registry.get_template_pptx_path("corp_light_v1")
        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = self.generator.generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        slide = presentation.slides[0]
        title_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and "Тестовый заголовок" in shape.text)
        meta_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and "Линия метаданных" in shape.text)

        title_sizes = {
            run.font.size.pt
            for paragraph in title_shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None
        }
        meta_sizes = {
            run.font.size.pt
            for paragraph in meta_shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.font.size is not None
        }
        self.assertEqual(title_shape.top, 900000)
        self.assertEqual(title_shape.left, 700000)
        self.assertEqual(meta_shape.top, 3400000)
        self.assertEqual(title_sizes, {40.0})
        self.assertEqual(meta_sizes, {19.0})

    def test_prototype_templates_keep_margin_contract_on_nonempty_text_shapes(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists() or manifest.generation_mode.value != "prototype":
                continue

            with self.subTest(template_id=manifest.template_id):
                plan = PresentationPlan(
                    template_id=manifest.template_id,
                    title=f"{manifest.display_name} Prototype Margin Contract",
                    slides=self._smoke_slides_for_manifest(manifest),
                )

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = self.generator.generate(
                        template_path=template_path,
                        manifest=manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))

                nonempty_text_shapes = [
                    shape
                    for slide in presentation.slides
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False) and (shape.text or "").strip()
                ]
                self.assertTrue(nonempty_text_shapes)
                for shape in nonempty_text_shapes:
                    with self.subTest(template_id=manifest.template_id, shape=getattr(shape, "name", "")):
                        text_frame = shape.text_frame
                        self.assertGreaterEqual(text_frame.margin_left, 0)
                        self.assertGreaterEqual(text_frame.margin_right, 0)
                        self.assertGreaterEqual(text_frame.margin_top, 0)
                        self.assertGreaterEqual(text_frame.margin_bottom, 0)

    def test_prototype_templates_render_nonempty_bound_content(self) -> None:
        manifests = self.registry.list_templates()

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists() or manifest.generation_mode.value != "prototype":
                continue

            with self.subTest(template_id=manifest.template_id):
                plan = PresentationPlan(
                    template_id=manifest.template_id,
                    title=f"{manifest.display_name} Prototype Content Contract",
                    slides=self._smoke_slides_for_manifest(manifest),
                )

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = self.generator.generate(
                        template_path=template_path,
                        manifest=manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))

                rendered_payload = "\n".join(
                    shape.text.strip()
                    for slide in presentation.slides
                    for shape in slide.shapes
                    if getattr(shape, "has_text_frame", False) and (shape.text or "").strip()
                )
                has_visual_payload = any(
                    getattr(shape, "has_table", False)
                    or getattr(shape, "has_chart", False)
                    or hasattr(shape, "image")
                    for slide in presentation.slides
                    for shape in slide.shapes
                )
                self.assertTrue(rendered_payload.strip() or has_visual_payload)

    def test_prototype_template_chart_image_binding_renders_chart_shape(self) -> None:
        manifests = [
            manifest
            for manifest in self.registry.list_templates()
            if manifest.generation_mode.value == "prototype"
            and any(any(token.binding == "chart_image" for token in prototype.tokens) for prototype in manifest.prototype_slides)
        ]
        self.assertTrue(manifests)

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists():
                continue

            with self.subTest(template_id=manifest.template_id):
                plan = PresentationPlan(
                    template_id=manifest.template_id,
                    title=f"{manifest.display_name} Chart Contract",
                    slides=[
                        SlideSpec(kind=SlideKind.TITLE, title=f"{manifest.display_name} Chart Contract", preferred_layout_key="cover"),
                        SlideSpec(
                            kind=SlideKind.CHART,
                            title="Выручка",
                            subtitle="Prototype chart_image binding должен стать реальным chart shape",
                            chart=ChartSpec(
                                chart_id="chart_prototype",
                                source_table_id="table_1",
                                chart_type=ChartType.COLUMN,
                                title="Выручка",
                                categories=["Q1", "Q2", "Q3"],
                                series=[ChartSeries(name="Выручка", values=[120.0, 200.0, 90.0])],
                                confidence=ChartConfidence.HIGH,
                            ),
                            preferred_layout_key="table",
                        ),
                    ],
                )

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = self.generator.generate(
                        template_path=template_path,
                        manifest=manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))

                chart_shapes = [
                    shape
                    for slide in presentation.slides
                    for shape in slide.shapes
                    if getattr(shape, "has_chart", False)
                ]
                self.assertEqual(len(chart_shapes), 1)

    def test_prototype_template_chart_image_binding_keeps_chart_audit_green(self) -> None:
        manifests = [
            manifest
            for manifest in self.registry.list_templates()
            if manifest.generation_mode.value == "prototype"
            and any(any(token.binding == "chart_image" for token in prototype.tokens) for prototype in manifest.prototype_slides)
        ]
        self.assertTrue(manifests)

        for manifest in manifests:
            template_path = self.settings.templates_dir / manifest.template_id / manifest.source_pptx
            if not template_path.exists():
                continue

            with self.subTest(template_id=manifest.template_id):
                plan = PresentationPlan(
                    template_id=manifest.template_id,
                    title=f"{manifest.display_name} Chart Audit Contract",
                    slides=[
                        SlideSpec(
                            kind=SlideKind.TITLE,
                            title=f"{manifest.display_name} Chart Audit Contract",
                            preferred_layout_key="cover",
                        ),
                        SlideSpec(
                            kind=SlideKind.CHART,
                            title="Выручка и маржа",
                            subtitle="Prototype chart_image binding должен проходить template-aware chart audit",
                            chart=ChartSpec(
                                chart_id="chart_prototype_audit",
                                source_table_id="table_1",
                                chart_type=ChartType.COMBO,
                                title="Выручка и маржа",
                                categories=["Q1", "Q2", "Q3"],
                                series=[
                                    ChartSeries(name="Выручка", values=[104_300_000.0, 111_300_000.0, 135_700_000.0], unit="RUB"),
                                    ChartSeries(name="Маржа", values=[18.0, 22.0, 27.0], unit="%"),
                                ],
                                confidence=ChartConfidence.HIGH,
                                value_format="number",
                            ),
                            preferred_layout_key="table",
                        ),
                    ],
                )

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = self.generator.generate(
                        template_path=template_path,
                        manifest=manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    audits = audit_generated_presentation(output_path, plan, manifest)
                    violations = find_capacity_violations(audits)

                chart_audit = next(audit for audit in audits if audit.kind == SlideKind.CHART.value)
                self.assertTrue(chart_audit.has_chart)
                self.assertEqual(chart_audit.expected_chart_type, ChartType.COMBO.value)
                self.assertEqual(chart_audit.rendered_chart_type, ChartType.COMBO.value)
                self.assertTrue(chart_audit.expected_chart_secondary_value_axis)
                self.assertTrue(chart_audit.rendered_chart_secondary_value_axis)
                self.assertEqual(chart_audit.expected_chart_value_axis_number_format, '0.0,," млн ₽"')
                self.assertEqual(chart_audit.rendered_chart_value_axis_number_format, '0.0,," млн ₽"')
                self.assertEqual(chart_audit.expected_chart_secondary_value_axis_number_format, '0"%"')
                self.assertEqual(chart_audit.rendered_chart_secondary_value_axis_number_format, '0"%"')
                self.assertFalse(any(item.rule == "narrow_chart_content" for item in violations))
                self.assertEqual(violations, [])

    def _smoke_slides_for_manifest(self, manifest) -> list[SlideSpec]:
        slides = [SlideSpec(kind=SlideKind.TITLE, title=f"{manifest.display_name} Smoke", preferred_layout_key="cover")]
        supported_kinds = self._supported_kinds(manifest)

        if SlideKind.TEXT.value in supported_kinds:
            slides.append(
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Smoke text",
                    text="Smoke generation validates runtime compatibility between manifest and generator.",
                )
            )
            return slides

        if SlideKind.BULLETS.value in supported_kinds:
            slides.append(
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Smoke bullets",
                    bullets=["Smoke generation", "Manifest compatibility", "Generator compatibility"],
                )
            )
            return slides

        if SlideKind.TABLE.value in supported_kinds:
            slides.append(
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Smoke table",
                    table=TableBlock(headers=["Metric", "Value"], rows=[["GMV", "125"]]),
                )
            )
            return slides

        slides.append(
            SlideSpec(
                kind=SlideKind.TEXT,
                title="Fallback smoke",
                text="Fallback slide for templates without explicit text support.",
            )
        )
        return slides

    def _supported_kinds(self, manifest) -> set[str]:
        supported: set[str] = set()
        for layout in manifest.layouts:
            supported.update(layout.supported_slide_kinds)
        for prototype in manifest.prototype_slides:
            supported.update(prototype.supported_slide_kinds)
        return supported

    def _has_geometry_metadata(self, spec) -> bool:
        return all(
            isinstance(value, int) and value > 0
            for value in (spec.left_emu, spec.top_emu, spec.width_emu, spec.height_emu)
        )

    def _has_text_margin_metadata(self, spec) -> bool:
        return all(
            isinstance(value, int) and value >= 0
            for value in (
                spec.margin_left_emu,
                spec.margin_right_emu,
                spec.margin_top_emu,
                spec.margin_bottom_emu,
            )
        )

    def _text_frame_font_sizes(self, text_frame) -> list[float]:
        sizes: list[float] = []
        for paragraph in text_frame.paragraphs:
            if paragraph.font.size is not None:
                sizes.append(paragraph.font.size.pt)
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
        return sizes


if __name__ == "__main__":
    unittest.main()
