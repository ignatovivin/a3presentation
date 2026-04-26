from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation
from pptx.chart.axis import ValueAxis
from pptx.shapes.autoshape import Shape
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_MARKER_STYLE
from pptx.enum.text import MSO_AUTO_SIZE

from a3presentation.domain.api import ChartOverride, DocumentBlock
from a3presentation.domain.chart import ChartConfidence, ChartSeries, ChartSpec, ChartType
from a3presentation.domain.presentation import (
    PresentationPlan,
    RenderTargetType,
    SlideContentBlock,
    SlideContentBlockKind,
    SlideKind,
    SlideSpec,
    TableBlock,
)
from a3presentation.services.layout_capacity import DENSE_TEXT_FULL_WIDTH_PROFILE, LIST_FULL_WIDTH_PROFILE, TEXT_FULL_WIDTH_PROFILE
from a3presentation.services.pptx_generator import PptxGenerator
from a3presentation.services.planner import ContinuationUnit, Section, TextToPlanService
from a3presentation.services.template_registry import TemplateRegistry
from a3presentation.settings import get_settings


class TextToPlanServiceTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        cls.manifest = registry.get_template("corp_light_v1")
        cls.template_path = registry.get_template_pptx_path("corp_light_v1")

    def test_build_plan_handles_large_text_sections_without_crashing(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="Большой аналитический раздел", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Первое длинное предложение объясняет контекст задачи и описывает ограничения модели. "
                    "Второе предложение добавляет параметры для расчета юнит-экономики по каждому сегменту. "
                    "Третье предложение вводит сравнение с конкурентами и связывает выводы с рекомендацией для CEO. "
                    "Четвертое предложение завершает блок и гарантирует, что раздел не помещается в один слайд."
                ),
            ),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Дополнительный абзац нужен, чтобы планировщик переходил к логике разбиения большого раздела. "
                    "Именно этот путь раньше падал из-за обращения к отсутствующей константе LIST_BULLET_MAX_CHARS."
                ),
            ),
            DocumentBlock(kind="heading", text="Вторая секция", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Эта секция нужна только для того, чтобы первая точно не была воспринята как титульная и skipped при сборке плана."
                ),
            ),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="A3 Presentation",
            blocks=blocks,
        )

        self.assertGreaterEqual(len(plan.slides), 2)
        self.assertTrue(all(slide.preferred_layout_key is None for slide in plan.slides))
        self.assertEqual(plan.slides[0].runtime_profile_key, "cover")
        self.assertTrue(any(slide.runtime_profile_key for slide in plan.slides[1:]))
        self.assertTrue(all(slide.render_target is not None for slide in plan.slides))
        self.assertTrue(all(slide.render_target.type == RenderTargetType.LAYOUT for slide in plan.slides if slide.render_target is not None))
        self.assertEqual(plan.slides[0].render_target.key, "cover")

    def test_build_plan_clears_legacy_template_layout_keys_from_output(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="Контакты", level=1),
            DocumentBlock(kind="paragraph", text="Иван Иванов"),
            DocumentBlock(kind="paragraph", text="+7 999 123-45-67"),
            DocumentBlock(kind="paragraph", text="ivan@example.com"),
            DocumentBlock(kind="heading", text="Иллюстрация", level=1),
            DocumentBlock(kind="paragraph", text="Подпись к изображению"),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="A3 Presentation",
            blocks=blocks,
        )

        self.assertTrue(plan.slides)
        self.assertTrue(all(slide.preferred_layout_key is None for slide in plan.slides))
        self.assertTrue(all(slide.runtime_profile_key for slide in plan.slides))
        self.assertTrue(all(slide.render_target is not None for slide in plan.slides))
        self.assertTrue(all((slide.render_target.key or "").strip() for slide in plan.slides if slide.render_target is not None))

    def test_safe_fallback_planner_marks_render_target_as_degraded(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Сводный операционный отчёт"),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Значение"], rows=[["GMV", "100"]])),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Значение"], rows=[["MAU", "200"]])),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Значение"], rows=[["NPS", "65"]])),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        degraded_targets = [slide.render_target for slide in plan.slides[1:] if slide.render_target is not None]
        self.assertTrue(degraded_targets)
        self.assertTrue(any("document_fallback" in target.degradation_reasons for target in degraded_targets))
        self.assertTrue(any(target.type == RenderTargetType.AUTO_LAYOUT for target in degraded_targets))

    def test_resume_fallback_planner_marks_render_target_as_degraded(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Иван Игнатов"),
            DocumentBlock(kind="paragraph", text="ivan@example.com"),
            DocumentBlock(kind="paragraph", text="+7 999 000-00-00"),
            DocumentBlock(kind="paragraph", text="ОПЫТ РАБОТЫ"),
            DocumentBlock(
                kind="paragraph",
                text="Руководил командой из десяти человек, запускал внутренние продукты, выстраивал процессы аналитики и координировал кросс-функциональные проекты.",
            ),
            DocumentBlock(kind="paragraph", text="ОБРАЗОВАНИЕ"),
            DocumentBlock(
                kind="paragraph",
                text="Высшее техническое образование, дополнительная программа по продуктовому менеджменту и регулярные курсы по аналитике данных.",
            ),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        degraded_targets = [slide.render_target for slide in plan.slides[1:] if slide.render_target is not None]
        self.assertTrue(degraded_targets)
        self.assertTrue(any("resume_fallback" in target.degradation_reasons for target in degraded_targets))
        self.assertTrue(any(target.type == RenderTargetType.AUTO_LAYOUT for target in degraded_targets))

    def test_planner_splits_dense_bullets_by_content_weight(self) -> None:
        service = TextToPlanService()
        bullets = [
            "Первый очень длинный пункт объясняет стратегическую развилку, перечисляет ограничения, содержит несколько зависимых условий для принятия решения, раскрывает влияние на экономику сегмента, описывает организационные последствия и дополняется уточнением по срокам внедрения.",
            "Второй очень длинный пункт продолжает тему, добавляет экономические критерии, риски реализации, ожидаемый эффект для нескольких сегментов, перечисляет технологические ограничения, требования к команде сопровождения и допущения по инфраструктуре.",
            "Третий очень длинный пункт суммирует логику, связывает рекомендации с KPI, инвестиционным циклом и влиянием на портфель поставщиков, а также отдельно фиксирует ограничения по данным, срокам и допустимой нагрузке на операционный контур.",
        ]
        chunks = service._chunk_bullets_for_slides(bullets)
        self.assertGreaterEqual(len(chunks), 2)

    def test_rebalance_continuation_group_balances_short_tail_bullet_slide(self) -> None:
        service = TextToPlanService()
        slides = [
            SlideSpec(
                kind=SlideKind.BULLETS,
                title="Партнерская стратегия",
                bullets=[
                    "Первый пункт подробно описывает рыночный контекст, ограничения онбординга, требования к SLA и связь с масштабом каталога для внешних партнеров.",
                    "Второй пункт связывает устойчивость интеграционного слоя с управляемой операционной моделью сопровождения и предсказуемым контуром изменений.",
                    "Третий пункт фиксирует роль каталога как массового охвата и отделяет его от глубокого интеграционного слоя для ключевых провайдеров.",
                ],
                preferred_layout_key="list_full_width",
            ),
            SlideSpec(
                kind=SlideKind.BULLETS,
                title="Партнерская стратегия (2)",
                bullets=[
                    "Четвертый пункт объясняет, как двухконтурная модель усиливает переговорную позицию в сделках с банками и экосистемами.",
                ],
                preferred_layout_key="list_full_width",
            ),
        ]

        rebalanced = service._rebalance_single_continuation_group(slides)

        self.assertLessEqual(len(rebalanced), 2)
        self.assertTrue(all(slide.kind == SlideKind.BULLETS for slide in rebalanced))
        total_bullets = sum(len(slide.bullets) for slide in rebalanced)
        self.assertEqual(total_bullets, 4)
        if len(rebalanced) == 2:
            self.assertGreaterEqual(len(rebalanced[1].bullets), 2)

    def test_rebalance_continuation_group_balances_mixed_bullets_and_text_tail(self) -> None:
        service = TextToPlanService()
        slides = [
            SlideSpec(
                kind=SlideKind.BULLETS,
                title='Как "5000" усиливает партнерства',
                bullets=[
                    "Если компания ограничивает управляемый слой поставщиков, она может удерживать качество интеграций и не размывать операционный контур.",
                    "Большой каталог при этом остается внешним сигналом масштаба для банков, супер-аппов и государственных каналов.",
                    "Такое разделение помогает объяснить рынку две цифры: охват каталога и количество глубоко интегрированных поставщиков.",
                ],
                preferred_layout_key="list_full_width",
            ),
            SlideSpec(
                kind=SlideKind.TEXT,
                title='Как "5000" усиливает партнерства (2)',
                text=(
                    "Надежность и предсказуемость онбординга превращаются в понятный продуктовый оффер, "
                    "а не в ручной консалтинг для каждого нового поставщика."
                ),
                preferred_layout_key="text_full_width",
            ),
        ]

        rebalanced = service._rebalance_single_continuation_group(slides)

        self.assertLessEqual(len(rebalanced), 2)
        tail_payload = " ".join(
            part
            for slide in rebalanced
            for part in [slide.text or "", slide.notes or "", *slide.bullets]
            if part
        )
        self.assertIn("Надежность и предсказуемость", tail_payload)
        self.assertTrue(any(slide.kind in {SlideKind.TEXT, SlideKind.BULLETS} for slide in rebalanced))

    def test_slide_to_continuation_units_rechunks_long_text_content_blocks(self) -> None:
        service = TextToPlanService()
        long_text = (
            "Первый абзац описывает стратегический контекст, ограничения интеграций и требования к SLA для high-touch слоя. "
            "Второй абзац переносит фокус на экономику сопровождения, объясняет рост нагрузки на команды и связь с unit-экономикой. "
            "Третий абзац фиксирует, почему слишком длинный paragraph нельзя оставлять атомарным при continuation rebalance, "
            "иначе последний слайд группы получается недозаполненным и визуально слабым."
        )
        slide = SlideSpec(
            kind=SlideKind.TEXT,
            title="Длинный narrative",
            content_blocks=[
                SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text=long_text),
            ],
            preferred_layout_key="text_full_width",
        )

        units = service._slide_to_continuation_units(slide)

        self.assertGreaterEqual(len(units), 2)
        self.assertTrue(all(unit.kind == "paragraph" for unit in units))
        self.assertEqual(" ".join(unit.text for unit in units), long_text)

    def test_split_large_text_section_rebalances_short_tail_via_finer_text_units(self) -> None:
        service = TextToPlanService()
        section = Section(
            title="Длинный narrative",
            paragraphs=[
                "A" * 299,
                "B" * 299,
                "C" * 446,
                "D" * 426,
                "E" * 420,
                "F" * 441,
                "G" * 108,
            ],
        )

        slides = service._split_large_section(section)

        self.assertGreaterEqual(len(slides), 2)
        self.assertEqual(slides[-1].kind, SlideKind.TEXT)
        self.assertLessEqual(len(slides), 7)

    def test_compact_continuation_slides_merges_short_leading_text_slide(self) -> None:
        service = TextToPlanService()
        slides = [
            SlideSpec(kind=SlideKind.TEXT, title="Раздел", text="A" * 140, preferred_layout_key="text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел (2)", text="B" * 320, preferred_layout_key="text_full_width"),
        ]

        compacted = service._compact_continuation_slides(slides, "Раздел", "")

        self.assertEqual(len(compacted), 1)
        self.assertEqual(compacted[0].title, "Раздел")
        self.assertGreaterEqual(len(compacted[0].text or "") + len(compacted[0].notes or ""), 460)

    def test_compact_continuation_slides_merges_short_single_bullet_tail(self) -> None:
        service = TextToPlanService()
        slides = [
            SlideSpec(
                kind=SlideKind.BULLETS,
                title="Раздел",
                bullets=["Очень короткий хвост"],
                content_blocks=[SlideContentBlock(kind=SlideContentBlockKind.BULLET_LIST, items=["Очень короткий хвост"])],
                preferred_layout_key="list_full_width",
            ),
            SlideSpec(
                kind=SlideKind.BULLETS,
                title="Раздел (2)",
                bullets=["Первый тезис", "Второй тезис", "Третий тезис"],
                content_blocks=[
                    SlideContentBlock(
                        kind=SlideContentBlockKind.BULLET_LIST,
                        items=["Первый тезис", "Второй тезис", "Третий тезис"],
                    )
                ],
                preferred_layout_key="list_full_width",
            ),
        ]

        compacted = service._compact_continuation_slides(slides, "Раздел", "")

        self.assertEqual(len(compacted), 1)
        self.assertEqual(compacted[0].title, "Раздел")
        self.assertGreaterEqual(len(compacted[0].bullets), 4)

    def test_compact_continuation_slides_keeps_pure_text_pair_when_group_has_subtitle(self) -> None:
        service = TextToPlanService()
        slides = [
            SlideSpec(kind=SlideKind.TEXT, title="Раздел", subtitle="Подзаголовок", text="A" * 320, preferred_layout_key="dense_text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел (2)", text="B" * 211, preferred_layout_key="dense_text_full_width"),
        ]

        compacted = service._compact_continuation_slides(slides, "Раздел", "Подзаголовок")

        self.assertEqual(len(compacted), 2)
        self.assertEqual(compacted[0].subtitle, "Подзаголовок")
        self.assertEqual(compacted[1].subtitle, "")

    def test_build_single_slide_keeps_numeric_metric_bullets_out_of_kpi_slide_layout(self) -> None:
        service = TextToPlanService()
        section = Section(
            title="A3 GIS",
            bullet_lists=[[
                "99,5 % успешных поисков",
                "53 млн активных начислений",
                "92 млн документов для поиска",
                "0,86 % среднее время поиска в секундах",
            ]],
        )

        slide = service._build_single_slide(section)

        self.assertEqual(slide.kind, SlideKind.BULLETS)
        self.assertNotEqual(slide.preferred_layout_key, "cards_kpi")

    def test_preferred_textual_layout_keeps_mixed_paragraph_dominant_slide_in_text_layout(self) -> None:
        service = TextToPlanService()
        layout_key = service._preferred_textual_layout_key(
            [
                SlideContentBlock(
                    kind=SlideContentBlockKind.PARAGRAPH,
                    text="Плотный narrative-блок объясняет логику сегмента, ограничения роста и экономику сопровождения.",
                ),
                SlideContentBlock(
                    kind=SlideContentBlockKind.PARAGRAPH,
                    text="Второй абзац продолжает narrative и должен оставаться в text layout, а не в широком bullet layout.",
                ),
                SlideContentBlock(
                    kind=SlideContentBlockKind.BULLET_LIST,
                    items=["Один уточняющий тезис"],
                ),
            ],
            total_chars=240,
        )

        self.assertEqual(layout_key, "text_full_width")

    def test_continuation_units_fit_single_slide_uses_text_capacity_for_mixed_text_render(self) -> None:
        service = TextToPlanService()
        units = [
            ContinuationUnit(
                kind="bullet",
                text="Предсказуемость процесса онбординга и сопровождения должна работать как продуктовый оффер.",
            ),
            ContinuationUnit(
                kind="paragraph",
                text="Ключевой внешний тезис объясняет, что ограничение управляемого слоя не уменьшает масштаб компании, если отдельно показывать охват каталога и число глубоко интегрированных поставщиков.",
            ),
            ContinuationUnit(
                kind="paragraph",
                text="Внутренний тезис связывает это с unit-экономикой, cost-to-serve и необходимостью считать маржинальный вклад каждого нового поставщика по сегментам.",
            ),
            ContinuationUnit(
                kind="paragraph",
                text="Завершающий абзац добавляет требования к данным, операционным метрикам и логике распределения затрат, из-за чего такой mixed bucket уже не должен считаться помещающимся в один text slide.",
            ),
        ]

        self.assertFalse(service._continuation_units_fit_single_slide(units))

    def test_rebalance_continuation_buckets_fills_short_leading_text_bucket(self) -> None:
        service = TextToPlanService()
        buckets = [
            [ContinuationUnit(kind="paragraph", text="A" * 128)],
            [
                ContinuationUnit(kind="paragraph", text="B" * 220),
                ContinuationUnit(kind="paragraph", text="C" * 220),
            ],
        ]

        rebalanced = service._rebalance_continuation_buckets(buckets, max_chars=900, max_weight=12.25)

        self.assertEqual(len(rebalanced[0]), 2)
        self.assertEqual(sum(len(unit.text) for unit in rebalanced[0]), 348)
        self.assertEqual(len(rebalanced[1]), 1)

    def test_rebalance_continuation_buckets_smooths_adjacent_payload_sizes(self) -> None:
        service = TextToPlanService()
        buckets = [
            [
                ContinuationUnit(kind="paragraph", text="A" * 220),
            ],
            [
                ContinuationUnit(kind="paragraph", text="C" * 220),
                ContinuationUnit(kind="paragraph", text="D" * 220),
                ContinuationUnit(kind="paragraph", text="E" * 220),
            ],
        ]

        rebalanced = service._rebalance_continuation_buckets(buckets, max_chars=900, max_weight=12.25)

        self.assertEqual([sum(len(unit.text) for unit in bucket) for bucket in rebalanced], [440, 440])

    def test_rebalance_continuation_groups_does_not_shift_later_section_ranges(self) -> None:
        service = TextToPlanService()
        slides = [
            SlideSpec(kind=SlideKind.TEXT, title="Раздел A", subtitle="Подзаголовок A", text="A" * 180, preferred_layout_key="text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел A (2)", text="B" * 120, preferred_layout_key="text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел A (3)", text="C" * 120, preferred_layout_key="text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел B", subtitle="Подзаголовок B", text="D" * 220, preferred_layout_key="text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел B (2)", text="E" * 220, preferred_layout_key="text_full_width"),
            SlideSpec(kind=SlideKind.TEXT, title="Раздел B (3)", text="F" * 220, preferred_layout_key="text_full_width"),
        ]

        rebalanced = service._rebalance_continuation_groups(slides)
        titles = [service._base_slide_title(slide.title) for slide in rebalanced]

        self.assertIn("Раздел A", titles)
        self.assertIn("Раздел B", titles)
        first_b_index = titles.index("Раздел B")
        self.assertTrue(all(title == "Раздел A" for title in titles[:first_b_index]))
        self.assertTrue(all(title == "Раздел B" for title in titles[first_b_index:]))

    def test_chunk_text_for_slides_keeps_short_inline_heading_with_following_paragraph(self) -> None:
        service = TextToPlanService()
        first_paragraph = (
            "А3 оперирует в уникальной нише на стыке инфраструктурных финтех-решений для регулярных платежей. "
            "Конкурентное поле формируют косвенные игроки, которые закрывают отдельные сегменты цепочки ценности, "
            "но не воспроизводят модель А3 целиком. Этот абзац специально сделан достаточно длинным, чтобы раздел "
            "ушел в continuation-режим и планировщик делил его на несколько текстовых слайдов."
        )
        second_paragraph = (
            "Ни один конкурент не воспроизводит модель А3 целиком: мультибанковская сеть, агрегация поставщиков "
            "и прямой доступ к инфраструктурным государственным сервисам формируют барьер входа. Главные риски "
            "связаны с крупными банками, которые строят собственную инфраструктуру, поэтому дальше нужен отдельный "
            "смысловой блок с преимуществами и митигацией."
        )
        chunks = service._refine_text_continuation_units(
            [first_paragraph, "Конкурентные преимущества А3", second_paragraph]
        )
        batches = service._chunk_text_for_slides(chunks)
        self.assertGreaterEqual(len(batches), 2)
        first_slide_payload = " ".join(batches[0])
        second_slide_payload = " ".join(batches[1])

        self.assertNotIn("Конкурентные преимущества А3", first_slide_payload)
        self.assertIn("Конкурентные преимущества А3", second_slide_payload)

    def test_choose_text_continuation_batches_can_select_dense_text_profile(self) -> None:
        service = TextToPlanService()
        chunks = ["A" * 430, "B" * 430, "C" * 430, "D" * 430]

        batches, profile = service._choose_text_continuation_batches(chunks, title="Плотный narrative", subtitle="")

        self.assertEqual(profile.layout_key, DENSE_TEXT_FULL_WIDTH_PROFILE.layout_key)
        self.assertEqual(len(batches), 2)
        self.assertTrue(all(sum(len(chunk) for chunk in batch) >= 800 for batch in batches))

    def test_split_large_section_uses_dense_text_layout_for_balanced_dense_narrative(self) -> None:
        service = TextToPlanService()
        section = Section(
            title="Плотный narrative",
            paragraphs=["A" * 430, "B" * 430, "C" * 430, "D" * 430],
        )

        slides = service._split_large_section(section)

        self.assertEqual(len(slides), 2)
        self.assertTrue(all(slide.preferred_layout_key == "dense_text_full_width" for slide in slides))

    def test_cover_meta_stays_compact_for_long_first_section(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Большая стратегическая презентация", level=0),
            DocumentBlock(kind="heading", text="Контекст задачи и что именно считать поставщиком", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    'В регулярных платежах метрика "количество поставщиков" почти всегда многозначна и требует отдельной '
                    "декомпозиции по юридическим лицам, договорам, биллинговым витринам и маршрутам списаний."
                ),
            ),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "На рынке обычно существуют как минимум три уровня поставщиков, и каждый по-разному влияет на экономику, "
                    "операционный контур и переговорную позицию компании."
                ),
            ),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий текст для второй секции."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="Большая стратегическая презентация",
            blocks=blocks,
        )

        cover = plan.slides[0]
        self.assertEqual(cover.preferred_layout_key, "cover")
        self.assertLessEqual(len((cover.notes or "").splitlines()), 2)
        self.assertNotIn("многозначна", cover.notes or "")
        self.assertNotIn("как минимум три уровня", cover.notes or "")

    def test_long_lead_paragraph_is_not_promoted_to_default_subtitle(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="heading", text="Раздел стратегии", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Это длинный вводный narrative-абзац, который объясняет контекст, ограничения, экономику и "
                    "операционные последствия, поэтому он не должен превращаться в subtitle."
                ),
            ),
            DocumentBlock(kind="paragraph", text="Второй абзац продолжает аналитическую мысль."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="Раздел стратегии",
            blocks=blocks,
        )

        self.assertEqual(plan.slides[1].title, "Раздел стратегии")
        self.assertFalse(plan.slides[1].subtitle)

    def test_first_content_section_does_not_repeat_cover_title_when_subheading_exists(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="heading", text="Большая стратегическая презентация", level=1),
            DocumentBlock(kind="subheading", text='Контекст задачи и что именно считать "поставщиком"', level=2),
            DocumentBlock(
                kind="paragraph",
                text=(
                    'В регулярных платежах метрика "количество поставщиков" почти всегда многозначна и требует '
                    "явного определения на уровне методологии."
                ),
            ),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Второй абзац нужен, чтобы первая содержательная секция точно рендерилась как самостоятельный "
                    "content slide, а не схлопывалась в пограничный короткий кейс."
                ),
            ),
            DocumentBlock(kind="subheading", text="Следующий раздел", level=2),
            DocumentBlock(kind="paragraph", text="Второй раздел нужен для проверки структуры."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="Большая стратегическая презентация",
            blocks=blocks,
        )

        self.assertEqual(plan.slides[0].title, "Большая стратегическая презентация")
        self.assertEqual(plan.slides[1].title, 'Контекст задачи и что именно считать "поставщиком"')
        self.assertNotEqual(plan.slides[1].title, plan.slides[0].title)

    def test_cover_keeps_short_meta_lines_when_present(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Стратегия развития А3", level=0),
            DocumentBlock(kind="paragraph", text="Горизонт планирования: 2026-2030"),
            DocumentBlock(kind="paragraph", text="Март 2026"),
            DocumentBlock(kind="paragraph", text="Конфиденциальный документ"),
            DocumentBlock(kind="heading", text="Первый раздел", level=1),
            DocumentBlock(kind="paragraph", text="Основной контент презентации."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="Стратегия развития А3",
            blocks=blocks,
        )

        cover = plan.slides[0]
        self.assertEqual(cover.notes, "Горизонт планирования: 2026-2030\nМарт 2026\nКонфиденциальный документ")

    def test_cover_with_five_leading_lines_does_not_create_duplicate_content_slide(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="А3"),
            DocumentBlock(kind="paragraph", text="Бизнес-стратегия 2026"),
            DocumentBlock(kind="paragraph", text="Горизонт планирования: 2026-2030"),
            DocumentBlock(kind="paragraph", text="Март 2026"),
            DocumentBlock(kind="paragraph", text="Конфиденциальный документ"),
            DocumentBlock(kind="heading", text="1. Vision и стратегические цели", level=1),
            DocumentBlock(kind="paragraph", text="Основной контент раздела."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        self.assertEqual(plan.slides[0].title, "А3 Бизнес-стратегия 2026")
        self.assertFalse(any(slide.title == "А3" for slide in plan.slides[1:]))

    def test_mixed_section_preserves_paragraphs_when_bullets_force_split(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="A3"),
            DocumentBlock(kind="heading", text="5.3 R&D новые продукты", level=1),
            DocumentBlock(
                kind="paragraph",
                text="Цель: построить системный пайплайн проверки гипотез и новых продуктов для роста выручки и диверсификации бизнеса А3.",
            ),
            DocumentBlock(
                kind="list",
                items=[
                    "Процесс Double Diamond запущен, первые инициативы в работе",
                    "Банк идей = единый бэклог, синхронизирован с Jira",
                    "ICE-приоритизация со стейкхолдерами раз в 2 недели",
                    "Юнит-экономика по поставщикам и партнёрам",
                    "Инициатива: Инфраструктура цифрового рубля",
                ],
            ),
            DocumentBlock(
                kind="paragraph",
                text="Q2 2026 (Discovery): Анализ архитектуры платформы цифрового рубля ЦБ и определение роли А3.",
            ),
            DocumentBlock(
                kind="list",
                items=[
                    "Q3 2026 (Proof of Concept): Техническое прототипирование с командой Payments",
                    "Q4 2026 — Пилот: Пилотный запуск с ограниченным объемом",
                ],
            ),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        bullets_text = "\n".join("\n".join(slide.bullets) for slide in plan.slides if slide.kind == SlideKind.BULLETS)
        self.assertIn("Цель: построить системный пайплайн проверки гипотез", bullets_text)
        self.assertIn("Q2 2026 (Discovery): Анализ архитектуры платформы цифрового рубля ЦБ", bullets_text)
        flattened = [
            item
            for slide in plan.slides
            if slide.kind == SlideKind.BULLETS and slide.title and slide.title.startswith("5.3 R&D новые продукты")
            for item in slide.bullets
        ]
        q2_index = next(index for index, item in enumerate(flattened) if item.startswith("Q2 2026 (Discovery)"))
        q3_index = next(index for index, item in enumerate(flattened) if item.startswith("Q3 2026 (Proof of Concept)"))
        self.assertLess(q2_index, q3_index)

    def test_mixed_section_preserves_paragraph_list_paragraph_order(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Стратегия партнерств"),
            DocumentBlock(kind="heading", text="Контур взаимодействия", level=1),
            DocumentBlock(kind="paragraph", text="Вводный абзац задает контекст и объясняет ограничение управляемого слоя."),
            DocumentBlock(kind="list", items=["Первый список фиксирует массовый каталог.", "Второй список фиксирует SLA-контур."]),
            DocumentBlock(kind="paragraph", text="Заключающий абзац должен остаться после списка, а не переехать перед ним."),
            DocumentBlock(kind="heading", text="Следующий блок", level=1),
            DocumentBlock(kind="paragraph", text="Граница следующего раздела."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        flattened = []
        for slide in plan.slides:
            if not (slide.title or "").startswith("Контур взаимодействия"):
                continue
            if slide.content_blocks:
                for block in slide.content_blocks:
                    if block.text and block.text.strip():
                        flattened.append(block.text.strip())
                    flattened.extend(item.strip() for item in block.items if item.strip())
                continue
            flattened.extend(item for item in slide.bullets if item.strip())
        intro_index = next(index for index, item in enumerate(flattened) if item.startswith("Вводный абзац"))
        list_index = next(index for index, item in enumerate(flattened) if item.startswith("Первый список"))
        outro_index = next(index for index, item in enumerate(flattened) if item.startswith("Заключающий абзац"))
        self.assertLess(intro_index, list_index)
        self.assertLess(list_index, outro_index)
        mixed_slides = [
            slide for slide in plan.slides if (slide.title or "").startswith("Контур взаимодействия") and slide.content_blocks
        ]
        self.assertTrue(mixed_slides)
        flattened_kinds = [block.kind for slide in mixed_slides for block in slide.content_blocks]
        self.assertIn(SlideContentBlockKind.PARAGRAPH, flattened_kinds)
        self.assertIn(SlideContentBlockKind.BULLET_LIST, flattened_kinds)

    def test_planner_marks_question_and_callout_content_blocks(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Стратегия платформы"),
            DocumentBlock(kind="heading", text="Ключевые развилки", level=1),
            DocumentBlock(kind="paragraph", text="Вопрос: почему нужен второй контур инфраструктуры?"),
            DocumentBlock(kind="paragraph", text="Важно: без резервного контура SLA для критичных платежей останется хрупким."),
            DocumentBlock(kind="paragraph", text="Итог: это решение уменьшает риск простоя в критическом платёжном контуре."),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий завершающий блок."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        target_slides = [slide for slide in plan.slides if (slide.title or "").startswith("Ключевые развилки")]
        self.assertTrue(target_slides)
        block_kinds = [block.kind for slide in target_slides for block in slide.content_blocks]
        self.assertIn(SlideContentBlockKind.QA_ITEM, block_kinds)
        self.assertIn(SlideContentBlockKind.CALLOUT, block_kinds)

    def test_semantic_question_and_callout_section_prefers_text_layout(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Стратегия платформы"),
            DocumentBlock(kind="heading", text="FAQ и выводы", level=1),
            DocumentBlock(kind="paragraph", text="Вопрос: почему нужен второй контур инфраструктуры?"),
            DocumentBlock(kind="paragraph", text="Важно: резервный контур удерживает SLA в часы пиковой нагрузки."),
            DocumentBlock(kind="paragraph", text="Итог: архитектура снижает риск простоя без усложнения операционной модели."),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий завершающий блок."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        target_slides = [slide for slide in plan.slides if (slide.title or "").startswith("FAQ и выводы")]
        self.assertEqual(len(target_slides), 1)
        slide = target_slides[0]
        self.assertEqual(slide.kind, SlideKind.TEXT)
        self.assertEqual(slide.preferred_layout_key, "text_full_width")
        self.assertEqual(
            [block.kind for block in slide.content_blocks],
            [
                SlideContentBlockKind.QA_ITEM,
                SlideContentBlockKind.CALLOUT,
                SlideContentBlockKind.CALLOUT,
            ],
        )

    def test_semantic_question_and_callout_with_bullets_prefers_list_layout(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Стратегия платформы"),
            DocumentBlock(kind="heading", text="Разбор сценария", level=1),
            DocumentBlock(kind="paragraph", text="Вопрос: где основной источник операционного риска?"),
            DocumentBlock(kind="list", items=["Пиковая нагрузка на платежный шлюз.", "Единая точка отказа в процессинге."]),
            DocumentBlock(kind="paragraph", text="Важно: резервирование нужно внедрять до масштабирования канала."),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий завершающий блок."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        target_slides = [slide for slide in plan.slides if (slide.title or "").startswith("Разбор сценария")]
        self.assertTrue(target_slides)
        self.assertTrue(all(slide.preferred_layout_key == "list_full_width" for slide in target_slides))
        self.assertTrue(any(block.kind == SlideContentBlockKind.QA_ITEM for slide in target_slides for block in slide.content_blocks))
        self.assertTrue(any(block.kind == SlideContentBlockKind.BULLET_LIST for slide in target_slides for block in slide.content_blocks))

    def test_mixed_continuation_with_paragraph_dominance_stays_text_slide(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Стратегия платформы"),
            DocumentBlock(kind="heading", text="Гибридный раздел", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Первый длинный абзац задает контекст, объясняет управленческую проблему и формирует narrative, "
                    "который не должен автоматически превращаться в bullet list."
                ),
            ),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Второй длинный абзац продолжает аналитическую мысль, связывает экономику поставщика с архитектурой "
                    "продукта и по смыслу должен остаться paragraph flow."
                ),
            ),
            DocumentBlock(kind="list", items=["Короткий тезис.", "Еще один тезис."]),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий завершающий блок."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        target_slides = [slide for slide in plan.slides if (slide.title or "").startswith("Гибридный раздел")]
        self.assertTrue(target_slides)
        self.assertEqual(target_slides[0].kind, SlideKind.TEXT)
        self.assertTrue(any(block.kind == SlideContentBlockKind.BULLET_LIST for block in target_slides[0].content_blocks))

    def test_preferred_textual_layout_uses_dense_text_for_paragraph_dominant_mixed_payload(self) -> None:
        service = TextToPlanService()
        layout_key = service._preferred_textual_layout_key(
            [
                SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text="A" * 260),
                SlideContentBlock(kind=SlideContentBlockKind.PARAGRAPH, text="B" * 260),
                SlideContentBlock(kind=SlideContentBlockKind.BULLET_LIST, items=["Первый короткий тезис.", "Второй короткий тезис."]),
            ],
            total_chars=560,
        )

        self.assertEqual(layout_key, "dense_text_full_width")

    def test_continuation_units_fit_single_slide_uses_dense_text_for_paragraph_dominant_mixed_payload(self) -> None:
        service = TextToPlanService()
        units = [
            ContinuationUnit(kind="paragraph", text="A" * 260),
            ContinuationUnit(kind="paragraph", text="B" * 260),
            ContinuationUnit(kind="bullet", text="Первый короткий тезис."),
            ContinuationUnit(kind="bullet", text="Второй короткий тезис."),
        ]

        self.assertTrue(service._continuation_units_fit_single_slide(units))

    def test_build_continuation_slide_uses_dense_text_for_paragraph_dominant_mixed_payload(self) -> None:
        service = TextToPlanService()
        units = [
            ContinuationUnit(kind="paragraph", text="A" * 260),
            ContinuationUnit(kind="paragraph", text="B" * 260),
            ContinuationUnit(kind="bullet", text="Первый короткий тезис."),
            ContinuationUnit(kind="bullet", text="Второй короткий тезис."),
        ]

        slide = service._build_continuation_slide("Гибридный раздел", "", units)

        self.assertEqual(slide.kind, SlideKind.TEXT)
        self.assertEqual(slide.preferred_layout_key, "dense_text_full_width")
        self.assertTrue(any(block.kind == SlideContentBlockKind.BULLET_LIST for block in slide.content_blocks))

    def test_build_continuation_slide_uses_dense_text_for_pure_narrative_payload(self) -> None:
        service = TextToPlanService()
        units = [
            ContinuationUnit(kind="paragraph", text="A" * 430),
            ContinuationUnit(kind="paragraph", text="B" * 430),
        ]

        slide = service._build_continuation_slide("Плотный narrative", "", units)

        self.assertEqual(slide.kind, SlideKind.TEXT)
        self.assertEqual(slide.preferred_layout_key, "dense_text_full_width")

    def test_continuation_units_fit_single_slide_uses_dense_text_for_pure_narrative_payload(self) -> None:
        service = TextToPlanService()
        units = [
            ContinuationUnit(kind="paragraph", text="A" * 430),
            ContinuationUnit(kind="paragraph", text="B" * 430),
        ]

        self.assertTrue(service._continuation_units_fit_single_slide(units))

    def test_first_section_with_tables_is_not_swallowed_by_cover(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="АНКЕТА КАНДИДАТА"),
            DocumentBlock(
                kind="table",
                table=TableBlock(
                    headers=["Поле", "Значение"],
                    rows=[["ФИО", "Игнатов Иван"], ["Телефон", "+7 900 000-00-00"]],
                ),
            ),
            DocumentBlock(kind="paragraph", text="ОБРАЗОВАНИЕ"),
            DocumentBlock(
                kind="table",
                table=TableBlock(
                    headers=["Период", "Учебное заведение", "Статус"],
                    rows=[["2018-2022", "Университет", "Высшее"]],
                ),
            ),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        self.assertGreaterEqual(len(plan.slides), 3)
        self.assertEqual(plan.slides[0].preferred_layout_key, "cover")
        self.assertTrue(any(slide.table is not None for slide in plan.slides[1:]))

    def test_planner_skips_appendix_like_source_section_from_main_deck(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="Стратегический обзор"),
            DocumentBlock(kind="heading", text="1. Контекст", level=1),
            DocumentBlock(kind="paragraph", text="Основной narrative раздел должен остаться в deck."),
            DocumentBlock(kind="heading", text="Какие источники данных использовать", level=1),
            DocumentBlock(kind="paragraph", text="https://example.com/source-1"),
            DocumentBlock(kind="paragraph", text="https://example.com/source-2"),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        titles = [slide.title or "" for slide in plan.slides]
        self.assertTrue(any(title.startswith("1. Контекст") for title in titles))
        self.assertFalse(any("Какие источники данных использовать" in title for title in titles))

    def test_non_narrative_document_uses_safe_fallback_planner(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="АНКЕТА КАНДИДАТА"),
            DocumentBlock(
                kind="table",
                table=TableBlock(
                    headers=["Поле", "Значение"],
                    rows=[["ФИО", "Игнатов Иван"], ["Телефон", "+7 900 000-00-00"]],
                ),
            ),
            DocumentBlock(kind="paragraph", text="ОБРАЗОВАНИЕ"),
            DocumentBlock(
                kind="table",
                table=TableBlock(
                    headers=["Период", "Учебное заведение", "Статус"],
                    rows=[["2018-2022", "Университет", "Высшее"]],
                ),
            ),
            DocumentBlock(kind="paragraph", text="ДОПОЛНИТЕЛЬНЫЕ СВЕДЕНИЯ"),
            DocumentBlock(
                kind="paragraph",
                text="Подписывая настоящую Анкету, кандидат выражает согласие на обработку персональных данных в целях трудоустройства.",
            ),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        self.assertGreaterEqual(len(plan.slides), 4)
        self.assertEqual(plan.slides[1].preferred_layout_key, "list_full_width")
        self.assertTrue(any(slide.preferred_layout_key == "text_full_width" for slide in plan.slides[1:]))
        self.assertTrue(any(slide.preferred_layout_key == "table" for slide in plan.slides[1:]))

    def test_resume_document_uses_resume_fallback(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Иван Игнатов"),
            DocumentBlock(kind="paragraph", text="ivan@example.com"),
            DocumentBlock(kind="paragraph", text="+7 999 000-00-00"),
            DocumentBlock(kind="paragraph", text="ОПЫТ РАБОТЫ"),
            DocumentBlock(
                kind="paragraph",
                text="Руководил командой из десяти человек, запускал внутренние продукты, выстраивал процессы аналитики и координировал кросс-функциональные проекты.",
            ),
            DocumentBlock(kind="paragraph", text="ОБРАЗОВАНИЕ"),
            DocumentBlock(
                kind="paragraph",
                text="Высшее техническое образование, дополнительная программа по продуктовому менеджменту и регулярные курсы по аналитике данных.",
            ),
            DocumentBlock(kind="paragraph", text="НАВЫКИ"),
            DocumentBlock(kind="list", items=["Product management", "SQL", "Презентации"]),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        self.assertGreaterEqual(len(plan.slides), 3)
        self.assertEqual(plan.slides[1].preferred_layout_key, "list_full_width")
        self.assertTrue(any(slide.preferred_layout_key == "text_full_width" for slide in plan.slides[1:]))
        self.assertFalse(any(slide.preferred_layout_key == "table" for slide in plan.slides[1:]))

    def test_table_heavy_document_adds_table_count_summary(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Сводный операционный отчёт"),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Значение"], rows=[["GMV", "100"]])),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Значение"], rows=[["MAU", "200"]])),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Значение"], rows=[["NPS", "65"]])),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title=None,
            tables=[],
            blocks=blocks,
        )

        self.assertGreaterEqual(len(plan.slides), 4)
        self.assertEqual(plan.slides[1].preferred_layout_key, "list_full_width")
        self.assertTrue(any("Таблиц в документе" in bullet for bullet in plan.slides[1].bullets))

    def test_planner_replaces_selected_table_with_chart_slide(self) -> None:
        service = TextToPlanService()
        table = TableBlock(
            headers=["Канал", "Лиды"],
            rows=[["SEO", "120"], ["Ads", "200"], ["Referral", "90"]],
        )

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="Отчет по каналам",
            title="Отчет по каналам",
            tables=[table],
            chart_overrides=[
                ChartOverride(
                    table_id="table_1",
                    mode="chart",
                    selected_chart=ChartSpec(
                        chart_id="chart_1",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Лиды по каналам",
                        categories=["SEO", "Ads", "Referral"],
                        series=[ChartSeries(name="Лиды", values=[120.0, 200.0, 90.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                )
            ],
        )

        chart_slides = [slide for slide in plan.slides if slide.kind == SlideKind.CHART]
        self.assertEqual(len(chart_slides), 1)
        self.assertEqual(chart_slides[0].source_table_id, "table_1")
        self.assertIsNotNone(chart_slides[0].chart)
        self.assertEqual(chart_slides[0].title, "Лиды по каналам")
        self.assertEqual(chart_slides[0].subtitle, "")
        self.assertIsNone(chart_slides[0].text)
        self.assertEqual(chart_slides[0].bullets, [])
        self.assertIsNone(chart_slides[0].table)

    def test_split_table_keeps_compact_two_column_ten_row_table_on_single_slide(self) -> None:
        service = TextToPlanService()
        table = TableBlock(
            headers=["Показатель", "Значение"],
            rows=[
                ["Выручка с НДС", "2 433 млн ₽"],
                ["Выручка без НДС", "2 122 млн ₽"],
                ["Нетранзакционный доход", "470 млн ₽ и выше"],
                ["Доля Альфа-Банка", "< 40% выручки"],
                ["НДС от выручки", "%"],
                ["РНКО", "Лицензия ЦБ и запуск у 1 партнёра"],
                ["А3 Лаб", "Создание продукта для формирования п/п"],
                ["Poseidon", "Промышленное использование у 6+ партнёров"],
                ["Банки-партнёры", "3-4"],
                ["Поставщики", "+1 000 новых подключений, итого 3 000+"],
            ],
        )

        chunks = service._split_table_for_slides(table)

        self.assertEqual(len(chunks), 1)

    def test_split_table_keeps_compact_five_column_market_map_on_single_slide(self) -> None:
        service = TextToPlanService()
        table = TableBlock(
            headers=["Игрок", "Тип", "Фокус", "Риск", "Комментарий"],
            rows=[
                ["Сбер", "Банк-экосистема", "ЖКХ, подписки, платежи", "Высокий", "Есть уникальный сервис, но слабый аккаунтинг и медлительность"],
                ["ВТБ", "Банк-экосистема", "Платежи, подписки", "Высокий", "Партнёр А3, пока нет мультибанковой сети"],
                ["Т Банк", "Банк-экосистема", "Платежи, персонализация", "Средний", "Нет мультибанковой сети и планов на её создание"],
                ["Монета", "Платёжные сервисы", "Онлайн-платежи, ЖКХ", "Средний", "Есть B2C и B2B модель, слабый аккаунтинг"],
                ["Элекснет", "Платёжный сервис", "ЖКХ, интернет, связь", "Низкий", "Есть мультибанковская сеть, нет доступа к ГИС"],
                ["Система Город", "Региональный агрегатор", "Оплата ЖКХ, локальные поставщики", "Низкий", "Локальный охват без масштаба и технологий"],
                ["Бурмистр", "SaaS для ЖКХ", "Управление УК и оплата ЖКХ", "Низкий", "Другой слой цепочки, нет банковской сети"],
            ],
        )

        chunks = service._split_table_for_slides(table)

        self.assertEqual(len(chunks), 1)

    def test_split_table_keeps_seven_row_three_column_journey_on_single_slide(self) -> None:
        service = TextToPlanService()
        table = TableBlock(
            headers=["Этап", "Статус", "Комментарий"],
            rows=[
                ["01. Лид", "Работает", "CRM / ГИС ЖКХ / холодный поиск / портал ЖКХ"],
                ["02. Договор", "Работает", ""],
                ["03. Интеграция", "Конкур. преим.", "AI: 5 минут, аналитики и PM"],
                [
                    "04. Открытие банками",
                    "Самостоятельное открытие банками",
                    "Банки сами решают срок открытия, внедряем систему для активного открытия и автоуведомлений",
                ],
                ["05. Выход на объём", "6 мес. лаг", "Около 6 месяцев до нормального платёжного потока"],
                ["06. Поддержка", "Улучшить", ""],
                ["07. Метрики", "В процессе", "NPS/CSAT внедряется"],
            ],
        )

        chunks = service._split_table_for_slides(table)

        self.assertEqual(len(chunks), 1)

    def test_split_table_reduces_over_fragmentation_for_risk_tables(self) -> None:
        service = TextToPlanService()
        table = TableBlock(
            headers=["Риск", "Вероятность", "Митигация"],
            rows=[
                ["Сбер/ВТБ строят собственную инфраструктуру", "Средняя", "Развитие уникальных технологий и партнёрств"],
                ["Регуляторные изменения ЦБ", "Средняя", "Проактивный мониторинг и РНКО как инструмент независимости"],
                ["Уход ключевых партнёров", "Низкая", "Диверсификация долей и рост новых банков"],
                ["Отмена реестра российского ПО", "Средняя", "Альтернативные модели налоговой оптимизации"],
                ["AI у конкурентов", "Средняя", "Инвестиции в AI и R&D фабрику гипотез"],
                ["Ужесточение налоговой политики", "Средняя", "Сценарий уже заложен в бюджет"],
                ["Задержка РНКО на год", "Средняя", "Скорректировать стратегию без смены курса"],
                ["Рост негатива из-за СМЭВ-согласий", "Средняя", "Коммуникационная стратегия и работа с брендом"],
                ["Потеря доступа к данным ГИС", "Низкая", "Прямые договоры и альтернативные источники данных"],
                [
                    "Макроэкономические факторы: высокая ключевая ставка, инфляция и санкционное давление",
                    "Высокая",
                    "Даже при сохранении текущего давления регулярные обязательные платежи остаются антицикличным сегментом и поддерживают рост GMV",
                ],
            ],
        )

        chunks = service._split_table_for_slides(table)

        self.assertLessEqual(len(chunks), 2)

    def test_planner_and_generator_keep_compact_table_on_single_slide(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="1.3 Цели 2026", level=1),
            DocumentBlock(
                kind="table",
                table=TableBlock(
                    headers=["Показатель", "Значение"],
                    rows=[
                        ["Выручка с НДС", "2 433 млн ₽"],
                        ["Выручка без НДС", "2 122 млн ₽"],
                        ["Нетранзакционный доход", "470 млн ₽ и выше"],
                        ["Доля Альфа-Банка", "< 40% выручки"],
                        ["НДС от выручки", "%"],
                        ["РНКО", "Лицензия ЦБ и запуск у 1 партнёра"],
                        ["А3 Лаб", "Создание продукта для формирования п/п"],
                        ["Poseidon", "Промышленное использование у 6+ партнёров"],
                        ["Банки-партнёры", "3-4"],
                        ["Поставщики", "+1 000 новых подключений, итого 3 000+"],
                    ],
                ),
            ),
            DocumentBlock(kind="heading", text="2. Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Текст для предотвращения fallback-сценария."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="A3 Presentation",
            blocks=blocks,
        )

        table_slides = [slide for slide in plan.slides if slide.kind == SlideKind.TABLE and slide.title == "1.3 Цели 2026"]
        self.assertEqual(len(table_slides), 1)
        self.assertEqual(len(table_slides[0].table.rows), 10)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), len(plan.slides))
            table_shapes = [shape for shape in presentation.slides[1].shapes if getattr(shape, "has_table", False)]
            self.assertEqual(len(table_shapes), 1)
            rendered_table = table_shapes[0].table
            self.assertEqual(len(rendered_table.rows), 11)
            self.assertEqual(rendered_table.cell(1, 0).text, "Выручка с НДС")

    def test_planner_and_generator_keep_competitor_map_in_single_rendered_table(self) -> None:
        service = TextToPlanService()
        section = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="3.3 Карта конкурентов", level=1),
            DocumentBlock(kind="paragraph", text="Контекст конкурентной карты с вводным описанием рынка."),
            DocumentBlock(
                kind="table",
                table=TableBlock(
                    headers=["Игрок", "Тип", "Фокус", "Риск", "Комментарий"],
                    rows=[
                        ["Сбер", "Банк-экосистема", "ЖКХ, подписки, платежи", "Высокий", "Есть уникальный сервис, но слабый аккаунтинг и медлительность"],
                        ["ВТБ", "Банк-экосистема", "Платежи, подписки", "Высокий", "Партнёр А3, пока нет мультибанковой сети"],
                        ["Т Банк", "Банк-экосистема", "Платежи, персонализация", "Средний", "Нет мультибанковой сети и планов на её создание"],
                        ["Монета", "Платёжные сервисы", "Онлайн-платежи, ЖКХ", "Средний", "Есть B2C и B2B модель, слабый аккаунтинг"],
                        ["Элекснет", "Платёжный сервис", "ЖКХ, интернет, связь", "Низкий", "Есть мультибанковская сеть, нет доступа к ГИС"],
                        ["Система Город", "Региональный агрегатор", "Оплата ЖКХ, локальные поставщики", "Низкий", "Локальный охват без масштаба и технологий"],
                        ["Бурмистр", "SaaS для ЖКХ", "Управление УК и оплата ЖКХ", "Низкий", "Другой слой цепочки, нет банковской сети"],
                    ],
                ),
            ),
            DocumentBlock(kind="heading", text="4. Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Текст для устойчивой narrative-классификации."),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in section),
            title="A3 Presentation",
            blocks=section,
        )

        table_slides = [slide for slide in plan.slides if slide.kind == SlideKind.TABLE and "3.3 Карта конкурентов" in (slide.title or "")]
        self.assertEqual(len(table_slides), 1)
        self.assertEqual(len(table_slides[0].table.rows), 7)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            matching_slides = []
            for slide in presentation.slides:
                slide_text = " ".join(shape.text for shape in slide.shapes if getattr(shape, "has_text_frame", False))
                if "3.3 Карта конкурентов" in slide_text:
                    matching_slides.append(slide)
            rendered_table_slides = [slide for slide in matching_slides if any(getattr(shape, "has_table", False) for shape in slide.shapes)]
            self.assertEqual(len(rendered_table_slides), 1)
            table_shape = next(shape for shape in rendered_table_slides[0].shapes if getattr(shape, "has_table", False))
            self.assertEqual(len(table_shape.table.rows), 8)

    def test_generator_expands_long_title_and_pushes_content_down(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Очень длинный заголовок слайда для проверки автоматического увеличения высоты заголовка и сдвига основного текста вниз без наложения на контент",
                    subtitle="Короткий подзаголовок",
                    text="Основной текст слайда должен оказаться ниже заголовка после перерасчета вертикального потока.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            slide = presentation.slides[1]
            placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}

            self.assertGreater(placeholders[0].height, 600000)
            self.assertEqual(placeholders[0].width, 11198224)
            self.assertGreater(placeholders[13].top, placeholders[0].top + placeholders[0].height)
            self.assertGreater(placeholders[14].top, 1791494)
            title_runs = [run for paragraph in placeholders[0].text_frame.paragraphs for run in paragraph.runs]
            self.assertTrue(title_runs)
            self.assertLessEqual(title_runs[0].font.size.pt, 30)

    def test_generator_keeps_short_title_compact_and_readable(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Короткий заголовок",
                    text="Основной текст без подзаголовка.",
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            slide = presentation.slides[1]
            placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}
            title_runs = [run for paragraph in placeholders[0].text_frame.paragraphs for run in paragraph.runs]

            self.assertTrue(title_runs)
            self.assertEqual(title_runs[0].font.size.pt, self.manifest.theme.master_text_styles["title"].font_size_pt)
            self.assertLess(placeholders[0].height, 1200000)
            self.assertGreater(placeholders[14].top, placeholders[0].top + placeholders[0].height)

    def test_generator_reduces_body_font_for_dense_text(self) -> None:
        dense_text = (
            "Первый длинный абзац описывает контекст, ограничения, допущения и критерии принятия решения по продуктовой стратегии. "
            "Второй длинный абзац добавляет финансовые ориентиры, риски реализации и требования к инфраструктуре. "
            "Третий длинный абзац связывает выводы с целевыми показателями, KPI и дорожной картой внедрения."
        )

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Плотный текстовый слайд",
                    text=dense_text,
                    preferred_layout_key="text_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            slide = presentation.slides[1]
            placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}
            body_runs = [run for paragraph in placeholders[14].text_frame.paragraphs for run in paragraph.runs]

            self.assertTrue(body_runs)
            self.assertIsNotNone(body_runs[0].font.size)
            self.assertEqual(body_runs[0].font.size.pt, 18.0)
            self.assertNotEqual(placeholders[14].text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_generator_applies_explicit_font_size_to_sparse_bullet_slide(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="R&D continuation",
                    bullets=[
                        "Юнит-экономика: по поставщикам и интеграциям",
                        "Инициатива: Инфраструктура цифрового рубля",
                        "Q3 2026: Техническое прототипирование",
                        "Q4 2026: Пилотный запуск и оценка unit-экономики",
                    ],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            placeholders = {shape.placeholder_format.idx: shape for shape in presentation.slides[1].placeholders}
            body = placeholders[14]
            body_runs = [run for paragraph in body.text_frame.paragraphs for run in paragraph.runs]

            self.assertTrue(body_runs)
            self.assertTrue(all(run.font.size is not None for run in body_runs))
            self.assertAlmostEqual(body_runs[0].font.size.pt, 18.0, places=1)
            self.assertNotEqual(body.text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_generator_shrinks_dense_bullet_container_to_avoid_overflow(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.BULLETS,
                    title="Плотный appendix slide",
                    bullets=[
                        "Q10: Зачем собственная инфраструктура Poseidon и второй ЦОД для критической финансовой платформы с жёсткими SLA и требованиями по контролю среды.",
                        "Q11: Почему целевыми считаются малые банки, несмотря на концентрацию выручки в верхнем сегменте и более скромный текущий оборот.",
                        "Q12: Что будет происходить после подключения 3000 поставщиков и какие операционные ограничения нужно закрыть заранее.",
                        "Контроль SLA: внешние решения не дают достаточного уровня гарантии для финансовой инфраструктуры и государственных интеграций.",
                        "Независимость от вендоров: критична для скорости изменений, контроля инцидентов и стоимости владения платформой.",
                        "Защита данных: госинтеграции требуют более высокого уровня контроля над инфраструктурой и доступом.",
                        "Инфраструктурное проникновение как барьер для конкурентов и основа для новых продуктов.",
                        "Дополнительный заработок на минимальных контрактах и последующем upsell поверх инфраструктурного контура.",
                    ],
                    preferred_layout_key="list_full_width",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            placeholders = {shape.placeholder_format.idx: shape for shape in presentation.slides[1].placeholders}
            body = placeholders[14]
            body_runs = [run for paragraph in body.text_frame.paragraphs for run in paragraph.runs]

            self.assertTrue(body_runs)
            self.assertLessEqual(body_runs[0].font.size.pt, 18.0)
            self.assertGreaterEqual(body_runs[0].font.size.pt, LIST_FULL_WIDTH_PROFILE.min_font_pt)
            self.assertNotEqual(body.text_frame.auto_size, MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE)

    def test_generator_styles_qa_and_callout_blocks_differently(self) -> None:
        generator = PptxGenerator()
        presentation = Presentation(str(self.template_path))
        slide = presentation.slides[0]
        shape: Shape = slide.shapes.add_textbox(0, 0, 6_000_000, 3_000_000)

        generator._set_content_blocks(
            shape,
            [
                SlideContentBlock(kind=SlideContentBlockKind.QA_ITEM, text="Почему нужен второй контур инфраструктуры?"),
                SlideContentBlock(
                    kind=SlideContentBlockKind.CALLOUT,
                    text="Важно: резервный контур снижает риск простоя в критическом платёжном контуре.",
                ),
                SlideContentBlock(
                    kind=SlideContentBlockKind.BULLET_LIST,
                    items=["Первый операционный тезис.", "Второй операционный тезис."],
                ),
            ],
            TEXT_FULL_WIDTH_PROFILE,
        )

        paragraphs = [paragraph for paragraph in shape.text_frame.paragraphs if paragraph.text.strip()]
        self.assertGreaterEqual(len(paragraphs), 4)
        qa_run = paragraphs[0].runs[0]
        callout_run = paragraphs[1].runs[0]
        bullet_run = paragraphs[2].runs[0]

        self.assertTrue(qa_run.font.bold)
        self.assertFalse(bool(qa_run.font.italic))
        self.assertTrue(callout_run.font.bold)
        self.assertTrue(callout_run.font.italic)
        self.assertFalse(bool(bullet_run.font.italic))

    def test_generator_keeps_footer_in_bottom_zone(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Очень длинное название презентации для проверки нижнего блока и корректного положения footer",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TEXT,
                    title="Тест footer",
                    text="Основной текст слайда.",
                    preferred_layout_key="text_full_width",
                ),
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Тест footer на табличном слайде",
                    subtitle="Подзаголовок",
                    table=TableBlock(headers=["Показатель", "Значение"], rows=[["GMV", "125"]]),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

            text_slide_placeholders = {shape.placeholder_format.idx: shape for shape in presentation.slides[1].placeholders}
            table_slide_placeholders = {shape.placeholder_format.idx: shape for shape in presentation.slides[2].placeholders}

            self.assertGreaterEqual(text_slide_placeholders[17].top, 6200000)
            self.assertGreaterEqual(table_slide_placeholders[15].top, 6200000)

    def test_generator_expands_table_footer_to_full_width_without_subtitle(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Таблица без подзаголовка",
                    table=TableBlock(headers=["Показатель", "Значение"], rows=[["GMV", "125"]]),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            placeholders = {shape.placeholder_format.idx: shape for shape in presentation.slides[1].placeholders}

            self.assertEqual(placeholders[15].left, 442913)
            self.assertEqual(placeholders[15].width, 11198224)
            self.assertGreaterEqual(placeholders[15].top, 6200000)

    def test_generator_keeps_table_subtitle_readable_instead_of_tiny_font(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.TABLE,
                    title="Раздел",
                    subtitle="6.1 РНКО — получение лицензии",
                    table=TableBlock(
                        headers=["Показатель", "Значение"],
                        rows=[["Статус", "В работе"], ["Срок", "Q4 2026"]],
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            subtitle = next(
                placeholder for placeholder in presentation.slides[1].placeholders if placeholder.placeholder_format.idx == 13
            )
            font_sizes = sorted(
                {
                    run.font.size.pt
                    for paragraph in subtitle.text_frame.paragraphs
                    for run in paragraph.runs
                    if run.font.size is not None
                }
            )

            self.assertTrue(font_sizes)
        self.assertEqual(font_sizes, [20.0])

    def test_generator_renders_chart_slide_into_pptx(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Лиды по каналам",
                    subtitle="Тест chart render",
                    chart=ChartSpec(
                        chart_id="chart_1",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Лиды",
                        categories=["SEO", "Ads", "Referral"],
                        series=[ChartSeries(name="Лиды", values=[120.0, 200.0, 90.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            slide = presentation.slides[1]
            chart_shapes = [shape for shape in slide.shapes if getattr(shape, "has_chart", False)]
            placeholders = {
                shape.placeholder_format.idx: shape
                for shape in slide.placeholders
                if getattr(shape, "is_placeholder", False)
            }

            self.assertEqual(len(chart_shapes), 1)
            title_sizes = sorted(
                {
                    run.font.size.pt
                    for paragraph in placeholders[0].text_frame.paragraphs
                    for run in paragraph.runs
                    if run.font.size is not None
                }
            )
            subtitle_sizes = sorted(
                {
                    run.font.size.pt
                    for paragraph in placeholders[13].text_frame.paragraphs
                    for run in paragraph.runs
                    if run.font.size is not None
                }
            )
            self.assertEqual(title_sizes, [self.manifest.theme.master_text_styles["title"].font_size_pt])
            self.assertEqual(subtitle_sizes, [20.0])
            self.assertEqual(chart_shapes[0].chart.series[0].name, "Лиды")
            self.assertEqual(chart_shapes[0].chart.series[0].format.fill.fore_color.rgb, RGBColor(0x67, 0x9A, 0xEA))
            self.assertEqual(
                chart_shapes[0].chart.chart_title.text_frame.paragraphs[0].runs[0].font.color.rgb,
                RGBColor(0x08, 0x1C, 0x4F),
            )

    def test_generator_styles_line_chart_markers_and_percent_format(self) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Конверсия",
                    chart=ChartSpec(
                        chart_id="chart_line",
                        source_table_id="table_1",
                        chart_type=ChartType.LINE,
                        title="Конверсия",
                        categories=["Q1", "Q2", "Q3"],
                        series=[ChartSeries(name="CR", values=[18.0, 22.0, 27.0])],
                        confidence=ChartConfidence.HIGH,
                        data_labels_visible=True,
                        value_format="percent",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))
            series = chart.series[0]

            self.assertEqual(series.marker.style, XL_MARKER_STYLE.CIRCLE)
            self.assertEqual(series.data_labels.number_format, '0"%"')

    def test_generator_renders_real_combo_chart_with_bar_and_line_plots(self) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Combo",
                    chart=ChartSpec(
                        chart_id="chart_combo",
                        source_table_id="table_1",
                        chart_type=ChartType.COMBO,
                        title="План и тренд",
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(name="План", values=[100.0, 130.0, 150.0]),
                            ChartSeries(name="Факт", values=[90.0, 120.0, 160.0]),
                            ChartSeries(name="Маржа", values=[18.0, 22.0, 27.0]),
                        ],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))

            bar_charts = chart._chartSpace.xpath(".//c:barChart")
            line_charts = chart._chartSpace.xpath(".//c:lineChart")

            self.assertEqual(len(bar_charts), 1)
            self.assertEqual(len(line_charts), 1)
            self.assertEqual(len(bar_charts[0].xpath("./c:ser")), 2)
            self.assertEqual(len(line_charts[0].xpath("./c:ser")), 1)
            self.assertEqual(chart.series[-1].marker.style, XL_MARKER_STYLE.CIRCLE)

    def test_generator_keeps_combo_as_column_when_line_series_is_hidden(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Combo hidden line",
                    chart=ChartSpec(
                        chart_id="chart_combo_hidden_line",
                        source_table_id="table_1",
                        chart_type=ChartType.COMBO,
                        title="План и тренд",
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(name="План", values=[100.0, 130.0, 150.0]),
                            ChartSeries(name="Факт", values=[90.0, 120.0, 160.0]),
                            ChartSeries(name="Маржа", values=[18.0, 22.0, 27.0], hidden=True),
                        ],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))

            bar_charts = chart._chartSpace.xpath(".//c:barChart")
            line_charts = chart._chartSpace.xpath(".//c:lineChart")

            self.assertEqual(len(bar_charts), 1)
            self.assertEqual(len(line_charts), 0)
            self.assertEqual(len(bar_charts[0].xpath("./c:ser")), 2)

    def test_generator_renders_secondary_value_axis_for_mixed_unit_combo(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Combo secondary axis",
                    chart=ChartSpec(
                        chart_id="chart_combo_secondary_axis",
                        source_table_id="table_1",
                        chart_type=ChartType.COMBO,
                        title="Выручка и маржа",
                        categories=["Q1", "Q2", "Q3"],
                        series=[
                            ChartSeries(name="Выручка", values=[104_300_000.0, 111_300_000.0, 135_700_000.0], unit="RUB"),
                            ChartSeries(name="Маржа", values=[18.0, 22.0, 27.0], unit="%"),
                        ],
                        confidence=ChartConfidence.HIGH,
                        value_format="number",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))
            chart_space = chart._chartSpace
            line_charts = chart_space.xpath(".//c:lineChart")
            value_axes = chart_space.xpath(".//c:valAx")

            self.assertEqual(len(line_charts), 1)
            self.assertEqual(len(value_axes), 2)
            self.assertEqual(ValueAxis(value_axes[0]).tick_labels.number_format, '0.0,," млн ₽"')
            self.assertEqual(ValueAxis(value_axes[1]).tick_labels.number_format, '0"%"')
            self.assertEqual(
                [element.get("val") for element in line_charts[0].xpath("./c:axId")],
                [chart_space.xpath(".//c:catAx/c:axId")[0].get("val"), value_axes[1].xpath("./c:axId")[0].get("val")],
            )

    def test_generator_formats_value_axis_in_millions_for_large_currency_values(self) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Доход",
                    chart=ChartSpec(
                        chart_id="chart_money",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Доход",
                        categories=["Q1", "Q2", "Q3"],
                        series=[ChartSeries(name="Доход", values=[104_300_000.0, 111_300_000.0, 135_700_000.0])],
                        confidence=ChartConfidence.HIGH,
                        value_format="currency",
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))

            self.assertEqual(chart.value_axis.tick_labels.number_format, '0.0,," млн ₽"')

    def test_generator_styles_pie_points_with_distinct_brand_colors(self) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Структура выручки",
                    chart=ChartSpec(
                        chart_id="chart_pie",
                        source_table_id="table_1",
                        chart_type=ChartType.PIE,
                        title="Структура",
                        categories=["A", "B", "C"],
                        series=[ChartSeries(name="Доли", values=[40.0, 35.0, 25.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))
            points = chart.series[0].points

            self.assertEqual(points[0].format.fill.fore_color.rgb, RGBColor(0x09, 0x1E, 0x38))
            self.assertEqual(points[1].format.fill.fore_color.rgb, RGBColor(0x34, 0x89, 0xF3))

    def test_generator_resolves_supported_chart_types_and_combo_fallback(self) -> None:
        generator = PptxGenerator()

        def make_spec(chart_type: ChartType) -> ChartSpec:
            return ChartSpec(
                chart_id=f"chart_{chart_type.value}",
                source_table_id="table_1",
                chart_type=chart_type,
                title="Тест",
                categories=["A", "B"],
                series=[ChartSeries(name="Series", values=[10.0, 20.0])],
                confidence=ChartConfidence.HIGH,
            )

        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.BAR)), XL_CHART_TYPE.BAR_CLUSTERED)
        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.COLUMN)), XL_CHART_TYPE.COLUMN_CLUSTERED)
        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.LINE)), XL_CHART_TYPE.LINE)
        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.STACKED_BAR)), XL_CHART_TYPE.BAR_STACKED)
        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.STACKED_COLUMN)), XL_CHART_TYPE.COLUMN_STACKED)
        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.PIE)), XL_CHART_TYPE.PIE)
        self.assertEqual(generator._resolve_chart_type(make_spec(ChartType.COMBO)), XL_CHART_TYPE.COLUMN_CLUSTERED)

    def test_generator_ranks_single_series_column_points_by_value_palette(self) -> None:
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="A3 Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.CHART,
                    title="Рейтинг значений",
                    chart=ChartSpec(
                        chart_id="chart_column_ranked",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Рейтинг",
                        categories=["A", "B", "C", "D"],
                        series=[ChartSeries(name="Значения", values=[40.0, 25.0, 35.0, 10.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                    preferred_layout_key="table",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            chart = next(shape.chart for shape in presentation.slides[1].shapes if getattr(shape, "has_chart", False))
            points = chart.series[0].points

            self.assertEqual(points[0].format.fill.fore_color.rgb, RGBColor(0x09, 0x1E, 0x38))
            self.assertEqual(points[2].format.fill.fore_color.rgb, RGBColor(0x34, 0x89, 0xF3))
            self.assertEqual(points[1].format.fill.fore_color.rgb, RGBColor(0x26, 0x45, 0x95))
            self.assertEqual(points[3].format.fill.fore_color.rgb, RGBColor(0xBF, 0xCE, 0xF5))

    def test_generator_renders_supported_chart_type_matrix(self) -> None:
        chart_specs = [
            ChartSpec(
                chart_id="chart_bar",
                source_table_id="table_1",
                chart_type=ChartType.BAR,
                title="Bar",
                categories=["A", "B", "C"],
                series=[ChartSeries(name="Series", values=[10.0, 20.0, 30.0])],
                confidence=ChartConfidence.HIGH,
            ),
            ChartSpec(
                chart_id="chart_column",
                source_table_id="table_2",
                chart_type=ChartType.COLUMN,
                title="Column",
                categories=["A", "B", "C"],
                series=[ChartSeries(name="Series", values=[10.0, 20.0, 30.0])],
                confidence=ChartConfidence.HIGH,
            ),
            ChartSpec(
                chart_id="chart_line",
                source_table_id="table_3",
                chart_type=ChartType.LINE,
                title="Line",
                categories=["A", "B", "C"],
                series=[ChartSeries(name="Series", values=[10.0, 20.0, 30.0])],
                confidence=ChartConfidence.HIGH,
            ),
            ChartSpec(
                chart_id="chart_stacked_bar",
                source_table_id="table_4",
                chart_type=ChartType.STACKED_BAR,
                title="Stacked bar",
                categories=["A", "B", "C"],
                series=[
                    ChartSeries(name="Series 1", values=[10.0, 20.0, 30.0]),
                    ChartSeries(name="Series 2", values=[8.0, 14.0, 18.0]),
                ],
                confidence=ChartConfidence.HIGH,
            ),
            ChartSpec(
                chart_id="chart_stacked_column",
                source_table_id="table_5",
                chart_type=ChartType.STACKED_COLUMN,
                title="Stacked column",
                categories=["A", "B", "C"],
                series=[
                    ChartSeries(name="Series 1", values=[10.0, 20.0, 30.0]),
                    ChartSeries(name="Series 2", values=[8.0, 14.0, 18.0]),
                ],
                confidence=ChartConfidence.HIGH,
            ),
            ChartSpec(
                chart_id="chart_pie",
                source_table_id="table_6",
                chart_type=ChartType.PIE,
                title="Pie",
                categories=["A", "B", "C"],
                series=[ChartSeries(name="Share", values=[40.0, 35.0, 25.0])],
                confidence=ChartConfidence.HIGH,
            ),
            ChartSpec(
                chart_id="chart_combo",
                source_table_id="table_7",
                chart_type=ChartType.COMBO,
                title="Combo",
                categories=["A", "B", "C"],
                series=[
                    ChartSeries(name="Bar 1", values=[100.0, 120.0, 140.0]),
                    ChartSeries(name="Bar 2", values=[90.0, 130.0, 150.0]),
                    ChartSeries(name="Line", values=[18.0, 22.0, 27.0]),
                ],
                confidence=ChartConfidence.HIGH,
            ),
        ]
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Chart Matrix",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Chart Matrix", preferred_layout_key="cover"),
                *[
                    SlideSpec(
                        kind=SlideKind.CHART,
                        title=spec.title,
                        chart=spec,
                        preferred_layout_key="table",
                    )
                    for spec in chart_specs
                ],
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))

        expected_xml = {
            ChartType.BAR: {"barChart": 1, "lineChart": 0, "pieChart": 0, "barDir": "bar", "grouping": "clustered"},
            ChartType.COLUMN: {"barChart": 1, "lineChart": 0, "pieChart": 0, "barDir": "col", "grouping": "clustered"},
            ChartType.LINE: {"barChart": 0, "lineChart": 1, "pieChart": 0},
            ChartType.STACKED_BAR: {"barChart": 1, "lineChart": 0, "pieChart": 0, "barDir": "bar", "grouping": "stacked"},
            ChartType.STACKED_COLUMN: {"barChart": 1, "lineChart": 0, "pieChart": 0, "barDir": "col", "grouping": "stacked"},
            ChartType.PIE: {"barChart": 0, "lineChart": 0, "pieChart": 1},
            ChartType.COMBO: {"barChart": 1, "lineChart": 1, "pieChart": 0},
        }
        for slide_index, spec in enumerate(chart_specs, start=2):
            with self.subTest(chart_type=spec.chart_type):
                chart = next(
                    shape.chart
                    for shape in presentation.slides[slide_index - 1].shapes
                    if getattr(shape, "has_chart", False)
                )
                chart_space = chart._chartSpace
                bar_charts = chart_space.xpath(".//c:barChart")
                line_charts = chart_space.xpath(".//c:lineChart")
                pie_charts = chart_space.xpath(".//c:pieChart")
                expected = expected_xml[spec.chart_type]
                self.assertEqual(len(bar_charts), expected["barChart"])
                self.assertEqual(len(line_charts), expected["lineChart"])
                self.assertEqual(len(pie_charts), expected["pieChart"])
                if bar_charts and "barDir" in expected:
                    self.assertEqual(next(element.get("val") for element in bar_charts[0].xpath("./c:barDir")), expected["barDir"])
                    self.assertEqual(next(element.get("val") for element in bar_charts[0].xpath("./c:grouping")), expected["grouping"])
                if spec.chart_type == ChartType.COMBO:
                    self.assertEqual(len(bar_charts[0].xpath("./c:ser")), 2)
                    self.assertEqual(len(line_charts[0].xpath("./c:ser")), 1)

    def test_generator_adapts_cover_title_height_and_meta_spacing(self) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="A3 Presentation",
            slides=[
                SlideSpec(
                    kind=SlideKind.TITLE,
                    title="Оптимальное количество поставщиков для А3 в регулярных платежах - методология, сценарии и рекомендации для CEO",
                    notes='Контекст задачи и что именно считать "поставщиком"',
                    preferred_layout_key="cover",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            slide = presentation.slides[0]
            placeholders = {shape.placeholder_format.idx: shape for shape in slide.placeholders}
            title = placeholders[0]
            meta = placeholders[15]
            title_runs = [run for paragraph in title.text_frame.paragraphs for run in paragraph.runs if run.text]

            self.assertTrue(title_runs)
            self.assertLessEqual(title_runs[0].font.size.pt, 46)
            self.assertGreater(meta.top, title.top + title.height)
            self.assertGreaterEqual(title.height, 1422646)

    def test_split_table_for_slides_preserves_fill_colors_across_chunks(self) -> None:
        service = TextToPlanService()
        table = TableBlock(
            headers=["Поставщик", "Статус"],
            header_fill_colors=["1F4E79", "1F4E79"],
            rows=[[f"Поставщик {index}", f"Статус {index}"] for index in range(1, 14)],
            row_fill_colors=[
                [None, "D9EAF7"] if index % 2 else ["FDE7D7", None]
                for index in range(1, 14)
            ],
        )

        chunks = service._split_table_for_slides(table)

        self.assertGreaterEqual(len(chunks), 2)
        self.assertTrue(all(chunk.header_fill_colors == ["1F4E79", "1F4E79"] for chunk in chunks))
        flattened_row_fill_colors = [row for chunk in chunks for row in chunk.row_fill_colors]
        self.assertEqual(flattened_row_fill_colors, table.row_fill_colors)


if __name__ == "__main__":
    unittest.main()
