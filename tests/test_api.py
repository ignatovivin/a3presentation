from __future__ import annotations

import asyncio
import importlib
import os
import shutil
import tempfile
import unittest
from io import BytesIO
from pathlib import Path
from unittest.mock import patch

from docx import Document
from fastapi import HTTPException
from pptx import Presentation
from starlette.datastructures import UploadFile

from a3presentation import main as main_module
from a3presentation import settings as settings_module
from a3presentation.api import routes as routes_module
from a3presentation.domain.template import TemplateManifest
from a3presentation.domain.api import TextPlanRequest
from a3presentation.domain.presentation import PresentationPlan, SlideKind, SlideSpec


class ApiContractTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls._env_backup = {key: os.environ.get(key) for key in ("TEMPLATES_DIR", "OUTPUTS_DIR", "STORAGE_DIR", "SEED_BUNDLED_TEMPLATES")}
        cls._temp_dir = tempfile.TemporaryDirectory()
        cls._root = Path(cls._temp_dir.name)
        cls._templates_dir = cls._root / "templates"
        cls._outputs_dir = cls._root / "outputs"
        cls._templates_dir.mkdir(parents=True, exist_ok=True)
        cls._outputs_dir.mkdir(parents=True, exist_ok=True)

        source_templates = Path(__file__).resolve().parents[1] / "storage" / "templates"
        for template_id in ("corp_light_v1",):
            shutil.copytree(source_templates / template_id, cls._templates_dir / template_id)
        manifest_path = cls._templates_dir / "missing_source" / "manifest.json"
        manifest_path.parent.mkdir(parents=True, exist_ok=True)
        manifest_path.write_text(
            TemplateManifest(
                template_id="missing_source",
                display_name="Missing Source",
                description="Manifest without pptx source file",
                source_pptx="missing.pptx",
                default_layout_key="cover",
                layouts=[],
            ).model_dump_json(),
            encoding="utf-8",
        )

        os.environ["TEMPLATES_DIR"] = str(cls._templates_dir)
        os.environ["OUTPUTS_DIR"] = str(cls._outputs_dir)
        os.environ["STORAGE_DIR"] = str(cls._root)

        importlib.reload(settings_module)
        importlib.reload(routes_module)
        importlib.reload(main_module)

    @classmethod
    def tearDownClass(cls) -> None:
        cls._temp_dir.cleanup()
        for key, value in cls._env_backup.items():
            if value is None:
                os.environ.pop(key, None)
            else:
                os.environ[key] = value
        importlib.reload(settings_module)
        importlib.reload(routes_module)
        importlib.reload(main_module)

    def test_health_endpoint_returns_ok(self) -> None:
        response = routes_module.healthcheck()
        self.assertEqual(response["status"], "ok")
        self.assertIn("commit", response)
        self.assertIn("branch", response)

    def test_create_app_does_not_seed_bundled_templates_by_default(self) -> None:
        unseeded_root = self._root / "unseeded"
        templates_dir = unseeded_root / "templates"
        outputs_dir = unseeded_root / "outputs"
        templates_dir.mkdir(parents=True, exist_ok=True)
        outputs_dir.mkdir(parents=True, exist_ok=True)
        os.environ["TEMPLATES_DIR"] = str(templates_dir)
        os.environ["OUTPUTS_DIR"] = str(outputs_dir)
        os.environ["STORAGE_DIR"] = str(unseeded_root)
        os.environ.pop("SEED_BUNDLED_TEMPLATES", None)

        importlib.reload(settings_module)
        importlib.reload(main_module)

        app = main_module.create_app()

        self.assertEqual(app.title, "A3 Presentation API")
        self.assertEqual(list(templates_dir.iterdir()), [])

        os.environ["TEMPLATES_DIR"] = str(self._templates_dir)
        os.environ["OUTPUTS_DIR"] = str(self._outputs_dir)
        os.environ["STORAGE_DIR"] = str(self._root)
        importlib.reload(settings_module)
        importlib.reload(main_module)

    def test_templates_endpoint_lists_available_templates(self) -> None:
        templates = routes_module.list_templates()
        template_ids = {item.template_id for item in templates}
        self.assertIn("corp_light_v1", template_ids)

    def test_template_details_expose_missing_template_file(self) -> None:
        response = routes_module.get_template("missing_source")
        self.assertFalse(response.has_template_file)
        self.assertEqual(response.inventory_summary.usability_status, "not_safely_editable")
        self.assertEqual(response.inventory_summary.generation_mode, "layout")
        self.assertEqual(response.inventory_summary.layout_target_count, 0)
        self.assertEqual(response.inventory_summary.prototype_target_count, 0)
        self.assertEqual(response.editable_targets, [])
        self.assertEqual(response.detected_components, [])

    def test_template_details_reject_path_traversal_template_id(self) -> None:
        with self.assertRaises(HTTPException) as error:
            routes_module.get_template("..\\..\\outside")

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn("escapes the storage root", error.exception.detail)

    def test_extract_text_endpoint_returns_blocks_tables_and_chart_assessments(self) -> None:
        document = Document()
        document.add_paragraph("Отчет по сегментам")
        document.add_heading("Рынок", level=1)
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Сегмент"
        table.cell(0, 1).text = "Выручка"
        table.cell(1, 0).text = "SMB"
        table.cell(1, 1).text = "120"
        table.cell(2, 0).text = "Enterprise"
        table.cell(2, 1).text = "250"
        buffer = BytesIO()
        document.save(buffer)
        upload = UploadFile(filename="report.docx", file=BytesIO(buffer.getvalue()))

        payload = asyncio.run(routes_module.extract_document_text(upload))
        self.assertEqual(payload.file_name, "report.docx")
        self.assertTrue(payload.text)
        self.assertEqual(len(payload.tables), 1)
        self.assertEqual(len(payload.chart_assessments), 1)

    def test_extract_text_endpoint_rejects_non_docx_uploads(self) -> None:
        upload = UploadFile(filename="report.txt", file=BytesIO(b"plain text"))

        with self.assertRaises(HTTPException) as error:
            asyncio.run(routes_module.extract_document_text(upload))

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn(".docx", error.exception.detail)

    def test_extract_text_endpoint_offers_safe_combo_variants_for_mixed_unit_table(self) -> None:
        document = Document()
        document.add_heading("Метрики", level=1)
        table = document.add_table(rows=4, cols=3)
        table.cell(0, 0).text = "Квартал"
        table.cell(0, 1).text = "Выручка"
        table.cell(0, 2).text = "Маржа"
        for row_index, values in enumerate(
            [
                ("Q1", "120 млн руб", "18%"),
                ("Q2", "150 млн руб", "22%"),
                ("Q3", "190 млн руб", "27%"),
            ],
            start=1,
        ):
            for col_index, value in enumerate(values):
                table.cell(row_index, col_index).text = value
        buffer = BytesIO()
        document.save(buffer)
        upload = UploadFile(filename="mixed-units.docx", file=BytesIO(buffer.getvalue()))

        payload = asyncio.run(routes_module.extract_document_text(upload))

        self.assertEqual(len(payload.chart_assessments), 1)
        chart_specs = payload.chart_assessments[0].candidate_specs
        self.assertEqual(
            [spec.chart_type.value for spec in chart_specs],
            ["combo", "combo", "column", "column", "line", "line"],
        )
        self.assertEqual(
            [spec.variant_label for spec in chart_specs],
            [
                "Комбинированный: столбцы Выручка; линия Маржа",
                "Комбинированный: столбцы Маржа; линия Выручка",
                "Единичный: Выручка",
                "Единичный: Маржа",
                "Единичный: Выручка",
                "Единичный: Маржа",
            ],
        )

    def test_extract_text_endpoint_rejects_too_ambiguous_mixed_unit_chart(self) -> None:
        document = Document()
        document.add_heading("Смешанные метрики", level=1)
        table = document.add_table(rows=4, cols=4)
        for col_index, value in enumerate(["Метрика", "Деньги", "Доля", "Количество"]):
            table.cell(0, col_index).text = value
        for row_index, values in enumerate(
            [
                ("A", "120 млн руб", "18%", "25"),
                ("B", "150 млн руб", "22%", "31"),
                ("C", "190 млн руб", "27%", "44"),
            ],
            start=1,
        ):
            for col_index, value in enumerate(values):
                table.cell(row_index, col_index).text = value
        buffer = BytesIO()
        document.save(buffer)
        upload = UploadFile(filename="too-mixed.docx", file=BytesIO(buffer.getvalue()))

        payload = asyncio.run(routes_module.extract_document_text(upload))

        self.assertEqual(len(payload.chart_assessments), 1)
        self.assertFalse(payload.chart_assessments[0].chartable)
        self.assertEqual(payload.chart_assessments[0].candidate_specs, [])

    def test_plan_from_text_returns_presentation_plan(self) -> None:
        payload = routes_module.plan_from_text(
            TextPlanRequest(
                template_id="corp_light_v1",
                title="Demo",
                raw_text="Основные выводы\n- Рост выручки\n- Снижение churn",
            )
        )
        self.assertEqual(payload.template_id, "corp_light_v1")
        self.assertGreaterEqual(len(payload.slides), 1)

    def test_plan_from_text_accepts_transient_uploaded_template_id(self) -> None:
        payload = routes_module.plan_from_text(
            TextPlanRequest(
                template_id="uploaded_customer_template",
                title="Demo",
                raw_text="Основные выводы\n- Рост выручки\n- Снижение churn",
            )
        )
        self.assertEqual(payload.template_id, "uploaded_customer_template")
        self.assertGreaterEqual(len(payload.slides), 1)

    def test_plan_from_text_with_uploaded_template_returns_plan_and_manifest(self) -> None:
        template_path = self._templates_dir / "corp_light_v1" / "template.pptx"
        upload = UploadFile(filename="customer-template.pptx", file=BytesIO(template_path.read_bytes()))
        payload = TextPlanRequest(
            template_id="ignored_template_id",
            title="Demo",
            raw_text="Основные выводы\n- Рост выручки\n- Снижение churn",
        )

        response = asyncio.run(
            routes_module.plan_from_text_with_template(
                payload_json=payload.model_dump_json(),
                template_file=upload,
            )
        )

        self.assertTrue(response.manifest.template_id.startswith("uploaded_customer-template"))
        self.assertEqual(response.plan.template_id, response.manifest.template_id)
        self.assertGreaterEqual(len(response.plan.slides), 1)
        self.assertTrue(response.inventory_summary.targets)
        self.assertIn(response.inventory_summary.usability_status, {"usable", "usable_with_degradation"})
        self.assertTrue(response.editable_targets)
        self.assertTrue(response.detected_components)
        self.assertEqual(response.inventory_summary.generation_mode, response.manifest.generation_mode.value)
        self.assertEqual(len(response.slide_layout_reviews), len(response.plan.slides))
        self.assertTrue(all(review.available_layouts for review in response.slide_layout_reviews))
        text_review = next((review for review in response.slide_layout_reviews if review.slide_index > 0), None)
        self.assertIsNotNone(text_review)
        self.assertEqual(text_review.current_target_key, response.plan.slides[text_review.slide_index].render_target.key)
        self.assertEqual(text_review.current_target_type, response.plan.slides[text_review.slide_index].render_target.type.value)
        self.assertEqual(text_review.current_runtime_profile_key, response.plan.slides[text_review.slide_index].runtime_profile_key)
        best_option = text_review.available_layouts[0]
        self.assertTrue(best_option.runtime_profile_key)
        self.assertTrue(best_option.source_label)
        self.assertTrue(best_option.match_summary)
        self.assertTrue(best_option.recommendation_label)
        self.assertTrue(best_option.recommendation_reasons)
        if best_option.supported_slide_kinds and "text" in best_option.supported_slide_kinds:
            self.assertIsNotNone(best_option.estimated_text_capacity_chars)

    def test_plan_from_text_with_arbitrary_uploaded_template_uses_synthesized_prototype_inventory(self) -> None:
        pptx = Presentation()
        slide = pptx.slides.add_slide(pptx.slide_layouts[6])
        title_shape = slide.shapes.add_textbox(600000, 400000, 8000000, 900000)
        title_shape.text_frame.text = "Произвольный заголовок"
        body_shape = slide.shapes.add_textbox(600000, 1900000, 7600000, 2400000)
        body_shape.text_frame.text = "Первый тезис"
        body_shape.text_frame.add_paragraph().text = "Второй тезис"
        body_shape.text_frame.add_paragraph().text = "Третий тезис"
        buffer = BytesIO()
        pptx.save(buffer)

        upload = UploadFile(filename="arbitrary-template.pptx", file=BytesIO(buffer.getvalue()))
        payload = TextPlanRequest(
            template_id="ignored_template_id",
            title="Demo",
            raw_text="Основные выводы\n- Рост выручки\n- Снижение churn",
        )

        response = asyncio.run(
            routes_module.plan_from_text_with_template(
                payload_json=payload.model_dump_json(),
                template_file=upload,
            )
        )

        self.assertEqual(response.manifest.generation_mode.value, "prototype")
        self.assertTrue(response.manifest.prototype_slides)
        self.assertTrue(response.manifest.inventory.components)
        self.assertTrue(response.manifest.inventory.slides)
        self.assertTrue(response.manifest.inventory.has_prototype_inventory)
        self.assertIn(response.manifest.inventory.degradation_mode, {None, "prototype_only"})
        self.assertTrue(response.inventory_summary.has_prototype_inventory)
        self.assertEqual(response.inventory_summary.usability_status, "usable_with_degradation")
        self.assertTrue(response.inventory_summary.prototype_target_count >= 1)
        self.assertTrue(any(target.source == "prototype" for target in response.inventory_summary.targets))
        self.assertTrue(response.editable_targets)
        self.assertTrue(response.detected_components)
        self.assertTrue(all(review.available_layouts for review in response.slide_layout_reviews))
        self.assertTrue(
            any(option.source == "prototype" for review in response.slide_layout_reviews for option in review.available_layouts)
        )
        self.assertTrue(
            any(
                slide.render_target is not None
                and slide.render_target.key
                and any(target.key == slide.render_target.key and target.source == "prototype" for target in response.inventory_summary.targets)
                for slide in response.plan.slides[1:]
            )
        )
        self.assertTrue(
            all(
                review.current_target_key == response.plan.slides[review.slide_index].render_target.key
                for review in response.slide_layout_reviews
                if response.plan.slides[review.slide_index].render_target is not None
            )
        )
        prototype_option = next(
            option
            for review in response.slide_layout_reviews
            for option in review.available_layouts
            if option.source == "prototype"
        )
        self.assertTrue(prototype_option.source_label)
        self.assertTrue(prototype_option.match_summary)
        self.assertTrue(prototype_option.recommendation_label)
        self.assertTrue(prototype_option.recommendation_reasons)

    def test_slide_layout_reviews_expose_stable_ranking_metadata(self) -> None:
        template_path = self._templates_dir / "corp_light_v1" / "template.pptx"
        upload = UploadFile(filename="customer-template.pptx", file=BytesIO(template_path.read_bytes()))
        payload = TextPlanRequest(
            template_id="ignored_template_id",
            title="Demo",
            raw_text="Основные выводы\n- Рост выручки\n- Снижение churn",
        )

        response = asyncio.run(
            routes_module.plan_from_text_with_template(
                payload_json=payload.model_dump_json(),
                template_file=upload,
            )
        )

        self.assertTrue(response.slide_layout_reviews)
        for review in response.slide_layout_reviews:
            self.assertGreaterEqual(review.slide_index, 0)
            self.assertTrue(review.available_layouts)
            self.assertEqual(review.current_runtime_profile_key, response.plan.slides[review.slide_index].runtime_profile_key)
            for option in review.available_layouts:
                self.assertIn(option.source, {"layout", "prototype"})
                self.assertTrue(option.runtime_profile_key)
                self.assertTrue(option.source_label)
                self.assertTrue(option.match_summary)
                self.assertIn(option.recommendation_label, {"Рекомендуем", "Подходит", "Запасной вариант"})
                self.assertTrue(option.recommendation_reasons)

    def test_generate_and_download_presentation_for_valid_template(self) -> None:
        payload = routes_module.generate_presentation(
            PresentationPlan(
                template_id="corp_light_v1",
                title="Smoke Test",
                slides=[
                    SlideSpec(kind=SlideKind.TITLE, title="Smoke Test", subtitle="API contract"),
                    SlideSpec(kind=SlideKind.TEXT, title="Итог", text="Проверка generate/download через API."),
                ],
            )
        )
        self.assertTrue((self._outputs_dir / payload.file_name).exists())

        download = routes_module.download_presentation(payload.file_name)
        self.assertEqual(
            download.media_type,
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        )
        self.assertEqual(Path(download.path).name, payload.file_name)

    def test_generate_with_uploaded_template_does_not_require_registry_template(self) -> None:
        template_path = self._templates_dir / "corp_light_v1" / "template.pptx"
        upload = UploadFile(filename="custom-template.pptx", file=BytesIO(template_path.read_bytes()))
        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Custom Template Smoke",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Custom Template Smoke", subtitle="Transient upload"),
                SlideSpec(kind=SlideKind.TEXT, title="Итог", text="Генерация через временно загруженный шаблон."),
            ],
        )

        payload = asyncio.run(
            routes_module.generate_presentation_with_template(
                plan_json=plan.model_dump_json(),
                template_file=upload,
            )
        )

        self.assertTrue((self._outputs_dir / payload.file_name).exists())
        self.assertNotIn("custom-template", {item.template_id for item in routes_module.list_templates()})

    def test_generate_returns_404_for_template_without_pptx(self) -> None:
        with self.assertRaises(HTTPException) as error:
            routes_module.generate_presentation(
                PresentationPlan(
                    template_id="missing_source",
                    title="Broken Template",
                    slides=[SlideSpec(kind=SlideKind.TITLE, title="Broken Template")],
                )
            )

        self.assertEqual(error.exception.status_code, 404)
        self.assertIn("Template PPTX not found", error.exception.detail)

    def test_upload_template_rejects_manifest_path_traversal(self) -> None:
        manifest = TemplateManifest(
            template_id="..\\..\\outside",
            display_name="Broken",
            source_pptx="template.pptx",
        )
        upload = UploadFile(
            filename="template.pptx",
            file=BytesIO(b"not-a-real-pptx"),
        )

        with self.assertRaises(HTTPException) as error:
            asyncio.run(
                routes_module.upload_template(
                    manifest_json=manifest.model_dump_json(),
                    template_file=upload,
                )
            )

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn("escapes the storage root", error.exception.detail)

    def test_upload_template_rejects_empty_pptx(self) -> None:
        manifest = TemplateManifest(
            template_id="empty_template",
            display_name="Empty",
            source_pptx="template.pptx",
        )
        upload = UploadFile(filename="template.pptx", file=BytesIO(b""))

        with self.assertRaises(HTTPException) as error:
            asyncio.run(
                routes_module.upload_template(
                    manifest_json=manifest.model_dump_json(),
                    template_file=upload,
                )
            )

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn("template_file is empty", error.exception.detail)

    def test_upload_template_auto_rejects_empty_pptx(self) -> None:
        upload = UploadFile(filename="empty-template.pptx", file=BytesIO(b""))

        with self.assertRaises(HTTPException) as error:
            asyncio.run(
                routes_module.upload_template_auto(
                    template_id="empty_auto",
                    display_name="Empty Auto",
                    description=None,
                    template_file=upload,
                )
            )

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn("template_file is empty", error.exception.detail)

    def test_upload_template_auto_returns_inventory_contract(self) -> None:
        template_path = self._templates_dir / "corp_light_v1" / "template.pptx"
        upload = UploadFile(filename="auto-template.pptx", file=BytesIO(template_path.read_bytes()))

        response = asyncio.run(
            routes_module.upload_template_auto(
                template_id="uploaded_api_auto",
                display_name="Uploaded API Auto",
                description="Contract check",
                template_file=upload,
            )
        )

        self.assertTrue(response.analyzed)
        self.assertIn(response.inventory_summary.usability_status, {"usable", "usable_with_degradation"})
        self.assertTrue(response.editable_targets)
        self.assertTrue(response.detected_components)

    def test_analyze_template_returns_inventory_contract(self) -> None:
        response = routes_module.analyze_template("corp_light_v1")

        self.assertEqual(response.template_id, "corp_light_v1")
        self.assertIn(response.inventory_summary.usability_status, {"usable", "usable_with_degradation"})
        self.assertTrue(response.editable_targets)
        self.assertTrue(response.detected_components)

    def test_plan_from_text_with_template_rejects_invalid_payload_json(self) -> None:
        template_path = self._templates_dir / "corp_light_v1" / "template.pptx"
        upload = UploadFile(filename="customer-template.pptx", file=BytesIO(template_path.read_bytes()))

        with self.assertRaises(HTTPException) as error:
            asyncio.run(
                routes_module.plan_from_text_with_template(
                    payload_json="{not valid json}",
                    template_file=upload,
                )
            )

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn("Invalid payload_json", error.exception.detail)

    def test_plan_from_text_with_template_returns_degradation_error_for_broken_template(self) -> None:
        upload = UploadFile(filename="broken-template.pptx", file=BytesIO(b"not-a-real-pptx"))
        payload = TextPlanRequest(
            template_id="ignored_template_id",
            title="Demo",
            raw_text="Основные выводы\n- Рост выручки\n- Снижение churn",
        )

        with patch.object(routes_module.analyzer, "analyze", side_effect=ValueError("broken template")):
            with self.assertRaises(HTTPException) as error:
                asyncio.run(
                    routes_module.plan_from_text_with_template(
                        payload_json=payload.model_dump_json(),
                        template_file=upload,
                    )
                )

        self.assertEqual(error.exception.status_code, 400)
        self.assertIn("Failed to analyze uploaded template", error.exception.detail)


if __name__ == "__main__":
    unittest.main()
