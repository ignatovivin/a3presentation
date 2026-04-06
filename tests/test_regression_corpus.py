from __future__ import annotations

import base64
import tempfile
import unittest
from io import BytesIO
from pathlib import Path

from docx import Document
from pptx import Presentation

from a3presentation.domain.api import ChartOverride, DocumentBlock
from a3presentation.domain.chart import ChartConfidence, ChartSeries, ChartSpec, ChartType
from a3presentation.domain.presentation import SlideKind
from a3presentation.services.deck_audit import audit_generated_presentation, find_capacity_violations
from a3presentation.services.document_text_extractor import DocumentTextExtractor
from a3presentation.services.planner import TextToPlanService
from a3presentation.services.pptx_generator import PptxGenerator
from a3presentation.services.template_registry import TemplateRegistry
from a3presentation.settings import get_settings


SMALL_PNG_BYTES = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9W6i8AAAAASUVORK5CYII="
)


class RegressionCorpusTests(unittest.TestCase):
    @classmethod
    def setUpClass(cls) -> None:
        cls.fixtures_dir = Path(__file__).parent / "fixtures" / "regression"
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        cls.manifest = registry.get_template("corp_light_v1")
        cls.template_path = registry.get_template_pptx_path("corp_light_v1")

    def test_fixture_text_documents_generate_non_empty_presentations(self) -> None:
        cases = [
            ("strategy_report.md", "Стратегический markdown-документ"),
            ("mixed_notes.txt", "Смешанный текстовый документ"),
        ]
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        for fixture_name, label in cases:
            with self.subTest(case=label):
                content = (self.fixtures_dir / fixture_name).read_bytes()
                text, tables, blocks = extractor.extract(fixture_name, content)
                plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

                self.assertGreaterEqual(len(plan.slides), 2)
                self.assertTrue(any((slide.title or "").strip() for slide in plan.slides[1:]))

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = PptxGenerator().generate(
                        template_path=self.template_path,
                        manifest=self.manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))
                    self.assertEqual(len(presentation.slides), len(plan.slides))

    def test_text_only_markdown_fixture_generates_deck_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        content = (self.fixtures_dir / "strategy_report.md").read_bytes()
        text, tables, blocks = extractor.extract("strategy_report.md", content)
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertTrue(any(audit.kind in {"text", "bullets"} for audit in audits))
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_mixed_text_fixture_generates_deck_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        content = (self.fixtures_dir / "mixed_notes.txt").read_bytes()
        text, tables, blocks = extractor.extract("mixed_notes.txt", content)
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertTrue(any(audit.kind in {"text", "bullets"} for audit in audits))
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_generic_mixed_section_rebalances_continuation_series_without_underfilled_tail(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Партнерская стратегия A3"),
            DocumentBlock(kind="heading", text="Как ограничение управляемого слоя усиливает партнерства", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Если компания ограничивает число глубоко интегрированных поставщиков, но сохраняет широкий каталог, "
                    "она может одновременно показывать рынку масштаб и удерживать качество операционного контура."
                ),
            ),
            DocumentBlock(
                kind="list",
                items=[
                    "Каталог закрывает массовый спрос и остается внешним маркером масштаба для партнеров.",
                    "Глубокий интеграционный слой удерживает SLA, безопасность и предсказуемость изменений.",
                    "Понятный процесс онбординга снижает зависимость от ручного сопровождения.",
                    "Разделение двух контуров делает коммерческий оффер для банков и экосистем понятнее.",
                ],
            ),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "CEO может объяснять модель двумя цифрами: сколько организаций доступно к оплате и сколько поставщиков "
                    "находятся в глубокой интеграции с полным сопровождением."
                ),
            ),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий завершающий блок нужен для фиксации границы между секциями."),
        ]
        raw_text = "\n".join(block.text or "" for block in blocks)
        plan = planner.build_plan("corp_light_v1", raw_text, None, [], blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        relevant = [
            violation
            for violation in find_capacity_violations(audits)
            if violation.rule in {
                "underfilled_continuation",
                "continuation_balance",
                "content_order_mismatch",
                "continuation_order_mismatch",
            }
            and "Как ограничение управляемого слоя усиливает партнерства" in violation.title
        ]
        self.assertEqual(relevant, [])

    def test_generic_mixed_section_keeps_block_order_through_planner_and_generator(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Архитектура партнерств"),
            DocumentBlock(kind="heading", text="Модель взаимодействия с поставщиками", level=1),
            DocumentBlock(
                kind="paragraph",
                text="Первый абзац описывает контекст и должен остаться перед операционными тезисами списка.",
            ),
            DocumentBlock(
                kind="list",
                items=[
                    "Первый тезис описывает каталог и охват массового клиента.",
                    "Второй тезис описывает глубокие интеграции и требования к SLA.",
                ],
            ),
            DocumentBlock(
                kind="paragraph",
                text="Финальный абзац завершает аргументацию и должен остаться после списка, а не перед ним.",
            ),
            DocumentBlock(kind="heading", text="Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткая отбивка следующего раздела."),
        ]
        raw_text = "\n".join(block.text or "" for block in blocks)
        plan = planner.build_plan("corp_light_v1", raw_text, None, [], blocks)

        flattened = []
        for slide in plan.slides:
            if not (slide.title or "").startswith("Модель взаимодействия с поставщиками"):
                continue
            if slide.content_blocks:
                for block in slide.content_blocks:
                    if block.text and block.text.strip():
                        flattened.append(block.text.strip())
                    flattened.extend(item.strip() for item in block.items if item.strip())
                continue
            flattened.extend(item for item in slide.bullets if item.strip())
        self.assertTrue(flattened)
        self.assertLess(
            next(index for index, item in enumerate(flattened) if item.startswith("Первый абзац")),
            next(index for index, item in enumerate(flattened) if item.startswith("Первый тезис")),
        )
        self.assertLess(
            next(index for index, item in enumerate(flattened) if item.startswith("Второй тезис")),
            next(index for index, item in enumerate(flattened) if item.startswith("Финальный абзац")),
        )

    def test_generated_docx_corpus_covers_report_form_resume_and_table_heavy(self) -> None:
        cases = [
            ("report.docx", self._build_report_docx()),
            ("form.docx", self._build_form_docx()),
            ("resume.docx", self._build_resume_docx()),
            ("table-heavy.docx", self._build_table_heavy_docx()),
        ]
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        for filename, content in cases:
            with self.subTest(case=filename):
                text, tables, blocks = extractor.extract(filename, content)
                plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

                self.assertGreaterEqual(len(plan.slides), 2)
                self.assertFalse(all(slide.kind.value == "title" for slide in plan.slides))

                with tempfile.TemporaryDirectory() as temp_dir:
                    output_path = PptxGenerator().generate(
                        template_path=self.template_path,
                        manifest=self.manifest,
                        plan=plan,
                        output_dir=Path(temp_dir),
                    )
                    presentation = Presentation(str(output_path))
                    self.assertEqual(len(presentation.slides), len(plan.slides))

    def test_form_like_docx_generates_deck_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("form.docx", self._build_form_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertGreaterEqual(len(plan.slides), 2)
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_resume_like_docx_generates_deck_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("resume.docx", self._build_resume_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertGreaterEqual(len(plan.slides), 2)
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_table_heavy_docx_generates_deck_without_text_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("table-heavy.docx", self._build_table_heavy_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertTrue(any(slide.kind == SlideKind.TABLE for slide in plan.slides))
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])
        table_audits = [audit for audit in audits if audit.kind == SlideKind.TABLE.value]
        self.assertTrue(table_audits)
        self.assertTrue(all(audit.has_table for audit in table_audits))
        self.assertTrue(all(audit.content_width_ratio >= 0.9 for audit in table_audits))
        self.assertTrue(all(audit.footer_width_ratio >= 0.9 for audit in table_audits))

    def test_fact_only_docx_generates_appendix_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("fact-only.docx", self._build_fact_only_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertTrue(any(slide.title and "Приложение" in slide.title for slide in plan.slides))
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_image_heavy_docx_generates_image_slide_and_preserves_text_capacity_contract(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Отчет по продукту"),
            DocumentBlock(kind="heading", text="1. Схема процесса", level=1),
            DocumentBlock(kind="paragraph", text="Ниже приведена ключевая схема целевого процесса."),
            DocumentBlock(
                kind="image",
                text="Схема целевого процесса",
                image_name="process.png",
                image_content_type="image/png",
                image_base64=base64.b64encode(SMALL_PNG_BYTES).decode("ascii"),
            ),
            DocumentBlock(kind="heading", text="2. Выводы", level=1),
            DocumentBlock(kind="paragraph", text="Иллюстрация подтверждает узкие места и точки автоматизации."),
        ]
        raw_text = "\n".join(block.text or "" for block in blocks)
        plan = planner.build_plan("corp_light_v1", raw_text, None, [], blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            audits = audit_generated_presentation(output_path, plan)

        self.assertTrue(any(slide.kind == SlideKind.IMAGE for slide in plan.slides))
        self.assertTrue(any(getattr(shape, "has_text_frame", False) for slide in presentation.slides for shape in slide.shapes))
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_cover_skip_does_not_swallow_numbered_first_section_with_image(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Отчет по продукту"),
            DocumentBlock(kind="heading", text="1. Схема процесса", level=1),
            DocumentBlock(kind="paragraph", text="Ниже приведена ключевая схема целевого процесса."),
            DocumentBlock(
                kind="image",
                text="Схема целевого процесса",
                image_name="process.png",
                image_content_type="image/png",
                image_base64=base64.b64encode(SMALL_PNG_BYTES).decode("ascii"),
            ),
            DocumentBlock(kind="heading", text="2. Выводы", level=1),
            DocumentBlock(kind="paragraph", text="Иллюстрация подтверждает узкие места и точки автоматизации."),
        ]

        plan = planner.build_plan("corp_light_v1", "\n".join(block.text or "" for block in blocks), None, [], blocks)

        self.assertTrue(any(slide.kind == SlideKind.IMAGE for slide in plan.slides))
        self.assertTrue(any((slide.title or "").startswith("1. Схема процесса") for slide in plan.slides))

    def test_strategy_like_docx_skips_cover_lines_as_first_content_section(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("report.docx", self._build_report_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        self.assertGreaterEqual(len(plan.slides), 2)
        self.assertNotEqual(plan.slides[1].title, "A3")
        self.assertTrue(any((slide.title or "").strip() for slide in plan.slides[1:]))

    def test_strategy_edge_case_docx_preserves_numbered_sections_and_mixed_content(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("strategy-edge.docx", self._build_strategy_edge_case_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        self.assertFalse(any(slide.title == "А3" for slide in plan.slides[1:]))
        self.assertTrue(any((slide.title or "").startswith("1.4 Показатели 2025") for slide in plan.slides))
        self.assertTrue(any("3.3 Карта конкурентов" in ((slide.title or "") + " " + (slide.subtitle or "")) for slide in plan.slides))

        slide_payload = "\n".join(
            " ".join(
                part
                for part in (
                    slide.title or "",
                    slide.subtitle or "",
                    slide.text or "",
                    slide.notes or "",
                    "\n".join(slide.bullets),
                )
                if part
            )
            for slide in plan.slides
        )
        self.assertIn("Структура: 5 продажников", slide_payload)
        self.assertIn("Ритм контактов: регулярные встречи", slide_payload)
        self.assertIn("Q2 2026 (Discovery): Анализ архитектуры платформы цифрового рубля ЦБ", slide_payload)

    def test_strategy_edge_case_docx_generates_deck_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("strategy-edge.docx", self._build_strategy_edge_case_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        self.assertGreaterEqual(len(audits), 3)
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_report_docx_generates_deck_without_capacity_violations(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("report.docx", self._build_report_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            audits = audit_generated_presentation(output_path, plan)

        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])

    def test_report_docx_prefers_text_flow_for_narrative_sections(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("narrative-report.docx", self._build_narrative_report_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        narrative_slides = [slide for slide in plan.slides if (slide.title or "").startswith("1. Контекст рынка")]
        self.assertTrue(narrative_slides)
        self.assertTrue(all(slide.kind == SlideKind.TEXT for slide in narrative_slides))
        self.assertGreaterEqual(len(narrative_slides), 2)

    def test_report_docx_skips_reference_tail_from_main_deck(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("report-with-refs.docx", self._build_report_with_reference_tail_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        payload = "\n".join(
            part
            for slide in plan.slides
            for part in [slide.title or "", slide.subtitle or "", slide.text or "", slide.notes or "", *slide.bullets]
            if part
        )
        self.assertNotIn("https://example.com/source-1", payload)
        self.assertNotIn("https://example.com/source-2", payload)
        self.assertTrue(any((slide.title or "").startswith("2. Выводы") for slide in plan.slides))

    def test_report_docx_does_not_add_appendix_from_false_semantic_facts(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        text, tables, blocks = extractor.extract("report-with-refs.docx", self._build_report_with_reference_tail_docx())
        plan = planner.build_plan("corp_light_v1", text, None, tables, blocks)

        self.assertFalse(any("Приложение" in (slide.title or "") for slide in plan.slides))

    def test_chart_heavy_docx_generates_chart_slide_and_preserves_text_capacity_contract(self) -> None:
        extractor = DocumentTextExtractor()
        planner = TextToPlanService()

        content = self._build_chart_heavy_docx()
        text, tables, blocks = extractor.extract("chart-heavy.docx", content)
        self.assertGreaterEqual(len(tables), 1)

        plan = planner.build_plan(
            "corp_light_v1",
            text,
            None,
            tables,
            blocks,
            chart_overrides=[
                ChartOverride(
                    table_id="table_1",
                    mode="chart",
                    selected_chart=ChartSpec(
                        chart_id="chart_1",
                        source_table_id="table_1",
                        chart_type=ChartType.COLUMN,
                        title="Выручка по каналам",
                        categories=["SEO", "Ads", "Referral"],
                        series=[ChartSeries(name="Выручка", values=[120.0, 200.0, 90.0])],
                        confidence=ChartConfidence.HIGH,
                    ),
                )
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=self.template_path,
                manifest=self.manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            audits = audit_generated_presentation(output_path, plan)

        self.assertTrue(any(slide.kind == SlideKind.CHART for slide in plan.slides))
        self.assertTrue(any(getattr(shape, "has_chart", False) for slide in presentation.slides for shape in slide.shapes))
        violations = find_capacity_violations(audits)
        self.assertEqual(violations, [])
        chart_audits = [audit for audit in audits if audit.kind == SlideKind.CHART.value]
        self.assertTrue(chart_audits)
        self.assertTrue(all(audit.has_chart for audit in chart_audits))
        self.assertTrue(all(audit.content_width_ratio >= 0.9 for audit in chart_audits))
        self.assertTrue(all(audit.footer_width_ratio >= 0.9 for audit in chart_audits))

    def test_planner_compresses_short_text_and_bullet_continuations(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="3.3 Карта конкурентов", level=1),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "А3 оперирует в уникальной нише на стыке инфраструктурных финтех-решений для регулярных платежей. "
                    "Конкурентное поле формируют косвенные игроки, которые закрывают отдельные сегменты цепочки ценности, "
                    "но не воспроизводят модель А3 целиком."
                ),
            ),
            DocumentBlock(kind="paragraph", text="Конкурентные преимущества А3"),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Ни один конкурент не воспроизводит модель А3 целиком: мультибанковская сеть и агрегация поставщиков. "
                    "Главные риски — Сбер и ВТБ, которые строят собственную инфраструктуру. "
                    "Митигация: развитие уникальных технологий, углубление экосистемы партнерств и запуск РНКО."
                ),
            ),
            DocumentBlock(kind="heading", text="4. Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Короткий текст для следующего раздела."),
        ]

        plan = planner.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="A3 Presentation",
            blocks=blocks,
        )

        competitor_slides = [slide for slide in plan.slides if (slide.title or "").startswith("3.3 Карта конкурентов")]
        self.assertLessEqual(len(competitor_slides), 2)
        self.assertTrue(all(slide.kind == SlideKind.TEXT for slide in competitor_slides))
        payload = " ".join(
            part
            for slide in competitor_slides
            for part in [slide.text or "", slide.notes or ""]
            if part
        )
        self.assertIn("Конкурентные преимущества А3", payload)
        self.assertIn("Главные риски", payload)

    def test_planner_keeps_short_quarter_plan_on_single_slide(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="Q2 2026", level=1),
            DocumentBlock(kind="list", items=[
                "Перевод всех партнеров на автономный режим работы",
                "Перенастройка сетей",
                "Выделение серверов под реплики БД",
                "k8s pci",
                "k8s invoice",
                "Перенос hadoop",
                "Решение проблемы с кросс-цодовой генерацией id транзакции",
                "Перевод всех партнеров на автономку",
            ]),
            DocumentBlock(kind="heading", text="Q3 2026", level=1),
            DocumentBlock(kind="paragraph", text="Следующий раздел."),
        ]

        plan = planner.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="A3 Presentation",
            blocks=blocks,
        )

        q2_slides = [slide for slide in plan.slides if (slide.title or "").startswith("Q2 2026")]
        self.assertEqual(len(q2_slides), 1)
        self.assertEqual(q2_slides[0].kind, SlideKind.BULLETS)
        self.assertEqual(len(q2_slides[0].bullets), 8)

    def test_planner_reduces_rd_section_to_two_slides_when_tail_is_short(self) -> None:
        planner = TextToPlanService()
        blocks = [
            DocumentBlock(kind="title", text="A3 Presentation", level=0),
            DocumentBlock(kind="heading", text="5.3 R&D новые продукты", level=1),
            DocumentBlock(
                kind="paragraph",
                text="Цель: построить системный пайплайн проверки гипотез и новых продуктов для роста выручки и диверсификации бизнеса А3.",
            ),
            DocumentBlock(
                kind="paragraph",
                text=(
                    "Q2 2026 (Discovery): Анализ архитектуры платформы цифрового рубля ЦБ. "
                    "Определение роли А3. Встречи с ЦБ и банками-партнерами."
                ),
            ),
            DocumentBlock(
                kind="list",
                items=[
                    "Процесс Double Diamond запущен, первые инициативы в работе",
                    "Банк идей = единый бэклог, синхронизирован с Jira",
                    "ICE-приоритизация со стейкхолдерами раз в 2 недели",
                    "Юнит-экономика: по поставщикам/интеграциям, по партнёрам, по продуктовым направлениям",
                    "Инициатива: Инфраструктура цифрового рубля",
                    "Q3 2026 (Proof of Concept): Техническое прототипирование с командой Payments",
                    "Q4 2026 — Пилот: Пилотный запуск с ограниченным объемом, тестирование с 1 банком-партнером, оценка unit-экономики и масштабируемости",
                ],
            ),
            DocumentBlock(kind="heading", text="6. Следующий раздел", level=1),
            DocumentBlock(kind="paragraph", text="Следующий раздел."),
        ]

        plan = planner.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            title="A3 Presentation",
            blocks=blocks,
        )

        rd_slides = [slide for slide in plan.slides if (slide.title or "").startswith("5.3 R&D новые продукты")]
        self.assertLessEqual(len(rd_slides), 2)
        self.assertGreaterEqual(len(rd_slides[0].bullets), 6)

    def _build_report_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("A3")
        document.add_paragraph("Бизнес-стратегия 2026")
        document.add_heading("1. Vision и стратегические цели", level=1)
        document.add_heading("1.1 Контекст роста", level=2)
        document.add_paragraph(
            "Платформа должна расти за счет новых сегментов, устойчивой экономики продукта и масштабируемой архитектуры."
        )
        document.add_paragraph(
            "Рынок требует масштабируемой инфраструктуры, диверсификации выручки и устойчивого продуктового позиционирования."
        )
        document.add_heading("2. Позиционирование", level=1)
        document.add_paragraph(
            "Компания должна восприниматься как инфраструктурный игрок, ускоряющий финансовые процессы клиентов."
        )
        return self._save_document(document)

    def _build_narrative_report_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("A3")
        document.add_paragraph("Стратегический обзор")
        document.add_heading("1. Контекст рынка", level=1)
        document.add_paragraph(
            "Рынок регулярных платежей растет за счет цифровых каналов, но одновременно становится чувствительным к "
            "устойчивости SLA, скорости вывода новых поставщиков и способности платформы масштабировать сопровождение "
            "без линейного роста ручных операций."
        )
        document.add_paragraph(
            "Для A3 это означает необходимость разделять публичный охват каталога, который важен для внешнего позиционирования, "
            "и управляемый слой глубоких интеграций, где каждая новая связь увеличивает нагрузку на процессы, безопасность и "
            "юридический контур."
        )
        document.add_paragraph(
            "Если не разделять эти два контура, управленческие решения начинают опираться на смешанную метрику, а длинный "
            "хвост малодеятельных поставщиков маскирует реальные ограничения по unit-экономике, качеству сервиса и "
            "операционной устойчивости."
        )
        document.add_heading("2. Выводы", level=1)
        document.add_paragraph("Нарративный отчет должен оставаться текстовым, а не распадаться на список из предложений.")
        return self._save_document(document)

    def _build_report_with_reference_tail_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("A3")
        document.add_paragraph("Стратегический обзор")
        document.add_heading("1. Контекст", level=1)
        document.add_paragraph(
            "Основной текст раздела должен попасть в презентацию, а технический хвост со ссылками не должен раздувать колоду."
        )
        document.add_heading("2. Выводы", level=1)
        document.add_paragraph("Нужно сохранить только смысловые выводы и убрать голые ссылки из основной части презентации.")
        document.add_paragraph("[1] https://example.com/source-1")
        document.add_paragraph("https://example.com/source-1")
        document.add_paragraph("[2] https://example.com/source-2")
        document.add_paragraph("https://example.com/source-2")
        return self._save_document(document)

    def _build_form_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("Анкета нового сотрудника")
        document.add_paragraph("ФИО: Иван Игнатов")
        document.add_paragraph("Дата рождения: 05.09.1994")
        document.add_paragraph("Телефон: +7 999 000-00-00")
        document.add_paragraph("Email: ivan@example.com")
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Поле"
        table.cell(0, 1).text = "Значение"
        table.cell(1, 0).text = "Город"
        table.cell(1, 1).text = "Екатеринбург"
        table.cell(2, 0).text = "Отдел"
        table.cell(2, 1).text = "Разработка"
        return self._save_document(document)

    def _build_resume_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("Иван Игнатов")
        document.add_paragraph("ivan@example.com")
        document.add_paragraph("+7 999 000-00-00")
        document.add_paragraph("ОПЫТ РАБОТЫ")
        document.add_paragraph("Руководил развитием внутренних платформ и продуктовой аналитики.")
        document.add_paragraph("НАВЫКИ")
        document.add_paragraph("Стратегия, аналитика, управление продуктом, интеграции.")
        return self._save_document(document)

    def _build_table_heavy_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("Сводка KPI по сегментам")
        table = document.add_table(rows=5, cols=3)
        headers = ["Сегмент", "Выручка", "Маржа"]
        for col_index, header in enumerate(headers):
            table.cell(0, col_index).text = header
        rows = [
            ("SMB", "120", "18%"),
            ("Enterprise", "260", "24%"),
            ("PSP", "180", "21%"),
            ("Acquiring", "140", "17%"),
        ]
        for row_index, row in enumerate(rows, start=1):
            for col_index, value in enumerate(row):
                table.cell(row_index, col_index).text = value
        document.add_picture(BytesIO(SMALL_PNG_BYTES))
        return self._save_document(document)

    def _build_chart_heavy_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("Отчёт по каналам продаж")
        document.add_heading("1. Итоги квартала", level=1)
        document.add_paragraph(
            "Основной рост пришёл из платных каналов и рекомендаций, при этом SEO сохранил устойчивую базу лидов."
        )
        table = document.add_table(rows=4, cols=2)
        table.cell(0, 0).text = "Канал"
        table.cell(0, 1).text = "Выручка"
        table.cell(1, 0).text = "SEO"
        table.cell(1, 1).text = "120"
        table.cell(2, 0).text = "Ads"
        table.cell(2, 1).text = "200"
        table.cell(3, 0).text = "Referral"
        table.cell(3, 1).text = "90"
        document.add_heading("2. Вывод", level=1)
        document.add_paragraph("Нужно перераспределить бюджет в пользу каналов с лучшей окупаемостью.")
        return self._save_document(document)

    def _build_fact_only_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("Паспорт клиента")
        document.add_paragraph("ФИО: Иван Игнатов")
        document.add_paragraph("Дата: 05.09.2025")
        document.add_paragraph("Email: ivan@example.com")
        document.add_paragraph("Телефон: +7 999 000-00-00")
        return self._save_document(document)

    def _build_image_heavy_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("Отчет по продукту")
        document.add_heading("1. Схема процесса", level=1)
        document.add_paragraph("Ниже приведена ключевая схема целевого процесса.")
        document.add_picture(BytesIO(SMALL_PNG_BYTES))
        document.add_heading("2. Выводы", level=1)
        document.add_paragraph("Иллюстрация подтверждает узкие места и точки автоматизации.")
        return self._save_document(document)

    def _build_strategy_edge_case_docx(self) -> bytes:
        document = Document()
        document.add_paragraph("А3")
        document.add_paragraph("Бизнес-стратегия 2026")
        document.add_paragraph("Горизонт планирования: 2026-2030")
        document.add_paragraph("Март 2026")
        document.add_paragraph("Конфиденциальный документ")

        document.add_heading("1. Vision и стратегические цели", level=1)
        document.add_heading("1.1 Vision 2030", level=2)
        document.add_paragraph("Стать одной из ведущих продуктовых финтех IT-компаний в России.")
        document.add_paragraph("1.4 Показатели 2025")
        table = document.add_table(rows=3, cols=2)
        table.cell(0, 0).text = "Показатель"
        table.cell(0, 1).text = "Значение"
        table.cell(1, 0).text = "Выручка"
        table.cell(1, 1).text = "1 678 млн ₽"
        table.cell(2, 0).text = "Партнёры"
        table.cell(2, 1).text = "49"

        document.add_heading("3. Анализ рынка и SWOT", level=1)
        document.add_paragraph("3.3 Карта конкурентов")
        document.add_paragraph(
            "А3 оперирует в уникальной нише на стыке инфраструктурных финтех-решений для регулярных платежей."
        )
        document.add_paragraph(
            "Митигация: развитие уникальных технологий, углубление экосистемы партнерств и запуск РНКО."
        )

        document.add_heading("4. Go-to-market и продажи", level=1)
        document.add_heading("4.3 Account Management", level=2)
        document.add_paragraph(
            "Структура: 5 продажников (новые сделки, upsell, стратегические встречи) + 5 аккаунт-менеджеров."
        )
        document.add_paragraph("Ритм контактов: регулярные встречи, мероприятия, рассылки.")
        document.add_paragraph(style="List Bullet").add_run(
            "Сегментация портфеля: Топ-10 (ядро), Средние 20-30, Хвост"
        )
        document.add_paragraph(style="List Bullet").add_run(
            "Upsell-план: для каждого банка хвоста - конкретный следующий продукт и квартал"
        )

        document.add_heading("5. Продуктовый роадмап 2026", level=1)
        document.add_heading("5.3 R&D новые продукты", level=2)
        document.add_paragraph(
            "Цель: построить системный пайплайн проверки гипотез и новых продуктов для роста выручки."
        )
        document.add_paragraph(style="List Bullet").add_run(
            "Процесс Double Diamond запущен, первые инициативы в работе"
        )
        document.add_paragraph(style="List Bullet").add_run(
            "Банк идей = единый бэклог, синхронизирован с Jira"
        )
        document.add_paragraph(
            "Q2 2026 (Discovery): Анализ архитектуры платформы цифрового рубля ЦБ и определение роли А3."
        )
        document.add_paragraph(style="List Bullet").add_run(
            "Q3 2026 (Proof of Concept): Техническое прототипирование с командой Payments"
        )

        return self._save_document(document)

    def _save_document(self, document: Document) -> bytes:
        buffer = BytesIO()
        document.save(buffer)
        return buffer.getvalue()


if __name__ == "__main__":
    unittest.main()
