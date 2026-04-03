from __future__ import annotations

import tempfile
import unittest
from pathlib import Path

from pptx import Presentation

from a3presentation.domain.api import DocumentBlock
from a3presentation.domain.presentation import PresentationPlan, SlideKind, SlideSpec, TableBlock
from a3presentation.domain.semantic import DocumentKind
from a3presentation.services.planner import TextToPlanService
from a3presentation.services.pptx_generator import PptxGenerator
from a3presentation.services.semantic_normalizer import SemanticDocumentNormalizer
from a3presentation.services.template_registry import TemplateRegistry
from a3presentation.settings import get_settings


SMALL_PNG_BASE64 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO9W6i8AAAAASUVORK5CYII="
)


class SemanticPipelineTests(unittest.TestCase):
    def test_normalizer_extracts_facts_contacts_dates_and_kind(self) -> None:
        normalizer = SemanticDocumentNormalizer()
        blocks = [
            DocumentBlock(kind="paragraph", text="Иван Игнатов"),
            DocumentBlock(kind="paragraph", text="ivan@example.com"),
            DocumentBlock(kind="paragraph", text="+7 999 000-00-00"),
            DocumentBlock(kind="paragraph", text="Дата: 05.09.2025"),
            DocumentBlock(kind="paragraph", text="ОПЫТ РАБОТЫ"),
            DocumentBlock(kind="paragraph", text="Образование: Высшее"),
            DocumentBlock(kind="paragraph", text="НАВЫКИ"),
        ]

        semantic_document = normalizer.normalize(
            raw_text="\n".join(block.text or "" for block in blocks),
            blocks=blocks,
            tables=[],
        )

        self.assertEqual(semantic_document.kind, DocumentKind.RESUME)
        self.assertTrue(any("ivan@example.com" in item for item in semantic_document.contacts))
        self.assertIn("05.09.2025", semantic_document.dates)
        self.assertTrue(any(fact.label == "Дата" for fact in semantic_document.facts))

    def test_normalizer_extracts_image_blocks(self) -> None:
        normalizer = SemanticDocumentNormalizer()
        blocks = [
            DocumentBlock(
                kind="image",
                text="Схема процесса",
                image_name="process.png",
                image_content_type="image/png",
                image_base64=SMALL_PNG_BASE64,
            )
        ]

        semantic_document = normalizer.normalize(raw_text="[Image] Схема процесса", blocks=blocks, tables=[])
        self.assertEqual(len(semantic_document.images), 1)
        self.assertEqual(semantic_document.images[0].content_type, "image/png")

    def test_strategy_document_with_tables_stays_report_not_form(self) -> None:
        normalizer = SemanticDocumentNormalizer()
        blocks = [
            DocumentBlock(kind="paragraph", text="A3"),
            DocumentBlock(kind="paragraph", text="Бизнес-стратегия 2026"),
            DocumentBlock(kind="heading", text="1. Vision и стратегические цели", level=1),
            DocumentBlock(kind="subheading", text="1.1 Vision 2029", level=2),
            DocumentBlock(
                kind="paragraph",
                text="Стать одной из ведущих продуктовых финтех IT-компаний в России за счет устойчивого роста, диверсификации выручки и масштабируемой инфраструктуры.",
            ),
            DocumentBlock(kind="subheading", text="1.2 Количественные параметры Vision", level=2),
            DocumentBlock(kind="table", table=TableBlock(headers=["Показатель", "Цель"], rows=[["Выручка", "x2"]])),
            DocumentBlock(kind="heading", text="2. Позиционирование и бренд", level=1),
            DocumentBlock(kind="subheading", text="2.1 Роль бренда", level=2),
            DocumentBlock(
                kind="paragraph",
                text="Инфраструктурный игрок, объединяющий участников финансового рынка и ускоряющий ключевые процессы за счет доступа к данным и платформенным возможностям.",
            ),
        ]

        semantic_document = normalizer.normalize(
            raw_text="\n".join(block.text or "" for block in blocks),
            blocks=blocks,
            tables=[TableBlock(headers=["Показатель", "Цель"], rows=[["Выручка", "x2"]])],
        )
        self.assertIn(semantic_document.kind, {DocumentKind.REPORT, DocumentKind.MIXED})

    def test_planner_creates_image_slide_from_semantic_content(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Отчет по продукту"),
            DocumentBlock(
                kind="image",
                text="Схема целевого процесса",
                image_name="process.png",
                image_content_type="image/png",
                image_base64=SMALL_PNG_BASE64,
            ),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            blocks=blocks,
        )

        self.assertTrue(any(slide.kind == SlideKind.IMAGE for slide in plan.slides))

    def test_generator_renders_image_slide_without_crashing(self) -> None:
        settings = get_settings()
        registry = TemplateRegistry(settings.templates_dir)
        manifest = registry.get_template("corp_light_v1")
        template_path = registry.get_template_pptx_path("corp_light_v1")

        plan = PresentationPlan(
            template_id="corp_light_v1",
            title="Image Presentation",
            slides=[
                SlideSpec(kind=SlideKind.TITLE, title="Image Presentation", preferred_layout_key="cover"),
                SlideSpec(
                    kind=SlideKind.IMAGE,
                    title="Схема процесса",
                    text="Ключевая иллюстрация из документа.",
                    preferred_layout_key="image_text",
                    image_base64=SMALL_PNG_BASE64,
                    image_content_type="image/png",
                ),
            ],
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            output_path = PptxGenerator().generate(
                template_path=template_path,
                manifest=manifest,
                plan=plan,
                output_dir=Path(temp_dir),
            )
            presentation = Presentation(str(output_path))
            self.assertEqual(len(presentation.slides), 2)

    def test_hard_safety_rules_add_appendix_for_fact_only_document(self) -> None:
        service = TextToPlanService()
        blocks = [
            DocumentBlock(kind="paragraph", text="Паспорт"),
            DocumentBlock(kind="paragraph", text="ФИО: Иван Игнатов"),
            DocumentBlock(kind="paragraph", text="Дата: 05.09.2025"),
            DocumentBlock(kind="paragraph", text="ivan@example.com"),
        ]

        plan = service.build_plan(
            template_id="corp_light_v1",
            raw_text="\n".join(block.text or "" for block in blocks),
            blocks=blocks,
        )

        self.assertTrue(any(slide.title and "Приложение" in slide.title for slide in plan.slides))


if __name__ == "__main__":
    unittest.main()
